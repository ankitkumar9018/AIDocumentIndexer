"""
PPTX Native Chart Generator
============================

Generates native PowerPoint charts using python-pptx.
Supports column, bar, line, pie, and area charts with theme colors.
"""

import re
from typing import List, Optional, Tuple, Dict, Any
from enum import Enum

import structlog
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, ChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR

logger = structlog.get_logger(__name__)


class ChartType(str, Enum):
    """Supported chart types."""
    COLUMN = "column"
    BAR = "bar"
    LINE = "line"
    PIE = "pie"
    AREA = "area"
    CLUSTERED_BAR = "clustered_bar"
    STACKED_COLUMN = "stacked_column"
    DOUGHNUT = "doughnut"
    SCATTER = "scatter"
    RADAR = "radar"
    LINE_MARKERS = "line_markers"
    STACKED_BAR = "stacked_bar"
    STACKED_AREA = "stacked_area"


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


class PPTXNativeChartGenerator:
    """
    Generator for native PowerPoint charts.

    Creates professional, theme-colored charts from extracted data.
    """

    # Mapping of chart type enum to python-pptx chart type
    CHART_TYPE_MAP = {
        ChartType.COLUMN: XL_CHART_TYPE.COLUMN_CLUSTERED,
        ChartType.BAR: XL_CHART_TYPE.BAR_CLUSTERED,
        ChartType.LINE: XL_CHART_TYPE.LINE,
        ChartType.PIE: XL_CHART_TYPE.PIE,
        ChartType.AREA: XL_CHART_TYPE.AREA,
        ChartType.CLUSTERED_BAR: XL_CHART_TYPE.BAR_CLUSTERED,
        ChartType.STACKED_COLUMN: XL_CHART_TYPE.COLUMN_STACKED,
        ChartType.DOUGHNUT: XL_CHART_TYPE.DOUGHNUT,
        ChartType.SCATTER: XL_CHART_TYPE.XY_SCATTER,
        ChartType.RADAR: XL_CHART_TYPE.RADAR,
        ChartType.LINE_MARKERS: XL_CHART_TYPE.LINE_MARKERS,
        ChartType.STACKED_BAR: XL_CHART_TYPE.BAR_STACKED,
        ChartType.STACKED_AREA: XL_CHART_TYPE.AREA_STACKED,
    }

    def __init__(self, theme_colors: Dict[str, str]):
        """
        Initialize chart generator with theme colors.

        Args:
            theme_colors: Dict with 'primary', 'secondary', 'accent', 'text' colors
        """
        self.theme_colors = theme_colors
        self.primary_color = hex_to_rgb(theme_colors.get('primary', '#3D5A80'))
        self.secondary_color = hex_to_rgb(theme_colors.get('secondary', '#98C1D9'))
        self.accent_color = hex_to_rgb(theme_colors.get('accent', '#EE6C4D'))
        self.text_color = hex_to_rgb(theme_colors.get('text', '#2D3748'))

        # Chart color palette derived from theme
        self.chart_colors = [
            self.primary_color,
            self.secondary_color,
            self.accent_color,
            hex_to_rgb('#66C2A5'),  # Teal
            hex_to_rgb('#FC8D62'),  # Orange
            hex_to_rgb('#8DA0CB'),  # Blue-gray
        ]

    def add_chart_to_slide(
        self,
        slide,
        chart_data: "ChartDataStructure",
        chart_type: ChartType = ChartType.COLUMN,
        left: Optional[int] = None,
        top: Optional[int] = None,
        width: Optional[int] = None,
        height: Optional[int] = None,
        title: Optional[str] = None,
        show_legend: bool = True,
        show_data_labels: bool = True,
    ):
        """
        Add a formatted chart to a slide.

        Args:
            slide: pptx.slide.Slide object
            chart_data: ChartDataStructure with categories and series
            chart_type: Type of chart to create
            left: Left position (defaults to centered)
            top: Top position
            width: Chart width
            height: Chart height
            title: Chart title
            show_legend: Whether to show legend
            show_data_labels: Whether to show data labels on chart elements

        Returns:
            The created chart shape
        """
        # Default positioning
        if left is None:
            left = Inches(0.75)
        if top is None:
            top = Inches(2.0)
        if width is None:
            width = Inches(8.5)
        if height is None:
            height = Inches(4.5)

        # Get chart type
        xl_chart_type = self.CHART_TYPE_MAP.get(
            chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED
        )

        # Create chart data - scatter charts use XyChartData, others use CategoryChartData
        if chart_type == ChartType.SCATTER:
            pptx_chart_data = XyChartData()
            # For scatter charts, categories are X values and series values are Y values
            for series in chart_data.series:
                xy_series = pptx_chart_data.add_series(series.name)
                # Pair up categories (as numeric X) with values (Y)
                for i, (x_val, y_val) in enumerate(zip(chart_data.categories, series.values)):
                    # Try to convert category to numeric for X axis
                    try:
                        x_numeric = float(x_val) if isinstance(x_val, (int, float)) else float(i + 1)
                    except (ValueError, TypeError):
                        x_numeric = float(i + 1)
                    xy_series.add_data_point(x_numeric, y_val)
        else:
            pptx_chart_data = CategoryChartData()
            pptx_chart_data.categories = chart_data.categories
            for series in chart_data.series:
                pptx_chart_data.add_series(series.name, series.values)

        # Add chart to slide
        chart_shape = slide.shapes.add_chart(
            xl_chart_type, left, top, width, height, pptx_chart_data
        )
        chart = chart_shape.chart

        # Style the chart
        self._style_chart(chart, title, show_legend, chart_type, show_data_labels)

        logger.info(
            "Created chart",
            type=chart_type.value,
            categories=len(chart_data.categories),
            series=len(chart_data.series),
            data_labels=show_data_labels,
        )
        return chart_shape

    def _style_chart(
        self,
        chart,
        title: Optional[str],
        show_legend: bool,
        chart_type: ChartType,
        show_data_labels: bool = True,
    ):
        """Apply theme styling to the chart."""
        from pptx.enum.chart import XL_DATA_LABEL_POSITION

        # Title
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
            chart.chart_title.text_frame.paragraphs[0].font.bold = True
        else:
            chart.has_title = False

        # Chart types that show legend on the right (circular charts)
        circular_charts = {ChartType.PIE, ChartType.DOUGHNUT}

        # Legend positioning based on chart type
        if chart_type in circular_charts:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
        elif chart_type == ChartType.RADAR:
            # Radar charts benefit from bottom legend
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        elif show_legend:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        else:
            chart.has_legend = False

        # Apply colors to series based on chart type
        try:
            for i, series in enumerate(chart.series):
                color = self.chart_colors[i % len(self.chart_colors)]

                if chart_type in {ChartType.LINE, ChartType.LINE_MARKERS}:
                    # Line charts - color the line
                    series.format.line.color.rgb = color
                    series.format.line.width = Pt(2.5)
                    # Add markers for LINE_MARKERS
                    if chart_type == ChartType.LINE_MARKERS:
                        try:
                            series.marker.style = 2  # Circle marker
                            series.marker.size = 8
                        except Exception:
                            pass
                elif chart_type == ChartType.SCATTER:
                    # Scatter charts - color the markers
                    try:
                        series.marker.format.fill.solid()
                        series.marker.format.fill.fore_color.rgb = color
                        series.marker.style = 2  # Circle marker
                        series.marker.size = 10
                    except Exception:
                        # Fallback to line formatting if markers don't work
                        series.format.line.color.rgb = color
                elif chart_type == ChartType.RADAR:
                    # Radar charts - color the line and fill with transparency
                    series.format.line.color.rgb = color
                    series.format.line.width = Pt(2)
                    try:
                        fill = series.format.fill
                        fill.solid()
                        fill.fore_color.rgb = color
                        # Radar fill should be semi-transparent (handled by PowerPoint)
                    except Exception:
                        pass
                else:
                    # Bar, column, pie, doughnut, area charts - solid fill
                    fill = series.format.fill
                    fill.solid()
                    fill.fore_color.rgb = color
        except Exception as e:
            logger.warning("Could not apply chart colors", error=str(e))

        # Add data labels if enabled
        if show_data_labels:
            try:
                self._add_data_labels(chart, chart_type)
            except Exception as e:
                logger.warning("Could not add data labels", error=str(e))

    def _add_data_labels(self, chart, chart_type: ChartType):
        """
        Add data labels to chart series.

        Different chart types require different label positioning:
        - Pie/Doughnut: Show percentage and category name
        - Bar: Outside end of bars
        - Column: Outside end of columns
        - Line: Above data points
        - Area: Above data points
        - Scatter: Next to markers
        - Radar: Above data points
        """
        from pptx.enum.chart import XL_DATA_LABEL_POSITION

        circular_charts = {ChartType.PIE, ChartType.DOUGHNUT}

        for series in chart.series:
            try:
                # Enable data labels for the series
                series.has_data_labels = True
                data_labels = series.data_labels

                # Common settings
                data_labels.font.size = Pt(9)
                data_labels.font.color.rgb = self.text_color

                # Chart-type specific positioning and content
                if chart_type in circular_charts:
                    # Pie/Doughnut: Show percentage, position outside
                    data_labels.show_percentage = True
                    data_labels.show_value = False
                    data_labels.show_category_name = False
                    try:
                        data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                    except Exception:
                        # Some chart types don't support OUTSIDE_END
                        pass
                    data_labels.font.size = Pt(10)

                elif chart_type in {ChartType.BAR, ChartType.CLUSTERED_BAR, ChartType.STACKED_BAR}:
                    # Bar charts: Show value at end of bar
                    data_labels.show_value = True
                    data_labels.show_percentage = False
                    try:
                        data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                    except Exception:
                        pass

                elif chart_type in {ChartType.COLUMN, ChartType.STACKED_COLUMN}:
                    # Column charts: Show value above column
                    data_labels.show_value = True
                    data_labels.show_percentage = False
                    try:
                        data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                    except Exception:
                        pass

                elif chart_type in {ChartType.LINE, ChartType.LINE_MARKERS}:
                    # Line charts: Show value above data points
                    data_labels.show_value = True
                    data_labels.show_percentage = False
                    try:
                        data_labels.position = XL_DATA_LABEL_POSITION.ABOVE
                    except Exception:
                        pass

                elif chart_type in {ChartType.AREA, ChartType.STACKED_AREA}:
                    # Area charts: Show value at data points
                    data_labels.show_value = True
                    data_labels.show_percentage = False
                    # Area charts may not support positioning well
                    try:
                        data_labels.position = XL_DATA_LABEL_POSITION.CENTER
                    except Exception:
                        pass

                elif chart_type == ChartType.SCATTER:
                    # Scatter charts: Show Y value
                    data_labels.show_value = True
                    data_labels.show_percentage = False
                    try:
                        data_labels.position = XL_DATA_LABEL_POSITION.RIGHT
                    except Exception:
                        pass

                elif chart_type == ChartType.RADAR:
                    # Radar charts: Show value at data points
                    data_labels.show_value = True
                    data_labels.show_percentage = False
                    # Radar doesn't support many position options

                else:
                    # Default: show value
                    data_labels.show_value = True

                # Number formatting for cleaner display
                try:
                    data_labels.number_format = '#,##0'
                except Exception:
                    pass

            except Exception as e:
                logger.debug(f"Could not add data labels to series: {e}")

    @staticmethod
    def _shorten_label(label: str, max_words: int = 3) -> str:
        """Shorten a chart label by stripping trailing function words and truncating.

        Examples:
            "Biodiversity loss has accelerated by" -> "Biodiversity Loss"
            "Ocean acidification has increased by" -> "Ocean Acidification"
            "Soil moisture" -> "Soil Moisture" (unchanged)
            "Total revenue for the quarter" -> "Total Revenue"
        """
        TRAILING_STOP = {
            'has', 'have', 'had', 'by', 'of', 'the', 'is', 'are', 'was', 'were',
            'been', 'in', 'at', 'to', 'for', 'with', 'from', 'on', 'a', 'an',
            'and', 'or', 'that', 'which', 'about', 'its', 'their', 'this',
            'increased', 'decreased', 'accelerated', 'reached', 'reported',
            'showed', 'experienced', 'affected', 'observed',
        }
        words = label.split()
        if not words:
            return label
        # Strip trailing function words
        while len(words) > 1 and words[-1].lower() in TRAILING_STOP:
            words.pop()
        # Truncate to max_words
        words = words[:max_words]
        # Second pass after truncation
        while len(words) > 1 and words[-1].lower() in TRAILING_STOP:
            words.pop()
        # Title-case, preserving all-caps acronyms (GDP, CO2)
        return ' '.join(
            w if (w.isupper() and len(w) >= 2) else w.capitalize()
            for w in words
        )

    @staticmethod
    def detect_chartable_data(text: str) -> Optional["ChartDataStructure"]:
        """
        Detect and extract chartable data from text.

        Supports:
        - Inline multi-point data (e.g., "Mumbai (40%), Bangkok (35%), Jakarta (32%)")
        - Numeric lists with labels
        - Percentage data
        - Time series data
        - Comparison data

        Args:
            text: Text content to analyze

        Returns:
            ChartDataStructure if chartable data found, None otherwise
        """
        lines = text.strip().split('\n')

        # 1. Try inline multi-point extraction first (cleanest labels from LLM output)
        inline_data = PPTXNativeChartGenerator._extract_inline_data_patterns(lines)
        if inline_data:
            return inline_data

        # 2. Try per-line numeric patterns (with label shortening)
        numeric_data = PPTXNativeChartGenerator._extract_numeric_patterns(lines)
        if numeric_data:
            return numeric_data

        # 3. Try per-line percentage patterns (with label shortening)
        percentage_data = PPTXNativeChartGenerator._extract_percentage_patterns(lines)
        if percentage_data:
            return percentage_data

        return None

    @staticmethod
    def _extract_inline_data_patterns(lines: List[str]) -> Optional["ChartDataStructure"]:
        """Extract multiple data points from a single line.

        Handles LLM output formats like:
          - "Mumbai (40%), Bangkok (35%), Jakarta (32%)"
          - "Q1 (500), Q2 (620), Q3 (710)"
          - "North America: 35%, Europe: 28%, Asia: 22%"
          - "Cities: Mumbai (40%), Bangkok (35%), Jakarta (32%)"
        """
        # Pattern A: Name (XX%) — inline percentage in parentheses
        pattern_pct_parens = re.compile(
            r'([A-Z][A-Za-z0-9]+(?:\s[A-Za-z0-9]+){0,3})\s*\(\s*(\d+(?:\.\d+)?)\s*%\s*\)'
        )
        # Pattern B: Name: XX% — inline percentage after colon
        pattern_colon_pct = re.compile(
            r'([A-Z][A-Za-z0-9]+(?:\s[A-Za-z0-9]+){0,3})\s*:\s*(\d+(?:\.\d+)?)\s*%'
        )
        # Pattern C: Name (number) — inline numeric in parentheses
        pattern_num_parens = re.compile(
            r'([A-Z][A-Za-z0-9]+(?:\s[A-Za-z0-9]+){0,3})\s*\(\s*(\d+(?:,\d{3})*(?:\.\d+)?)\s*\)'
        )
        # Pattern D: Name: number — inline numeric after colon
        pattern_colon_num = re.compile(
            r'([A-Z][A-Za-z0-9]+(?:\s[A-Za-z0-9]+){0,3})\s*:\s*\$?\s*(\d+(?:,\d{3})*(?:\.\d+)?)\s*(?=[,;)\s]|$)'
        )

        best_result = None
        best_count = 0

        for line in lines:
            clean_line = line.strip().lstrip('-*•● ').strip()
            if len(clean_line) < 15:
                continue

            for pattern, is_percentage in [
                (pattern_pct_parens, True),
                (pattern_colon_pct, True),
                (pattern_num_parens, False),
                (pattern_colon_num, False),
            ]:
                matches = pattern.findall(clean_line)
                if len(matches) >= 3 and len(matches) > best_count:
                    data_points = []
                    for label_raw, value_raw in matches:
                        label = PPTXNativeChartGenerator._shorten_label(
                            label_raw.strip(), max_words=3
                        )
                        value_str = value_raw.replace(',', '')
                        try:
                            value = float(value_str)
                            data_points.append((label, value))
                        except ValueError:
                            continue

                    if len(data_points) >= 3:
                        # Guard against citation years (non-percentage only)
                        if not is_percentage and all(1900 <= v <= 2099 for _, v in data_points):
                            continue

                        categories = [dp[0] for dp in data_points]
                        values = [dp[1] for dp in data_points]

                        suggested_type = None
                        if is_percentage:
                            total = sum(values)
                            suggested_type = ChartType.PIE if 80 <= total <= 120 else None

                        best_result = ChartDataStructure(
                            categories=categories,
                            series=[SeriesData(
                                name="Percentage" if is_percentage else "Values",
                                values=values,
                            )],
                            suggested_type=suggested_type,
                        )
                        best_count = len(data_points)

        return best_result

    @staticmethod
    def _extract_numeric_patterns(lines: List[str]) -> Optional["ChartDataStructure"]:
        """
        Extract numeric patterns like "Category: 100" or "Category - 100".

        Uses label shortening to produce chart-friendly labels from prose-style bullets.
        """
        data_points = []

        # Pattern: "Label: number" or "Label - number" or "Label (number)"
        # The 4th pattern is tighter than before: max 60-char label, no sentence punctuation
        patterns = [
            re.compile(r'^([^:\d]+):\s*([\d,\.]+)\s*$'),
            re.compile(r'^([^-\d]+)\s*-\s*([\d,\.]+)\s*$'),
            re.compile(r'^([^\(\d]+)\s*\(([\d,\.]+)\)\s*$'),
            re.compile(r'^([^.!?]{1,60})[\s:]+\$?([\d,\.]+)%?\s*$'),
        ]

        for line in lines:
            line = line.strip()
            # Strip bullet markers before matching
            clean_line = line.lstrip('-*•● ').strip()
            for pattern in patterns:
                match = pattern.match(clean_line)
                if match:
                    label = match.group(1).strip().rstrip(':- ')
                    # Reject labels that are clearly full sentences (>10 words)
                    if len(label.split()) > 10:
                        continue
                    # Reject labels longer than 60 chars (raw)
                    if len(label) > 60:
                        continue
                    # Shorten long labels to chart-friendly form
                    if len(label.split()) > 3:
                        label = PPTXNativeChartGenerator._shorten_label(label, max_words=3)
                    value_str = match.group(2).replace(',', '')
                    try:
                        value = float(value_str)
                        data_points.append((label, value))
                        break
                    except ValueError:
                        continue

        if len(data_points) >= 3:
            categories = [dp[0] for dp in data_points]
            values = [dp[1] for dp in data_points]
            return ChartDataStructure(
                categories=categories,
                series=[SeriesData(name="Values", values=values)],
            )
        return None

    @staticmethod
    def _extract_percentage_patterns(lines: List[str]) -> Optional["ChartDataStructure"]:
        """Extract percentage data for pie charts.

        Uses label shortening to produce chart-friendly labels from prose-style bullets.
        Only suggests PIE chart if values plausibly sum to ~100.
        """
        data_points = []

        # Pattern: "Label: XX%" or "Label - XX%"
        pattern = re.compile(r'^([^:\d%]+)[:\-\s]+([\d\.]+)\s*%')

        for line in lines:
            clean_line = line.strip().lstrip('-*•● ').strip()
            match = pattern.match(clean_line)
            if match:
                label = match.group(1).strip()
                # Reject labels that are clearly full sentences (>10 words)
                if len(label.split()) > 10 or len(label) > 60:
                    continue
                # Shorten long labels to chart-friendly form
                if len(label.split()) > 3:
                    label = PPTXNativeChartGenerator._shorten_label(label, max_words=3)
                try:
                    value = float(match.group(2))
                    data_points.append((label, value))
                except ValueError:
                    continue

        if len(data_points) >= 3:
            categories = [dp[0] for dp in data_points]
            values = [dp[1] for dp in data_points]
            # Only suggest PIE if values plausibly sum to ~100
            total = sum(values)
            suggested_type = ChartType.PIE if 80 <= total <= 120 else None
            return ChartDataStructure(
                categories=categories,
                series=[SeriesData(name="Percentage", values=values)],
                suggested_type=suggested_type,
            )
        return None

    @staticmethod
    def suggest_chart_type(data: "ChartDataStructure") -> ChartType:
        """
        Suggest the best chart type based on data characteristics.

        Args:
            data: Chart data structure

        Returns:
            Recommended chart type
        """
        if data.suggested_type:
            return data.suggested_type

        num_categories = len(data.categories)
        num_series = len(data.series)
        values = data.series[0].values if data.series else []

        # Check if values look like percentages that sum to ~100
        is_percentage_data = False
        if values and all(0 <= v <= 100 for v in values):
            total = sum(values)
            is_percentage_data = 90 <= total <= 110  # Allow some tolerance

        # Pie/Doughnut chart for percentage data with single series
        if is_percentage_data and num_series == 1:
            # Use doughnut for more categories (looks better), pie for fewer
            if num_categories > 5:
                return ChartType.DOUGHNUT
            return ChartType.PIE

        # Radar chart for multi-dimensional comparison (3-8 categories, multiple series)
        if 3 <= num_categories <= 8 and num_series >= 2:
            # Radar is good for comparing attributes across entities
            return ChartType.RADAR

        # Scatter chart when categories look numeric (for correlation data)
        categories_numeric = all(
            isinstance(cat, (int, float)) or
            (isinstance(cat, str) and cat.replace('.', '').replace('-', '').isdigit())
            for cat in data.categories
        )
        if categories_numeric and num_series == 1:
            return ChartType.SCATTER

        # Stacked charts for part-of-whole with multiple series
        if num_series >= 2 and num_categories <= 6:
            # Check if values could represent parts of a whole
            if all(v >= 0 for s in data.series for v in s.values):
                return ChartType.STACKED_COLUMN

        # Bar chart for many categories (horizontal is easier to read)
        if num_categories > 6:
            return ChartType.BAR

        # Line chart with markers for time-series-like data with multiple series
        if num_series > 1 or num_categories > 8:
            return ChartType.LINE_MARKERS if num_series > 1 else ChartType.LINE

        # Default to column chart
        return ChartType.COLUMN

    @staticmethod
    def should_render_as_chart(text: str) -> bool:
        """
        Determine if text content should be rendered as a chart.

        Args:
            text: Text content to analyze

        Returns:
            True if content appears to be chartable numeric data
        """
        if not text or len(text) < 20:
            return False

        lines = text.strip().split('\n')

        # Count lines with numeric data
        numeric_pattern = re.compile(r'[\d,\.]+\s*%?')
        numeric_lines = sum(
            1 for line in lines
            if numeric_pattern.search(line) and len(line) > 5
        )

        # Should have at least 3 data points
        return numeric_lines >= 3

    @staticmethod
    def validate_chart_data(data: "ChartDataStructure") -> bool:
        """Validate that chart data is meaningful, not extracted from prose.

        Rejects data where category labels look like full sentences
        (indicating the chart detector matched bullet prose, not real data).
        """
        if not data or not data.categories or not data.series:
            return False
        # Reject if fewer than 3 data points
        if len(data.categories) < 3:
            return False
        # Reject if average category label is too long (prose, not labels)
        avg_label_len = sum(len(c) for c in data.categories) / len(data.categories)
        if avg_label_len > 40:
            return False
        # Reject if any category has too many words (full sentence)
        for cat in data.categories:
            if len(cat.split()) > 8:
                return False
        return True


class SeriesData:
    """Data for a single chart series."""

    def __init__(self, name: str, values: List[float]):
        self.name = name
        self.values = values


class ChartDataStructure:
    """Structured data for creating a chart."""

    def __init__(
        self,
        categories: List[str],
        series: List[SeriesData],
        suggested_type: Optional[ChartType] = None,
    ):
        self.categories = categories
        self.series = series
        self.suggested_type = suggested_type


# Helper function for easy integration
def create_chart_generator(theme_colors: Dict[str, str]) -> PPTXNativeChartGenerator:
    """Create a new chart generator with the given theme."""
    return PPTXNativeChartGenerator(theme_colors)

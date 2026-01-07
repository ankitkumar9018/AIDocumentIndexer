"""
Document Preview Service
========================

Generates previews for generated documents in various formats:
- PDF: Page images using pdf2image
- PPTX: Slide images using python-pptx + Pillow
- DOCX: HTML conversion using mammoth
"""

import base64
import io
import os
from pathlib import Path
from typing import List, Optional, Tuple, Dict, Any
import tempfile

import structlog
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont

logger = structlog.get_logger(__name__)

# Optional imports - gracefully handle missing dependencies
try:
    import pdf2image
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False
    logger.warning("pdf2image not installed - PDF preview disabled")

try:
    import mammoth
    HAS_MAMMOTH = True
except ImportError:
    HAS_MAMMOTH = False
    logger.warning("mammoth not installed - DOCX preview disabled")


class PreviewError(Exception):
    """Error generating preview."""
    pass


class DocumentPreviewService:
    """Service for generating document previews."""

    def __init__(self, storage_path: str = "generated_docs"):
        """
        Initialize the preview service.

        Args:
            storage_path: Path to generated documents storage
        """
        self.storage_path = Path(storage_path)

    def get_document_path(self, job_id: str, format: str) -> Optional[Path]:
        """
        Get the path to a generated document.

        Args:
            job_id: Generation job ID
            format: Document format (pptx, docx, pdf, etc.)

        Returns:
            Path to document if exists, None otherwise
        """
        # Check for file with job_id prefix
        for file in self.storage_path.glob(f"{job_id}*"):
            if file.suffix.lower() == f".{format.lower()}":
                return file
        return None

    async def get_preview_metadata(
        self,
        job_id: str,
        format: str,
    ) -> Dict[str, Any]:
        """
        Get metadata about a document's preview capabilities.

        Args:
            job_id: Generation job ID
            format: Document format

        Returns:
            Dict with preview info (page_count, type, supported)
        """
        doc_path = self.get_document_path(job_id, format)
        if not doc_path or not doc_path.exists():
            return {"supported": False, "error": "Document not found"}

        format_lower = format.lower()

        if format_lower == "pdf":
            if not HAS_PDF2IMAGE:
                return {"supported": False, "error": "PDF preview not available"}
            try:
                from pdf2image import pdfinfo_from_path
                info = pdfinfo_from_path(str(doc_path))
                return {
                    "supported": True,
                    "type": "image",
                    "page_count": info.get("Pages", 1),
                    "format": "pdf",
                }
            except Exception as e:
                logger.error("Failed to get PDF info", error=str(e))
                return {"supported": False, "error": str(e)}

        elif format_lower == "pptx":
            try:
                prs = Presentation(str(doc_path))
                return {
                    "supported": True,
                    "type": "slides",
                    "page_count": len(prs.slides),
                    "format": "pptx",
                }
            except Exception as e:
                logger.error("Failed to get PPTX info", error=str(e))
                return {"supported": False, "error": str(e)}

        elif format_lower == "docx":
            if not HAS_MAMMOTH:
                return {"supported": False, "error": "DOCX preview not available"}
            return {
                "supported": True,
                "type": "html",
                "page_count": 1,  # DOCX doesn't have pages in preview
                "format": "docx",
            }

        elif format_lower in ["md", "markdown", "html", "txt"]:
            return {
                "supported": True,
                "type": "text",
                "page_count": 1,
                "format": format_lower,
            }

        return {"supported": False, "error": f"Unsupported format: {format}"}

    async def generate_pdf_preview(
        self,
        job_id: str,
        page: int = 1,
        dpi: int = 150,
    ) -> Tuple[bytes, str]:
        """
        Generate a preview image for a PDF page.

        Args:
            job_id: Generation job ID
            page: Page number (1-indexed)
            dpi: Image resolution

        Returns:
            Tuple of (image_bytes, content_type)
        """
        if not HAS_PDF2IMAGE:
            raise PreviewError("PDF preview not available - pdf2image not installed")

        doc_path = self.get_document_path(job_id, "pdf")
        if not doc_path:
            raise PreviewError("PDF document not found")

        try:
            from pdf2image import convert_from_path
            images = convert_from_path(
                str(doc_path),
                dpi=dpi,
                first_page=page,
                last_page=page,
            )
            if not images:
                raise PreviewError(f"No page {page} in document")

            img = images[0]
            buffer = io.BytesIO()
            img.save(buffer, format="PNG")
            return buffer.getvalue(), "image/png"
        except Exception as e:
            logger.error("Failed to generate PDF preview", error=str(e), page=page)
            raise PreviewError(f"Failed to generate PDF preview: {e}")

    async def generate_pptx_slide_preview(
        self,
        job_id: str,
        slide_num: int = 1,
        width: int = 1280,
        height: int = 720,
    ) -> Tuple[bytes, str]:
        """
        Generate a preview image for a PPTX slide.

        This creates a simplified rendering of the slide since python-pptx
        doesn't have native export capabilities.

        Args:
            job_id: Generation job ID
            slide_num: Slide number (1-indexed)
            width: Image width
            height: Image height

        Returns:
            Tuple of (image_bytes, content_type)
        """
        doc_path = self.get_document_path(job_id, "pptx")
        if not doc_path:
            raise PreviewError("PPTX document not found")

        try:
            prs = Presentation(str(doc_path))
            if slide_num < 1 or slide_num > len(prs.slides):
                raise PreviewError(f"Slide {slide_num} not found")

            slide = prs.slides[slide_num - 1]

            # Create image canvas
            img = Image.new("RGB", (width, height), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)

            # Try to load a font (fallback to default)
            try:
                font_title = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 32)
                font_body = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 18)
            except:
                font_title = ImageFont.load_default()
                font_body = ImageFont.load_default()

            # Extract text from slide shapes
            y_offset = 50
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            # Determine if title (first significant text or larger text)
                            is_title = y_offset < 100
                            font = font_title if is_title else font_body
                            color = (30, 30, 30) if is_title else (60, 60, 60)

                            # Word wrap
                            max_width = width - 100
                            words = text.split()
                            lines = []
                            current_line = []
                            for word in words:
                                current_line.append(word)
                                test_line = " ".join(current_line)
                                bbox = draw.textbbox((0, 0), test_line, font=font)
                                if bbox[2] > max_width:
                                    if len(current_line) > 1:
                                        current_line.pop()
                                        lines.append(" ".join(current_line))
                                        current_line = [word]
                                    else:
                                        lines.append(test_line)
                                        current_line = []
                            if current_line:
                                lines.append(" ".join(current_line))

                            for line in lines:
                                if y_offset < height - 50:
                                    draw.text((50, y_offset), line, fill=color, font=font)
                                    y_offset += 35 if is_title else 25

                            y_offset += 10  # Paragraph spacing

            # Add slide number
            slide_text = f"Slide {slide_num} of {len(prs.slides)}"
            draw.text((width - 150, height - 30), slide_text, fill=(150, 150, 150), font=font_body)

            buffer = io.BytesIO()
            img.save(buffer, format="PNG")
            return buffer.getvalue(), "image/png"

        except PreviewError:
            raise
        except Exception as e:
            logger.error("Failed to generate PPTX preview", error=str(e), slide=slide_num)
            raise PreviewError(f"Failed to generate PPTX preview: {e}")

    async def generate_pptx_all_slides(
        self,
        job_id: str,
        width: int = 640,
        height: int = 360,
    ) -> List[str]:
        """
        Generate preview images for all slides in a PPTX.

        Args:
            job_id: Generation job ID
            width: Thumbnail width
            height: Thumbnail height

        Returns:
            List of base64-encoded PNG images
        """
        doc_path = self.get_document_path(job_id, "pptx")
        if not doc_path:
            raise PreviewError("PPTX document not found")

        try:
            prs = Presentation(str(doc_path))
            images = []

            for i in range(len(prs.slides)):
                img_bytes, _ = await self.generate_pptx_slide_preview(
                    job_id, slide_num=i + 1, width=width, height=height
                )
                images.append(base64.b64encode(img_bytes).decode())

            return images
        except Exception as e:
            logger.error("Failed to generate all PPTX slides", error=str(e))
            raise PreviewError(f"Failed to generate slides: {e}")

    async def generate_docx_preview(self, job_id: str) -> Tuple[str, str]:
        """
        Generate an HTML preview of a DOCX document.

        Args:
            job_id: Generation job ID

        Returns:
            Tuple of (html_content, content_type)
        """
        if not HAS_MAMMOTH:
            raise PreviewError("DOCX preview not available - mammoth not installed")

        doc_path = self.get_document_path(job_id, "docx")
        if not doc_path:
            raise PreviewError("DOCX document not found")

        try:
            with open(doc_path, "rb") as f:
                result = mammoth.convert_to_html(f)
                html = result.value

                # Wrap in basic styling
                styled_html = f"""
                <!DOCTYPE html>
                <html>
                <head>
                    <style>
                        body {{
                            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
                            max-width: 800px;
                            margin: 0 auto;
                            padding: 20px;
                            line-height: 1.6;
                            color: #333;
                        }}
                        h1, h2, h3, h4, h5, h6 {{
                            color: #111;
                            margin-top: 1.5em;
                        }}
                        table {{
                            border-collapse: collapse;
                            width: 100%;
                            margin: 1em 0;
                        }}
                        th, td {{
                            border: 1px solid #ddd;
                            padding: 8px;
                            text-align: left;
                        }}
                        th {{
                            background-color: #f5f5f5;
                        }}
                        img {{
                            max-width: 100%;
                            height: auto;
                        }}
                    </style>
                </head>
                <body>
                    {html}
                </body>
                </html>
                """
                return styled_html, "text/html"
        except Exception as e:
            logger.error("Failed to generate DOCX preview", error=str(e))
            raise PreviewError(f"Failed to generate DOCX preview: {e}")

    async def generate_text_preview(
        self,
        job_id: str,
        format: str,
    ) -> Tuple[str, str]:
        """
        Generate a text preview for text-based formats.

        Args:
            job_id: Generation job ID
            format: Document format (md, html, txt)

        Returns:
            Tuple of (content, content_type)
        """
        doc_path = self.get_document_path(job_id, format)
        if not doc_path:
            raise PreviewError(f"{format.upper()} document not found")

        try:
            with open(doc_path, "r", encoding="utf-8") as f:
                content = f.read()

            content_type = {
                "md": "text/markdown",
                "markdown": "text/markdown",
                "html": "text/html",
                "txt": "text/plain",
            }.get(format.lower(), "text/plain")

            return content, content_type
        except Exception as e:
            logger.error("Failed to read text file", error=str(e))
            raise PreviewError(f"Failed to read document: {e}")

    async def generate_thumbnail(
        self,
        job_id: str,
        format: str,
        width: int = 300,
        height: int = 200,
    ) -> Tuple[bytes, str]:
        """
        Generate a thumbnail image for a document.

        Args:
            job_id: Generation job ID
            format: Document format
            width: Thumbnail width
            height: Thumbnail height

        Returns:
            Tuple of (image_bytes, content_type)
        """
        format_lower = format.lower()

        if format_lower == "pdf" and HAS_PDF2IMAGE:
            # Get first page at lower resolution
            img_bytes, _ = await self.generate_pdf_preview(job_id, page=1, dpi=72)
            img = Image.open(io.BytesIO(img_bytes))
            img.thumbnail((width, height), Image.Resampling.LANCZOS)
            buffer = io.BytesIO()
            img.save(buffer, format="PNG")
            return buffer.getvalue(), "image/png"

        elif format_lower == "pptx":
            # Get first slide as thumbnail
            return await self.generate_pptx_slide_preview(
                job_id, slide_num=1, width=width, height=height
            )

        else:
            # Generate a placeholder thumbnail with format icon
            img = Image.new("RGB", (width, height), color=(240, 240, 240))
            draw = ImageDraw.Draw(img)

            # Draw format text
            try:
                font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
            except:
                font = ImageFont.load_default()

            text = format.upper()
            bbox = draw.textbbox((0, 0), text, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
            x = (width - text_width) // 2
            y = (height - text_height) // 2
            draw.text((x, y), text, fill=(100, 100, 100), font=font)

            buffer = io.BytesIO()
            img.save(buffer, format="PNG")
            return buffer.getvalue(), "image/png"


# Global instance
preview_service = DocumentPreviewService()

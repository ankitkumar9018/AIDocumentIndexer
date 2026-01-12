"""
AIDocumentIndexer - Document Generation Service
================================================

Human-in-the-loop document generation using LangGraph.
Supports PPTX, DOCX, PDF, and other output formats.

LLM provider and model are configured via Admin UI (Operation-Level Config).
Configure the "content_generation" operation in Admin > Settings > LLM Configuration.
"""

import os
import re
import unicodedata
import uuid
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
from io import BytesIO
from pathlib import Path
from typing import Optional, List, Dict, Any, AsyncGenerator

import structlog

from backend.services.image_generator import (
    ImageGeneratorService,
    ImageGeneratorConfig,
    ImageBackend,
    get_image_generator,
)
from backend.services.content_quality import ContentQualityScorer, QualityReport

logger = structlog.get_logger(__name__)


# =============================================================================
# Language Configuration
# =============================================================================

LANGUAGE_NAMES = {
    "en": "English",
    "de": "German",
    "es": "Spanish",
    "fr": "French",
    "it": "Italian",
    "pt": "Portuguese",
    "nl": "Dutch",
    "pl": "Polish",
    "ru": "Russian",
    "zh": "Chinese",
    "ja": "Japanese",
    "ko": "Korean",
    "ar": "Arabic",
    "hi": "Hindi",
}


# =============================================================================
# Theme Configuration
# =============================================================================

THEMES = {
    # === EXISTING THEMES (Enhanced with distinctive visual properties) ===
    "business": {
        "name": "Business Professional",
        "primary": "#1E3A5F",
        "secondary": "#3D5A80",
        "accent": "#E0E1DD",
        "text": "#2D3A45",
        "light_gray": "#888888",
        "description": "Clean, corporate look ideal for business presentations",
        "slide_background": "solid",
        "header_style": "underline",
        "bullet_style": "circle",
        "accent_position": "top",
    },
    "creative": {
        "name": "Creative & Bold",
        "primary": "#6B4C9A",
        "secondary": "#9B6B9E",
        "accent": "#F4E4BA",
        "text": "#333333",
        "light_gray": "#666666",
        "description": "Vibrant colors for marketing and creative content",
        "slide_background": "gradient",
        "header_style": "bar",
        "bullet_style": "arrow",
        "accent_position": "side",
    },
    "modern": {
        "name": "Modern Minimal",
        "primary": "#212529",
        "secondary": "#495057",
        "accent": "#00B4D8",
        "text": "#212529",
        "light_gray": "#6C757D",
        "description": "Sleek, contemporary design with bold accents",
        "slide_background": "solid",
        "header_style": "none",
        "bullet_style": "dash",
        "accent_position": "bottom",
    },
    "nature": {
        "name": "Nature & Organic",
        "primary": "#2D5016",
        "secondary": "#5A7D3A",
        "accent": "#F5F0E1",
        "text": "#2D3A2E",
        "light_gray": "#7A8B6E",
        "description": "Earthy tones for sustainability and wellness topics",
        "slide_background": "textured",
        "header_style": "leaf",
        "bullet_style": "leaf",
        "accent_position": "corner",
    },
    "elegant": {
        "name": "Elegant & Refined",
        "primary": "#2C3E50",
        "secondary": "#7F8C8D",
        "accent": "#BDC3C7",
        "text": "#2C3E50",
        "light_gray": "#95A5A6",
        "description": "Sophisticated look for executive presentations",
        "slide_background": "solid",
        "header_style": "serif",
        "bullet_style": "square",
        "accent_position": "top",
    },
    "vibrant": {
        "name": "Vibrant & Energetic",
        "primary": "#E74C3C",
        "secondary": "#F39C12",
        "accent": "#FDF2E9",
        "text": "#2D3436",
        "light_gray": "#BDC3C7",
        "description": "Bold colors for high-energy content",
        "slide_background": "gradient",
        "header_style": "colorblock",
        "bullet_style": "circle-filled",
        "accent_position": "diagonal",
    },
    "tech": {
        "name": "Tech & Digital",
        "primary": "#0984E3",
        "secondary": "#6C5CE7",
        "accent": "#DFE6E9",
        "text": "#2D3436",
        "light_gray": "#B2BEC3",
        "description": "Modern tech aesthetic for digital topics",
        "slide_background": "gradient",
        "header_style": "glow",
        "bullet_style": "chevron",
        "accent_position": "side",
    },
    "warm": {
        "name": "Warm & Inviting",
        "primary": "#D35400",
        "secondary": "#E67E22",
        "accent": "#FDEBD0",
        "text": "#2C3E50",
        "light_gray": "#A6ACAF",
        "description": "Cozy colors for community and wellness",
        "slide_background": "warm-gradient",
        "header_style": "rounded",
        "bullet_style": "circle",
        "accent_position": "corner",
    },
    # === NEW THEMES ===
    "minimalist": {
        "name": "Ultra Minimalist",
        "primary": "#333333",
        "secondary": "#666666",
        "accent": "#F5F5F5",
        "text": "#222222",
        "light_gray": "#AAAAAA",
        "description": "Ultra-clean design with maximum whitespace and focus on content",
        "slide_background": "white",
        "header_style": "none",
        "bullet_style": "dash",
        "accent_position": "none",
    },
    "dark": {
        "name": "Dark Mode",
        "primary": "#1A1A2E",
        "secondary": "#16213E",
        "accent": "#0F3460",
        "text": "#E4E4E4",
        "light_gray": "#7A7A8C",
        "description": "Elegant dark theme for low-light viewing and modern aesthetics",
        "slide_background": "dark",
        "header_style": "glow",
        "bullet_style": "square",
        "accent_position": "border",
    },
    "colorful": {
        "name": "Colorful & Fun",
        "primary": "#FF6B6B",
        "secondary": "#4ECDC4",
        "accent": "#FFE66D",
        "text": "#2C3E50",
        "light_gray": "#95A5A6",
        "description": "Bold, multi-color theme for engaging and memorable presentations",
        "slide_background": "gradient-multi",
        "header_style": "colorblock",
        "bullet_style": "circle-filled",
        "accent_position": "corners",
    },
    "academic": {
        "name": "Academic & Scholarly",
        "primary": "#2C3E50",
        "secondary": "#8E44AD",
        "accent": "#ECF0F1",
        "text": "#2C3E50",
        "light_gray": "#7F8C8D",
        "description": "Traditional academic style ideal for research and educational presentations",
        "slide_background": "solid",
        "header_style": "underline",
        "bullet_style": "number",
        "accent_position": "footer",
    },
}

# Font family configurations
FONT_FAMILIES = {
    "modern": {
        "name": "Modern",
        "heading": "Calibri",
        "body": "Calibri",
        "description": "Clean, contemporary sans-serif"
    },
    "classic": {
        "name": "Classic",
        "heading": "Georgia",
        "body": "Times New Roman",
        "description": "Traditional serif fonts for formal documents"
    },
    "professional": {
        "name": "Professional",
        "heading": "Arial",
        "body": "Arial",
        "description": "Universal business-standard fonts"
    },
    "technical": {
        "name": "Technical",
        "heading": "Consolas",
        "body": "Courier New",
        "description": "Monospace fonts for technical content"
    },
}

# Layout templates for PPTX
LAYOUT_TEMPLATES = {
    "standard": {
        "name": "Standard",
        "description": "Title with bullet points",
        "content_width": 0.85,
        "image_position": "right"
    },
    "two_column": {
        "name": "Two Column",
        "description": "Split content into two columns",
        "content_width": 0.45,
        "image_position": "side"
    },
    "image_focused": {
        "name": "Image Focused",
        "description": "Large images with minimal text",
        "content_width": 0.4,
        "image_position": "center"
    },
    "minimal": {
        "name": "Minimal",
        "description": "Clean layout with lots of whitespace",
        "content_width": 0.7,
        "image_position": "bottom"
    },
}


# =============================================================================
# Utility Functions
# =============================================================================

def strip_markdown(text: str) -> str:
    """Remove markdown formatting, returning clean text.

    Used by PPTX, XLSX, and other generators that don't support markdown.
    """
    if not text:
        return ""
    # Remove headers (# ## ### #### etc.)
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
    # Remove bold **text** or __text__
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'__([^_]+)__', r'\1', text)
    # Remove italic *text* or _text_ (but not bullet points)
    text = re.sub(r'(?<!\s)\*([^*\n]+)\*(?!\s)', r'\1', text)
    text = re.sub(r'(?<!\s)_([^_\n]+)_(?!\s)', r'\1', text)
    # Remove code backticks
    text = re.sub(r'`([^`]+)`', r'\1', text)
    # Remove link syntax [text](url)
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    # Clean up multiple spaces
    text = re.sub(r'  +', ' ', text)
    return text.strip()


def filter_llm_metatext(text: str) -> str:
    """Remove common LLM conversational artifacts from generated content.

    Filters out:
    - Preamble text like "Here are the bullet points for..."
    - Closing remarks like "Let me know if you need any adjustments"
    - Conversational artifacts like "Certainly, here's...", "Sure, here are..."
    """
    if not text:
        return ""

    # Patterns to remove
    patterns = [
        # Preamble patterns (at start of content)
        r'^.*?[Hh]ere (?:are|is) (?:the )?.*?:\s*\n?',           # "Here are the bullet points:"
        r'^.*?[Ll]et me (?:provide|create|write|explain).*?:\s*\n?',  # "Let me provide..."
        r'^.*?[Cc]ertainly[,!]?\s*[Hh]ere.*?:\s*\n?',              # "Certainly, here's..."
        r'^.*?[Ss]ure[,!]?\s*[Hh]ere.*?:\s*\n?',                   # "Sure, here are..."
        r'^.*?[Ii]\'ll (?:provide|create|write|give).*?:\s*\n?',  # "I'll provide..."
        r'^.*?[Bb]elow (?:are|is).*?:\s*\n?',                      # "Below are the..."
        # Closing patterns (at end of content)
        r'\n?[Ll]et me know if you (?:need|want|have).*$',        # "Let me know if you need..."
        r'\n?[Hh]ope this helps.*$',                              # "Hope this helps!"
        r'\n?[Ff]eel free to (?:ask|reach|contact).*$',           # "Feel free to ask..."
        r'\n?[Ii]f you (?:have|need) any (?:questions|changes).*$',  # "If you have any questions..."
        r'\n?[Pp]lease let me know.*$',                           # "Please let me know..."
        # Section metadata echoed as content
        r'^\s*[-•*▪◦▸]?\s*\(Section \d+ of \d+\)\s*$',           # "(Section 1 of 8)" standalone
        # Preamble with "we will" / "to measure"
        r'^.*?[Ww]e (?:will|\'ll) (?:track|provide|create|use|measure).*?:\s*\n?',  # "We will track..."
        r'^.*?[Tt]o (?:measure|track|monitor) (?:the )?success.*?:\s*\n?',  # "To measure the success..."
        # Style instructions appearing as content (LLM outputting the style guide)
        r'^\s*[-•*▪◦▸]?\s*(?:Writing )?[Ss]tyle [Rr]equirements?.*$',  # "Writing Style Requirements"
        r'^\s*[-•*▪◦▸]?\s*[Mm]aintain a professional tone.*$',    # "Maintain a professional tone..."
        r'^\s*[-•*▪◦▸]?\s*[Uu]se (?:simple|clear|concise) language.*$',  # "Use simple language..."
        r'^\s*[-•*▪◦▸]?\s*[Tt]he (?:new )?content should (?:use|have|be|match).*$',  # "The content should use..."
        r'^\s*[-•*▪◦▸]?\s*[Kk]ey characteristics of the desired writing style.*$',  # "Key characteristics..."
        r'^\s*[-•*▪◦▸]?\s*[Uu]sing medium-length sentences.*$',   # "Using medium-length sentences..."
        r'^\s*[-•*▪◦▸]?\s*[Ii]ncorporating action verbs.*$',      # "Incorporating action verbs..."
        r'^\s*[-•*▪◦▸]?\s*[Bb]y following these style requirements.*$',  # "By following these style requirements..."
        # Preamble "To [action]..." patterns that describe what will be done
        r'^\s*[-•*▪◦▸]?\s*[Tt]o (?:create|establish|build|develop|implement|ensure|achieve) (?:a |the )?(?:buzz|strong|solid|effective|successful).*?,\s*we\s+.*$',  # "To create buzz..., we..."
        r'^\s*[-•*▪◦▸]?\s*[Tt]o (?:create|establish|build|develop|implement|ensure|achieve).*?:\s*$',  # "To create buzz:" followed by colon
        # In order to patterns
        r'^\s*[-•*▪◦▸]?\s*[Ii]n order to (?:create|establish|build|develop|implement|ensure|achieve).*?:\s*$',  # "In order to..."
        # For the purpose patterns
        r'^\s*[-•*▪◦▸]?\s*[Ff]or the purpose of (?:creating|establishing|building|developing).*?:\s*$',  # "For the purpose of..."
    ]

    result = text
    for pattern in patterns:
        result = re.sub(pattern, '', result, flags=re.MULTILINE | re.IGNORECASE)

    return result.strip()


def filter_title_echo(content: str, section_title: str) -> str:
    """Remove bullet points that just echo the section title.

    Filters out lines where the bullet text matches the section title,
    including variants with section count suffix like "(Section 1 of 8)".
    """
    if not content or not section_title:
        return content

    # Normalize section title for comparison
    title_normalized = section_title.upper().strip()
    # Also handle roman numeral prefixes (I., II., III., etc.)
    title_no_roman = re.sub(r'^[IVXLCDM]+\.\s*', '', title_normalized)

    lines = content.split('\n')
    filtered = []

    for line in lines:
        # Strip bullet markers for comparison
        text = re.sub(r'^[-•*▪◦▸]\s*', '', line.strip())
        # Remove section count suffix like "(Section 1 of 8)"
        text = re.sub(r'\s*\(Section \d+ of \d+\)\s*$', '', text, flags=re.IGNORECASE)
        text_normalized = text.upper().strip()

        # Also remove roman numeral prefix from text
        text_no_roman = re.sub(r'^[IVXLCDM]+\.\s*', '', text_normalized)

        # Skip if it matches the title (with or without roman numerals)
        if text_normalized == title_normalized or text_no_roman == title_no_roman:
            continue
        if text_normalized == title_no_roman or text_no_roman == title_normalized:
            continue

        # Keep the line
        filtered.append(line)

    return '\n'.join(filtered)


def smart_truncate(text: str, max_chars: int) -> str:
    """Truncate text at word boundaries to avoid mid-word cuts."""
    if len(text) <= max_chars:
        return text
    truncated = text[:max_chars]
    last_space = truncated.rfind(' ')
    if last_space > max_chars * 0.7:  # Don't cut too much
        return truncated[:last_space] + '...'
    return truncated + '...'


def sentence_truncate(text: str, max_chars: int) -> str:
    """Truncate text at sentence boundaries, avoiding mid-sentence cuts.

    Handles common abbreviations and ensures meaningful content is preserved.
    Falls back to clause boundaries, then word boundaries.

    Args:
        text: Text to truncate
        max_chars: Maximum character limit

    Returns:
        Truncated text ending at a complete sentence, clause, or word boundary
    """
    if len(text) <= max_chars:
        return text

    truncated = text[:max_chars]

    # Common abbreviations that shouldn't be treated as sentence ends
    # Use negative lookbehind to exclude these
    abbreviations = r'(?<![Mm]r)(?<![Mm]rs)(?<![Mm]s)(?<![Dd]r)(?<![Pp]rof)(?<![Ii]nc)(?<![Ll]td)(?<![Jj]r)(?<![Ss]r)(?<![Vv]s)(?<!etc)(?<!e\.g)(?<!i\.e)'

    # Find sentence endings: period/exclamation/question followed by space or end
    # Exclude common abbreviations
    pattern = abbreviations + r'[.!?](?:\s|$)'
    sentence_ends = list(re.finditer(pattern, truncated))

    if sentence_ends:
        last_end = sentence_ends[-1].end()
        # Use 50% threshold to preserve more complete sentences
        if last_end > max_chars * 0.5:
            return text[:last_end].strip()

    # Fallback: find last complete clause (before comma, semicolon, colon, or dash)
    clause_ends = list(re.finditer(r'[,;:—–-]\s', truncated))
    if clause_ends and clause_ends[-1].start() > max_chars * 0.5:
        return text[:clause_ends[-1].start()].strip() + '...'

    # Final fallback to word boundary
    return smart_truncate(text, max_chars)


async def llm_condense_text(
    text: str,
    max_chars: int,
    fallback_truncate: bool = True,
    preserve_numbers: bool = True,
    context_type: str = "bullet_point",
) -> str:
    """Use LLM to condense text while preserving meaning and critical data.

    Instead of truncating with '...', this function uses LLM to intelligently
    rephrase the text to fit within the character limit while preserving
    numbers, percentages, statistics, and key facts.

    Args:
        text: The text to condense
        max_chars: Maximum allowed characters
        fallback_truncate: If True, fall back to sentence_truncate on LLM failure
        preserve_numbers: If True, explicitly preserve numerical data
        context_type: Type of content ("bullet_point", "paragraph", "title")

    Returns:
        Condensed text that fits within max_chars
    """
    if len(text) <= max_chars:
        return text

    try:
        from backend.services.llm import EnhancedLLMFactory

        llm, _ = await EnhancedLLMFactory.get_chat_model_for_operation(
            operation="content_generation",
            user_id=None,
        )

        # Extract numerical data to preserve
        numerical_data = []
        if preserve_numbers:
            # Find all numbers, percentages, currency values
            number_patterns = [
                r'\$[\d,]+(?:\.\d+)?(?:\s*(?:million|billion|trillion|M|B|K))?',  # Currency
                r'[\d,]+(?:\.\d+)?%',  # Percentages
                r'\b\d{1,3}(?:,\d{3})+(?:\.\d+)?\b',  # Large numbers with commas
                r'\b\d+(?:\.\d+)?\s*(?:million|billion|trillion|M|B|K)\b',  # Numbers with magnitude
                r'\b(?:Q[1-4]|FY)\s*\d{2,4}\b',  # Fiscal quarters/years
                r'\b\d{4}\b',  # Years
            ]
            for pattern in number_patterns:
                matches = re.findall(pattern, text, re.IGNORECASE)
                numerical_data.extend(matches)

        # Build context-specific instructions
        context_instructions = {
            "bullet_point": "Keep it punchy and actionable. Start with a verb when possible.",
            "paragraph": "Maintain readability and flow. Keep complete sentences.",
            "title": "Be concise but descriptive. Capture the main theme.",
        }
        context_hint = context_instructions.get(context_type, context_instructions["bullet_point"])

        # Build preservation instructions
        preserve_instructions = ""
        if numerical_data:
            unique_numbers = list(set(numerical_data))[:10]  # Limit to first 10 unique
            preserve_instructions = f"""
CRITICAL: You MUST preserve these exact numerical values in your condensed version:
{', '.join(unique_numbers)}

Do not round, approximate, or omit any of these numbers."""

        prompt = f"""Condense this text to under {max_chars} characters while preserving the key meaning.

RULES:
1. Preserve ALL numbers, percentages, and statistics exactly as written
2. Keep the main point and most important facts
3. Use concise, professional language
4. {context_hint}
5. Do NOT add any prefixes like "Here is..." or "The condensed version is..."
6. Output ONLY the condensed text - nothing else
{preserve_instructions}

Original ({len(text)} chars): {text}

Condensed version (under {max_chars} chars):"""

        response = await llm.ainvoke(prompt)
        condensed = response.content.strip() if hasattr(response, 'content') else str(response).strip()

        # Clean up any accidental prefixes the LLM might add
        prefixes_to_remove = [
            "Here is the condensed version:",
            "Condensed version:",
            "Here's the condensed text:",
            "The condensed text is:",
        ]
        for prefix in prefixes_to_remove:
            if condensed.lower().startswith(prefix.lower()):
                condensed = condensed[len(prefix):].strip()

        # Verify it fits
        if len(condensed) <= max_chars:
            # Verify numerical data was preserved (warning only)
            if preserve_numbers and numerical_data:
                preserved_count = sum(1 for num in numerical_data if num in condensed)
                if preserved_count < len(numerical_data) * 0.5:  # Less than 50% preserved
                    logger.debug(
                        "Smart condense: some numerical data may have been lost",
                        original_numbers=len(numerical_data),
                        preserved=preserved_count,
                    )
            return condensed

        # If still too long, truncate the LLM output
        return sentence_truncate(condensed, max_chars)

    except Exception as e:
        logger.warning(f"LLM condense failed, using fallback: {e}")
        if fallback_truncate:
            return sentence_truncate(text, max_chars)
        return text[:max_chars]


async def smart_condense_content(
    content: str,
    max_length: int,
    content_type: str = "bullet_point",
    preserve_numbers: bool = True,
) -> str:
    """Smart content condensation for document generation.

    High-level wrapper around llm_condense_text that provides
    content-type-specific condensation with fallback strategies.

    Args:
        content: Content to condense
        max_length: Maximum character length
        content_type: Type of content (bullet_point, paragraph, title, subtitle)
        preserve_numbers: Whether to preserve numerical data

    Returns:
        Condensed content fitting within max_length
    """
    if len(content) <= max_length:
        return content

    # Map content types to condensation strategies
    type_mapping = {
        "bullet_point": "bullet_point",
        "paragraph": "paragraph",
        "title": "title",
        "subtitle": "title",
        "heading": "title",
        "body": "paragraph",
    }
    context_type = type_mapping.get(content_type, "bullet_point")

    return await llm_condense_text(
        text=content,
        max_chars=max_length,
        fallback_truncate=True,
        preserve_numbers=preserve_numbers,
        context_type=context_type,
    )


# =============================================================================
# Inline Citation Support
# =============================================================================

@dataclass
class CitationMapping:
    """Maps citation markers to source references."""
    marker: str  # e.g., "[1]"
    source_document: str
    source_page: Optional[int]
    source_snippet: str
    relevance_score: float = 0.0


@dataclass
class ContentWithCitations:
    """Content with inline citations and mapping."""
    content: str
    citations: List[CitationMapping]
    citation_style: str = "numbered"  # "numbered", "superscript", "author_year"


async def generate_content_with_citations(
    prompt: str,
    sources: List["SourceReference"],
    citation_style: str = "numbered",
    max_citations: int = 10,
) -> ContentWithCitations:
    """Generate content with inline citation markers.

    Uses LLM to generate content that includes citation markers [1], [2], etc.
    for facts sourced from the provided documents.

    Args:
        prompt: The content generation prompt
        sources: List of SourceReference objects to cite from
        citation_style: Style of citations ("numbered", "superscript", "author_year")
        max_citations: Maximum number of distinct citations to include

    Returns:
        ContentWithCitations with content and citation mapping
    """
    if not sources:
        # No sources to cite - generate without citations
        try:
            from backend.services.llm import EnhancedLLMFactory
            llm, _ = await EnhancedLLMFactory.get_chat_model_for_operation(
                operation="content_generation",
                user_id=None,
            )
            response = await llm.ainvoke(prompt)
            content = response.content if hasattr(response, 'content') else str(response)
            return ContentWithCitations(
                content=content.strip(),
                citations=[],
                citation_style=citation_style,
            )
        except Exception as e:
            logger.error(f"Content generation failed: {e}")
            return ContentWithCitations(content="", citations=[], citation_style=citation_style)

    # Prepare sources for citation
    source_context = _format_sources_for_citation(sources[:max_citations])

    # Build citation-aware prompt
    citation_prompt = f"""{prompt}

---CITATION REQUIREMENTS---
You have access to the following source documents. Include inline citations [1], [2], etc. after statements that use information from these sources.

SOURCES:
{source_context}

CITATION RULES:
1. Add citation markers [1], [2], etc. AFTER statements that use specific information from the sources
2. Multiple citations can be combined: [1][2] or [1, 2]
3. Only cite when using specific facts, statistics, or claims from the sources
4. General knowledge does not need citations
5. Each source number corresponds to the numbered sources above
6. Do NOT cite information that isn't in the sources - only cite what you actually use
7. Place citation markers at the end of the sentence, before the period

Example: The company's revenue increased by 45% in Q3 [1]. This growth was driven by new product launches [2].

Generate the content with appropriate citations:"""

    try:
        from backend.services.llm import EnhancedLLMFactory
        llm, _ = await EnhancedLLMFactory.get_chat_model_for_operation(
            operation="content_generation",
            user_id=None,
        )
        response = await llm.ainvoke(citation_prompt)
        content = response.content if hasattr(response, 'content') else str(response)
        content = content.strip()

        # Extract citation mappings
        citations = _extract_citation_mappings(content, sources[:max_citations])

        return ContentWithCitations(
            content=content,
            citations=citations,
            citation_style=citation_style,
        )

    except Exception as e:
        logger.error(f"Content with citations generation failed: {e}")
        return ContentWithCitations(content="", citations=[], citation_style=citation_style)


def _format_sources_for_citation(sources: List["SourceReference"]) -> str:
    """Format sources for LLM prompt with numbered references."""
    formatted = []
    for i, source in enumerate(sources, 1):
        page_info = f", page {source.page_number}" if source.page_number else ""
        formatted.append(
            f"[{i}] {source.document_name}{page_info}:\n"
            f"    {source.snippet[:500]}..."
        )
    return "\n\n".join(formatted)


def _extract_citation_mappings(
    content: str,
    sources: List["SourceReference"],
) -> List[CitationMapping]:
    """Extract citation markers from content and map to sources."""
    citations = []

    # Find all citation markers in the content
    citation_pattern = r'\[(\d+(?:,\s*\d+)*)\]'
    found_numbers = set()

    for match in re.finditer(citation_pattern, content):
        # Handle both [1] and [1, 2] formats
        numbers_str = match.group(1)
        numbers = [int(n.strip()) for n in numbers_str.split(',')]
        found_numbers.update(numbers)

    # Create mappings for found citations
    for num in sorted(found_numbers):
        if 1 <= num <= len(sources):
            source = sources[num - 1]
            citations.append(CitationMapping(
                marker=f"[{num}]",
                source_document=source.document_name,
                source_page=source.page_number,
                source_snippet=source.snippet[:200],
                relevance_score=source.relevance_score,
            ))

    return citations


def format_citations_for_footnotes(citations: List[CitationMapping]) -> str:
    """Format citations as footnotes for DOCX documents."""
    footnotes = []
    for i, citation in enumerate(citations, 1):
        page_ref = f", p. {citation.source_page}" if citation.source_page else ""
        footnotes.append(f"{citation.marker} {citation.source_document}{page_ref}")
    return "\n".join(footnotes)


def format_citations_for_speaker_notes(citations: List[CitationMapping]) -> str:
    """Format citations for PPTX speaker notes."""
    if not citations:
        return ""

    notes = ["Sources referenced in this slide:"]
    for citation in citations:
        page_ref = f" (p. {citation.source_page})" if citation.source_page else ""
        notes.append(f"  {citation.marker} {citation.source_document}{page_ref}")

    return "\n".join(notes)


def strip_citation_markers(content: str) -> str:
    """Remove citation markers from content.

    Useful when outputting to formats that don't support citations
    or when citations should be handled separately.
    """
    return re.sub(r'\s*\[\d+(?:,\s*\d+)*\]\s*', ' ', content).strip()


def convert_citations_to_superscript(content: str) -> str:
    """Convert [1] style citations to superscript numbers.

    Note: This returns Unicode superscript characters.
    For actual superscript formatting, use the document library's
    formatting capabilities.
    """
    superscript_map = {
        '0': '⁰', '1': '¹', '2': '²', '3': '³', '4': '⁴',
        '5': '⁵', '6': '⁶', '7': '⁷', '8': '⁸', '9': '⁹',
    }

    def replace_citation(match):
        numbers = match.group(1)
        superscript = ''.join(superscript_map.get(c, c) for c in numbers if c.isdigit())
        return superscript

    return re.sub(r'\[(\d+(?:,\s*\d+)*)\]', replace_citation, content)


async def add_citations_to_section(
    section: "Section",
    include_citations: bool = True,
    citation_style: str = "numbered",
) -> "Section":
    """Add inline citations to a section's content.

    Regenerates content with citation markers if include_citations is True.

    Args:
        section: Section to add citations to
        include_citations: Whether to include inline citations
        citation_style: Style of citations

    Returns:
        Updated section with citations
    """
    if not include_citations or not section.sources:
        return section

    # Generate content with citations
    prompt = f"""Rewrite this content with inline citation markers [1], [2], etc. where appropriate.

Original content:
{section.content}

Add citation markers after statements that use information from the sources.
Keep the content largely the same - just add citation markers."""

    result = await generate_content_with_citations(
        prompt=prompt,
        sources=section.sources,
        citation_style=citation_style,
    )

    if result.content:
        section.content = result.content
        # Store citations in metadata for later use
        if section.metadata is None:
            section.metadata = {}
        section.metadata["citations"] = [
            {
                "marker": c.marker,
                "source_document": c.source_document,
                "source_page": c.source_page,
            }
            for c in result.citations
        ]

    return section


# =============================================================================
# Style Learning from Existing Documents
# =============================================================================

@dataclass
class StyleProfile:
    """Learned writing style from user documents.

    Captures detailed style characteristics that can be applied
    to new document generation for consistency.
    """
    # Core style attributes
    tone: str = "professional"  # formal, casual, technical, friendly, academic
    vocabulary_level: str = "moderate"  # simple, moderate, advanced, technical
    formality: str = "moderate"  # casual, moderate, formal, highly_formal

    # Structural preferences
    avg_sentence_length: float = 15.0  # Average words per sentence
    avg_paragraph_length: float = 4.0  # Average sentences per paragraph
    structure_pattern: str = "mixed"  # bullet-lists, paragraphs, mixed, headers-heavy
    bullet_preference: bool = False  # Prefers bullet points over prose

    # Language patterns
    uses_passive_voice: float = 0.3  # 0-1 ratio of passive voice usage
    uses_first_person: bool = False  # Uses "I", "we"
    uses_contractions: bool = False  # Uses "don't", "can't", etc.
    heading_style: str = "title_case"  # title_case, sentence_case, all_caps

    # Key phrases and terminology
    key_phrases: List[str] = field(default_factory=list)
    domain_terms: List[str] = field(default_factory=list)

    # Formatting preferences
    uses_bold_emphasis: bool = True
    uses_italic_emphasis: bool = True
    uses_numbered_lists: bool = False

    # Source information
    source_documents: List[str] = field(default_factory=list)
    confidence_score: float = 0.5  # 0-1, how confident we are in the analysis


async def learn_style_from_documents(
    document_contents: List[str],
    document_names: Optional[List[str]] = None,
    use_llm_analysis: bool = True,
) -> StyleProfile:
    """Learn writing style from a collection of documents.

    Analyzes document content to extract comprehensive style patterns
    that can be applied to generate consistently styled new documents.

    Args:
        document_contents: List of document text contents
        document_names: Optional list of source document names
        use_llm_analysis: Whether to use LLM for deeper style analysis

    Returns:
        StyleProfile with learned style characteristics
    """
    if not document_contents:
        return StyleProfile()

    # Statistical analysis of documents
    all_sentences = []
    all_paragraphs = []
    word_counts = []
    total_words = 0
    uses_bullets = False
    uses_numbered = False
    uses_bold = False
    uses_italic = False
    first_person_count = 0
    contraction_count = 0
    passive_patterns = 0
    total_sentences = 0

    for content in document_contents:
        # Split into sentences
        sentences = re.split(r'[.!?]+', content)
        sentences = [s.strip() for s in sentences if s.strip()]
        all_sentences.extend(sentences)

        # Split into paragraphs
        paragraphs = re.split(r'\n\s*\n', content)
        paragraphs = [p.strip() for p in paragraphs if p.strip()]
        all_paragraphs.extend(paragraphs)

        # Word count per sentence
        for sentence in sentences:
            words = sentence.split()
            word_counts.append(len(words))
            total_words += len(words)
            total_sentences += 1

            # Check for first person
            if re.search(r'\b(I|we|our|my)\b', sentence, re.IGNORECASE):
                first_person_count += 1

            # Check for contractions
            if re.search(r"'(t|s|re|ll|ve|d)\b|n't\b", sentence, re.IGNORECASE):
                contraction_count += 1

            # Check for passive voice patterns
            if re.search(r'\b(is|are|was|were|been|being)\s+\w+ed\b', sentence, re.IGNORECASE):
                passive_patterns += 1

        # Check for formatting patterns
        if '•' in content or re.search(r'^[-*]\s', content, re.MULTILINE):
            uses_bullets = True
        if re.search(r'^\d+[.)]\s', content, re.MULTILINE):
            uses_numbered = True
        if '**' in content or '<b>' in content.lower():
            uses_bold = True
        if '_' in content or '*' in content or '<i>' in content.lower():
            uses_italic = True

    # Calculate averages
    avg_sentence_length = sum(word_counts) / len(word_counts) if word_counts else 15.0
    sentences_per_para = []
    for para in all_paragraphs:
        para_sentences = re.split(r'[.!?]+', para)
        para_sentences = [s for s in para_sentences if s.strip()]
        sentences_per_para.append(len(para_sentences))
    avg_paragraph_length = sum(sentences_per_para) / len(sentences_per_para) if sentences_per_para else 4.0

    # Calculate ratios
    first_person_ratio = first_person_count / total_sentences if total_sentences > 0 else 0
    contraction_ratio = contraction_count / total_sentences if total_sentences > 0 else 0
    passive_ratio = passive_patterns / total_sentences if total_sentences > 0 else 0

    # Determine vocabulary level based on average word length
    all_words = ' '.join(document_contents).split()
    avg_word_length = sum(len(w) for w in all_words) / len(all_words) if all_words else 5
    if avg_word_length > 7:
        vocabulary_level = "advanced"
    elif avg_word_length > 5.5:
        vocabulary_level = "moderate"
    else:
        vocabulary_level = "simple"

    # Determine structure pattern
    if uses_bullets and not uses_numbered:
        structure_pattern = "bullet-lists"
    elif uses_numbered:
        structure_pattern = "headers-heavy"
    elif avg_paragraph_length > 5:
        structure_pattern = "paragraphs"
    else:
        structure_pattern = "mixed"

    # Determine formality
    if contraction_ratio > 0.2 or first_person_ratio > 0.3:
        formality = "casual"
    elif contraction_ratio < 0.05 and first_person_ratio < 0.1:
        formality = "formal"
    else:
        formality = "moderate"

    # Build initial profile
    profile = StyleProfile(
        vocabulary_level=vocabulary_level,
        formality=formality,
        avg_sentence_length=round(avg_sentence_length, 1),
        avg_paragraph_length=round(avg_paragraph_length, 1),
        structure_pattern=structure_pattern,
        bullet_preference=uses_bullets,
        uses_passive_voice=round(passive_ratio, 2),
        uses_first_person=first_person_ratio > 0.1,
        uses_contractions=contraction_ratio > 0.1,
        uses_bold_emphasis=uses_bold,
        uses_italic_emphasis=uses_italic,
        uses_numbered_lists=uses_numbered,
        source_documents=document_names or [],
        confidence_score=min(1.0, len(document_contents) * 0.2),  # More docs = higher confidence
    )

    # Use LLM for deeper analysis if enabled
    if use_llm_analysis and document_contents:
        try:
            from backend.services.llm import EnhancedLLMFactory

            llm, _ = await EnhancedLLMFactory.get_chat_model_for_operation(
                operation="content_generation",
                user_id=None,
            )

            # Sample text for LLM analysis
            sample_text = "\n\n---\n\n".join([
                content[:1500] for content in document_contents[:5]
            ])

            analysis_prompt = f"""Analyze this text sample and identify writing style characteristics.

TEXT SAMPLES:
{sample_text}

Provide a JSON response with:
{{
    "tone": "formal" | "casual" | "technical" | "friendly" | "academic",
    "key_phrases": ["list of 3-5 common phrases or patterns"],
    "domain_terms": ["list of 3-5 domain-specific terms used"],
    "heading_style": "title_case" | "sentence_case" | "all_caps",
    "recommended_tone": "brief description of how to write in this style"
}}

Return ONLY valid JSON:"""

            response = await llm.ainvoke(analysis_prompt)
            response_text = response.content if hasattr(response, 'content') else str(response)

            # Parse LLM response
            import json
            json_match = re.search(r'\{[^{}]*\}', response_text, re.DOTALL)
            if json_match:
                llm_analysis = json.loads(json_match.group())
                profile.tone = llm_analysis.get("tone", profile.tone)
                profile.key_phrases = llm_analysis.get("key_phrases", [])
                profile.domain_terms = llm_analysis.get("domain_terms", [])
                profile.heading_style = llm_analysis.get("heading_style", "title_case")
                profile.confidence_score = min(1.0, profile.confidence_score + 0.2)

        except Exception as e:
            logger.warning(f"LLM style analysis failed: {e}")

    return profile


def apply_style_to_prompt(prompt: str, style: StyleProfile) -> str:
    """Apply style profile instructions to a generation prompt.

    Args:
        prompt: Original generation prompt
        style: StyleProfile to apply

    Returns:
        Prompt with style instructions appended
    """
    style_instructions = f"""
---WRITING STYLE GUIDANCE (internal - follow but do not output)---
Based on existing documents, match this style:

TONE & FORMALITY:
- Tone: {style.tone}
- Formality: {style.formality}
- Vocabulary: {style.vocabulary_level} complexity

SENTENCE STRUCTURE:
- Target sentence length: ~{style.avg_sentence_length:.0f} words
- Target paragraph length: ~{style.avg_paragraph_length:.0f} sentences
{"- Use first person (I, we) when appropriate" if style.uses_first_person else "- Avoid first person pronouns"}
{"- Contractions are OK (don't, can't)" if style.uses_contractions else "- Avoid contractions (use 'do not', 'cannot')"}

FORMATTING:
- Structure: {style.structure_pattern}
{"- Prefer bullet points for lists" if style.bullet_preference else "- Prefer prose over bullet points"}
{"- Use **bold** for emphasis" if style.uses_bold_emphasis else ""}
{"- Use _italic_ for emphasis" if style.uses_italic_emphasis else ""}
- Heading style: {style.heading_style}
"""

    if style.key_phrases:
        style_instructions += f"\nKEY PHRASES to use when relevant:\n"
        for phrase in style.key_phrases[:5]:
            style_instructions += f"- \"{phrase}\"\n"

    if style.domain_terms:
        style_instructions += f"\nDOMAIN TERMS to incorporate:\n"
        for term in style.domain_terms[:5]:
            style_instructions += f"- {term}\n"

    style_instructions += "---END STYLE GUIDANCE---\n"

    return prompt + "\n" + style_instructions


def hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color string to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def sanitize_filename(title: str, max_length: int = 50) -> str:
    """Sanitize title for use as filename.

    Removes/replaces invalid characters for cross-platform compatibility.
    Handles special characters like /, :, |, ?, * that cause issues on various OS.

    Args:
        title: Document title to sanitize
        max_length: Maximum filename length (default 50)

    Returns:
        Safe filename string with only alphanumeric, underscores, and hyphens
    """
    # Normalize unicode characters (e.g., convert é to e)
    title = unicodedata.normalize('NFKD', title)
    title = title.encode('ascii', 'ignore').decode('ascii')
    # Keep only alphanumeric, spaces, hyphens, underscores
    title = re.sub(r'[^\w\s-]', '', title)
    # Replace spaces and multiple hyphens/underscores with single underscore
    title = re.sub(r'[-\s]+', '_', title)
    # Remove leading/trailing underscores
    title = title.strip('_')
    # Limit length
    return title[:max_length] if title else 'document'


def get_theme_colors(theme_key: str = "business", custom_colors: dict = None) -> dict:
    """Get theme colors, with fallback to business theme.

    If custom_colors is provided, those values override the theme colors.
    Custom colors can include: primary, secondary, accent, text, background.
    """
    theme = THEMES.get(theme_key, THEMES["business"]).copy()

    # Apply custom color overrides if provided
    if custom_colors:
        if "primary" in custom_colors:
            theme["primary"] = custom_colors["primary"]
        if "secondary" in custom_colors:
            theme["secondary"] = custom_colors["secondary"]
        if "accent" in custom_colors:
            theme["accent"] = custom_colors["accent"]
        if "text" in custom_colors:
            theme["text"] = custom_colors["text"]
        if "background" in custom_colors:
            theme["background"] = custom_colors["background"]

    return theme


def check_spelling(text: str) -> dict:
    """Check spelling and return suggestions for review.

    This flags potential spelling issues for user approval rather than
    auto-correcting, allowing users to decide what to fix.

    Returns:
        dict with:
        - has_issues: bool indicating if spelling issues were found
        - issues: list of dicts with word, position, suggestion, context
        - original: the original text
    """
    if not text:
        return {"has_issues": False, "issues": [], "original": text}

    try:
        from spellchecker import SpellChecker

        spell = SpellChecker()
        words = text.split()
        issues = []

        for i, word in enumerate(words):
            # Clean punctuation from word
            clean_word = word.strip('.,!?:;"\'-()[]{}')
            if not clean_word or len(clean_word) < 3:
                continue

            # Skip common patterns that aren't words
            if clean_word.isdigit():
                continue
            if any(c.isdigit() for c in clean_word):  # Skip alphanumeric codes
                continue
            if clean_word.startswith(('@', '#', 'http', 'www')):
                continue

            # Check if word is known
            if clean_word.lower() not in spell:
                correction = spell.correction(clean_word.lower())
                # Only suggest if we have a different correction
                if correction and correction != clean_word.lower():
                    # Get context (surrounding words)
                    context_start = max(0, i - 2)
                    context_end = min(len(words), i + 3)
                    context = ' '.join(words[context_start:context_end])

                    issues.append({
                        "word": clean_word,
                        "position": i,
                        "suggestion": correction,
                        "context": context
                    })

        return {
            "has_issues": len(issues) > 0,
            "issues": issues[:20],  # Limit to top 20 issues
            "original": text
        }
    except ImportError:
        logger.warning("pyspellchecker not installed, skipping spell check")
        return {"has_issues": False, "issues": [], "original": text}
    except Exception as e:
        logger.warning(f"Spell check failed: {e}")
        return {"has_issues": False, "issues": [], "original": text}


# =============================================================================
# Types
# =============================================================================

class OutputFormat(str, Enum):
    """Supported output formats."""
    PPTX = "pptx"
    DOCX = "docx"
    PDF = "pdf"
    XLSX = "xlsx"
    MARKDOWN = "markdown"
    HTML = "html"
    TXT = "txt"


class GenerationStatus(str, Enum):
    """Document generation workflow status."""
    DRAFT = "draft"
    OUTLINE_PENDING = "outline_pending"
    OUTLINE_APPROVED = "outline_approved"
    GENERATING = "generating"
    SECTION_REVIEW = "section_review"
    REVISION = "revision"
    COMPLETED = "completed"
    CANCELLED = "cancelled"
    FAILED = "failed"


@dataclass
class SourceReference:
    """Reference to a source document used in generation."""
    document_id: str
    document_name: str
    chunk_id: Optional[str] = None
    page_number: Optional[int] = None
    relevance_score: float = 0.0
    snippet: str = ""


@dataclass
class Section:
    """A section of the generated document."""
    id: str
    title: str
    content: str
    order: int
    sources: List[SourceReference] = field(default_factory=list)
    approved: bool = False
    feedback: Optional[str] = None
    revised_content: Optional[str] = None
    rendered_content: Optional[str] = None  # Final content after processing (what appears in output)
    metadata: Optional[Dict[str, Any]] = None  # Quality scores, etc.


@dataclass
class DocumentOutline:
    """Outline for document generation."""
    title: str
    description: str
    sections: List[Dict[str, str]]  # List of {title, description}
    target_audience: Optional[str] = None
    tone: Optional[str] = None
    word_count_target: Optional[int] = None


@dataclass
class GenerationJob:
    """A document generation job."""
    id: str
    user_id: str
    title: str
    description: str
    output_format: OutputFormat
    status: GenerationStatus
    outline: Optional[DocumentOutline] = None
    sections: List[Section] = field(default_factory=list)
    sources_used: List[SourceReference] = field(default_factory=list)
    created_at: datetime = field(default_factory=datetime.utcnow)
    updated_at: datetime = field(default_factory=datetime.utcnow)
    completed_at: Optional[datetime] = None
    output_path: Optional[str] = None
    error_message: Optional[str] = None
    metadata: Dict[str, Any] = field(default_factory=dict)


def _get_generation_setting(key: str, env_key: str, default: Any) -> Any:
    """
    Get a generation setting with fallback chain:
    1. Database settings (via settings service defaults)
    2. Environment variable
    3. Hardcoded default
    """
    from backend.services.settings import get_settings_service
    settings = get_settings_service()

    # Try settings service first
    value = settings.get_default_value(key)
    if value is not None:
        return value

    # Fall back to environment variable
    env_value = os.getenv(env_key)
    if env_value is not None:
        # Handle boolean conversion
        if isinstance(default, bool):
            return env_value.lower() in ("true", "1", "yes", "on")
        return env_value

    return default


@dataclass
class GenerationConfig:
    """Configuration for document generation."""
    # LLM settings
    model: str = "gpt-4"
    temperature: float = 0.7
    max_tokens: int = 4000

    # RAG settings
    use_rag: bool = True
    max_sources: int = 10
    min_relevance_score: float = 0.01  # RRF scores are typically 0-0.1, not 0-1

    # Output settings - loaded from database settings or environment
    # Use persistent storage instead of /tmp
    output_dir: str = str(Path(__file__).resolve().parents[2] / "data" / "generated_docs")
    include_sources: bool = field(
        default_factory=lambda: _get_generation_setting(
            "generation.include_sources", "GENERATION_INCLUDE_SOURCES", True
        )
    )
    include_toc: bool = True

    # Image generation settings - loaded from database settings
    # Configure via Admin UI: Settings > Document Generation
    include_images: bool = field(
        default_factory=lambda: _get_generation_setting(
            "generation.include_images", "GENERATION_INCLUDE_IMAGES", True
        )
    )
    # Image backend: "picsum" (free, no API key), "unsplash" (requires API key),
    # "pexels" (requires API key), "openai" (DALL-E, requires API key),
    # "stability" (Stable Diffusion API), "automatic1111" (local SD), or "disabled"
    image_backend: str = field(
        default_factory=lambda: _get_generation_setting(
            "generation.image_backend", "GENERATION_IMAGE_BACKEND", "picsum"
        )
    )

    # Style settings - loaded from database settings
    default_tone: str = field(
        default_factory=lambda: _get_generation_setting(
            "generation.default_tone", "GENERATION_DEFAULT_TONE", "professional"
        )
    )
    default_style: str = field(
        default_factory=lambda: _get_generation_setting(
            "generation.default_style", "GENERATION_DEFAULT_STYLE", "business"
        )
    )

    # Chart generation settings
    auto_charts: bool = field(
        default_factory=lambda: _get_generation_setting(
            "generation.auto_charts", "GENERATION_AUTO_CHARTS", False
        )
    )
    chart_style: str = field(
        default_factory=lambda: _get_generation_setting(
            "generation.chart_style", "GENERATION_CHART_STYLE", "business"
        )
    )
    chart_dpi: int = field(
        default_factory=lambda: _get_generation_setting(
            "generation.chart_dpi", "GENERATION_CHART_DPI", 150
        )
    )

    # Workflow settings
    require_outline_approval: bool = True
    require_section_approval: bool = False
    auto_generate_on_approval: bool = True

    # Quality scoring settings
    enable_quality_review: bool = field(
        default_factory=lambda: _get_generation_setting(
            "generation.enable_quality_review", "GENERATION_ENABLE_QUALITY_REVIEW", False
        )
    )
    min_quality_score: float = field(
        default_factory=lambda: _get_generation_setting(
            "generation.min_quality_score", "GENERATION_MIN_QUALITY_SCORE", 0.7
        )
    )
    auto_regenerate_low_quality: bool = field(
        default_factory=lambda: _get_generation_setting(
            "generation.auto_regenerate_low_quality", "GENERATION_AUTO_REGENERATE", True
        )
    )
    max_regeneration_attempts: int = 2  # Max times to regenerate a low-quality section


# =============================================================================
# Document Generation Service
# =============================================================================

class DocumentGenerationService:
    """
    Service for generating documents with human-in-the-loop workflow.

    Features:
    - RAG-powered content generation
    - Outline review and approval
    - Section-by-section generation
    - Source attribution
    - Multiple output formats
    """

    def __init__(self, config: Optional[GenerationConfig] = None):
        """Initialize the document generation service."""
        self.config = config or GenerationConfig()
        self._jobs: Dict[str, GenerationJob] = {}

        # Ensure output directory exists
        Path(self.config.output_dir).mkdir(parents=True, exist_ok=True)

    async def reload_settings(self) -> None:
        """
        Reload generation settings from database.

        Call this to pick up admin setting changes without service restart.
        Settings are read from Admin UI > Settings > Document Generation.
        """
        from backend.services.settings import get_settings_service
        settings_service = get_settings_service()

        # Get current database settings (async)
        include_images = await settings_service.get_setting("generation.include_images")
        image_backend = await settings_service.get_setting("generation.image_backend")
        include_sources = await settings_service.get_setting("generation.include_sources")
        default_tone = await settings_service.get_setting("generation.default_tone")
        default_style = await settings_service.get_setting("generation.default_style")
        auto_charts = await settings_service.get_setting("generation.auto_charts")
        chart_style = await settings_service.get_setting("generation.chart_style")
        chart_dpi = await settings_service.get_setting("generation.chart_dpi")

        # Update config with database values
        if include_images is not None:
            self.config.include_images = include_images
        if image_backend is not None:
            self.config.image_backend = image_backend
        if include_sources is not None:
            self.config.include_sources = include_sources
        if default_tone is not None:
            self.config.default_tone = default_tone
        if default_style is not None:
            self.config.default_style = default_style
        if auto_charts is not None:
            self.config.auto_charts = auto_charts
        if chart_style is not None:
            self.config.chart_style = chart_style
        if chart_dpi is not None:
            self.config.chart_dpi = int(chart_dpi)

        logger.info(
            "Generation settings reloaded",
            include_images=self.config.include_images,
            image_backend=self.config.image_backend,
            default_tone=self.config.default_tone,
            default_style=self.config.default_style,
            auto_charts=self.config.auto_charts,
        )

    async def suggest_theme(
        self,
        title: str,
        description: str,
        document_type: str = "pptx",
    ) -> Dict[str, Any]:
        """
        Use LLM to suggest optimal theming for a document based on its content.

        Args:
            title: Document title
            description: Document description/topic
            document_type: Type of document (pptx, docx, pdf, etc.)

        Returns:
            Dictionary with recommended theme, font_family, layout, and reason
        """
        try:
            from backend.services.llm import EnhancedLLMFactory

            llm, _ = await EnhancedLLMFactory.get_chat_model_for_operation(
                operation="content_generation",
                user_id=None,
            )

            # Build theme options for the prompt
            theme_options = "\n".join([
                f"- {key}: {theme['name']} - {theme['description']}"
                for key, theme in THEMES.items()
            ])

            font_options = "\n".join([
                f"- {key}: {font['name']} - {font['description']}"
                for key, font in FONT_FAMILIES.items()
            ])

            layout_options = "\n".join([
                f"- {key}: {layout['name']} - {layout['description']}"
                for key, layout in LAYOUT_TEMPLATES.items()
            ])

            prompt = f"""Analyze this document and suggest optimal theming.

Document Title: {title}
Document Description: {description}
Document Type: {document_type}

Available Themes:
{theme_options}

Available Font Families:
{font_options}

Available Layouts (for presentations):
{layout_options}

Based on the document topic, audience, and purpose, recommend the best options.
Consider: industry conventions, emotional tone, readability, and visual impact.

Return ONLY a JSON object with this exact structure:
{{
    "theme": "theme_key",
    "font_family": "font_key",
    "layout": "layout_key",
    "animations": true/false,
    "reason": "Brief explanation of why these choices suit this document"
}}"""

            response = await llm.ainvoke(prompt)
            response_text = response.content if hasattr(response, 'content') else str(response)

            # Parse JSON response
            import json
            import re

            # Extract JSON from response
            json_match = re.search(r'\{[^{}]*\}', response_text, re.DOTALL)
            if json_match:
                suggestion = json.loads(json_match.group())
            else:
                suggestion = json.loads(response_text)

            # Validate keys exist
            if suggestion.get("theme") not in THEMES:
                suggestion["theme"] = "business"
            if suggestion.get("font_family") not in FONT_FAMILIES:
                suggestion["font_family"] = "modern"
            if suggestion.get("layout") not in LAYOUT_TEMPLATES:
                suggestion["layout"] = "standard"

            # Add full theme/font/layout details
            suggestion["theme_details"] = THEMES.get(suggestion["theme"])
            suggestion["font_details"] = FONT_FAMILIES.get(suggestion["font_family"])
            suggestion["layout_details"] = LAYOUT_TEMPLATES.get(suggestion["layout"])

            logger.info(
                "Theme suggestion generated",
                theme=suggestion.get("theme"),
                font_family=suggestion.get("font_family"),
                layout=suggestion.get("layout"),
            )

            return suggestion

        except Exception as e:
            logger.warning(f"Theme suggestion failed, using defaults: {e}")
            # Return sensible defaults
            return {
                "theme": "business",
                "font_family": "modern",
                "layout": "standard",
                "animations": False,
                "reason": "Default professional theme selected",
                "theme_details": THEMES["business"],
                "font_details": FONT_FAMILIES["modern"],
                "layout_details": LAYOUT_TEMPLATES["standard"],
            }

    async def create_job(
        self,
        user_id: str,
        title: str,
        description: str,
        output_format: OutputFormat = OutputFormat.DOCX,
        collection_filter: Optional[str] = None,
        folder_id: Optional[str] = None,
        include_subfolders: bool = True,
        metadata: Optional[Dict[str, Any]] = None,
        include_images: Optional[bool] = None,  # Override from request
    ) -> GenerationJob:
        """
        Create a new document generation job.

        Args:
            user_id: ID of the user creating the job
            title: Title for the document
            description: Description of what to generate
            output_format: Desired output format
            collection_filter: Optional collection(s) to search (comma-separated for multiple)
            folder_id: Optional folder ID to scope search to
            include_subfolders: Whether to include subfolders when folder_id is set
            metadata: Additional metadata
            include_images: Override image generation setting for this job

        Returns:
            New GenerationJob instance
        """
        # Reload settings from database to pick up admin changes
        await self.reload_settings()

        job_id = str(uuid.uuid4())

        job = GenerationJob(
            id=job_id,
            user_id=user_id,
            title=title,
            description=description,
            output_format=output_format,
            status=GenerationStatus.DRAFT,
            metadata=metadata or {},
        )

        if collection_filter:
            job.metadata["collection_filter"] = collection_filter
        if folder_id:
            job.metadata["folder_id"] = folder_id
            job.metadata["include_subfolders"] = include_subfolders

        # Store image generation setting - use request override or config setting
        job.metadata["include_images"] = (
            include_images if include_images is not None else self.config.include_images
        )
        job.metadata["image_backend"] = self.config.image_backend
        job.metadata["default_tone"] = self.config.default_tone
        job.metadata["default_style"] = self.config.default_style
        # Store LLM model for notes/metadata - fetch actual model from LLM config
        try:
            from backend.services.llm import LLMConfigManager
            llm_config = await LLMConfigManager.get_config_for_operation("generation")
            job.metadata["llm_model"] = f"{llm_config.provider_type}/{llm_config.model}"
        except Exception as e:
            logger.warning(f"Could not fetch LLM config for notes: {e}")
            job.metadata["llm_model"] = self.config.model

        self._jobs[job_id] = job

        logger.info(
            "Document generation job created",
            job_id=job_id,
            user_id=user_id,
            title=title,
            include_images=job.metadata["include_images"],
            image_backend=job.metadata["image_backend"],
        )

        return job

    async def generate_outline(
        self,
        job_id: str,
        num_sections: Optional[int] = None,
    ) -> DocumentOutline:
        """
        Generate an outline for the document.

        Args:
            job_id: Job ID
            num_sections: Number of sections to generate. None = auto mode (LLM decides)

        Returns:
            Generated outline
        """
        job = self._get_job(job_id)

        if job.status not in [GenerationStatus.DRAFT, GenerationStatus.OUTLINE_PENDING]:
            raise ValueError(f"Cannot generate outline in status: {job.status}")

        job.status = GenerationStatus.OUTLINE_PENDING
        job.updated_at = datetime.utcnow()

        logger.info(
            "Generating outline",
            job_id=job_id,
            num_sections=num_sections,
            mode="auto" if num_sections is None else "manual",
        )

        # Analyze style from existing documents if enabled
        style_guide = None
        if job.metadata.get("use_existing_docs"):
            logger.info(
                "Analyzing document styles for style-aware generation",
                job_id=job_id,
                style_collections=job.metadata.get("style_collection_filters"),
                style_folder=job.metadata.get("style_folder_id"),
            )
            style_guide = await self._analyze_document_styles(
                collection_filters=job.metadata.get("style_collection_filters"),
                folder_id=job.metadata.get("style_folder_id"),
                include_subfolders=job.metadata.get("include_style_subfolders", True),
            )
            if style_guide:
                # Store for later use in content generation
                job.metadata["style_guide"] = style_guide
                logger.info(
                    "Style analysis complete",
                    job_id=job_id,
                    tone=style_guide.get("tone"),
                    vocabulary=style_guide.get("vocabulary_level"),
                    source_docs=len(style_guide.get("source_documents", [])),
                )
            else:
                logger.warning(
                    "No documents found for style analysis",
                    job_id=job_id,
                )

        # Use RAG to find relevant sources
        sources = []

        # Add style source documents to sources for references slide/page
        # These are the documents used for style, language, and pattern learning
        if style_guide and style_guide.get("source_documents"):
            style_sources = style_guide.get("source_documents", [])
            logger.info(
                "Adding style source documents to references",
                job_id=job_id,
                count=len(style_sources),
            )
            for doc_name in style_sources:
                sources.append(SourceReference(
                    document_id="",  # Style sources identified by name
                    document_name=doc_name,
                    chunk_id="",
                    page_number=None,
                    relevance_score=1.0,
                    snippet="Used for style, language, and pattern learning",
                ))
        if self.config.use_rag:
            rag_sources = await self._search_sources(
                query=f"{job.title}: {job.description}",
                collection_filter=job.metadata.get("collection_filter"),
                max_results=self.config.max_sources,
                enhance_query=job.metadata.get("enhance_query"),
            )
            sources.extend(rag_sources)  # Add RAG sources to style sources
            logger.info(
                "RAG sources found for outline generation",
                job_id=job_id,
                rag_sources_count=len(rag_sources),
                total_sources=len(sources),
                query=f"{job.title}: {job.description}"[:100],
                enhance_query=job.metadata.get("enhance_query"),
            )
        else:
            logger.warning(
                "RAG disabled for outline generation - no document sources",
                job_id=job_id,
                use_rag=self.config.use_rag,
            )

        if not sources:
            logger.warning(
                "No sources found for outline generation - outline may be generic",
                job_id=job_id,
                title=job.title,
            )

        # Generate outline using LLM
        outline = await self._generate_outline_with_llm(
            title=job.title,
            description=job.description,
            sources=sources,
            num_sections=num_sections,
            output_format=job.output_format.value,
            style_guide=style_guide,
            output_language=job.metadata.get("output_language", "en"),
        )

        job.outline = outline
        job.sources_used = sources
        job.updated_at = datetime.utcnow()

        logger.info(
            "Outline generated",
            job_id=job_id,
            sections=len(outline.sections),
        )

        return outline

    async def approve_outline(
        self,
        job_id: str,
        modifications: Optional[Dict[str, Any]] = None,
    ) -> GenerationJob:
        """
        Approve the outline and optionally apply modifications.

        Args:
            job_id: Job ID
            modifications: Optional changes to the outline

        Returns:
            Updated job
        """
        job = self._get_job(job_id)

        if job.status != GenerationStatus.OUTLINE_PENDING:
            raise ValueError(f"Cannot approve outline in status: {job.status}")

        if not job.outline:
            raise ValueError("No outline to approve")

        # Apply modifications if provided
        if modifications:
            if "title" in modifications:
                job.outline.title = modifications["title"]
                # Also update the job title so it's used in document generation (PPTX, etc.)
                job.title = modifications["title"]
            if "sections" in modifications:
                job.outline.sections = modifications["sections"]
            if "tone" in modifications:
                job.outline.tone = modifications["tone"]
            # Allow theme change after outline generation
            if "theme" in modifications and modifications["theme"]:
                job.metadata["theme"] = modifications["theme"]
                logger.info(f"Theme changed to: {modifications['theme']}")

        job.status = GenerationStatus.OUTLINE_APPROVED
        job.updated_at = datetime.utcnow()

        logger.info("Outline approved", job_id=job_id)

        # Auto-generate if configured
        if self.config.auto_generate_on_approval:
            await self.generate_content(job_id)

        return job

    async def generate_content(
        self,
        job_id: str,
    ) -> GenerationJob:
        """
        Generate the document content based on approved outline.

        Args:
            job_id: Job ID

        Returns:
            Updated job with generated content
        """
        job = self._get_job(job_id)

        if job.status != GenerationStatus.OUTLINE_APPROVED:
            raise ValueError(f"Cannot generate content in status: {job.status}")

        if not job.outline:
            raise ValueError("No outline available")

        job.status = GenerationStatus.GENERATING
        job.updated_at = datetime.utcnow()

        logger.info("Generating content", job_id=job_id)

        try:
            # Generate each section
            sections = []
            for idx, section_def in enumerate(job.outline.sections):
                section = await self._generate_section(
                    job=job,
                    section_title=section_def.get("title", f"Section {idx + 1}"),
                    section_description=section_def.get("description", ""),
                    order=idx,
                )
                sections.append(section)

                logger.info(
                    "Section generated",
                    job_id=job_id,
                    section_title=section.title,
                    order=idx,
                )

            job.sections = sections

            # Aggregate section sources into job.sources_used for reference slide
            # This ensures sources found during section generation are included
            if not job.sources_used:
                job.sources_used = []
            existing_source_ids = {s.document_id for s in job.sources_used}
            for section in sections:
                if section.sources:
                    for source in section.sources:
                        if source.document_id not in existing_source_ids:
                            job.sources_used.append(source)
                            existing_source_ids.add(source.document_id)

            logger.info(
                "Sources aggregated from sections",
                job_id=job_id,
                total_sources=len(job.sources_used),
            )

            # Backfill section sources from job.sources_used for sections with no sources
            # This ensures speaker notes have sources to display
            backfilled_count = 0
            for section in job.sections:
                if not section.sources and job.sources_used:
                    section.sources = job.sources_used[:5]
                    backfilled_count += 1
            if backfilled_count > 0:
                logger.info(
                    "Backfilled section sources from job sources",
                    job_id=job_id,
                    backfilled_sections=backfilled_count,
                    sources_per_section=min(5, len(job.sources_used)),
                )

            # Move to review or completed based on config
            if self.config.require_section_approval:
                job.status = GenerationStatus.SECTION_REVIEW
            else:
                job.status = GenerationStatus.COMPLETED
                job.completed_at = datetime.utcnow()

                # Generate output file
                await self._generate_output_file(job)

            job.updated_at = datetime.utcnow()

        except Exception as e:
            logger.error("Content generation failed", job_id=job_id, error=str(e))
            job.status = GenerationStatus.FAILED
            job.error_message = str(e)
            job.updated_at = datetime.utcnow()
            raise

        return job

    async def approve_section(
        self,
        job_id: str,
        section_id: str,
        feedback: Optional[str] = None,
        approved: bool = True,
    ) -> Section:
        """
        Approve or request revision for a section.

        Args:
            job_id: Job ID
            section_id: Section ID
            feedback: Optional feedback for revision
            approved: Whether the section is approved

        Returns:
            Updated section
        """
        job = self._get_job(job_id)

        if job.status not in [GenerationStatus.SECTION_REVIEW, GenerationStatus.REVISION]:
            raise ValueError(f"Cannot approve section in status: {job.status}")

        section = next((s for s in job.sections if s.id == section_id), None)
        if not section:
            raise ValueError(f"Section not found: {section_id}")

        section.approved = approved
        section.feedback = feedback if not approved else None

        if not approved:
            job.status = GenerationStatus.REVISION
        else:
            # Check if all sections are approved
            all_approved = all(s.approved for s in job.sections)
            if all_approved:
                job.status = GenerationStatus.COMPLETED
                job.completed_at = datetime.utcnow()

                # Generate output file
                await self._generate_output_file(job)

        job.updated_at = datetime.utcnow()

        return section

    async def revise_section(
        self,
        job_id: str,
        section_id: str,
    ) -> Section:
        """
        Revise a section based on feedback.

        Args:
            job_id: Job ID
            section_id: Section ID

        Returns:
            Revised section
        """
        job = self._get_job(job_id)

        if job.status != GenerationStatus.REVISION:
            raise ValueError(f"Cannot revise section in status: {job.status}")

        section = next((s for s in job.sections if s.id == section_id), None)
        if not section:
            raise ValueError(f"Section not found: {section_id}")

        if not section.feedback:
            raise ValueError("No feedback provided for revision")

        # Regenerate with feedback
        revised_section = await self._regenerate_section_with_feedback(
            job=job,
            section=section,
        )

        # Update section
        section.revised_content = revised_section.content
        section.sources = revised_section.sources
        section.approved = False
        section.feedback = None

        job.status = GenerationStatus.SECTION_REVIEW
        job.updated_at = datetime.utcnow()

        return section

    async def get_job(self, job_id: str) -> GenerationJob:
        """Get a generation job by ID."""
        return self._get_job(job_id)

    async def list_jobs(
        self,
        user_id: str,
        status: Optional[GenerationStatus] = None,
    ) -> List[GenerationJob]:
        """List jobs for a user."""
        jobs = [j for j in self._jobs.values() if j.user_id == user_id]

        if status:
            jobs = [j for j in jobs if j.status == status]

        return sorted(jobs, key=lambda j: j.created_at, reverse=True)

    async def cancel_job(self, job_id: str) -> GenerationJob:
        """Cancel a generation job."""
        job = self._get_job(job_id)

        if job.status in [GenerationStatus.COMPLETED, GenerationStatus.CANCELLED]:
            raise ValueError(f"Cannot cancel job in status: {job.status}")

        job.status = GenerationStatus.CANCELLED
        job.updated_at = datetime.utcnow()

        logger.info("Job cancelled", job_id=job_id)

        return job

    async def get_output_file(self, job_id: str) -> tuple[bytes, str, str]:
        """
        Get the generated output file.

        Args:
            job_id: Job ID

        Returns:
            Tuple of (file_bytes, filename, content_type)
        """
        job = self._get_job(job_id)

        if job.status != GenerationStatus.COMPLETED:
            raise ValueError("Job is not completed")

        if not job.output_path:
            raise ValueError("No output file available")

        # Read file
        with open(job.output_path, "rb") as f:
            file_bytes = f.read()

        filename = Path(job.output_path).name
        content_type = self._get_content_type(job.output_format)

        return file_bytes, filename, content_type

    # -------------------------------------------------------------------------
    # Private Methods
    # -------------------------------------------------------------------------

    def _get_job(self, job_id: str) -> GenerationJob:
        """Get job by ID or raise error."""
        job = self._jobs.get(job_id)
        if not job:
            raise ValueError(f"Job not found: {job_id}")
        return job

    async def _search_sources(
        self,
        query: str,
        collection_filter: Optional[str] = None,
        max_results: int = 10,
        enhance_query: Optional[bool] = None,
    ) -> List[SourceReference]:
        """Search for relevant sources using RAG with semantic search.

        Args:
            query: The search query
            collection_filter: Optional collection filter
            max_results: Maximum number of results to return
            enhance_query: Enable query enhancement (expansion + HyDE).
                          None = use admin setting, True/False = override.
        """
        try:
            from backend.services.vectorstore import get_vector_store, SearchType
            from backend.services.llm import generate_embedding
            from backend.services.settings import SettingsService

            vector_store = get_vector_store()

            # Determine if query enhancement is enabled
            settings = SettingsService()
            use_enhancement = enhance_query
            if use_enhancement is None:
                # Check admin settings
                use_enhancement = settings.get_default_value("rag.query_expansion_enabled")
                if use_enhancement is None:
                    use_enhancement = True  # Default to enabled

            # Apply query enhancement if enabled
            enhanced_query = query
            all_queries = [query]

            if use_enhancement:
                try:
                    from backend.services.query_expander import get_query_expander, QueryExpansionConfig
                    from backend.services.hyde import get_hyde_expander

                    # Query expansion - generate variations
                    expander_config = QueryExpansionConfig(expansion_count=3)
                    expander = get_query_expander(expander_config)
                    if expander:
                        expanded = await expander.expand(query)
                        if expanded and expanded.expanded_queries:
                            all_queries.extend(expanded.expanded_queries[:2])  # Add top 2 variations
                            logger.debug(
                                "Query expanded for source search",
                                original=query[:50],
                                variations=len(expanded.expanded_queries),
                            )

                    # HyDE - Hypothetical Document Embedding (optional, can be expensive)
                    hyde_enabled = settings.get_default_value("rag.hyde_enabled")
                    if hyde_enabled:
                        hyde_expander = get_hyde_expander()
                        if hyde_expander:
                            from backend.services.llm import LLMService
                            llm_service = LLMService()
                            hyde_result = await hyde_expander.expand(query, llm_service)
                            if hyde_result and hyde_result.hypothetical_doc:
                                # Use HyDE doc for embedding instead of original query
                                enhanced_query = hyde_result.hypothetical_doc[:1000]
                                logger.debug(
                                    "HyDE applied for source search",
                                    original=query[:50],
                                )
                except Exception as enhance_err:
                    logger.warning(
                        "Query enhancement failed, using original query",
                        error=str(enhance_err),
                        query=query[:50],
                    )

            # Generate embedding for semantic search (much better than keyword-only)
            query_embedding = None
            try:
                query_embedding = await generate_embedding(enhanced_query)
                logger.debug(
                    "Generated query embedding for source search",
                    query=query[:50],
                    enhanced=enhanced_query != query,
                    embedding_dims=len(query_embedding) if query_embedding else 0,
                )
            except Exception as embed_err:
                logger.warning(
                    "Failed to generate query embedding, falling back to keyword search",
                    error=str(embed_err),
                    query=query[:50],
                )

            logger.info(
                "Executing vector store search",
                query=query[:100],
                has_embedding=query_embedding is not None,
                search_type="hybrid" if query_embedding else "keyword",
                max_results=max_results,
            )

            results = await vector_store.search(
                query=query,
                query_embedding=query_embedding,
                search_type=SearchType.HYBRID if query_embedding else SearchType.KEYWORD,
                top_k=max_results,
            )

            logger.info(
                "Vector store search completed",
                query=query[:50],
                results_count=len(results) if results else 0,
                min_score=self.config.min_relevance_score,
                search_type="hybrid" if query_embedding else "keyword",
                top_scores=[r.score for r in (results or [])[:3]],
            )

            sources = []
            for result in results:
                if result.score >= self.config.min_relevance_score:
                    # Get document name from multiple possible fields
                    doc_name = (
                        getattr(result, 'document_title', '') or
                        getattr(result, 'document_filename', '') or
                        str(result.document_id)
                    )
                    sources.append(
                        SourceReference(
                            document_id=str(result.document_id),
                            document_name=doc_name,
                            chunk_id=str(result.chunk_id),
                            page_number=result.page_number,
                            relevance_score=result.score,
                            snippet=result.content[:200],
                        )
                    )

            if sources:
                logger.info(
                    "Section sources found",
                    query=query[:50],
                    sources_count=len(sources),
                    source_names=[s.document_name for s in sources[:3]],
                )
            else:
                logger.warning(
                    "No sources found after filtering",
                    query=query[:50],
                    raw_results_count=len(results) if results else 0,
                    min_score=self.config.min_relevance_score,
                    hint="Check if documents are indexed or min_relevance_score is too high",
                )

            return sources

        except Exception as e:
            logger.error(
                "Failed to search sources - sources will be empty for this section",
                error=str(e),
                error_type=type(e).__name__,
                query=query[:50],
            )
            import traceback
            logger.debug("Source search exception traceback", traceback=traceback.format_exc())
            return []

    async def _analyze_document_styles(
        self,
        collection_filters: Optional[List[str]] = None,
        folder_id: Optional[str] = None,
        include_subfolders: bool = True,
        sample_size: Optional[int] = None,  # None = use ALL documents
    ) -> Optional[dict]:
        """Analyze existing documents to extract style patterns for new document generation.

        This method retrieves documents from the specified collections/folders
        and uses LLM to analyze their writing style, tone, vocabulary, and structure.
        All document names are recorded for references, but content is sampled for LLM analysis.

        Args:
            collection_filters: List of collections to sample from
            folder_id: Folder ID to scope the search
            include_subfolders: Whether to include subfolders
            sample_size: Number of unique documents to sample for content analysis (None = all)

        Returns:
            Style analysis dict with tone, vocabulary_level, structure_pattern, etc.
            Returns None if no documents found or analysis fails.
        """
        try:
            import json
            from backend.services.vectorstore import get_vector_store, SearchType
            from backend.db.database import async_session_context
            from backend.db.models import Document as DBDocument, Folder
            from sqlalchemy import select, cast, String, literal, or_

            vector_store = get_vector_store()

            # Build document ID filter based on collection/folder filters
            document_ids = None

            # Get document IDs matching collection filters
            if collection_filters:
                async with async_session_context() as db:
                    all_collection_doc_ids = set()
                    for collection in collection_filters:
                        if collection == "(Untagged)":
                            query_stmt = select(DBDocument.id).where(
                                or_(
                                    DBDocument.tags.is_(None),
                                    DBDocument.tags == [],
                                )
                            )
                        else:
                            # Escape LIKE special characters
                            safe_filter = collection.replace("\\", "\\\\")
                            safe_filter = safe_filter.replace("%", "\\%")
                            safe_filter = safe_filter.replace("_", "\\_")
                            safe_filter = safe_filter.replace('"', '\\"')
                            pattern = f'%"{safe_filter}"%'
                            query_stmt = select(DBDocument.id).where(
                                cast(DBDocument.tags, String).like(literal(pattern))
                            )
                        result = await db.execute(query_stmt)
                        all_collection_doc_ids.update(str(row[0]) for row in result.fetchall())

                    if all_collection_doc_ids:
                        document_ids = list(all_collection_doc_ids)
                        logger.info(
                            "Style analysis filtered by collections",
                            collections=collection_filters,
                            document_count=len(document_ids),
                        )
                    else:
                        logger.info("No documents match collection filters for style analysis")
                        return None

            # Get document IDs from folder filter
            if folder_id:
                async with async_session_context() as db:
                    folder_doc_ids = set()

                    if include_subfolders:
                        # Get folder and all subfolders
                        folder_ids = [folder_id]
                        # Get subfolders recursively
                        result = await db.execute(
                            select(Folder.id).where(Folder.parent_id == folder_id)
                        )
                        subfolder_ids = [str(row[0]) for row in result.fetchall()]
                        folder_ids.extend(subfolder_ids)

                        # Get documents in all folders
                        for fid in folder_ids:
                            result = await db.execute(
                                select(DBDocument.id).where(DBDocument.folder_id == fid)
                            )
                            folder_doc_ids.update(str(row[0]) for row in result.fetchall())
                    else:
                        # Just the specified folder
                        result = await db.execute(
                            select(DBDocument.id).where(DBDocument.folder_id == folder_id)
                        )
                        folder_doc_ids.update(str(row[0]) for row in result.fetchall())

                    if folder_doc_ids:
                        if document_ids:
                            # Intersect with collection filter results
                            document_ids = list(set(document_ids) & folder_doc_ids)
                        else:
                            document_ids = list(folder_doc_ids)
                        logger.info(
                            "Style analysis filtered by folder",
                            folder_id=folder_id,
                            include_subfolders=include_subfolders,
                            document_count=len(document_ids),
                        )
                    else:
                        logger.info("No documents in folder for style analysis")
                        return None

            # Log if searching all documents (no filters)
            if document_ids is None:
                logger.info("Style analysis searching all documents (no filters specified)")

            # Get document samples directly from database (more reliable than search)
            # This approach ensures we get samples even when search terms don't match
            from backend.db.models import Chunk as DBChunk

            samples = []
            all_doc_names = set()  # Track ALL document names for references
            async with async_session_context() as db:
                # Build query for chunks with content
                chunk_query = select(
                    DBChunk.content,
                    DBChunk.document_id,
                    DBDocument.original_filename.label("document_name"),
                    DBDocument.tags.label("collection"),
                ).join(DBDocument, DBChunk.document_id == DBDocument.id)

                # Apply document_ids filter if specified
                if document_ids:
                    chunk_query = chunk_query.where(
                        cast(DBChunk.document_id, String).in_(document_ids)
                    )

                # Order by document to get variety
                chunk_query = chunk_query.order_by(DBDocument.id)

                # Only limit if sample_size is specified
                if sample_size is not None:
                    chunk_query = chunk_query.limit(sample_size * 5)

                result = await db.execute(chunk_query)
                rows = result.fetchall()

                # Extract unique document samples (for LLM analysis, limit content samples)
                seen_docs = set()
                max_content_samples = sample_size if sample_size else 10  # Cap content samples at 10 for LLM
                for row in rows:
                    doc_id = str(row.document_id) if row.document_id else None
                    doc_name = row.document_name or ""

                    # Track ALL document names for references
                    if doc_name:
                        all_doc_names.add(doc_name)

                    # Sample content for LLM analysis (limited)
                    if doc_id and doc_id not in seen_docs and len(samples) < max_content_samples:
                        seen_docs.add(doc_id)
                        collection_tags = row.collection if row.collection else []
                        samples.append({
                            "content": row.content or "",
                            "document_name": doc_name,
                            "collection": collection_tags[0] if collection_tags else "",
                        })
                        logger.debug(
                            "Style sample added",
                            doc_id=doc_id,
                            doc_name=doc_name,
                        )

            if not samples:
                logger.info("No document samples found for style analysis")
                return None

            logger.info(
                "Style analysis samples collected",
                sample_count=len(samples),
                total_documents=len(all_doc_names),
                document_names=list(all_doc_names),
            )

            # Use LLM to analyze style patterns
            sample_text = "\n\n".join([
                f"--- Sample from {s['document_name'] or 'Unknown'} ---\n{s['content'][:1000]}"
                for s in samples
            ])

            analysis_prompt = f"""Analyze the following document excerpts and extract their writing style patterns.

DOCUMENT SAMPLES:
{sample_text}

Analyze and return a JSON object with:
{{
    "tone": "formal" | "casual" | "technical" | "friendly" | "academic",
    "vocabulary_level": "simple" | "moderate" | "advanced",
    "structure_pattern": "bullet-lists" | "paragraphs" | "mixed" | "headers-heavy",
    "sentence_style": "short-concise" | "medium" | "long-detailed",
    "key_phrases": ["common phrases or terms used"],
    "formatting_notes": "any notable formatting patterns observed",
    "recommended_approach": "brief recommendation for matching this style"
}}

Return ONLY valid JSON, no other text."""

            # Use EnhancedLLMFactory to get LLM instance (same pattern as other methods)
            from backend.services.llm import EnhancedLLMFactory

            llm, llm_config = await EnhancedLLMFactory.get_chat_model_for_operation(
                operation="content_generation",
                user_id=None,  # System-level operation
            )
            response = await llm.ainvoke(analysis_prompt)
            response_text = response.content if hasattr(response, 'content') else str(response)

            # Parse the JSON response
            import json
            # Try to extract JSON from the response
            json_match = re.search(r'\{[^{}]*\}', response_text, re.DOTALL)
            if json_match:
                style_analysis = json.loads(json_match.group())
            else:
                style_analysis = json.loads(response_text)

            # Use ALL document names for references, not just the sampled ones
            style_analysis["source_documents"] = list(all_doc_names)

            logger.info(
                "Style analysis completed",
                tone=style_analysis.get("tone"),
                vocabulary=style_analysis.get("vocabulary_level"),
                num_sources=len(style_analysis.get("source_documents", [])),
            )

            return style_analysis

        except json.JSONDecodeError as e:
            logger.warning(f"Failed to parse style analysis JSON: {e}")
            # Return default style guide with all document names for references
            return {
                "tone": "professional",
                "vocabulary_level": "moderate",
                "structure_pattern": "mixed",
                "sentence_style": "medium",
                "source_documents": list(all_doc_names) if all_doc_names else [],
            }
        except Exception as e:
            logger.warning(f"Failed to analyze document styles: {e}")
            return None

    async def _generate_outline_with_llm(
        self,
        title: str,
        description: str,
        sources: List[SourceReference],
        num_sections: Optional[int] = None,
        output_format: str = "docx",
        style_guide: Optional[Dict[str, Any]] = None,
        output_language: str = "en",
    ) -> DocumentOutline:
        """Generate outline using LLM.

        Args:
            title: Document title
            description: Document description
            sources: Relevant sources from RAG
            num_sections: Number of sections. None = auto mode (LLM decides)
            output_format: Target output format (affects section count guidance)
            style_guide: Optional style analysis from existing documents
            output_language: Language code for generated content (default: en)
        """
        # Build context from sources
        context = ""
        if sources:
            context = "Relevant information from the knowledge base:\n\n"
            for source in sources[:5]:
                context += f"- {source.snippet}...\n\n"
            logger.debug(
                "Outline context built from sources",
                sources_used=len(sources[:5]),
                context_length=len(context),
                first_snippet=sources[0].snippet[:100] if sources else "none",
            )
        else:
            logger.warning(
                "No sources available for outline context - LLM will generate without document knowledge"
            )

        # Build style instructions if available
        style_instructions = ""
        if style_guide:
            style_instructions = f"""
---INTERNAL STYLE GUIDANCE (DO NOT OUTPUT THIS AS CONTENT)---
Use these style hints when writing, but DO NOT include them as sections:
- Tone: {style_guide.get('tone', 'professional')}
- Vocabulary: {style_guide.get('vocabulary_level', 'moderate')}
Match the style of existing documents. This is internal guidance only.
---END INTERNAL GUIDANCE---
"""

        # Build language instruction
        language_instruction = ""
        if output_language == "auto":
            # Auto-detect: prioritize the TOPIC/TITLE language, then source documents
            language_instruction = """
LANGUAGE REQUIREMENT:
1. FIRST, detect the language of the document TITLE/TOPIC provided by the user.
2. Generate ALL section titles and descriptions in the SAME LANGUAGE as the title/topic.
3. If the title is in Hinglish (Hindi+English mix), respond in Hinglish/Hindi.
4. If the title is in German, respond in German.
5. If the title is in English, respond in English.
6. Source documents may be in ANY language - translate relevant information to match the title language.
7. Do NOT default to English just because source documents are in English.

Example: If title is "maketing startegy gaadi ke bare mean" - this is Hinglish, so ALL sections should be in Hinglish/Hindi.
"""
        elif output_language != "en":
            language_name = LANGUAGE_NAMES.get(output_language, "English")
            language_instruction = f"""
LANGUAGE REQUIREMENT:
- Generate ALL section titles and descriptions in {language_name}.
- If source material is in a different language, translate concepts to {language_name}.
- Technical terms may remain in English if commonly used that way.
"""

        # Build section instruction based on mode
        if num_sections is None:
            # Auto mode - LLM decides optimal count
            section_instruction = f"""Analyze the topic complexity and determine the optimal number of sections.

Consider:
- Topic depth and complexity
- Output format: {output_format} (presentations typically need 5-10 slides, documents 3-8 pages)
- Amount of source material available
- Target audience expectations

First, output your recommended section count on a line by itself like:
SECTION_COUNT: N

Where N is between 3 and 15.

Then generate exactly N sections with specific, descriptive titles."""
        else:
            section_instruction = f"Generate exactly {num_sections} sections with specific, descriptive titles."

        # Create prompt for outline generation
        prompt = f"""Create a professional document outline for the following:
{style_instructions}
{language_instruction}
Title: {title}
Description: {description}

{context}

{section_instruction}

CRITICAL RULES:
1. Each section title MUST contain specific keywords from the topic "{title}"
2. NEVER use generic template titles - every title must be unique to THIS specific topic
3. Include specific nouns, brands, products, or concepts from the topic in each title
4. Descriptions must explain WHAT SPECIFIC CONTENT will be in that section

ABSOLUTELY FORBIDDEN GENERIC TITLES (never use these or similar):
- "Background and Context", "Key Analysis", "Strategic Recommendations"
- "Implementation Details", "Conclusion and Next Steps", "Overview"
- "Introduction", "Summary", "Key Findings", "Analysis and Findings"
- "Action Items", "Next Steps", "Recommendations"
- Any title that could apply to ANY document topic

Format each section EXACTLY like this:
## [Title with specific topic keywords]
Description: [Specific content about THIS topic, not generic filler]

GOOD EXAMPLE for "marketing strategy for upcoming shoe launch":
## Target Demographics for Athletic Footwear
Description: Analysis of key customer segments including runners, casual athletes, and fashion-conscious buyers aged 18-35.

## Social Media Campaign for Shoe Release
Description: Instagram, TikTok, and influencer partnership strategies to build pre-launch excitement.

## Retail and E-commerce Distribution Plan
Description: Store placement strategy, online launch timing, and inventory allocation across channels.

## Competitive Pricing Analysis vs Nike and Adidas
Description: Price point positioning relative to competitors, value proposition, and promotional pricing strategy.

## Launch Event and PR Timeline
Description: Press release schedule, influencer seeding dates, and launch day activation plans.

BAD EXAMPLE (NEVER DO THIS):
## Background and Context ❌
## Key Analysis and Findings ❌
## Strategic Recommendations ❌
## Implementation Details ❌
## Conclusion and Next Steps ❌

Generate the outline for "{title}" with SPECIFIC, NON-GENERIC titles now:"""

        # Use LLM to generate (database-driven configuration)
        try:
            from backend.services.llm import EnhancedLLMFactory

            llm, config = await EnhancedLLMFactory.get_chat_model_for_operation(
                operation="content_generation",
                user_id=None,  # System-level operation
            )
            response = await llm.ainvoke(prompt)

            # Log the raw LLM response for debugging
            logger.info(
                "LLM outline response received",
                response_length=len(response.content) if response.content else 0,
                response_preview=response.content[:500] if response.content else "empty",
            )

            # Parse response into sections with improved parsing
            import re
            sections = []
            lines = response.content.split("\n")
            current_section = None

            # In auto mode, extract the section count from LLM response
            target_sections = num_sections
            if num_sections is None:
                # Look for SECTION_COUNT: N pattern
                section_count_match = re.search(r'SECTION_COUNT:\s*(\d+)', response.content)
                if section_count_match:
                    target_sections = int(section_count_match.group(1))
                    # Clamp to valid range
                    target_sections = max(3, min(15, target_sections))
                    logger.info(f"Auto mode: LLM suggested {target_sections} sections")
                else:
                    # Fallback: count sections in response, or default to 5
                    target_sections = 5
                    logger.warning("Auto mode: Could not extract section count, defaulting to 5")

            # Generic title patterns to reject
            generic_patterns = [
                r'^section\s*\d*$',
                r'^introduction$',
                r'^overview$',
                r'^summary$',
                r'^conclusion$',
                r'^part\s*\d+$',
                r'^section_count',  # Metadata line, not a section
                r'^content\s+covering',  # Description of metadata, not a section
            ]

            for line in lines:
                line = line.strip()
                if not line:
                    continue

                # Skip SECTION_COUNT metadata line (already extracted above)
                # Check for SECTION_COUNT anywhere in the line (LLM might format it differently)
                if re.search(r'SECTION_COUNT\s*:', line, re.IGNORECASE):
                    continue

                # Skip conversational LLM artifacts (User:, Assistant:, Human:, etc.)
                if re.match(r'^#*\s*(User|Assistant|Human|AI|System|Here is|Below is|I\'ll|Let me|Sure|Certainly)[:.]?\s', line, re.IGNORECASE):
                    continue

                # Check for section header - multiple formats supported:
                # - ## Title (markdown)
                # - 1. Title (numbered)
                # - **Title** (bold markdown)
                # - - Title (bullet)
                # - Title: (colon-ending titles)
                is_section_header = (
                    line.startswith("##") or
                    line.startswith("**") or
                    line.startswith("- ") or
                    (len(line) > 0 and line[0].isdigit() and ("." in line[:4] or ")" in line[:4])) or
                    (line.endswith(":") and len(line) < 100 and not line.lower().startswith("description"))
                )

                if is_section_header:
                    if current_section and current_section["title"]:
                        sections.append(current_section)

                    # Extract section title, removing markdown formatting and numbering
                    # Use section_title to avoid shadowing the document title parameter
                    section_title = re.sub(r'^[#\d.\s\-\*\)]+', '', line).strip()
                    section_title = section_title.rstrip(':').strip()  # Remove trailing colon
                    section_title = re.sub(r'\*+', '', section_title).strip()  # Remove bold markers

                    # Check if title is generic
                    is_generic = any(
                        re.match(pattern, section_title.lower().strip())
                        for pattern in generic_patterns
                    )

                    if is_generic or not section_title:
                        # Generate a better title based on document topic
                        section_title = f"Key Aspect {len(sections) + 1} of {description[:30].split()[0].title() if description else 'Topic'}"

                    current_section = {"title": section_title, "description": ""}

                elif current_section and line:
                    # Handle description lines
                    if line.lower().startswith("description:"):
                        desc_content = line[12:].strip()
                        if desc_content:
                            current_section["description"] = desc_content + " "
                    elif not current_section["description"]:
                        # First non-header, non-description-prefixed line after section title
                        current_section["description"] = line.strip() + " "
                    elif current_section["description"] and not current_section["description"].strip():
                        # Description was set but is empty/whitespace, use this line
                        current_section["description"] = line.strip() + " "

            if current_section and current_section["title"]:
                sections.append(current_section)

            # Log parsed sections count
            logger.info(
                "Outline sections parsed from LLM response",
                parsed_sections=len(sections),
                target_sections=target_sections,
                section_titles=[s.get("title", "")[:50] for s in sections[:5]],
            )

            # Ensure all sections have non-empty descriptions
            for section in sections:
                if not section.get("description") or not section["description"].strip():
                    # Generate a fallback description based on the section title
                    section["description"] = f"Content covering {section['title'].lower()}. "

            # Ensure we have the requested number of sections with meaningful titles
            topic_words = title.split()[:3] if title else ["Document"]
            topic_prefix = " ".join(topic_words)

            while len(sections) < target_sections:
                section_num = len(sections) + 1
                # Generate contextual fallback titles
                fallback_titles = [
                    f"Key Findings and Insights",
                    f"Analysis and Recommendations",
                    f"Implementation Approach",
                    f"Strategic Considerations",
                    f"Supporting Details",
                    f"Additional Context",
                ]
                fallback_title = fallback_titles[min(section_num - 1, len(fallback_titles) - 1)]
                sections.append({
                    "title": f"{fallback_title} for {topic_prefix}",
                    "description": f"Detailed content covering {fallback_title.lower()}",
                })

            return DocumentOutline(
                title=title,
                description=description,
                sections=sections[:target_sections],
            )

        except Exception as e:
            logger.error("Failed to generate outline with LLM", error=str(e))
            # Return a more descriptive fallback outline
            # Use num_sections if specified, otherwise default to 5
            fallback_count = num_sections if num_sections is not None else 5
            fallback_section_templates = [
                ("Background and Context", f"Overview and background information about {title}"),
                ("Key Analysis and Findings", f"Main analysis and findings related to {title}"),
                ("Strategic Recommendations", f"Recommendations and action items for {title}"),
                ("Implementation Details", f"How to implement the strategies for {title}"),
                ("Conclusion and Next Steps", f"Summary and suggested next steps for {title}"),
                ("Supporting Information", f"Additional details and references for {title}"),
            ]
            return DocumentOutline(
                title=title,
                description=description,
                sections=[
                    {
                        "title": fallback_section_templates[min(i, len(fallback_section_templates) - 1)][0],
                        "description": fallback_section_templates[min(i, len(fallback_section_templates) - 1)][1]
                    }
                    for i in range(fallback_count)
                ],
            )

    async def _generate_section(
        self,
        job: GenerationJob,
        section_title: str,
        section_description: str,
        order: int,
    ) -> Section:
        """Generate content for a section."""
        section_id = str(uuid.uuid4())

        # Search for relevant sources for this section
        # Include job title/description for better context - section titles alone are often too generic
        section_query = f"{job.title} - {section_title}"
        if section_description:
            section_query += f": {section_description}"

        sources = await self._search_sources(
            query=section_query,
            collection_filter=job.metadata.get("collection_filter"),
            max_results=5,
            enhance_query=job.metadata.get("enhance_query"),
        )

        # Log sources found for debugging
        logger.info(
            "Section sources search completed",
            section_title=section_title[:50],
            sources_found=len(sources),
            source_names=[s.document_name for s in sources[:3]] if sources else [],
        )

        # Build context
        context = ""
        if sources:
            context = "Use the following information:\n\n"
            for source in sources:
                context += f"- {source.snippet}\n\n"

        # Generate format-specific prompts
        if job.output_format == OutputFormat.PPTX:
            format_instructions = """FORMAT REQUIREMENTS FOR PRESENTATION SLIDES:
- Write 6-10 bullet points maximum
- Each bullet MUST be a COMPLETE sentence under 90 characters
- If a point is too long, split it into two shorter points
- Use simple, concise language - avoid complex sentences
- NO markdown formatting (no **, no ##, no _)
- Focus on key takeaways, not detailed explanations
- Start each point with an action verb or key noun
- NEVER leave a sentence incomplete or cut off
- Use "• " (bullet) for main points, "  ◦ " (2-space indent + open circle) for sub-points

CRITICAL OUTPUT RULES:
1. Start DIRECTLY with the first bullet point - NO introductory text
2. Do NOT include phrases like "Here are the bullet points", "Let me provide", "I'll create"
3. Do NOT include closing remarks like "Let me know if you need adjustments"
4. ONLY output the bullet points themselves - nothing else

Every bullet point must be a complete, standalone thought.
If you cannot express an idea in under 90 characters, break it into multiple shorter points.

Example format (each is a complete sentence):
• Revenue increased 25% year-over-year across all regions.
  ◦ North America led with 32% growth.
  ◦ Europe showed steady 18% improvement.
• Customer acquisition cost reduced by 15% through optimization.
• New market expansion in Q3 shows promising early results."""
        elif job.output_format in (OutputFormat.DOCX, OutputFormat.PDF):
            format_instructions = """FORMAT REQUIREMENTS FOR DOCUMENT:
- Write 3-5 well-structured paragraphs
- Use clear topic sentences for each paragraph
- Include relevant details and examples
- Maintain professional, formal tone
- Use **bold** for key terms (will be styled)
- Use _italic_ for emphasis (will be styled)
- Avoid excessive bullet points - prefer prose
- Target 200-400 words per section"""
        elif job.output_format == OutputFormat.XLSX:
            format_instructions = """FORMAT REQUIREMENTS FOR SPREADSHEET:
- Write concise, data-oriented content
- Use short sentences that fit in cells
- NO markdown formatting (no **, no ##, no _)
- Focus on quantifiable information
- Use clear, structured points
- Each point on a new line
- Target 5-10 key points
- Avoid long paragraphs"""
        else:
            format_instructions = """Write clear, well-structured content.
Include relevant details and maintain a professional tone."""

        # Calculate section position context
        total_sections = len(job.outline.sections) if job.outline else 5
        current_section_num = order + 1  # order is 0-indexed

        if current_section_num <= 2:
            position_context = "This is an EARLY section - set the stage and introduce key concepts."
        elif current_section_num >= total_sections - 1:
            position_context = "This is a CONCLUDING section - summarize key points and provide actionable takeaways."
        else:
            position_context = "This is a MIDDLE section - dive into details, analysis, and supporting information."

        # Build style context if available from style analysis
        style_context = ""
        style_guide = job.metadata.get("style_guide")
        if style_guide:
            style_context = f"""
---INTERNAL STYLE GUIDANCE (follow these rules but do NOT include them in your output)---
Tone: {style_guide.get('tone', 'professional')}
Vocabulary: {style_guide.get('vocabulary_level', 'moderate')}
Structure: {style_guide.get('structure_pattern', 'mixed')}
Sentence style: {style_guide.get('sentence_style', 'medium')}
{f"Approach: {style_guide.get('recommended_approach')}" if style_guide.get('recommended_approach') else ""}
Match the style and tone of existing documents. Do NOT output these instructions as content.
---END INTERNAL GUIDANCE---
"""

        # Build language instruction
        output_language = job.metadata.get("output_language", "en")
        language_instruction = ""
        if output_language == "auto":
            # Auto-detect: programmatically detect title language with confidence threshold
            detected_lang = "en"  # default fallback
            detected_lang_name = "English"
            detection_confidence = 0.0
            try:
                from langdetect import detect_langs, DetectorFactory
                DetectorFactory.seed = 0  # Make deterministic
                # Combine title and description for better detection
                text_to_detect = f"{job.title} {job.description or ''}"

                # detect_langs returns list of languages with probabilities
                detected_results = detect_langs(text_to_detect)
                if detected_results:
                    top_result = detected_results[0]
                    detected_lang = top_result.lang
                    detection_confidence = top_result.prob

                    # Map detected language to name
                    lang_map = {
                        "en": "English", "de": "German", "es": "Spanish", "fr": "French",
                        "it": "Italian", "pt": "Portuguese", "zh-cn": "Chinese", "zh-tw": "Chinese",
                        "ja": "Japanese", "ko": "Korean", "ar": "Arabic", "hi": "Hindi",
                        "ru": "Russian", "nl": "Dutch", "pl": "Polish", "tr": "Turkish"
                    }
                    detected_lang_name = lang_map.get(detected_lang, detected_lang.upper())

                    logger.info(
                        "Auto-detected language from title",
                        detected_lang=detected_lang,
                        detected_lang_name=detected_lang_name,
                        confidence=f"{detection_confidence:.2%}",
                        text_sample=text_to_detect[:100]
                    )

                    # If confidence is low (< 80%), fall back to English
                    if detection_confidence < 0.80:
                        logger.warning(
                            "Low confidence language detection, defaulting to English",
                            detected_lang=detected_lang,
                            confidence=f"{detection_confidence:.2%}"
                        )
                        detected_lang = "en"
                        detected_lang_name = "English"
            except Exception as e:
                logger.warning(f"Language detection failed, using English: {e}")

            # Only trigger Hinglish mode if:
            # 1. Hindi was detected with HIGH confidence (>= 0.7)
            # 2. AND the text clearly contains Hinglish markers
            title_lower = job.title.lower()
            hinglish_markers = ["ke", "ka", "ki", "hai", "ko", "se", "mein", "par", "aur", "liye", "kaise"]
            hinglish_word_count = sum(1 for marker in hinglish_markers if marker in title_lower.split())
            has_hinglish = hinglish_word_count >= 2  # Require at least 2 markers

            if has_hinglish and detected_lang == "hi" and detection_confidence >= 0.70:
                detected_lang_name = "Hinglish (Hindi mixed with English)"
                language_instruction = f"""
---CRITICAL LANGUAGE REQUIREMENT---
DETECTED LANGUAGE: Hinglish (Hindi/English mix)
The document title "{job.title}" is in Hinglish.

YOU MUST:
1. Write ALL content in Hinglish (Hindi words written in Roman script mixed with English)
2. Example: "Marketing strategy ko implement karna" not "Implementing marketing strategy"
3. DO NOT write in pure English!
4. DO NOT write in Devanagari script - use Roman letters only
5. If source documents are in German/English/other, TRANSLATE to Hinglish
6. Every slide/section MUST be in Hinglish - no exceptions!
---END LANGUAGE REQUIREMENT---
"""
            else:
                language_instruction = f"""
---LANGUAGE REQUIREMENT---
OUTPUT LANGUAGE: {detected_lang_name}

Generate ALL content in {detected_lang_name}.
If source documents are in a different language, translate them to {detected_lang_name}.
Keep technical terms and proper nouns in their original form if commonly used that way.
---END LANGUAGE REQUIREMENT---
"""
        elif output_language != "en":
            language_name = LANGUAGE_NAMES.get(output_language, "English")
            language_instruction = f"""
---LANGUAGE REQUIREMENT---
Generate ALL content in {language_name}.
If source material is in a different language, translate it to {language_name}.
Do NOT mix languages - all text must be in {language_name}.
Technical terms may remain in English if commonly used that way.
---END LANGUAGE REQUIREMENT---
"""

        # Generate content
        prompt = f"""Write content for the following section:

Document Title: {job.title}
Section Title: {section_title}
Description: {section_description}

{position_context}

{context}

{format_instructions}
{style_context}
{language_instruction}"""

        try:
            from backend.services.llm import EnhancedLLMFactory

            llm, config = await EnhancedLLMFactory.get_chat_model_for_operation(
                operation="content_generation",
                user_id=None,  # System-level operation
            )
            response = await llm.ainvoke(prompt)
            content = response.content

            # Filter out LLM conversational artifacts (meta-text)
            content = filter_llm_metatext(content)

            # Filter out bullets that just echo the section title
            content = filter_title_echo(content, section_title)

        except Exception as e:
            logger.error("Failed to generate section", error=str(e))
            content = f"[Content for {section_title} - generation failed]"

        # Quality scoring and optional auto-regeneration
        quality_report = None
        regeneration_attempts = 0

        # Check if quality review is enabled (from config or job metadata)
        enable_quality = job.metadata.get("enable_quality_review", self.config.enable_quality_review)

        if enable_quality and content and not content.startswith("[Content for"):
            quality_scorer = ContentQualityScorer(min_score=self.config.min_quality_score)

            # Get other sections' content for consistency check
            other_sections_content = []
            for section in getattr(job, 'sections', []):
                if section.order != order and section.content:
                    other_sections_content.append(section.content)

            # Build context for quality scoring
            quality_context = {
                "title": job.title,
                "description": job.description,
                "output_format": job.output_format.value if hasattr(job.output_format, 'value') else str(job.output_format),
            }

            # Convert sources to dict format for quality scorer
            sources_dicts = [{"snippet": s.snippet} for s in sources] if sources else None

            # Score the content
            quality_report = await quality_scorer.score_section(
                content=content,
                title=section_title,
                sources=sources_dicts,
                context=quality_context,
                other_sections=other_sections_content if other_sections_content else None,
            )

            logger.info(
                "Section quality scored",
                section_title=section_title,
                score=quality_report.overall_score,
                needs_revision=quality_report.needs_revision,
            )

            # Auto-regenerate if quality is too low
            if (quality_report.needs_revision and
                self.config.auto_regenerate_low_quality and
                regeneration_attempts < self.config.max_regeneration_attempts):

                logger.info(
                    "Auto-regenerating low quality section",
                    section_title=section_title,
                    score=quality_report.overall_score,
                    issues=quality_report.critical_issues[:3],
                )

                # Create feedback from quality report
                feedback_items = quality_report.critical_issues + quality_report.improvements[:3]
                quality_feedback = "Quality issues to address:\n" + "\n".join(f"- {item}" for item in feedback_items)

                # Build regeneration prompt (include language instruction!)
                regen_prompt = f"""Revise this content to improve quality:
{language_instruction}

ORIGINAL CONTENT:
{content}

QUALITY FEEDBACK:
{quality_feedback}

REQUIREMENTS:
- Address all the quality issues listed above
- Keep the same topic and key information
- Maintain the required format ({format_instructions[:200]}...)
- CRITICAL: Keep the content in the SAME LANGUAGE as the original!

Write the improved content:"""

                try:
                    regeneration_attempts += 1
                    response = await llm.ainvoke(regen_prompt)
                    improved_content = response.content

                    # Re-score the improved content
                    new_report = await quality_scorer.score_section(
                        content=improved_content,
                        title=section_title,
                        sources=sources_dicts,
                        context=quality_context,
                        other_sections=other_sections_content if other_sections_content else None,
                    )

                    if new_report.overall_score > quality_report.overall_score:
                        content = improved_content
                        quality_report = new_report
                        logger.info(
                            "Section quality improved after regeneration",
                            section_title=section_title,
                            new_score=new_report.overall_score,
                        )

                except Exception as e:
                    logger.warning("Failed to regenerate section", error=str(e))

        # Initialize section metadata
        section_metadata = {}

        # CriticAgent proofreading - runs after basic quality scoring
        enable_critic = job.metadata.get("enable_critic_review", False)
        if enable_critic and content and not content.startswith("[Content for"):
            content, critic_metadata = await self._review_with_critic(
                content=content,
                section_title=section_title,
                job=job,
                format_instructions=format_instructions,
                output_language=output_language,  # Pass language to prevent English reversion
            )
            # Merge critic metadata
            if critic_metadata:
                section_metadata.update(critic_metadata)

        # Store quality report in section metadata
        if quality_report:
            section_metadata["quality_score"] = quality_report.overall_score
            section_metadata["quality_summary"] = quality_report.summary
            section_metadata["needs_revision"] = quality_report.needs_revision

        return Section(
            id=section_id,
            title=section_title,
            content=content,
            order=order,
            sources=sources,
            approved=not self.config.require_section_approval,
            metadata=section_metadata if section_metadata else None,
        )

    async def _regenerate_section_with_feedback(
        self,
        job: GenerationJob,
        section: Section,
    ) -> Section:
        """Regenerate a section based on feedback with enhanced context."""

        # Build cross-section context for consistency
        other_sections_context = self._build_cross_section_context(job, section)

        # Build source material context for accuracy
        source_context = self._build_source_context(section)

        # Determine format-specific guidelines
        format_guidelines = self._get_format_guidelines(job.output_format)

        # Build enhanced revision prompt
        prompt = f"""You are revising a section of a {job.output_format.upper()} document.

# DOCUMENT CONTEXT
Title: {job.title}
Description: {job.description}
Target Audience: {job.outline.target_audience if job.outline else 'General professional audience'}
Tone: {job.outline.tone if job.outline else 'Professional'}
Output Format: {job.output_format.upper()}

# CURRENT SECTION TO REVISE
Section Title: {section.title}
Section Order: {section.order + 1} of {len(job.sections)}

Current Content:
---
{section.content}
---

# USER FEEDBACK
{section.feedback}

{other_sections_context}

{source_context}

# FORMAT-SPECIFIC GUIDELINES
{format_guidelines}

# QUALITY REQUIREMENTS
Please revise the content ensuring:
1. **Address the feedback directly** - Make specific changes requested
2. **Maintain consistency** - Use similar terminology and tone as other sections
3. **Preserve accuracy** - Keep source-backed claims accurate
4. **Improve clarity** - Ensure content is clear and well-structured
5. **Keep appropriate length** - Match the original section length unless feedback requests changes
6. **Smooth transitions** - Ensure logical flow from previous section

# OUTPUT
Provide ONLY the revised section content. Do not include the section title or any meta-commentary.
Write the content ready for direct use in the document."""

        try:
            from backend.services.llm import EnhancedLLMFactory

            llm, config = await EnhancedLLMFactory.get_chat_model_for_operation(
                operation="content_generation",
                user_id=None,  # System-level operation
            )
            response = await llm.ainvoke(prompt)

            return Section(
                id=section.id,
                title=section.title,
                content=response.content,
                order=section.order,
                sources=section.sources,
                approved=False,
            )

        except Exception as e:
            logger.error("Failed to revise section", error=str(e))
            return section

    def _build_cross_section_context(
        self,
        job: GenerationJob,
        current_section: Section,
    ) -> str:
        """Build context from other sections for consistency."""
        context_parts = []

        # Get previous section for flow
        prev_section = None
        next_section = None
        for i, s in enumerate(job.sections):
            if s.id == current_section.id:
                if i > 0:
                    prev_section = job.sections[i - 1]
                if i < len(job.sections) - 1:
                    next_section = job.sections[i + 1]
                break

        if prev_section or next_section:
            context_parts.append("# ADJACENT SECTIONS (for context and consistency)")

            if prev_section:
                prev_content = prev_section.revised_content or prev_section.content
                # Truncate to last 300 chars for context
                preview = prev_content[-300:] if len(prev_content) > 300 else prev_content
                context_parts.append(f"Previous Section ({prev_section.title}) ends with:")
                context_parts.append(f"...{preview}")

            if next_section:
                next_content = next_section.revised_content or next_section.content
                # Truncate to first 300 chars for context
                preview = next_content[:300] if len(next_content) > 300 else next_content
                context_parts.append(f"\nNext Section ({next_section.title}) begins with:")
                context_parts.append(f"{preview}...")

        # Extract key terminology from other sections for consistency
        all_content = []
        for s in job.sections:
            if s.id != current_section.id:
                content = s.revised_content or s.content
                all_content.append(content)

        if all_content:
            combined = " ".join(all_content)
            # Simple term extraction - just note length
            context_parts.append(f"\n# Document has {len(job.sections)} total sections with ~{len(combined.split())} words overall.")

        return "\n".join(context_parts) if context_parts else ""

    def _build_source_context(self, section: Section) -> str:
        """Build source material context for accuracy checking."""
        if not section.sources:
            return ""

        context_parts = ["# SOURCE MATERIAL (maintain accuracy with these)"]

        for i, source in enumerate(section.sources[:5], 1):  # Top 5 sources
            snippet = source.snippet[:200] if len(source.snippet) > 200 else source.snippet
            context_parts.append(f"[{i}] {source.document_name}: {snippet}")

        return "\n".join(context_parts)

    async def _review_with_critic(
        self,
        content: str,
        section_title: str,
        job: GenerationJob,
        format_instructions: str,
        output_language: str = "en",
    ) -> tuple[str, dict]:
        """
        Use CriticAgent to review and auto-fix content issues.

        Args:
            content: The generated content to review
            section_title: Title of the section being reviewed
            job: The generation job (for settings and context)
            format_instructions: Format-specific instructions for regeneration
            output_language: Language code for output (en, de, es, auto, etc.)

        Returns:
            Tuple of (possibly improved content, metadata dict with review info)
        """
        import uuid as uuid_module
        from backend.services.agents.worker_agents import CriticAgent
        from backend.services.agents.agent_base import AgentConfig, AgentTask

        metadata = {
            "critic_reviewed": True,
            "critic_score": None,
            "critic_feedback": None,
            "was_revised_by_critic": False,
        }

        # Get settings from job metadata
        quality_threshold = job.metadata.get("quality_threshold", 0.7)
        fix_styling = job.metadata.get("fix_styling", True)
        fix_incomplete = job.metadata.get("fix_incomplete", True)

        try:
            # Create CriticAgent
            critic_config = AgentConfig(
                agent_id=str(uuid_module.uuid4()),
                name="Document Critic",
                description="Reviews generated content for quality, styling, and completeness",
            )
            critic = CriticAgent(critic_config)

            # Build evaluation criteria based on settings
            criteria = ["accuracy", "clarity", "relevance"]
            if fix_styling:
                criteria.append("formatting")
            if fix_incomplete:
                criteria.append("completeness")

            # Use the evaluate convenience method
            evaluation = await critic.evaluate(
                content=content,
                original_request=f"Section '{section_title}' for document '{job.title}': {job.description}",
                criteria=criteria,
            )

            # Normalize score to 0-1 range (critic returns 1-5)
            normalized_score = evaluation.overall_score / 5.0
            metadata["critic_score"] = normalized_score
            metadata["critic_feedback"] = evaluation.feedback

            logger.info(
                "CriticAgent reviewed section",
                section_title=section_title,
                score=normalized_score,
                passed=evaluation.passed,
                threshold=quality_threshold,
            )

            # Auto-fix if below threshold
            if normalized_score < quality_threshold and evaluation.improvements_needed:
                logger.info(
                    "CriticAgent auto-fixing low quality section",
                    section_title=section_title,
                    score=normalized_score,
                    improvements=evaluation.improvements_needed[:3],
                )

                # Build fix prompt with critic feedback
                feedback_text = "\n".join(f"- {item}" for item in evaluation.improvements_needed[:5])

                # Build language instruction for critic fix - use same strong instruction as main generation
                critic_language_instruction = ""
                if output_language == "auto":
                    # Check for Hinglish markers in title
                    title_lower = job.title.lower()
                    hinglish_markers = ["ke", "ka", "ki", "hai", "ko", "se", "mein", "par", "aur", "liye", "kaise"]
                    has_hinglish = any(marker in title_lower.split() for marker in hinglish_markers)

                    if has_hinglish:
                        critic_language_instruction = f"""
---CRITICAL LANGUAGE REQUIREMENT---
The document title "{job.title}" is in Hinglish.
YOU MUST keep the improved content in Hinglish (Hindi+English mix, Roman script).
DO NOT translate to pure English or German!
Every sentence must be in Hinglish - no exceptions!
---END LANGUAGE REQUIREMENT---
"""
                    else:
                        critic_language_instruction = f"""
---CRITICAL LANGUAGE REQUIREMENT---
The document title is: "{job.title}"
Keep the improved content in the EXACT SAME LANGUAGE as the original content.
DO NOT translate to English or any other language!
Maintain language consistency throughout.
---END LANGUAGE REQUIREMENT---
"""
                elif output_language != "en":
                    language_name = LANGUAGE_NAMES.get(output_language, "English")
                    critic_language_instruction = f"""
---CRITICAL LANGUAGE REQUIREMENT---
Keep ALL content in {language_name}.
DO NOT translate to English - maintain {language_name} throughout.
Every sentence must be in {language_name} - no exceptions!
---END LANGUAGE REQUIREMENT---
"""

                fix_prompt = f"""Improve this content based on the following feedback:
{critic_language_instruction}

ORIGINAL CONTENT:
{content}

QUALITY ISSUES TO ADDRESS:
{feedback_text}

{f"FORMATTING ISSUES: Fix any styling or formatting problems" if fix_styling else ""}
{f"COMPLETENESS ISSUES: Complete any incomplete sentences or bullet points" if fix_incomplete else ""}

REQUIREMENTS:
- Address all the quality issues listed above
- Maintain the same topic and key information
- Follow format requirements: {format_instructions[:300]}
- CRITICAL: Keep the content in the SAME LANGUAGE as the original!

Write the improved content:"""

                try:
                    from backend.services.llm import EnhancedLLMFactory

                    llm, _ = await EnhancedLLMFactory.get_chat_model_for_operation(
                        operation="content_generation",
                        user_id=None,
                    )
                    response = await llm.ainvoke(fix_prompt)
                    improved_content = response.content

                    # Re-evaluate to confirm improvement
                    new_evaluation = await critic.evaluate(
                        content=improved_content,
                        original_request=f"Section '{section_title}' for document '{job.title}'",
                        criteria=criteria,
                    )
                    new_score = new_evaluation.overall_score / 5.0

                    if new_score > normalized_score:
                        content = improved_content
                        metadata["critic_score"] = new_score
                        metadata["was_revised_by_critic"] = True
                        logger.info(
                            "CriticAgent improved section quality",
                            section_title=section_title,
                            old_score=normalized_score,
                            new_score=new_score,
                        )

                except Exception as e:
                    logger.warning("CriticAgent failed to fix content", error=str(e))

        except Exception as e:
            logger.warning("CriticAgent review failed", error=str(e))
            metadata["critic_reviewed"] = False

        return content, metadata

    def _get_format_guidelines(self, output_format: OutputFormat) -> str:
        """Get format-specific writing guidelines."""
        guidelines = {
            OutputFormat.PPTX: """For PowerPoint slides:
- Keep bullet points concise (max 8-10 words each)
- Use 3-6 bullet points per slide typically
- Avoid long paragraphs
- Use action-oriented language
- Include speaker notes context if appropriate""",

            OutputFormat.DOCX: """For Word documents:
- Use clear paragraph structure
- Include smooth transitions between ideas
- Maintain consistent heading hierarchy
- Balance detail with readability""",

            OutputFormat.PDF: """For PDF reports:
- Structure content clearly with logical flow
- Use professional language
- Include appropriate detail for formal documents
- Consider visual hierarchy in text structure""",

            OutputFormat.MARKDOWN: """For Markdown:
- Use proper markdown formatting (headers, lists, code blocks)
- Keep structure clean and readable
- Use bullet points and numbered lists appropriately""",

            OutputFormat.HTML: """For HTML:
- Structure content semantically
- Use appropriate heading levels
- Keep paragraphs well-organized""",

            OutputFormat.XLSX: """For Excel:
- Structure data clearly in rows/columns
- Include headers and labels
- Keep text concise for cell content""",

            OutputFormat.TXT: """For plain text:
- Use clear structure and spacing
- Avoid relying on formatting
- Keep content well-organized without markup""",
        }

        return guidelines.get(output_format, "Write clear, professional content.")

    async def _generate_output_file(self, job: GenerationJob) -> str:
        """Generate the output file in the requested format."""
        safe_title = sanitize_filename(job.title)
        filename = f"{job.id}_{safe_title}"

        if job.output_format == OutputFormat.PPTX:
            output_path = await self._generate_pptx(job, filename)
        elif job.output_format == OutputFormat.DOCX:
            output_path = await self._generate_docx(job, filename)
        elif job.output_format == OutputFormat.PDF:
            output_path = await self._generate_pdf(job, filename)
        elif job.output_format == OutputFormat.XLSX:
            output_path = await self._generate_xlsx(job, filename)
        elif job.output_format == OutputFormat.MARKDOWN:
            output_path = await self._generate_markdown(job, filename)
        elif job.output_format == OutputFormat.HTML:
            output_path = await self._generate_html(job, filename)
        else:
            output_path = await self._generate_txt(job, filename)

        job.output_path = output_path
        return output_path

    async def _generate_pptx(self, job: GenerationJob, filename: str) -> str:
        """Generate professional PowerPoint file with modern styling."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE

            # Check if using a template PPTX
            template_path = job.metadata.get("template_pptx_path")
            use_template_styling = False  # When True, skip hardcoded styling to preserve template

            # Standard dimensions for 16:9 (used as reference for scaling)
            STANDARD_WIDTH = 13.333
            STANDARD_HEIGHT = 7.5

            if template_path and os.path.exists(template_path):
                # Use the template - this preserves slide masters, themes, colors
                prs = Presentation(template_path)
                # Remove all existing slides from template (keep only masters)
                while len(prs.slides) > 0:
                    slide_id = prs.slides._sldIdLst[0].rId
                    prs.part.drop_rel(slide_id)
                    del prs.slides._sldIdLst[0]
                use_template_styling = True

                # Calculate scaling factors for content positioning
                # Template may have different dimensions than our standard 16:9
                template_width = prs.slide_width.inches
                template_height = prs.slide_height.inches
                width_scale = template_width / STANDARD_WIDTH
                height_scale = template_height / STANDARD_HEIGHT

                logger.info(
                    "Using PPTX template - preserving template styling",
                    template_path=template_path,
                    masters_count=len(prs.slide_masters),
                    layouts_count=len(prs.slide_layouts),
                    template_width=f"{template_width:.2f}in",
                    template_height=f"{template_height:.2f}in",
                    width_scale=f"{width_scale:.2f}",
                    height_scale=f"{height_scale:.2f}",
                )
            else:
                prs = Presentation()
                prs.slide_width = Inches(STANDARD_WIDTH)  # 16:9 aspect ratio
                prs.slide_height = Inches(STANDARD_HEIGHT)
                width_scale = 1.0
                height_scale = 1.0

            # Helper function to scale dimensions when using templates
            def scale_inches(value: float, axis: str = "w") -> float:
                """Scale an Inches value based on template dimensions.

                Args:
                    value: The original value in inches (designed for 13.333x7.5)
                    axis: "w" for width scaling, "h" for height scaling

                Returns:
                    Scaled value in inches
                """
                if axis == "w":
                    return value * width_scale
                else:
                    return value * height_scale

            def scaled_inches(value: float, axis: str = "w"):
                """Return Inches object with scaled value."""
                return Inches(scale_inches(value, axis))

            # Get theme colors from job metadata or use default (with custom color overrides)
            # Only apply custom theme if NOT using template styling
            theme_key = job.metadata.get("theme", "business")
            custom_colors = job.metadata.get("custom_colors")
            theme = get_theme_colors(theme_key, custom_colors)

            # Get font family from job metadata or use default
            # Only apply custom fonts if NOT using template styling
            font_family_key = job.metadata.get("font_family", "modern")
            font_config = FONT_FAMILIES.get(font_family_key, FONT_FAMILIES["modern"])
            heading_font = font_config["heading"]
            body_font = font_config["body"]

            # Get layout from job metadata (for future use)
            layout_key = job.metadata.get("layout", "standard")
            layout_config = LAYOUT_TEMPLATES.get(layout_key, LAYOUT_TEMPLATES["standard"])

            # Check if animations are enabled and get speed
            enable_animations = job.metadata.get("animations", False)
            animation_speed = job.metadata.get("animation_speed", "med")  # very_slow, slow, med, fast, very_fast, custom
            animation_duration_ms = job.metadata.get("animation_duration_ms")  # Custom duration in ms

            # Duration mapping for preset speeds (in milliseconds)
            TRANSITION_DURATIONS = {
                "very_slow": 2000,  # 2 seconds
                "slow": 1500,       # 1.5 seconds
                "med": 750,         # 0.75 seconds
                "fast": 400,        # 0.4 seconds
                "very_fast": 200,   # 0.2 seconds
            }

            # Calculate the actual transition duration to use
            def get_transition_duration():
                """Get transition duration based on settings."""
                if animation_speed == "custom" and animation_duration_ms:
                    return animation_duration_ms
                return TRANSITION_DURATIONS.get(animation_speed, 750)

            transition_duration = get_transition_duration()

            def get_slide_layout(layout_type: str = "content"):
                """
                Get appropriate slide layout based on template or default.

                When using template, tries to use template's layouts:
                - "title": Layout 0 (Title Slide)
                - "content": Layout 1 (Title and Content)
                - "blank": Layout 6 (Blank) or last layout

                When not using template, always uses blank layout.

                Args:
                    layout_type: "title", "content", or "blank"

                Returns:
                    SlideLayout object
                """
                if use_template_styling:
                    try:
                        if layout_type == "title" and len(prs.slide_layouts) > 0:
                            return prs.slide_layouts[0]
                        elif layout_type == "content" and len(prs.slide_layouts) > 1:
                            return prs.slide_layouts[1]
                        elif len(prs.slide_layouts) > 6:
                            return prs.slide_layouts[6]
                        else:
                            return prs.slide_layouts[-1]  # Use last available layout
                    except Exception as e:
                        logger.warning(f"Failed to get template layout, using default: {e}")
                        return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
                else:
                    return prs.slide_layouts[6]  # Blank layout for non-template

            def add_slide_notes(slide, notes_text: str):
                """
                Add speaker notes to a slide.

                Args:
                    slide: The slide to add notes to
                    notes_text: The text to add as notes
                """
                try:
                    notes_slide = slide.notes_slide
                    notes_frame = notes_slide.notes_text_frame
                    notes_frame.text = notes_text
                except Exception as e:
                    logger.warning(f"Failed to add slide notes: {e}")

            def add_slide_transition(slide, transition_type="fade", duration=500, speed="med"):
                """
                Add slide transition using XML manipulation.

                Transition types:
                - fade: Smooth fade transition
                - push: Push from right
                - wipe: Wipe from left
                - dissolve: Dissolve effect
                - blinds: Vertical blinds effect
                """
                try:
                    from lxml import etree

                    # PowerPoint XML namespaces
                    nsmap = {
                        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                    }

                    # Get the slide's XML element
                    slide_elem = slide._element

                    # Remove existing transition if any
                    for trans in slide_elem.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition'):
                        slide_elem.remove(trans)

                    # Create transition element
                    # Duration in milliseconds (e.g., 500 = 0.5 seconds)
                    trans_elem = etree.Element(
                        '{http://schemas.openxmlformats.org/presentationml/2006/main}transition',
                        nsmap={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
                    )

                    # Map extended speed names to PowerPoint's native values for compatibility
                    ppt_speed = "slow" if speed in ("very_slow", "slow") else ("fast" if speed in ("fast", "very_fast") else "med")
                    trans_elem.set('spd', ppt_speed)  # Speed: slow, med, fast (native PowerPoint values)

                    # Use p14:dur for custom transition duration (PowerPoint 2010+)
                    # This provides more precise control than the native spd attribute
                    trans_elem.set('{http://schemas.microsoft.com/office/powerpoint/2010/main}dur', str(duration))

                    # Also set advTm for auto-advance (optional)
                    trans_elem.set('advTm', str(duration * 3))  # Auto-advance after 3x transition duration

                    # Add transition type element
                    if transition_type == "fade":
                        effect = etree.SubElement(trans_elem, '{http://schemas.openxmlformats.org/presentationml/2006/main}fade')
                        effect.set('thruBlk', 'true')  # Fade through black
                    elif transition_type == "push":
                        effect = etree.SubElement(trans_elem, '{http://schemas.openxmlformats.org/presentationml/2006/main}push')
                        effect.set('dir', 'r')  # From right
                    elif transition_type == "wipe":
                        effect = etree.SubElement(trans_elem, '{http://schemas.openxmlformats.org/presentationml/2006/main}wipe')
                        effect.set('dir', 'r')  # From right
                    elif transition_type == "dissolve":
                        etree.SubElement(trans_elem, '{http://schemas.openxmlformats.org/presentationml/2006/main}dissolve')
                    elif transition_type == "blinds":
                        effect = etree.SubElement(trans_elem, '{http://schemas.openxmlformats.org/presentationml/2006/main}blinds')
                        effect.set('dir', 'vert')  # Vertical blinds
                    else:
                        # Default to fade
                        effect = etree.SubElement(trans_elem, '{http://schemas.openxmlformats.org/presentationml/2006/main}fade')

                    # Insert transition as first child of slide element (after cSld)
                    cSld = slide_elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
                    if cSld is not None:
                        cSld_index = list(slide_elem).index(cSld)
                        slide_elem.insert(cSld_index + 1, trans_elem)
                    else:
                        slide_elem.append(trans_elem)

                except Exception as e:
                    logger.warning(f"Failed to add slide transition: {e}")

            logger.info(
                "PPTX generation settings",
                theme=theme_key,
                font_family=font_family_key,
                heading_font=heading_font,
                body_font=body_font,
                layout=layout_key,
                animations=enable_animations,
            )

            # Apply theme color scheme
            primary_rgb = hex_to_rgb(theme["primary"])
            secondary_rgb = hex_to_rgb(theme["secondary"])
            accent_rgb = hex_to_rgb(theme["accent"])
            text_rgb = hex_to_rgb(theme["text"])
            light_gray_rgb = hex_to_rgb(theme["light_gray"])

            PRIMARY_COLOR = RGBColor(*primary_rgb)
            SECONDARY_COLOR = RGBColor(*secondary_rgb)
            ACCENT_COLOR = RGBColor(*accent_rgb)
            TEXT_COLOR = RGBColor(*text_rgb)
            LIGHT_GRAY = RGBColor(*light_gray_rgb)
            WHITE = RGBColor(0xFF, 0xFF, 0xFF)

            def apply_title_style(shape, font_size=44, bold=True, color=WHITE):
                """Apply consistent title styling using configured heading font."""
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = heading_font
                    paragraph.font.size = Pt(font_size)
                    paragraph.font.bold = bold
                    paragraph.font.color.rgb = color
                    paragraph.alignment = PP_ALIGN.LEFT

            def apply_body_style(shape, font_size=18, color=TEXT_COLOR):
                """Apply consistent body text styling using configured body font."""
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = body_font
                    paragraph.font.size = Pt(font_size)
                    paragraph.font.color.rgb = color
                    paragraph.line_spacing = 1.5

            def add_header_bar(slide, text=""):
                """Add a colored header bar to slide."""
                header = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(0),
                    prs.slide_width, Inches(1.2)
                )
                header.fill.solid()
                header.fill.fore_color.rgb = PRIMARY_COLOR
                header.line.fill.background()

            def sanitize_text(text: str) -> str:
                """Sanitize text for PPTX XML compatibility."""
                if not text:
                    return ""
                # Remove control characters that can corrupt XML
                import re
                # Remove ASCII control chars except tab, newline, carriage return
                text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
                # Replace common problematic characters
                text = text.replace('\x0b', ' ')  # Vertical tab
                text = text.replace('\x0c', ' ')  # Form feed
                return text

            # Note: strip_markdown and filter_llm_metatext are now module-level functions

            def add_footer(slide, page_num, total_pages):
                """Add footer with page number."""
                footer_text = f"{page_num} / {total_pages}"
                footer = slide.shapes.add_textbox(
                    Inches(12.3), Inches(7.1),
                    Inches(0.8), Inches(0.3)
                )
                tf = footer.text_frame
                tf.paragraphs[0].text = footer_text
                tf.paragraphs[0].font.size = Pt(10)
                tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
                tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

            # ========== THEME-AWARE STYLING FUNCTIONS ==========
            # Get theme visual properties
            slide_background_style = theme.get("slide_background", "solid")
            header_style = theme.get("header_style", "none")
            bullet_style = theme.get("bullet_style", "circle")
            accent_position = theme.get("accent_position", "top")

            # Bullet character mapping based on theme
            BULLET_CHARS = {
                "circle": ["•", "◦", "▪"],
                "circle-filled": ["●", "○", "▪"],
                "arrow": ["▸", "▹", "▫"],
                "chevron": ["»", "›", "·"],
                "dash": ["—", "-", "·"],
                "square": ["■", "□", "▪"],
                "number": None,  # Use numbering instead
                "leaf": ["❧", "✿", "·"],
            }

            def get_bullet_chars():
                """Get bullet characters based on theme."""
                return BULLET_CHARS.get(bullet_style, BULLET_CHARS["circle"])

            def apply_slide_background(slide, is_title_slide=False):
                """Apply theme-specific background styling to slide."""
                bg_style = slide_background_style

                if bg_style == "solid" or bg_style == "white":
                    # Solid color background (or white for minimalist)
                    bg_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        prs.slide_width, prs.slide_height
                    )
                    if bg_style == "white":
                        bg_shape.fill.solid()
                        bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    else:
                        bg_shape.fill.solid()
                        if is_title_slide:
                            bg_shape.fill.fore_color.rgb = PRIMARY_COLOR
                        else:
                            bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    bg_shape.line.fill.background()
                    return bg_shape

                elif bg_style == "gradient" or bg_style == "warm-gradient":
                    # Create gradient effect with two shapes
                    # Top portion with primary color
                    top_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        prs.slide_width, Inches(3.75) if not is_title_slide else prs.slide_height
                    )
                    top_shape.fill.solid()
                    top_shape.fill.fore_color.rgb = PRIMARY_COLOR if is_title_slide else SECONDARY_COLOR
                    top_shape.line.fill.background()

                    if not is_title_slide:
                        # Bottom portion lighter
                        bottom_shape = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            Inches(0), Inches(3.75),
                            prs.slide_width, Inches(3.75)
                        )
                        bottom_shape.fill.solid()
                        bottom_shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
                        bottom_shape.line.fill.background()
                    return top_shape

                elif bg_style == "gradient-multi":
                    # Multi-color gradient for colorful theme
                    top_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        prs.slide_width, Inches(1.5) if not is_title_slide else Inches(3)
                    )
                    top_shape.fill.solid()
                    top_shape.fill.fore_color.rgb = PRIMARY_COLOR
                    top_shape.line.fill.background()

                    if not is_title_slide:
                        # Main content area white
                        main_shape = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            Inches(0), Inches(1.5),
                            prs.slide_width, Inches(6)
                        )
                        main_shape.fill.solid()
                        main_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        main_shape.line.fill.background()
                    return top_shape

                elif bg_style == "dark":
                    # Dark mode background
                    bg_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        prs.slide_width, prs.slide_height
                    )
                    bg_shape.fill.solid()
                    # Use a very dark color
                    bg_shape.fill.fore_color.rgb = RGBColor(26, 26, 46)
                    bg_shape.line.fill.background()
                    return bg_shape

                elif bg_style == "textured":
                    # Textured look with subtle pattern
                    bg_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        prs.slide_width, prs.slide_height
                    )
                    bg_shape.fill.solid()
                    bg_shape.fill.fore_color.rgb = ACCENT_COLOR if not is_title_slide else PRIMARY_COLOR
                    bg_shape.line.fill.background()
                    return bg_shape

                else:
                    # Default solid
                    bg_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        prs.slide_width, prs.slide_height
                    )
                    bg_shape.fill.solid()
                    bg_shape.fill.fore_color.rgb = PRIMARY_COLOR if is_title_slide else RGBColor(255, 255, 255)
                    bg_shape.line.fill.background()
                    return bg_shape

            def add_header_accent(slide, title_text=""):
                """Add theme-specific header accent based on header_style."""
                h_style = header_style

                if h_style == "bar" or h_style == "colorblock":
                    # Full-width colored bar at top
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        prs.slide_width, Inches(1.2)
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = PRIMARY_COLOR
                    bar.line.fill.background()
                    return bar

                elif h_style == "underline":
                    # Header bar + underline accent
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        prs.slide_width, Inches(1.2)
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = PRIMARY_COLOR
                    bar.line.fill.background()

                    # Add underline accent
                    underline = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0.8), Inches(1.15),
                        Inches(3), Inches(0.05)
                    )
                    underline.fill.solid()
                    underline.fill.fore_color.rgb = SECONDARY_COLOR
                    underline.line.fill.background()
                    return bar

                elif h_style == "glow":
                    # Header with glow effect (gradient bar)
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        prs.slide_width, Inches(1.3)
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = PRIMARY_COLOR
                    bar.line.fill.background()

                    # Add subtle accent line
                    glow_line = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(1.25),
                        prs.slide_width, Inches(0.08)
                    )
                    glow_line.fill.solid()
                    glow_line.fill.fore_color.rgb = ACCENT_COLOR
                    glow_line.line.fill.background()
                    return bar

                elif h_style == "rounded":
                    # Rounded header bar
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.3), Inches(0.2),
                        Inches(12.7), Inches(1.0)
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = PRIMARY_COLOR
                    bar.line.fill.background()
                    return bar

                elif h_style == "serif" or h_style == "leaf":
                    # Elegant header with accent
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        prs.slide_width, Inches(1.2)
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = PRIMARY_COLOR
                    bar.line.fill.background()
                    return bar

                elif h_style == "none":
                    # No header bar - just return None
                    return None

                else:
                    # Default header bar
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        prs.slide_width, Inches(1.2)
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = PRIMARY_COLOR
                    bar.line.fill.background()
                    return bar

            def add_accent_elements(slide, position="top"):
                """Add decorative accent elements based on theme."""
                pos = accent_position

                if pos == "top":
                    # Accent line at top
                    accent = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        prs.slide_width, Inches(0.08)
                    )
                    accent.fill.solid()
                    accent.fill.fore_color.rgb = SECONDARY_COLOR
                    accent.line.fill.background()

                elif pos == "bottom":
                    # Accent line at bottom
                    accent = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(7.42),
                        prs.slide_width, Inches(0.08)
                    )
                    accent.fill.solid()
                    accent.fill.fore_color.rgb = SECONDARY_COLOR
                    accent.line.fill.background()

                elif pos == "side":
                    # Accent bar on left side
                    accent = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        Inches(0.15), prs.slide_height
                    )
                    accent.fill.solid()
                    accent.fill.fore_color.rgb = SECONDARY_COLOR
                    accent.line.fill.background()

                elif pos == "corner" or pos == "corners":
                    # Corner accents
                    # Top-left corner
                    corner1 = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        Inches(0.5), Inches(0.08)
                    )
                    corner1.fill.solid()
                    corner1.fill.fore_color.rgb = SECONDARY_COLOR
                    corner1.line.fill.background()

                    corner2 = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        Inches(0.08), Inches(0.5)
                    )
                    corner2.fill.solid()
                    corner2.fill.fore_color.rgb = SECONDARY_COLOR
                    corner2.line.fill.background()

                elif pos == "diagonal":
                    # Diagonal accent in corner
                    accent = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_TRIANGLE,
                        Inches(12), Inches(6),
                        Inches(1.333), Inches(1.5)
                    )
                    accent.fill.solid()
                    accent.fill.fore_color.rgb = SECONDARY_COLOR
                    accent.line.fill.background()

                elif pos == "border":
                    # Full border accent
                    # Top
                    slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        prs.slide_width, Inches(0.05)
                    ).fill.solid()
                    # Bottom
                    slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(7.45),
                        prs.slide_width, Inches(0.05)
                    ).fill.solid()

                elif pos == "footer":
                    # Footer accent bar
                    accent = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(7.2),
                        prs.slide_width, Inches(0.3)
                    )
                    accent.fill.solid()
                    accent.fill.fore_color.rgb = PRIMARY_COLOR
                    accent.line.fill.background()

                # "none" position - no accent

            def calculate_luminance(color: RGBColor) -> float:
                """Calculate relative luminance of a color (0-1 scale)."""
                # Convert to linear values
                r = color[0] / 255.0
                g = color[1] / 255.0
                b = color[2] / 255.0

                # Apply gamma correction
                r = r / 12.92 if r <= 0.03928 else ((r + 0.055) / 1.055) ** 2.4
                g = g / 12.92 if g <= 0.03928 else ((g + 0.055) / 1.055) ** 2.4
                b = b / 12.92 if b <= 0.03928 else ((b + 0.055) / 1.055) ** 2.4

                return 0.2126 * r + 0.7152 * g + 0.0722 * b

            def get_contrast_ratio(color1: RGBColor, color2: RGBColor) -> float:
                """Calculate WCAG contrast ratio between two colors."""
                lum1 = calculate_luminance(color1)
                lum2 = calculate_luminance(color2)
                lighter = max(lum1, lum2)
                darker = min(lum1, lum2)
                return (lighter + 0.05) / (darker + 0.05)

            def get_best_text_color(bg_color: RGBColor) -> RGBColor:
                """Get text color with best contrast against background."""
                white_contrast = get_contrast_ratio(bg_color, WHITE)
                black_contrast = get_contrast_ratio(bg_color, RGBColor(0, 0, 0))

                # WCAG AA requires 4.5:1 for normal text, 3:1 for large text
                # Title text is large, so we use 3:1 minimum
                if white_contrast >= black_contrast:
                    return WHITE
                else:
                    return RGBColor(0, 0, 0)

            def get_text_color_for_background(is_header=False, is_dark_bg=False, bg_color=None):
                """Get appropriate text color based on background."""
                if slide_background_style == "dark":
                    return RGBColor(228, 228, 228)  # Light text for dark mode

                # If background color provided, calculate best contrast
                if bg_color is not None:
                    return get_best_text_color(bg_color)

                if is_header or is_dark_bg:
                    # For title slide, check actual primary color luminance
                    return get_best_text_color(PRIMARY_COLOR)
                else:
                    return TEXT_COLOR

            # Determine include_sources from job metadata or fall back to config
            include_sources_override = job.metadata.get("include_sources")
            if include_sources_override is not None:
                include_sources = include_sources_override
            else:
                include_sources = self.config.include_sources

            logger.info(f"PPTX include_sources check - override: {include_sources_override}, config: {self.config.include_sources}, resolved: {include_sources}")
            logger.info(f"PPTX sources_used count: {len(job.sources_used) if job.sources_used else 0}")

            # Calculate total slides accurately
            total_slides = 1  # Title slide
            if self.config.include_toc:
                total_slides += 1  # TOC slide
            total_slides += len(job.sections)  # Content slides
            if include_sources and job.sources_used:
                total_slides += 1  # Sources slide only if both conditions met

            current_slide = 0

            # Initialize image generator if images are enabled
            include_images = job.metadata.get("include_images", self.config.include_images)
            section_images = {}
            if include_images:
                try:
                    image_config = ImageGeneratorConfig(
                        enabled=True,
                        backend=ImageBackend(self.config.image_backend),
                        default_width=400,
                        default_height=300,
                    )
                    image_service = get_image_generator(image_config)

                    # Generate images for all sections in parallel
                    sections_data = [
                        (section.title, section.revised_content or section.content)
                        for section in job.sections
                    ]
                    images = await image_service.generate_batch(
                        sections=sections_data,
                        document_title=job.title,
                    )

                    # Map images to sections by index
                    for idx, img in enumerate(images):
                        if img and img.success and img.path:
                            section_images[idx] = img.path

                    logger.info(
                        "Generated images for PPTX",
                        total_sections=len(job.sections),
                        images_generated=len(section_images),
                    )
                except Exception as e:
                    logger.warning(f"Image generation failed, continuing without images: {e}")

            # ========== TITLE SLIDE ==========
            current_slide += 1
            slide = prs.slides.add_slide(get_slide_layout("title"))

            # Add transition if animations are enabled
            if enable_animations:
                add_slide_transition(slide, "fade", duration=transition_duration, speed=animation_speed)

            # Apply theme-specific background and styling ONLY if not using template
            if not use_template_styling:
                apply_slide_background(slide, is_title_slide=True)

                # Accent bar at bottom (theme-aware)
                accent_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(6.5),
                    prs.slide_width, Inches(1)
                )
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = SECONDARY_COLOR
                accent_bar.line.fill.background()

            # Title text color adapts to background
            title_text_color = get_text_color_for_background(is_dark_bg=True)

            # Title text
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2.5),
                Inches(11), Inches(1.5)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sanitize_text(job.title) or "Untitled"
            p.font.name = heading_font
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = title_text_color

            # Subtitle/description
            desc_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(4.2),
                Inches(11), Inches(1)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            desc_text = job.outline.description if job.outline else job.description
            p.text = sanitize_text(desc_text) or ""
            p.font.name = body_font
            p.font.size = Pt(20)
            p.font.color.rgb = ACCENT_COLOR

            # Date
            from datetime import datetime
            generation_date = datetime.now()
            date_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(6.7),
                Inches(4), Inches(0.4)
            )
            tf = date_box.text_frame
            p = tf.paragraphs[0]
            p.text = generation_date.strftime("%B %d, %Y")
            p.font.name = body_font
            p.font.size = Pt(14)
            p.font.color.rgb = title_text_color

            # Add notes to title slide (always - contains generation info)
            output_language = job.metadata.get("output_language", "en")
            # Get user info - prefer email from metadata, fallback to job.user_id
            user_info = job.metadata.get("user_email") or job.user_id
            # Get LLM model - prefer from metadata (set during generation), fallback to config
            llm_model = job.metadata.get("llm_model") or self.config.model
            title_notes = f"""Generated by AIDocumentIndexer

Model: {llm_model}
User: {user_info}
Date: {generation_date.strftime('%Y-%m-%d %H:%M')}
Language: {output_language}
Theme: {theme_key}
Font: {font_family_key}
Total sections: {len(job.sections)}
"""
            add_slide_notes(slide, title_notes)

            # ========== TABLE OF CONTENTS SLIDE ==========
            if self.config.include_toc:
                current_slide += 1
                slide = prs.slides.add_slide(get_slide_layout("content"))

                # Add transition if animations are enabled
                if enable_animations:
                    add_slide_transition(slide, "push", duration=transition_duration, speed=animation_speed)

                # Apply theme-specific header accent ONLY if not using template
                if not use_template_styling:
                    add_header_accent(slide, "Contents")

                # TOC Title - position depends on header style
                title_y = Inches(0.3) if header_style != "none" else Inches(0.5)
                title_color = WHITE if header_style != "none" else TEXT_COLOR
                toc_title = slide.shapes.add_textbox(
                    Inches(0.8), title_y,
                    Inches(8), Inches(0.8)
                )
                tf = toc_title.text_frame
                p = tf.paragraphs[0]
                p.text = "Contents"
                p.font.name = heading_font
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = title_color

                # Add accent elements based on theme (skip if using template)
                if header_style == "none" and not use_template_styling:
                    add_accent_elements(slide)

                # TOC items - get theme-specific bullet characters
                toc_bullets = get_bullet_chars()
                toc_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(1.8),
                    Inches(11), Inches(5)
                )
                tf = toc_box.text_frame
                tf.word_wrap = True

                # Determine text color based on background
                content_text_color = get_text_color_for_background()

                first_toc_used = False
                for idx, section in enumerate(job.sections):
                    if first_toc_used:
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                        first_toc_used = True
                    section_title = sanitize_text(section.title) or f"Section {idx + 1}"

                    # Check if title already has Roman numeral prefix (I., II., III., etc.)
                    # If so, don't add Arabic number to avoid double numbering
                    roman_pattern = r'^[IVXLCDM]+\.\s+'
                    if re.match(roman_pattern, section_title):
                        p.text = f"{idx + 1}.  {re.sub(roman_pattern, '', section_title)}"
                    else:
                        p.text = f"{idx + 1}.  {section_title}"
                    p.font.name = body_font
                    p.font.size = Pt(20)
                    p.font.color.rgb = content_text_color
                    p.space_after = Pt(12)

                # Ensure first paragraph is initialized even if no sections
                if not first_toc_used:
                    p = tf.paragraphs[0]
                    p.text = "No sections"
                    p.font.name = body_font
                    p.font.size = Pt(20)

                add_footer(slide, current_slide, total_slides)

            # ========== CONTENT SLIDES ==========
            # Transition types to cycle through for variety
            content_transitions = ["wipe", "fade", "push", "dissolve"]

            # Get theme-specific bullet characters for content
            content_bullets = get_bullet_chars()

            for section_idx, section in enumerate(job.sections):
                current_slide += 1
                slide = prs.slides.add_slide(get_slide_layout("content"))

                # Add transition if animations are enabled (cycle through types)
                if enable_animations:
                    trans_type = content_transitions[section_idx % len(content_transitions)]
                    add_slide_transition(slide, trans_type, duration=transition_duration, speed=animation_speed)

                # Apply theme-specific header accent (skip if using template)
                if not use_template_styling:
                    add_header_accent(slide, section.title)

                # Check if we have an image for this section
                has_image = section_idx in section_images

                # Section title - position depends on header style
                title_y = Inches(0.3) if header_style != "none" else Inches(0.5)
                title_color = WHITE if header_style != "none" else TEXT_COLOR
                section_title = slide.shapes.add_textbox(
                    Inches(0.8), title_y,
                    Inches(11), Inches(0.8)
                )
                tf = section_title.text_frame
                p = tf.paragraphs[0]
                p.text = sanitize_text(section.title) or f"Section {section_idx + 1}"
                p.font.name = heading_font
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = title_color

                # Add accent elements for themes without header bar (skip if using template)
                if header_style == "none" and not use_template_styling:
                    add_accent_elements(slide)

                # Adjust content area based on whether we have an image and layout
                content = sanitize_text(section.revised_content or section.content) or ""
                # Strip markdown formatting for clean slide content
                content = strip_markdown(content)

                # Limit content for presentation readability
                # ~450 words = ~3000 chars - use LLM to condense if needed
                if len(content) > 3000:
                    content = await llm_condense_text(content, 3000)

                # Store the processed content for review/preview consistency
                section.rendered_content = content

                # Calculate content dimensions based on layout template
                slide_content_width = prs.slide_width.inches - 1.6  # Total usable width
                content_width_ratio = layout_config.get("content_width", 0.85)

                # Initialize image positioning variables (used later if has_image)
                image_left = Inches(8.5)  # Default
                image_width = Inches(4.3)  # Default

                if has_image:
                    # Layout-aware image positioning
                    image_pos = layout_config.get("image_position", "right")

                    # Default content height - will be reduced for some layouts
                    # Scale all dimensions based on template size
                    content_height = scaled_inches(5.2, "h")
                    content_top = scaled_inches(1.6, "h")

                    if layout_key == "two_column":
                        # Split screen: content left, image right (side by side)
                        content_width = Inches(slide_content_width * 0.48)
                        content_left = scaled_inches(0.8, "w")
                        image_left = scaled_inches(7.2, "w")
                        image_width = scaled_inches(5.5, "w")
                        # Side-by-side layout can use full height
                        content_height = scaled_inches(5.0, "h")
                    elif layout_key == "image_focused":
                        # Large image BELOW text - reduce content height to prevent overlap
                        content_width = Inches(slide_content_width)
                        content_left = scaled_inches(0.8, "w")
                        # Content area limited to top portion (image starts at Y=4.2")
                        content_height = scaled_inches(2.4, "h")  # Reduced from 5.2" to prevent overlap
                        content_top = scaled_inches(1.6, "h")
                        # Image positioned below content
                        image_left = scaled_inches(2.5, "w")
                        image_width = scaled_inches(8.0, "w")
                    elif layout_key == "minimal":
                        # Clean layout with image on the side
                        content_width = Inches(slide_content_width * 0.6)
                        content_left = scaled_inches(0.8, "w")
                        image_left = scaled_inches(8.5, "w")
                        image_width = scaled_inches(4.0, "w")
                        # Side-by-side can use more height
                        content_height = scaled_inches(5.0, "h")
                    else:  # standard layout
                        content_width = scaled_inches(7.5, "w")
                        content_left = scaled_inches(0.8, "w")
                        image_left = scaled_inches(8.5, "w")
                        image_width = scaled_inches(4.3, "w")
                        # Side-by-side can use more height
                        content_height = scaled_inches(5.0, "h")

                    content_box = slide.shapes.add_textbox(
                        content_left, content_top,
                        content_width, content_height
                    )
                else:
                    # Full width content (no image)
                    if layout_key == "minimal":
                        # Centered narrow content for minimal layout
                        content_width = Inches(slide_content_width * 0.7)
                        content_left = Inches((prs.slide_width.inches - content_width.inches) / 2)
                    elif layout_key == "two_column":
                        # Use full width for two-column without image
                        content_width = Inches(slide_content_width)
                        content_left = scaled_inches(0.8, "w")
                    else:
                        # Standard full width
                        content_width = scaled_inches(11.5, "w")
                        content_left = scaled_inches(0.8, "w")

                    content_box = slide.shapes.add_textbox(
                        content_left, scaled_inches(1.6, "h"),
                        content_width, scaled_inches(5.2, "h")
                    )

                tf = content_box.text_frame
                tf.word_wrap = True

                # Split content into bullet points for cleaner slides
                paragraphs = content.split('\n')
                first_para_used = False
                # Adjust max paragraphs based on layout and image presence
                if has_image and layout_key == "image_focused":
                    max_paras = 5  # Very limited space above image
                elif has_image:
                    max_paras = 10  # Side-by-side layouts have more room
                else:
                    max_paras = 12  # Full content area
                para_count = 0

                # Parse bullet hierarchy - collect (text, level) tuples
                def parse_bullet_hierarchy(lines: list) -> list:
                    """Parse content lines into (text, level) tuples preserving hierarchy."""
                    import re
                    result = []

                    # Patterns to skip - LLM meta-text lines (safety net)
                    skip_patterns = [
                        r'^[Hh]ere (are|is) ',
                        r'^[Ll]et me ',
                        r'^[Cc]ertainly',
                        r'^[Ss]ure[,!]',
                        r'^[Ii]\'ll ',
                        r'^[Bb]elow (are|is)',
                        r'[Ll]et me know',
                        r'[Hh]ope this helps',
                        r'[Ff]eel free to',
                    ]

                    for line in lines:
                        if not line:
                            continue
                        # Count leading whitespace for indentation level
                        stripped = line.lstrip()
                        if not stripped:
                            continue

                        # Skip meta-text lines
                        if any(re.match(pattern, stripped, re.IGNORECASE) for pattern in skip_patterns):
                            continue

                        indent = len(line) - len(stripped)

                        # Detect markdown nested lists (  - item,    - item)
                        # Each 2 spaces = 1 level of nesting
                        level = min(indent // 2, 3)  # Max 3 levels deep (0-3)

                        # Strip bullet markers - include ◦ (open circle) for sub-points
                        if stripped.startswith(('- ', '• ', '* ', '◦ ', '○ ', '▪ ', '▸ ')):
                            text = stripped[2:].strip()
                            # If it's a sub-point marker, ensure level >= 1
                            if stripped.startswith(('◦ ', '○ ')) and level == 0:
                                level = 1
                        elif stripped.startswith(tuple(f'{i}.' for i in range(1, 10))):
                            # Numbered list: "1. text" -> "text"
                            text = stripped.split('.', 1)[1].strip() if '.' in stripped else stripped
                        else:
                            text = stripped

                        if text:
                            result.append((text, level))
                    return result

                # Collect valid paragraphs with hierarchy info
                valid_paragraphs = parse_bullet_hierarchy(paragraphs)

                # Dynamic font sizing based on content volume
                # valid_paragraphs is now list of (text, level) tuples
                # Increased max_chars limits to reduce truncation
                total_bullets = min(len(valid_paragraphs), max_paras)
                total_chars = sum(len(text) for text, _ in valid_paragraphs[:max_paras])

                if total_bullets <= 5 and total_chars < 500:
                    base_font_size = Pt(20)
                    max_chars = 180  # Was 120
                    line_spacing = Pt(10)
                elif total_bullets <= 8 and total_chars < 900:
                    base_font_size = Pt(18)
                    max_chars = 150  # Was 100
                    line_spacing = Pt(8)
                elif total_bullets <= 10 and total_chars < 1200:
                    base_font_size = Pt(16)
                    max_chars = 130  # Was 90
                    line_spacing = Pt(6)
                else:
                    base_font_size = Pt(14)
                    max_chars = 110  # Was 80
                    line_spacing = Pt(4)

                # Adjust max chars if we have an image (narrower content area)
                if has_image:
                    max_chars = int(max_chars * 0.75)  # Was 0.7

                for bullet_text, bullet_level in valid_paragraphs:
                    if para_count >= max_paras:
                        break

                    if first_para_used:
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                        first_para_used = True
                    para_count += 1

                    # Use LLM to condense text intelligently instead of truncating with '...'
                    if len(bullet_text) > max_chars:
                        bullet_text = await llm_condense_text(bullet_text, max_chars)

                    # Set paragraph level for hierarchy (0=main, 1=sub, 2=sub-sub, etc)
                    p.level = bullet_level

                    # Use theme-specific bullet characters
                    # content_bullets is from get_bullet_chars() which returns theme-appropriate chars
                    theme_bullets = content_bullets if content_bullets else ['•', '◦', '▪']
                    # Add fallback for deep nesting
                    bullet_char = theme_bullets[min(bullet_level, len(theme_bullets) - 1)]
                    p.text = f"{bullet_char} {bullet_text}"
                    p.font.name = body_font
                    # Decrease font size slightly for nested levels
                    level_font_size = Pt(max(base_font_size.pt - (bullet_level * 2), 10))
                    p.font.size = level_font_size
                    # Use theme-aware text color
                    content_text_color = get_text_color_for_background()
                    p.font.color.rgb = content_text_color
                    p.space_after = line_spacing

                # Ensure first paragraph is initialized even if content was empty
                if not first_para_used:
                    p = tf.paragraphs[0]
                    p.text = ""  # Empty but properly initialized
                    p.font.name = body_font
                    p.font.size = Pt(16)

                # Add image if available (use layout-based positioning)
                if has_image:
                    try:
                        image_path = section_images[section_idx]

                        # Position image based on layout template (scaled for template dimensions)
                        if layout_key == "image_focused":
                            # Large centered image BELOW content (content ends at ~4.0")
                            slide.shapes.add_picture(
                                image_path,
                                image_left, scaled_inches(4.2, "h"),  # Start below content area
                                width=image_width,
                                height=scaled_inches(3.0, "h"),  # Slightly smaller to fit
                            )
                        elif layout_key == "two_column":
                            # Image on right side, vertically centered
                            slide.shapes.add_picture(
                                image_path,
                                image_left, scaled_inches(2.0, "h"),
                                width=image_width,
                                height=scaled_inches(4.0, "h"),
                            )
                        else:
                            # Standard/minimal: image on right (side-by-side, no vertical overlap)
                            slide.shapes.add_picture(
                                image_path,
                                image_left, scaled_inches(1.8, "h"),
                                width=image_width,
                                height=scaled_inches(4.5, "h"),  # Can be taller since side-by-side
                            )

                        logger.debug(
                            "Added image to PPTX slide",
                            section=section.title,
                            image_path=image_path,
                            layout=layout_key,
                        )
                    except Exception as e:
                        logger.warning(f"Failed to add image to slide: {e}")

                add_footer(slide, current_slide, total_slides)

                # Add notes with optional summary and per-section sources
                # Check if user enabled notes explanation (which includes summary/chain of thought)
                include_notes_explanation = job.metadata.get("include_notes_explanation", False)

                # Build notes - always include section title and sources
                notes_parts = []
                section_title = section.title or f"Section {section_idx + 1}"
                notes_parts.append(f"SECTION: {section_title}")
                notes_parts.append("")

                logger.debug(
                    "Building notes for section",
                    section_idx=section_idx,
                    section_title=section_title,
                    include_notes_explanation=include_notes_explanation,
                )

                # Only add summary if include_notes_explanation is enabled
                if include_notes_explanation:
                    section_content_text = section.revised_content or section.content or ""

                    # Extract first 2-3 sentences or first 300 chars for summary
                    def extract_summary(text: str, max_chars: int = 300) -> str:
                        """Extract a brief summary from content."""
                        if not text:
                            return ""
                        # Remove bullet points and clean up
                        clean_text = re.sub(r'^[•●○◆-]\s*', '', text, flags=re.MULTILINE)
                        clean_text = clean_text.replace('\n', ' ').strip()
                        # Get first few sentences
                        sentences = re.split(r'(?<=[.!?])\s+', clean_text)
                        summary = ""
                        for sent in sentences[:3]:
                            if len(summary) + len(sent) < max_chars:
                                summary += sent + " "
                            else:
                                break
                        return summary.strip() or clean_text[:max_chars] + "..."

                    content_summary = extract_summary(section_content_text)
                    if content_summary:
                        notes_parts.append("SUMMARY:")
                        notes_parts.append(content_summary)
                        notes_parts.append("")

                # Add sources - fallback to job-level sources if section has none
                section_sources = section.sources or []
                used_fallback = False
                if not section_sources and job.sources_used:
                    # Use top sources from job as fallback
                    section_sources = job.sources_used[:5]
                    used_fallback = True
                    logger.info(
                        "Using job-level sources as fallback for section notes",
                        section_title=section_title[:30],
                        fallback_sources=len(section_sources),
                    )
                logger.debug(
                    "PPTX notes - section sources check",
                    section_title=section_title[:30],
                    sources_count=len(section_sources),
                    used_fallback=used_fallback,
                    source_names=[s.document_name for s in section_sources[:3]] if section_sources else [],
                )
                if section_sources:
                    source_lines = []
                    for s in section_sources[:5]:
                        source_name = s.document_name or s.document_id
                        if s.page_number:
                            # Detect if source was a PPTX (slide) or other doc (page)
                            if source_name.lower().endswith('.pptx'):
                                source_lines.append(f"• {source_name} (Slide {s.page_number})")
                            else:
                                source_lines.append(f"• {source_name} (Page {s.page_number})")
                        else:
                            source_lines.append(f"• {source_name}")
                    notes_parts.append("SOURCES:")
                    notes_parts.extend(source_lines)
                else:
                    notes_parts.append("SOURCES: No specific sources for this section.")

                section_notes = "\n".join(notes_parts)

                # Add LLM explanation if enabled and available (reuse flag from above)
                if include_notes_explanation:
                    llm_reasoning = section.metadata.get("llm_reasoning") if section.metadata else None
                    if llm_reasoning:
                        section_notes += f"\n\nAI Explanation:\n{llm_reasoning}"

                # Add metadata footer (user, model, date)
                section_notes += "\n\n---"
                if job.metadata.get("user_email"):
                    section_notes += f"\nUser: {job.metadata.get('user_email')}"
                section_notes += f"\nModel: {job.metadata.get('llm_model', 'N/A')}"
                section_notes += f"\nGenerated: {job.created_at.strftime('%Y-%m-%d %H:%M') if job.created_at else 'N/A'}"
                section_notes += f"\nSlide: {current_slide} / {total_slides}"

                add_slide_notes(slide, section_notes)
                logger.debug(
                    "Added notes to slide",
                    slide_number=current_slide,
                    notes_length=len(section_notes),
                    notes_preview=section_notes[:100] if section_notes else "EMPTY",
                )

            # ========== SOURCES SLIDE ==========
            if not include_sources:
                logger.warning(f"Sources slide SKIPPED - include_sources is False")
            elif not job.sources_used:
                logger.warning(f"Sources slide SKIPPED - no sources_used (RAG may have returned no results)")
            if include_sources and job.sources_used:
                current_slide += 1
                slide = prs.slides.add_slide(get_slide_layout("content"))

                # Add transition if animations are enabled
                if enable_animations:
                    add_slide_transition(slide, "blinds", duration=transition_duration, speed=animation_speed)

                # Apply theme-specific header accent (skip if using template)
                if not use_template_styling:
                    add_header_accent(slide, "Sources & References")

                # Sources title - position depends on header style
                title_y = Inches(0.3) if header_style != "none" else Inches(0.5)
                title_color = WHITE if header_style != "none" else TEXT_COLOR
                sources_title = slide.shapes.add_textbox(
                    Inches(0.8), title_y,
                    Inches(8), Inches(0.8)
                )
                tf = sources_title.text_frame
                p = tf.paragraphs[0]
                p.text = "Sources & References"
                p.font.name = heading_font
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = title_color

                # Add accent elements for themes without header bar (skip if using template)
                if header_style == "none" and not use_template_styling:
                    add_accent_elements(slide)

                # Sources list
                sources_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(1.6),
                    Inches(11), Inches(5)
                )
                tf = sources_box.text_frame
                tf.word_wrap = True

                # Get theme-aware text color
                sources_text_color = get_text_color_for_background()
                # Get theme-specific bullet character
                source_bullet = content_bullets[0] if content_bullets else "•"

                first_source_used = False
                for source in job.sources_used[:10]:
                    if first_source_used:
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                        first_source_used = True
                    doc_name = sanitize_text(source.document_name or source.document_id[:20])
                    # Add page/slide number if available
                    location_info = ""
                    if source.page_number:
                        if doc_name.lower().endswith('.pptx'):
                            location_info = f" (Slide {source.page_number})"
                        else:
                            location_info = f" (Page {source.page_number})"
                    p.text = f"{source_bullet}  {doc_name}{location_info}"
                    p.font.name = body_font
                    p.font.size = Pt(14)
                    p.font.color.rgb = sources_text_color
                    p.space_after = Pt(6)

                # Ensure first paragraph is initialized even if no sources
                if not first_source_used:
                    p = tf.paragraphs[0]
                    p.text = "No sources available"
                    p.font.name = body_font
                    p.font.size = Pt(14)

                add_footer(slide, current_slide, total_slides)

                # Add notes to sources slide with metadata
                sources_notes_lines = [
                    "Sources & References",
                    "",
                    f"Generated: {job.created_at.strftime('%Y-%m-%d %H:%M') if job.created_at else 'N/A'}",
                    f"Model: {job.metadata.get('llm_model', 'N/A')}",
                ]
                if job.metadata.get("user_email"):
                    sources_notes_lines.insert(2, f"User: {job.metadata.get('user_email')}")
                sources_notes_lines.append("")
                sources_notes_lines.append(f"Total sources referenced: {len(job.sources_used)}")
                if job.sources_used:
                    sources_notes_lines.append("")
                    sources_notes_lines.append("Documents used:")
                    for source in job.sources_used[:10]:
                        doc_name = source.document_name or source.document_id[:20]
                        sources_notes_lines.append(f"  - {doc_name}")
                    if len(job.sources_used) > 10:
                        sources_notes_lines.append(f"  ... and {len(job.sources_used) - 10} more")
                add_slide_notes(slide, "\n".join(sources_notes_lines))

            output_path = os.path.join(self.config.output_dir, f"{filename}.pptx")
            prs.save(output_path)

            logger.info("Professional PPTX generated", path=output_path)
            return output_path

        except ImportError as e:
            logger.error(f"python-pptx import error: {e}")
            return await self._generate_txt(job, filename)

    async def _generate_docx(self, job: GenerationJob, filename: str) -> str:
        """Generate professional Word document with cover page and proper styling."""
        try:
            from docx import Document
            from docx.shared import Pt, Inches, RGBColor, Cm
            from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
            from docx.enum.style import WD_STYLE_TYPE
            from docx.enum.table import WD_TABLE_ALIGNMENT

            doc = Document()

            # Get theme colors from job metadata or use default (with custom color overrides)
            theme_key = job.metadata.get("theme", "business")
            custom_colors = job.metadata.get("custom_colors")
            theme = get_theme_colors(theme_key, custom_colors)

            # Get font family from job metadata or use default
            font_family_key = job.metadata.get("font_family", "modern")
            font_config = FONT_FAMILIES.get(font_family_key, FONT_FAMILIES["modern"])
            heading_font = font_config["heading"]
            body_font = font_config["body"]

            logger.info(
                "DOCX generation settings",
                theme=theme_key,
                font_family=font_family_key,
                heading_font=heading_font,
                body_font=body_font,
            )

            # Apply theme color scheme
            primary_rgb = hex_to_rgb(theme["primary"])
            secondary_rgb = hex_to_rgb(theme["secondary"])
            text_rgb = hex_to_rgb(theme["text"])
            light_gray_rgb = hex_to_rgb(theme["light_gray"])

            PRIMARY_COLOR = RGBColor(*primary_rgb)
            SECONDARY_COLOR = RGBColor(*secondary_rgb)
            TEXT_COLOR = RGBColor(*text_rgb)
            LIGHT_GRAY = RGBColor(*light_gray_rgb)

            # Determine include_sources from job metadata or fall back to config
            include_sources = job.metadata.get("include_sources", self.config.include_sources)

            # Initialize image generator if images are enabled
            include_images = job.metadata.get("include_images", self.config.include_images)
            section_images = {}
            if include_images:
                try:
                    image_config = ImageGeneratorConfig(
                        enabled=True,
                        backend=ImageBackend(self.config.image_backend),
                        default_width=600,
                        default_height=400,
                    )
                    image_service = get_image_generator(image_config)

                    # Generate images for all sections in parallel
                    sections_data = [
                        (section.title, section.revised_content or section.content)
                        for section in job.sections
                    ]
                    images = await image_service.generate_batch(
                        sections=sections_data,
                        document_title=job.title,
                    )

                    # Map images to sections by index
                    for idx, img in enumerate(images):
                        if img and img.success and img.path:
                            section_images[idx] = img.path

                    logger.info(
                        "Generated images for DOCX",
                        total_sections=len(job.sections),
                        images_generated=len(section_images),
                    )
                except Exception as e:
                    logger.warning(f"Image generation failed, continuing without images: {e}")

            # Configure document margins
            for section in doc.sections:
                section.top_margin = Inches(1)
                section.bottom_margin = Inches(1)
                section.left_margin = Inches(1.25)
                section.right_margin = Inches(1.25)

            # ========== COVER PAGE ==========
            # Add spacing before title
            for _ in range(8):
                doc.add_paragraph()

            # Document title
            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run = title_para.add_run(job.title)
            title_run.font.name = heading_font
            title_run.font.size = Pt(36)
            title_run.font.bold = True
            title_run.font.color.rgb = PRIMARY_COLOR

            # Subtitle / description
            doc.add_paragraph()
            if job.outline:
                desc_para = doc.add_paragraph()
                desc_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                desc_run = desc_para.add_run(job.outline.description)
                desc_run.font.name = body_font
                desc_run.font.size = Pt(14)
                desc_run.font.color.rgb = SECONDARY_COLOR
                desc_run.font.italic = True

            # Add more spacing
            for _ in range(6):
                doc.add_paragraph()

            # Decorative line
            line_para = doc.add_paragraph()
            line_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            line_run = line_para.add_run("─" * 40)
            line_run.font.color.rgb = SECONDARY_COLOR

            # Date
            from datetime import datetime
            doc.add_paragraph()
            date_para = doc.add_paragraph()
            date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            date_run = date_para.add_run(datetime.now().strftime("%B %d, %Y"))
            date_run.font.name = body_font
            date_run.font.size = Pt(12)
            date_run.font.color.rgb = LIGHT_GRAY

            # Page break after cover
            doc.add_page_break()

            # ========== TABLE OF CONTENTS ==========
            if self.config.include_toc:
                toc_heading = doc.add_heading("Table of Contents", level=1)
                toc_heading.runs[0].font.color.rgb = PRIMARY_COLOR
                toc_heading.runs[0].font.size = Pt(24)

                doc.add_paragraph()

                for idx, section in enumerate(job.sections):
                    toc_entry = doc.add_paragraph()
                    toc_entry.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
                    toc_entry.paragraph_format.space_after = Pt(6)

                    # Section number
                    num_run = toc_entry.add_run(f"{idx + 1}.  ")
                    num_run.font.name = body_font
                    num_run.font.size = Pt(12)
                    num_run.font.bold = True
                    num_run.font.color.rgb = SECONDARY_COLOR

                    # Section title - strip Roman numeral prefix to avoid double numbering
                    section_title = section.title
                    roman_pattern = r'^[IVXLCDM]+\.\s+'
                    if re.match(roman_pattern, section_title):
                        section_title = re.sub(roman_pattern, '', section_title)

                    title_run = toc_entry.add_run(section_title)
                    title_run.font.name = body_font
                    title_run.font.size = Pt(12)
                    title_run.font.color.rgb = TEXT_COLOR

                doc.add_page_break()

            # ========== CONTENT SECTIONS ==========
            for idx, section in enumerate(job.sections):
                # Section heading
                heading = doc.add_heading(section.title, level=1)
                heading.runs[0].font.color.rgb = PRIMARY_COLOR
                heading.runs[0].font.size = Pt(20)
                heading.paragraph_format.space_before = Pt(12)
                heading.paragraph_format.space_after = Pt(12)

                content = section.revised_content or section.content

                # Split content into paragraphs and format
                paragraphs = content.split('\n\n')
                for para_text in paragraphs:
                    para_text = para_text.strip()
                    if not para_text:
                        continue

                    # Check if it's a subheading (starts with ## or ###)
                    if para_text.startswith('###'):
                        sub = doc.add_heading(para_text.lstrip('#').strip(), level=3)
                        sub.runs[0].font.color.rgb = SECONDARY_COLOR
                        sub.runs[0].font.size = Pt(13)
                    elif para_text.startswith('##'):
                        sub = doc.add_heading(para_text.lstrip('#').strip(), level=2)
                        sub.runs[0].font.color.rgb = SECONDARY_COLOR
                        sub.runs[0].font.size = Pt(15)
                    # Check if it's a bullet list
                    elif para_text.startswith(('- ', '• ', '* ')) or '  -' in para_text:
                        lines = para_text.split('\n')
                        for line in lines:
                            # Detect indentation level
                            stripped_line = line.lstrip()
                            indent_spaces = len(line) - len(stripped_line)
                            indent_level = min(indent_spaces // 2, 2)  # Max 2 levels (List Bullet 3)

                            if stripped_line.startswith(('- ', '• ', '* ')):
                                # Use List Bullet, List Bullet 2, or List Bullet 3 based on indent
                                bullet_style = 'List Bullet' if indent_level == 0 else f'List Bullet {indent_level + 1}'
                                try:
                                    bullet_para = doc.add_paragraph(style=bullet_style)
                                except KeyError:
                                    # Fallback to regular bullet if style doesn't exist
                                    bullet_para = doc.add_paragraph(style='List Bullet')
                                    # Manually indent for deeper levels
                                    bullet_para.paragraph_format.left_indent = Inches(0.25 * indent_level)

                                bullet_run = bullet_para.add_run(stripped_line.lstrip('-•* ').strip())
                                bullet_run.font.name = body_font
                                # Slightly smaller font for nested bullets
                                bullet_run.font.size = Pt(11 - indent_level)
                                bullet_run.font.color.rgb = TEXT_COLOR
                    # Regular paragraph with inline formatting support
                    else:
                        para = doc.add_paragraph()
                        para.paragraph_format.line_spacing = 1.5
                        para.paragraph_format.space_after = Pt(8)

                        # Parse inline bold (**text**) and italic (*text*) formatting
                        import re
                        # Pattern to match **bold**, *italic*, or plain text
                        pattern = r'(\*\*[^*]+\*\*|\*[^*]+\*|[^*]+)'
                        parts = re.findall(pattern, para_text)

                        for part in parts:
                            if part.startswith('**') and part.endswith('**'):
                                # Bold text
                                run = para.add_run(part[2:-2])
                                run.bold = True
                            elif part.startswith('*') and part.endswith('*') and len(part) > 2:
                                # Italic text
                                run = para.add_run(part[1:-1])
                                run.italic = True
                            else:
                                # Plain text
                                run = para.add_run(part)

                            run.font.name = body_font
                            run.font.size = Pt(11)
                            run.font.color.rgb = TEXT_COLOR

                # Add image for this section if available
                if idx in section_images:
                    try:
                        image_path = section_images[idx]
                        doc.add_paragraph()  # Add spacing before image
                        img_para = doc.add_paragraph()
                        img_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run = img_para.add_run()
                        run.add_picture(image_path, width=Inches(5.0))
                        doc.add_paragraph()  # Add spacing after image
                        logger.debug(
                            "Added image to DOCX section",
                            section=section.title,
                            image_path=image_path,
                        )
                    except Exception as e:
                        logger.warning(f"Failed to add image to DOCX: {e}")

                # Add section sources
                if include_sources and section.sources:
                    doc.add_paragraph()
                    source_label = doc.add_paragraph()
                    label_run = source_label.add_run("Sources for this section:")
                    label_run.font.name = body_font
                    label_run.font.size = Pt(9)
                    label_run.font.italic = True
                    label_run.font.color.rgb = LIGHT_GRAY

                    for source in section.sources[:3]:
                        src_para = doc.add_paragraph()
                        src_para.paragraph_format.left_indent = Inches(0.25)
                        # Add page/slide number if available
                        doc_name = source.document_name or source.document_id
                        location_info = ""
                        if source.page_number:
                            if doc_name.lower().endswith('.pptx'):
                                location_info = f" (Slide {source.page_number})"
                            else:
                                location_info = f" (Page {source.page_number})"
                        src_run = src_para.add_run(f"• {doc_name}{location_info}")
                        src_run.font.name = body_font
                        src_run.font.size = Pt(9)
                        src_run.font.color.rgb = LIGHT_GRAY

                # Add page break between sections (except last)
                if idx < len(job.sections) - 1:
                    doc.add_page_break()

            # ========== REFERENCES SECTION ==========
            if include_sources and job.sources_used:
                doc.add_page_break()

                ref_heading = doc.add_heading("References", level=1)
                ref_heading.runs[0].font.color.rgb = PRIMARY_COLOR
                ref_heading.runs[0].font.size = Pt(20)

                doc.add_paragraph()

                for source in job.sources_used:
                    ref_para = doc.add_paragraph()
                    ref_para.paragraph_format.space_after = Pt(4)
                    # Add page/slide number if available
                    doc_name = source.document_name or source.document_id
                    location_info = ""
                    if source.page_number:
                        if doc_name.lower().endswith('.pptx'):
                            location_info = f" (Slide {source.page_number})"
                        else:
                            location_info = f" (Page {source.page_number})"
                    ref_run = ref_para.add_run(f"• {doc_name}{location_info}")
                    ref_run.font.name = body_font
                    ref_run.font.size = Pt(10)
                    ref_run.font.color.rgb = TEXT_COLOR

            output_path = os.path.join(self.config.output_dir, f"{filename}.docx")
            doc.save(output_path)

            logger.info("Professional DOCX generated", path=output_path)
            return output_path

        except ImportError:
            logger.error("python-docx not installed")
            return await self._generate_txt(job, filename)

    async def _generate_pdf(self, job: GenerationJob, filename: str) -> str:
        """Generate professional PDF document with cover page and styling."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.colors import HexColor, Color
            from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
            from reportlab.platypus import (
                SimpleDocTemplate, Paragraph, Spacer, PageBreak,
                Table, TableStyle, ListFlowable, ListItem, Image
            )
            from reportlab.lib.units import inch, cm

            # Determine include_sources from job metadata or fall back to config
            include_sources = job.metadata.get("include_sources", self.config.include_sources)

            # Initialize image generator if images are enabled
            include_images = job.metadata.get("include_images", self.config.include_images)
            section_images = {}
            if include_images:
                try:
                    image_config = ImageGeneratorConfig(
                        enabled=True,
                        backend=ImageBackend(self.config.image_backend),
                        default_width=600,
                        default_height=400,
                    )
                    image_service = get_image_generator(image_config)

                    # Generate images for all sections in parallel
                    sections_data = [
                        (section.title, section.revised_content or section.content)
                        for section in job.sections
                    ]
                    images = await image_service.generate_batch(
                        sections=sections_data,
                        document_title=job.title,
                    )

                    # Map images to sections by index
                    for idx, img in enumerate(images):
                        if img and img.success and img.path:
                            section_images[idx] = img.path

                    logger.info(
                        "Generated images for PDF",
                        total_sections=len(job.sections),
                        images_generated=len(section_images),
                    )
                except Exception as e:
                    logger.warning(f"Image generation failed, continuing without images: {e}")

            # Get theme colors from job metadata or use default (with custom color overrides)
            theme_key = job.metadata.get("theme", "business")
            custom_colors = job.metadata.get("custom_colors")
            theme = get_theme_colors(theme_key, custom_colors)

            # Get font family from job metadata or use default
            # ReportLab built-in fonts mapping
            PDF_FONT_MAP = {
                "modern": {
                    "heading": "Helvetica-Bold",
                    "heading_oblique": "Helvetica-BoldOblique",
                    "body": "Helvetica",
                    "body_bold": "Helvetica-Bold",
                    "body_italic": "Helvetica-Oblique",
                },
                "classic": {
                    "heading": "Times-Bold",
                    "heading_oblique": "Times-BoldItalic",
                    "body": "Times-Roman",
                    "body_bold": "Times-Bold",
                    "body_italic": "Times-Italic",
                },
                "professional": {
                    "heading": "Helvetica-Bold",
                    "heading_oblique": "Helvetica-BoldOblique",
                    "body": "Helvetica",
                    "body_bold": "Helvetica-Bold",
                    "body_italic": "Helvetica-Oblique",
                },
                "technical": {
                    "heading": "Courier-Bold",
                    "heading_oblique": "Courier-BoldOblique",
                    "body": "Courier",
                    "body_bold": "Courier-Bold",
                    "body_italic": "Courier-Oblique",
                },
            }

            font_family_key = job.metadata.get("font_family", "modern")
            pdf_fonts = PDF_FONT_MAP.get(font_family_key, PDF_FONT_MAP["modern"])
            heading_font = pdf_fonts["heading"]
            heading_oblique = pdf_fonts["heading_oblique"]
            body_font = pdf_fonts["body"]
            body_bold = pdf_fonts["body_bold"]
            body_italic = pdf_fonts["body_italic"]

            logger.info(
                "PDF generation settings",
                theme=theme_key,
                font_family=font_family_key,
                heading_font=heading_font,
                body_font=body_font,
            )

            # Apply theme color scheme
            PRIMARY_COLOR = HexColor(theme["primary"])
            SECONDARY_COLOR = HexColor(theme["secondary"])
            TEXT_COLOR = HexColor(theme["text"])
            LIGHT_GRAY = HexColor(theme["light_gray"])
            ACCENT_BG = HexColor('#F5F5F5')  # Light background (keep neutral)

            output_path = os.path.join(self.config.output_dir, f"{filename}.pdf")

            # Custom page template for headers/footers
            def add_page_number(canvas, doc, font=body_font):
                """Add page number to footer."""
                canvas.saveState()
                canvas.setFont(font, 9)
                canvas.setFillColor(LIGHT_GRAY)
                page_num = canvas.getPageNumber()
                text = f"Page {page_num}"
                canvas.drawRightString(letter[0] - 0.75*inch, 0.5*inch, text)
                canvas.restoreState()

            doc = SimpleDocTemplate(
                output_path,
                pagesize=letter,
                rightMargin=0.75*inch,
                leftMargin=0.75*inch,
                topMargin=0.75*inch,
                bottomMargin=0.75*inch,
            )

            # Create custom styles
            styles = getSampleStyleSheet()

            # Cover title style
            cover_title_style = ParagraphStyle(
                'CoverTitle',
                parent=styles['Title'],
                fontSize=36,
                textColor=PRIMARY_COLOR,
                alignment=TA_CENTER,
                spaceAfter=20,
                fontName=heading_font,
            )

            # Cover subtitle style
            cover_subtitle_style = ParagraphStyle(
                'CoverSubtitle',
                parent=styles['Normal'],
                fontSize=14,
                textColor=SECONDARY_COLOR,
                alignment=TA_CENTER,
                fontName=body_italic,
                spaceAfter=12,
            )

            # Section heading style
            heading_style = ParagraphStyle(
                'SectionHeading',
                parent=styles['Heading1'],
                fontSize=20,
                textColor=PRIMARY_COLOR,
                spaceBefore=16,
                spaceAfter=12,
                fontName=heading_font,
            )

            # Subheading style
            subheading_style = ParagraphStyle(
                'SubHeading',
                parent=styles['Heading2'],
                fontSize=14,
                textColor=SECONDARY_COLOR,
                spaceBefore=12,
                spaceAfter=8,
                fontName=body_bold,
            )

            # Body text style
            body_style = ParagraphStyle(
                'BodyText',
                parent=styles['Normal'],
                fontSize=11,
                textColor=TEXT_COLOR,
                alignment=TA_JUSTIFY,
                spaceBefore=4,
                spaceAfter=8,
                leading=16,
                fontName=body_font,
            )

            # TOC style
            toc_style = ParagraphStyle(
                'TOCEntry',
                parent=styles['Normal'],
                fontSize=12,
                textColor=TEXT_COLOR,
                spaceBefore=6,
                spaceAfter=6,
                leftIndent=20,
                fontName=body_font,
            )

            # Source/reference style
            source_style = ParagraphStyle(
                'SourceStyle',
                parent=styles['Normal'],
                fontSize=9,
                textColor=LIGHT_GRAY,
                fontName=body_italic,
                spaceBefore=4,
                spaceAfter=2,
            )

            story = []

            # ========== COVER PAGE ==========
            story.append(Spacer(1, 2.5*inch))

            # Title
            story.append(Paragraph(job.title, cover_title_style))
            story.append(Spacer(1, 0.3*inch))

            # Decorative line
            line_data = [['─' * 50]]
            line_table = Table(line_data, colWidths=[5*inch])
            line_table.setStyle(TableStyle([
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('TEXTCOLOR', (0, 0), (-1, -1), SECONDARY_COLOR),
            ]))
            story.append(line_table)
            story.append(Spacer(1, 0.3*inch))

            # Description
            if job.outline:
                story.append(Paragraph(job.outline.description, cover_subtitle_style))

            story.append(Spacer(1, 2*inch))

            # Date
            from datetime import datetime
            date_style = ParagraphStyle(
                'DateStyle',
                parent=styles['Normal'],
                fontSize=12,
                textColor=LIGHT_GRAY,
                alignment=TA_CENTER,
                fontName=body_font,
            )
            story.append(Paragraph(datetime.now().strftime("%B %d, %Y"), date_style))

            story.append(PageBreak())

            # ========== TABLE OF CONTENTS ==========
            if self.config.include_toc:
                toc_title_style = ParagraphStyle(
                    'TOCTitle',
                    parent=styles['Heading1'],
                    fontSize=24,
                    textColor=PRIMARY_COLOR,
                    spaceAfter=20,
                    fontName=heading_font,
                )
                story.append(Paragraph("Table of Contents", toc_title_style))
                story.append(Spacer(1, 0.3*inch))

                for idx, section in enumerate(job.sections):
                    # Strip Roman numeral prefix to avoid double numbering
                    section_title = section.title
                    roman_pattern = r'^[IVXLCDM]+\.\s+'
                    if re.match(roman_pattern, section_title):
                        section_title = re.sub(roman_pattern, '', section_title)

                    toc_entry = f"<b><font color='#3D5A80'>{idx + 1}.</font></b>  {section_title}"
                    story.append(Paragraph(toc_entry, toc_style))

                story.append(PageBreak())

            # ========== CONTENT SECTIONS ==========
            for idx, section in enumerate(job.sections):
                # Section heading with number
                heading_text = f"{idx + 1}. {section.title}"
                story.append(Paragraph(heading_text, heading_style))

                content = section.revised_content or section.content

                # Split content and format
                paragraphs = content.split('\n')
                current_list = []
                current_level = 0  # Track bullet nesting level

                def flush_list(lst):
                    """Flush list items to story with proper nesting."""
                    if lst:
                        story.append(ListFlowable(
                            lst,
                            bulletType='bullet',
                            leftIndent=20,
                        ))
                    return []

                for para_text in paragraphs:
                    original_text = para_text  # Keep original for indentation detection
                    para_text = para_text.strip()
                    if not para_text:
                        # Flush any pending list
                        current_list = flush_list(current_list)
                        current_level = 0
                        continue

                    # Handle markdown-style headers
                    if para_text.startswith('###'):
                        current_list = flush_list(current_list)
                        current_level = 0
                        story.append(Paragraph(para_text.lstrip('#').strip(), subheading_style))
                    elif para_text.startswith('##'):
                        current_list = flush_list(current_list)
                        current_level = 0
                        story.append(Paragraph(para_text.lstrip('#').strip(), subheading_style))
                    # Handle bullet points with hierarchy
                    elif para_text.startswith(('- ', '• ', '* ')):
                        bullet_text = para_text.lstrip('-•* ').strip()
                        # Detect indentation level from original text
                        indent = len(original_text) - len(original_text.lstrip())
                        level = min(indent // 2, 3)  # Max 3 levels (0-3)

                        # Calculate left indent based on nesting level
                        # Base indent is 20, add 15 for each level
                        left_indent = 20 + (level * 15)

                        # Create a style with appropriate indent for this level
                        indented_style = ParagraphStyle(
                            f'BulletLevel{level}',
                            parent=body_style,
                            leftIndent=left_indent,
                            bulletIndent=left_indent - 10,
                        )
                        current_list.append(ListItem(Paragraph(bullet_text, indented_style), leftIndent=left_indent))
                    # Handle numbered lists
                    elif para_text[:2].replace('.', '').isdigit():
                        current_list = flush_list(current_list)
                        current_level = 0
                        import re
                        num_match = re.match(r'^(\d+\.)\s*(.+)', para_text)
                        if num_match:
                            story.append(Paragraph(f"<b>{num_match.group(1)}</b> {num_match.group(2)}", body_style))
                        else:
                            story.append(Paragraph(para_text, body_style))
                    # Regular paragraph
                    else:
                        current_list = flush_list(current_list)
                        current_level = 0
                        # Escape XML special characters first
                        formatted = para_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        # Convert markdown bold (**text**) to HTML bold for ReportLab
                        import re
                        formatted = re.sub(r'\*\*([^*]+)\*\*', r'<b>\1</b>', formatted)
                        # Convert markdown italic (*text*) to HTML italic
                        formatted = re.sub(r'\*([^*]+)\*', r'<i>\1</i>', formatted)
                        story.append(Paragraph(formatted, body_style))

                # Flush remaining list items
                current_list = flush_list(current_list)

                # Add image for this section if available
                if idx in section_images:
                    try:
                        image_path = section_images[idx]
                        story.append(Spacer(1, 0.2*inch))
                        # Create image with max width of 5 inches
                        img = Image(image_path, width=5*inch, height=3.5*inch)
                        story.append(img)
                        story.append(Spacer(1, 0.2*inch))
                        logger.debug(
                            "Added image to PDF section",
                            section=section.title,
                            image_path=image_path,
                        )
                    except Exception as e:
                        logger.warning(f"Failed to add image to PDF: {e}")

                # Section sources
                if include_sources and section.sources:
                    story.append(Spacer(1, 0.2*inch))
                    story.append(Paragraph("Sources for this section:", source_style))
                    for source in section.sources[:3]:
                        doc_name = source.document_name or source.document_id
                        # Add page/slide number if available
                        location_info = ""
                        if source.page_number:
                            if doc_name.lower().endswith('.pptx'):
                                location_info = f" (Slide {source.page_number})"
                            else:
                                location_info = f" (Page {source.page_number})"
                        story.append(Paragraph(f"• {doc_name}{location_info}", source_style))

                # Add page break after each section (except the last one)
                if idx < len(job.sections) - 1:
                    story.append(PageBreak())
                else:
                    story.append(Spacer(1, 0.3*inch))

            # ========== REFERENCES ==========
            if include_sources and job.sources_used:
                story.append(PageBreak())
                story.append(Paragraph("References", heading_style))
                story.append(Spacer(1, 0.2*inch))

                for source in job.sources_used:
                    doc_name = source.document_name or source.document_id
                    # Add page/slide number if available
                    location_info = ""
                    if source.page_number:
                        if doc_name.lower().endswith('.pptx'):
                            location_info = f" (Slide {source.page_number})"
                        else:
                            location_info = f" (Page {source.page_number})"
                    ref_style = ParagraphStyle(
                        'RefStyle',
                        parent=styles['Normal'],
                        fontSize=10,
                        textColor=TEXT_COLOR,
                        spaceBefore=4,
                        spaceAfter=4,
                        leftIndent=15,
                        fontName=body_font,
                    )
                    story.append(Paragraph(f"• {doc_name}{location_info}", ref_style))

            # Build PDF with page numbers
            doc.build(story, onFirstPage=add_page_number, onLaterPages=add_page_number)

            logger.info("Professional PDF generated", path=output_path)
            return output_path

        except ImportError:
            logger.error("reportlab not installed")
            return await self._generate_txt(job, filename)

    async def _generate_markdown(self, job: GenerationJob, filename: str) -> str:
        """Generate Markdown document."""
        # Determine include_sources from job metadata or fall back to config
        include_sources = job.metadata.get("include_sources", self.config.include_sources)

        lines = []

        # Title
        lines.append(f"# {job.title}")
        lines.append("")

        # Description
        if job.outline:
            lines.append(job.outline.description)
            lines.append("")

        # Table of contents
        if self.config.include_toc:
            lines.append("## Table of Contents")
            lines.append("")
            for section in job.sections:
                anchor = section.title.lower().replace(" ", "-")
                lines.append(f"- [{section.title}](#{anchor})")
            lines.append("")

        # Content
        for section in job.sections:
            lines.append(f"## {section.title}")
            lines.append("")
            content = section.revised_content or section.content
            lines.append(content)
            lines.append("")

            # Sources
            if include_sources and section.sources:
                lines.append("**Sources:**")
                for source in section.sources[:3]:
                    doc_name = source.document_name or source.document_id
                    # Add page/slide number if available
                    location_info = ""
                    if source.page_number:
                        if doc_name.lower().endswith('.pptx'):
                            location_info = f" (Slide {source.page_number})"
                        else:
                            location_info = f" (Page {source.page_number})"
                    lines.append(f"- {doc_name}{location_info}")
                lines.append("")

        # References
        if include_sources and job.sources_used:
            lines.append("## References")
            lines.append("")
            for source in job.sources_used:
                doc_name = source.document_name or source.document_id
                # Add page/slide number if available
                location_info = ""
                if source.page_number:
                    if doc_name.lower().endswith('.pptx'):
                        location_info = f" (Slide {source.page_number})"
                    else:
                        location_info = f" (Page {source.page_number})"
                lines.append(f"- {doc_name}{location_info}")

        output_path = os.path.join(self.config.output_dir, f"{filename}.md")
        with open(output_path, "w") as f:
            f.write("\n".join(lines))

        logger.info("Markdown generated", path=output_path)
        return output_path

    async def _generate_html(self, job: GenerationJob, filename: str) -> str:
        """Generate HTML document."""
        # Determine include_sources from job metadata or fall back to config
        include_sources = job.metadata.get("include_sources", self.config.include_sources)

        # Get theme colors (with custom color overrides)
        theme_key = job.metadata.get("theme", "business")
        custom_colors = job.metadata.get("custom_colors")
        theme = get_theme_colors(theme_key, custom_colors)

        # HTML font family mapping
        HTML_FONT_MAP = {
            "modern": "'Segoe UI', Roboto, 'Helvetica Neue', Arial, sans-serif",
            "classic": "Georgia, 'Times New Roman', Times, serif",
            "professional": "Arial, Helvetica, sans-serif",
            "technical": "'Courier New', Consolas, monospace",
        }
        font_family_key = job.metadata.get("font_family", "modern")
        font_family = HTML_FONT_MAP.get(font_family_key, HTML_FONT_MAP["modern"])

        html = f"""<!DOCTYPE html>
<html>
<head>
    <title>{job.title}</title>
    <style>
        body {{ font-family: {font_family}; max-width: 800px; margin: 0 auto; padding: 20px; line-height: 1.6; }}
        h1 {{ color: {theme["primary"]}; margin-bottom: 0.5em; }}
        h2 {{ color: {theme["secondary"]}; border-bottom: 2px solid {theme["secondary"]}; padding-bottom: 0.3em; }}
        .section {{ margin-bottom: 30px; }}
        .description {{ color: {theme["secondary"]}; font-style: italic; margin-bottom: 2em; }}
        .sources {{ font-size: 0.9em; color: {theme["light_gray"]}; margin-top: 1em; padding-top: 0.5em; border-top: 1px dashed {theme["light_gray"]}; }}
        .references {{ margin-top: 3em; padding-top: 1em; border-top: 2px solid {theme["primary"]}; }}
        .references h2 {{ border-bottom: none; }}
        .references ul {{ list-style-type: disc; padding-left: 1.5em; }}
        .references li {{ color: {theme["text"]}; margin-bottom: 0.5em; }}
    </style>
</head>
<body>
    <h1>{job.title}</h1>
"""

        if job.outline:
            html += f'    <p class="description">{job.outline.description}</p>\n'

        for section in job.sections:
            content = section.revised_content or section.content
            html += f"""    <div class="section">
        <h2>{section.title}</h2>
        <p>{content.replace(chr(10), '</p><p>')}</p>
"""
            if include_sources and section.sources:
                html += '        <div class="sources">Sources: '
                source_items = []
                for s in section.sources[:3]:
                    doc_name = s.document_name or s.document_id
                    # Add page/slide number if available
                    if s.page_number:
                        if doc_name.lower().endswith('.pptx'):
                            source_items.append(f"{doc_name} (Slide {s.page_number})")
                        else:
                            source_items.append(f"{doc_name} (Page {s.page_number})")
                    else:
                        source_items.append(doc_name)
                html += ", ".join(source_items)
                html += "</div>\n"

            html += "    </div>\n"

        # Add references section if sources are available
        if include_sources and job.sources_used:
            html += '    <div class="references">\n'
            html += "        <h2>References</h2>\n"
            html += "        <ul>\n"
            seen_docs = set()
            for source in job.sources_used:
                doc_name = source.document_name or source.document_id
                if doc_name and doc_name not in seen_docs:
                    seen_docs.add(doc_name)
                    # Add page/slide number if available
                    location_info = ""
                    if source.page_number:
                        if doc_name.lower().endswith('.pptx'):
                            location_info = f" (Slide {source.page_number})"
                        else:
                            location_info = f" (Page {source.page_number})"
                    html += f"            <li>{doc_name}{location_info}</li>\n"
            html += "        </ul>\n"
            html += "    </div>\n"

        html += """</body>
</html>"""

        output_path = os.path.join(self.config.output_dir, f"{filename}.html")
        with open(output_path, "w") as f:
            f.write(html)

        logger.info("HTML generated", path=output_path)
        return output_path

    async def _generate_xlsx(self, job: GenerationJob, filename: str) -> str:
        """
        Generate Excel spreadsheet with document content.

        Creates a structured Excel file with:
        - Summary sheet with document overview
        - Content sheet with all sections
        - Sources sheet with references
        """
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
            from openpyxl.utils import get_column_letter

            # Determine include_sources from job metadata or fall back to config
            include_sources = job.metadata.get("include_sources", self.config.include_sources)

            wb = Workbook()

            # Get theme colors from job metadata or use default (with custom color overrides)
            theme_key = job.metadata.get("theme", "business")
            custom_colors = job.metadata.get("custom_colors")
            theme = get_theme_colors(theme_key, custom_colors)

            # Convert theme color to Excel format (without #)
            primary_hex = theme["primary"].lstrip('#')

            # Excel font family mapping
            XLSX_FONT_MAP = {
                "modern": "Calibri",
                "classic": "Times New Roman",
                "professional": "Arial",
                "technical": "Consolas",
            }
            font_family_key = job.metadata.get("font_family", "modern")
            excel_font_name = XLSX_FONT_MAP.get(font_family_key, XLSX_FONT_MAP["modern"])

            # Define styles with theme colors and selected font
            header_font = Font(name=excel_font_name, bold=True, size=12, color="FFFFFF")
            header_fill = PatternFill(start_color=primary_hex, end_color=primary_hex, fill_type="solid")
            header_alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            cell_alignment = Alignment(vertical="top", wrap_text=True)
            thin_border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )

            # Sheet 1: Summary
            ws_summary = wb.active
            ws_summary.title = "Summary"

            # Title row
            ws_summary.merge_cells('A1:B1')
            ws_summary['A1'] = job.title
            ws_summary['A1'].font = Font(name=excel_font_name, bold=True, size=16)
            ws_summary['A1'].alignment = Alignment(horizontal="center")

            # Summary data
            summary_data = [
                ("Description", job.outline.description if job.outline else ""),
                ("Status", job.status.value),
                ("Created", job.created_at.strftime("%Y-%m-%d %H:%M")),
                ("Completed", job.completed_at.strftime("%Y-%m-%d %H:%M") if job.completed_at else "In Progress"),
                ("Total Sections", str(len(job.sections))),
                ("Sources Used", str(len(job.sources_used))),
            ]

            for i, (label, value) in enumerate(summary_data, start=3):
                ws_summary[f'A{i}'] = label
                ws_summary[f'A{i}'].font = Font(name=excel_font_name, bold=True)
                ws_summary[f'B{i}'] = value
                ws_summary[f'A{i}'].border = thin_border
                ws_summary[f'B{i}'].border = thin_border

            # Adjust column widths
            ws_summary.column_dimensions['A'].width = 20
            ws_summary.column_dimensions['B'].width = 60

            # Sheet 2: Content
            ws_content = wb.create_sheet("Content")

            # Headers
            headers = ["Section #", "Title", "Content", "Approved", "Feedback"]
            for col, header in enumerate(headers, start=1):
                cell = ws_content.cell(row=1, column=col, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_alignment
                cell.border = thin_border

            # Content rows
            for i, section in enumerate(job.sections, start=2):
                content = section.revised_content or section.content
                # Strip markdown formatting for clean Excel content
                content = strip_markdown(content) if content else ""

                ws_content.cell(row=i, column=1, value=section.order).border = thin_border
                ws_content.cell(row=i, column=2, value=section.title).border = thin_border
                ws_content.cell(row=i, column=3, value=content).border = thin_border
                ws_content.cell(row=i, column=3).alignment = cell_alignment
                ws_content.cell(row=i, column=4, value="Yes" if section.approved else "No").border = thin_border
                ws_content.cell(row=i, column=5, value=section.feedback or "").border = thin_border

            # Adjust column widths
            ws_content.column_dimensions['A'].width = 12
            ws_content.column_dimensions['B'].width = 30
            ws_content.column_dimensions['C'].width = 80
            ws_content.column_dimensions['D'].width = 12
            ws_content.column_dimensions['E'].width = 40

            # Sheet 3: Sources (only if include_sources is enabled)
            if include_sources and job.sources_used:
                ws_sources = wb.create_sheet("Sources")

                # Headers
                source_headers = ["#", "Document Name", "Location", "Relevance", "Snippet"]
                for col, header in enumerate(source_headers, start=1):
                    cell = ws_sources.cell(row=1, column=col, value=header)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = header_alignment
                    cell.border = thin_border

                # Source rows
                for i, source in enumerate(job.sources_used, start=2):
                    ws_sources.cell(row=i, column=1, value=i-1).border = thin_border
                    ws_sources.cell(row=i, column=2, value=source.document_name).border = thin_border
                    # Show Page or Slide based on document type
                    doc_name = source.document_name or ""
                    if source.page_number:
                        location_label = f"Slide {source.page_number}" if doc_name.lower().endswith('.pptx') else f"Page {source.page_number}"
                    else:
                        location_label = "N/A"
                    ws_sources.cell(row=i, column=3, value=location_label).border = thin_border
                    ws_sources.cell(row=i, column=4, value=f"{source.relevance_score:.2f}").border = thin_border
                    ws_sources.cell(row=i, column=5, value=source.snippet[:200] + "..." if len(source.snippet) > 200 else source.snippet).border = thin_border
                    ws_sources.cell(row=i, column=5).alignment = cell_alignment

                # Adjust column widths
                ws_sources.column_dimensions['A'].width = 8
                ws_sources.column_dimensions['B'].width = 40
                ws_sources.column_dimensions['C'].width = 10
                ws_sources.column_dimensions['D'].width = 12
                ws_sources.column_dimensions['E'].width = 60

            # Save workbook
            output_path = os.path.join(self.config.output_dir, f"{filename}.xlsx")
            wb.save(output_path)

            logger.info("XLSX generated", path=output_path)
            return output_path

        except ImportError:
            logger.error("openpyxl not installed, falling back to TXT")
            return await self._generate_txt(job, filename)
        except Exception as e:
            logger.error("XLSX generation failed", error=str(e))
            raise

    async def _generate_txt(self, job: GenerationJob, filename: str) -> str:
        """Generate plain text document."""
        lines = []

        lines.append("=" * 60)
        lines.append(job.title.upper())
        lines.append("=" * 60)
        lines.append("")

        if job.outline:
            lines.append(job.outline.description)
            lines.append("")

        for section in job.sections:
            lines.append("-" * 40)
            lines.append(section.title)
            lines.append("-" * 40)
            lines.append("")
            content = section.revised_content or section.content
            lines.append(content)
            lines.append("")

        output_path = os.path.join(self.config.output_dir, f"{filename}.txt")
        with open(output_path, "w") as f:
            f.write("\n".join(lines))

        logger.info("TXT generated", path=output_path)
        return output_path

    def _get_content_type(self, format: OutputFormat) -> str:
        """Get MIME content type for format."""
        content_types = {
            OutputFormat.PPTX: "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            OutputFormat.DOCX: "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            OutputFormat.PDF: "application/pdf",
            OutputFormat.XLSX: "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            OutputFormat.MARKDOWN: "text/markdown",
            OutputFormat.HTML: "text/html",
            OutputFormat.TXT: "text/plain",
        }
        return content_types.get(format, "application/octet-stream")


# =============================================================================
# Singleton Instance
# =============================================================================

_generation_service: Optional[DocumentGenerationService] = None


def get_generation_service() -> DocumentGenerationService:
    """Get or create document generation service instance."""
    global _generation_service

    if _generation_service is None:
        _generation_service = DocumentGenerationService()

    return _generation_service

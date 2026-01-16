"""
Theme Applier

Applies themes to documents during generation.
Handles format-specific theme application for PPTX, DOCX, XLSX, and PDF.
"""

from typing import Optional, Dict, Any, Union

import structlog

from .models import ThemeProfile, PPTXTheme, DOCXTheme, XLSXTheme, PDFTheme

logger = structlog.get_logger(__name__)


class ThemeApplier:
    """Applies themes to documents during generation.

    This class provides format-specific theme application methods
    that work with python-pptx, python-docx, openpyxl, etc.
    """

    def __init__(self, theme: Union[ThemeProfile, Dict[str, Any]]):
        """Initialize the theme applier.

        Args:
            theme: ThemeProfile object or theme dict with color/font settings
        """
        if isinstance(theme, dict):
            self.theme = ThemeProfile(
                name=theme.get("name", "Custom"),
                primary=theme.get("primary", "#333333"),
                secondary=theme.get("secondary", "#666666"),
                accent=theme.get("accent", "#0066CC"),
                background=theme.get("background", "#FFFFFF"),
                text=theme.get("text", "#333333"),
                font_heading=theme.get("font_heading", "Arial"),
                font_body=theme.get("font_body", "Arial"),
            )
        else:
            self.theme = theme

    def get_rgb_color(self, color_key: str):
        """Get RGBColor object for a theme color key.

        Args:
            color_key: One of primary, secondary, accent, background, text

        Returns:
            RGBColor object for use with python-pptx/docx
        """
        from pptx.dml.color import RGBColor

        color_map = {
            "primary": self.theme.primary,
            "secondary": self.theme.secondary,
            "accent": self.theme.accent,
            "background": self.theme.background,
            "text": self.theme.text,
        }

        hex_color = color_map.get(color_key, self.theme.text)
        hex_color = hex_color.lstrip("#")

        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )

    # =========================================================================
    # PPTX Theme Application
    # =========================================================================

    def apply_to_pptx_slide(self, slide, is_title_slide: bool = False):
        """Apply theme colors and fonts to a PPTX slide.

        Args:
            slide: python-pptx slide object
            is_title_slide: Whether this is a title slide (affects styling)
        """
        from pptx.util import Pt

        primary_color = self.get_rgb_color("primary")
        text_color = self.get_rgb_color("text")

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # Apply font
                    if is_title_slide or shape == slide.shapes.title:
                        run.font.name = self.theme.font_heading
                        run.font.color.rgb = primary_color
                    else:
                        run.font.name = self.theme.font_body
                        run.font.color.rgb = text_color

    def apply_to_pptx_presentation(self, presentation,
                                   skip_template_slides: bool = True):
        """Apply theme to entire PPTX presentation.

        Args:
            presentation: python-pptx Presentation object
            skip_template_slides: If True, don't modify slides from template
        """
        for i, slide in enumerate(presentation.slides):
            is_title = i == 0
            self.apply_to_pptx_slide(slide, is_title_slide=is_title)

    def get_pptx_background_fill(self, is_title_slide: bool = False):
        """Get background fill settings for PPTX slides.

        Args:
            is_title_slide: Whether this is a title slide

        Returns:
            Dict with fill settings
        """
        return {
            "type": "solid",
            "color": self.theme.background,
        }

    def style_pptx_chart(self, chart, use_accent_colors: bool = True):
        """Apply theme colors to a PPTX chart.

        Args:
            chart: python-pptx chart object
            use_accent_colors: Whether to use accent colors for data series
        """
        from pptx.dml.color import RGBColor

        # Apply primary color to chart title
        if chart.has_title:
            title_frame = chart.chart_title.text_frame
            for para in title_frame.paragraphs:
                for run in para.runs:
                    run.font.name = self.theme.font_heading
                    run.font.color.rgb = self.get_rgb_color("primary")

        # Apply colors to data series
        if use_accent_colors:
            accent_color = self.get_rgb_color("accent")
            for series in chart.plots[0].series:
                if hasattr(series, 'format') and hasattr(series.format, 'fill'):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = accent_color

    # =========================================================================
    # DOCX Theme Application
    # =========================================================================

    def apply_to_docx_paragraph(self, paragraph, is_heading: bool = False):
        """Apply theme styling to a DOCX paragraph.

        Args:
            paragraph: python-docx paragraph object
            is_heading: Whether this is a heading paragraph
        """
        from docx.shared import RGBColor as DocxRGBColor

        for run in paragraph.runs:
            if is_heading:
                run.font.name = self.theme.font_heading
                hex_color = self.theme.primary.lstrip("#")
            else:
                run.font.name = self.theme.font_body
                hex_color = self.theme.text.lstrip("#")

            run.font.color.rgb = DocxRGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16),
            )

    def apply_to_docx_document(self, document):
        """Apply theme to entire DOCX document.

        Args:
            document: python-docx Document object
        """
        for para in document.paragraphs:
            is_heading = para.style.name.startswith("Heading")
            self.apply_to_docx_paragraph(para, is_heading=is_heading)

    def get_docx_style_dict(self) -> Dict[str, Any]:
        """Get style dictionary for DOCX document creation.

        Returns:
            Dict with style settings for document generation
        """
        return {
            "heading_font": self.theme.font_heading,
            "body_font": self.theme.font_body,
            "heading_color": self.theme.primary,
            "body_color": self.theme.text,
            "accent_color": self.theme.accent,
        }

    # =========================================================================
    # XLSX Theme Application
    # =========================================================================

    def apply_to_xlsx_cell(self, cell, is_header: bool = False):
        """Apply theme styling to an XLSX cell.

        Args:
            cell: openpyxl cell object
            is_header: Whether this is a header cell
        """
        from openpyxl.styles import Font, PatternFill

        if is_header:
            hex_bg = self.theme.primary.lstrip("#")
            cell.fill = PatternFill(
                start_color=hex_bg,
                end_color=hex_bg,
                fill_type="solid",
            )
            cell.font = Font(
                name=self.theme.font_heading,
                bold=True,
                color="FFFFFF",
            )
        else:
            hex_text = self.theme.text.lstrip("#")
            cell.font = Font(
                name=self.theme.font_body,
                color=hex_text,
            )

    def apply_to_xlsx_worksheet(self, worksheet, header_row: int = 1):
        """Apply theme to an XLSX worksheet.

        Args:
            worksheet: openpyxl worksheet object
            header_row: Row number of the header (1-indexed)
        """
        for row_idx, row in enumerate(worksheet.iter_rows(), start=1):
            is_header = row_idx == header_row
            for cell in row:
                if cell.value:  # Only style non-empty cells
                    self.apply_to_xlsx_cell(cell, is_header=is_header)

    def get_xlsx_style_dict(self) -> Dict[str, Any]:
        """Get style dictionary for XLSX workbook creation.

        Returns:
            Dict with style settings for spreadsheet generation
        """
        return {
            "header_font": self.theme.font_heading,
            "data_font": self.theme.font_body,
            "header_bg_color": self.theme.primary,
            "header_text_color": "#FFFFFF",
            "data_text_color": self.theme.text,
            "accent_color": self.theme.accent,
        }

    # =========================================================================
    # PDF/HTML Theme Application
    # =========================================================================

    def get_pdf_css(self) -> str:
        """Generate CSS styles for PDF generation.

        Returns:
            CSS string for use with WeasyPrint or similar
        """
        return f"""
:root {{
    --primary-color: {self.theme.primary};
    --secondary-color: {self.theme.secondary};
    --accent-color: {self.theme.accent};
    --background-color: {self.theme.background};
    --text-color: {self.theme.text};
    --font-heading: "{self.theme.font_heading}", sans-serif;
    --font-body: "{self.theme.font_body}", sans-serif;
}}

body {{
    font-family: var(--font-body);
    color: var(--text-color);
    background-color: var(--background-color);
}}

h1, h2, h3, h4, h5, h6 {{
    font-family: var(--font-heading);
    color: var(--primary-color);
}}

a {{
    color: var(--accent-color);
}}

.accent {{
    color: var(--accent-color);
}}

.highlight {{
    background-color: var(--accent-color);
    color: white;
}}
"""

    def get_html_style_tag(self) -> str:
        """Get HTML style tag with theme CSS.

        Returns:
            HTML <style> tag with theme CSS
        """
        return f"<style>\n{self.get_pdf_css()}\n</style>"

    def get_pdf_style_dict(self) -> Dict[str, Any]:
        """Get style dictionary for PDF generation.

        Returns:
            Dict with style settings for PDF generation
        """
        return {
            "heading_font": self.theme.font_heading,
            "body_font": self.theme.font_body,
            "heading_color": self.theme.primary,
            "body_color": self.theme.text,
            "accent_color": self.theme.accent,
            "background_color": self.theme.background,
        }


def create_applier_from_dict(theme_dict: Dict[str, Any]) -> ThemeApplier:
    """Create a ThemeApplier from a dictionary.

    Args:
        theme_dict: Dictionary with theme settings

    Returns:
        ThemeApplier instance
    """
    return ThemeApplier(theme_dict)


def create_applier_from_profile(theme: ThemeProfile) -> ThemeApplier:
    """Create a ThemeApplier from a ThemeProfile.

    Args:
        theme: ThemeProfile object

    Returns:
        ThemeApplier instance
    """
    return ThemeApplier(theme)

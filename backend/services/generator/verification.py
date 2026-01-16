"""
Document Verification and Repair Module

Provides functionality for verifying generated documents are complete
and properly formatted, with automatic repair capabilities.
"""

import os
import re
import json
from typing import Dict, Any, List, Optional, Tuple, TYPE_CHECKING

import structlog

if TYPE_CHECKING:
    from .models import GenerationJob, OutputFormat

logger = structlog.get_logger(__name__)


class DocumentVerifier:
    """Handles document verification and repair operations.

    This class provides methods to verify generated documents for completeness
    and proper formatting, with automatic repair capabilities for common issues.
    """

    def __init__(self):
        self._llm = None

    async def _get_llm(self):
        """Get the LLM instance for AI-assisted operations."""
        if self._llm is None:
            from backend.services.llm import EnhancedLLMFactory
            self._llm, _ = await EnhancedLLMFactory.get_chat_model_for_operation(
                operation="content_generation",
                user_id=None,
            )
        return self._llm

    async def suggest_theme(
        self,
        title: str,
        description: str,
        document_type: str = "pptx",
    ) -> Dict[str, Any]:
        """
        Use LLM to suggest optimal theming for a document based on its content.

        Args:
            title: Document title
            description: Document description/topic
            document_type: Type of document (pptx, docx, pdf, etc.)

        Returns:
            Dictionary with recommended theme, font_family, layout, and reason
        """
        from .config import THEMES, FONT_FAMILIES, LAYOUT_TEMPLATES

        try:
            llm = await self._get_llm()

            # Build theme options for the prompt
            theme_options = "\n".join([
                f"- {key}: {theme['name']} - {theme['description']}"
                for key, theme in THEMES.items()
            ])

            font_options = "\n".join([
                f"- {key}: {font['name']} - {font['description']}"
                for key, font in FONT_FAMILIES.items()
            ])

            layout_options = "\n".join([
                f"- {key}: {layout['name']} - {layout['description']}"
                for key, layout in LAYOUT_TEMPLATES.items()
            ])

            prompt = f"""Analyze this document and suggest optimal theming.

Document Title: {title}
Document Description: {description}
Document Type: {document_type}

Available Themes:
{theme_options}

Available Font Families:
{font_options}

Available Layouts (for presentations):
{layout_options}

Based on the document topic, audience, and purpose, recommend the best options.
Consider: industry conventions, emotional tone, readability, and visual impact.

Return ONLY a JSON object with this exact structure:
{{
    "theme": "theme_key",
    "font_family": "font_key",
    "layout": "layout_key",
    "animations": true/false,
    "reason": "Brief explanation of why these choices suit this document"
}}"""

            response = await llm.ainvoke(prompt)
            response_text = response.content if hasattr(response, 'content') else str(response)

            # Extract JSON from response
            json_match = re.search(r'\{[^{}]*\}', response_text, re.DOTALL)
            if json_match:
                suggestion = json.loads(json_match.group())
            else:
                suggestion = json.loads(response_text)

            # Validate keys exist
            if suggestion.get("theme") not in THEMES:
                suggestion["theme"] = "business"
            if suggestion.get("font_family") not in FONT_FAMILIES:
                suggestion["font_family"] = "modern"
            if suggestion.get("layout") not in LAYOUT_TEMPLATES:
                suggestion["layout"] = "standard"

            # Add full theme/font/layout details
            suggestion["theme_details"] = THEMES.get(suggestion["theme"])
            suggestion["font_details"] = FONT_FAMILIES.get(suggestion["font_family"])
            suggestion["layout_details"] = LAYOUT_TEMPLATES.get(suggestion["layout"])

            logger.info(
                "Theme suggestion generated",
                theme=suggestion.get("theme"),
                font_family=suggestion.get("font_family"),
                layout=suggestion.get("layout"),
            )

            return suggestion

        except Exception as e:
            logger.warning(f"Theme suggestion failed, using defaults: {e}")
            # Return sensible defaults
            return {
                "theme": "business",
                "font_family": "modern",
                "layout": "standard",
                "animations": False,
                "reason": "Default professional theme selected",
                "theme_details": THEMES["business"],
                "font_details": FONT_FAMILIES["modern"],
                "layout_details": LAYOUT_TEMPLATES["standard"],
            }

    async def analyze_section_layout(
        self,
        section_title: str,
        section_content: str,
        has_chart_data: bool = False,
        has_image: bool = False,
        document_type: str = "pptx",
    ) -> Dict[str, Any]:
        """
        Use AI to analyze section content and recommend optimal layout.

        Args:
            section_title: Title of the section
            section_content: Content of the section
            has_chart_data: Whether chartable data was detected
            has_image: Whether an image is available for this section
            document_type: Type of document being generated

        Returns:
            Dictionary with layout recommendations
        """
        try:
            llm = await self._get_llm()

            # Truncate content for analysis (first 500 chars)
            content_preview = section_content[:500] + "..." if len(section_content) > 500 else section_content

            prompt = f"""Analyze this section content and recommend the optimal layout strategy.

Section Title: {section_title}
Content Preview: {content_preview}
Has Chartable Data: {has_chart_data}
Has Image Available: {has_image}
Document Type: {document_type}

Analyze the content and recommend:
1. Layout type (bullets, numbered_list, paragraph, chart, image_focused, two_column)
2. Content emphasis (text_heavy, visual_heavy, balanced)
3. Image position if applicable (left, right, top, bottom, full_width, none)
4. Font size recommendation (small, medium, large) based on content amount
5. Whether to split into multiple slides/pages if content is long

Return ONLY a JSON object:
{{
    "layout_type": "bullets|numbered_list|paragraph|chart|image_focused|two_column",
    "content_emphasis": "text_heavy|visual_heavy|balanced",
    "image_position": "left|right|top|bottom|full_width|none",
    "font_size": "small|medium|large",
    "split_content": true|false,
    "bullet_count_per_slide": 5-8,
    "reason": "Brief explanation"
}}"""

            response = await llm.ainvoke(prompt)
            response_text = response.content if hasattr(response, 'content') else str(response)

            # Parse JSON response
            json_match = re.search(r'\{[^{}]*\}', response_text, re.DOTALL)
            if json_match:
                recommendation = json.loads(json_match.group())
            else:
                recommendation = json.loads(response_text)

            # Validate and set defaults
            valid_layouts = ["bullets", "numbered_list", "paragraph", "chart", "image_focused", "two_column"]
            if recommendation.get("layout_type") not in valid_layouts:
                recommendation["layout_type"] = "bullets"

            valid_positions = ["left", "right", "top", "bottom", "full_width", "none"]
            if recommendation.get("image_position") not in valid_positions:
                recommendation["image_position"] = "right" if has_image else "none"

            logger.debug(
                "AI layout analysis complete",
                section=section_title,
                layout=recommendation.get("layout_type"),
                image_position=recommendation.get("image_position"),
            )

            return recommendation

        except Exception as e:
            logger.warning(f"AI layout analysis failed, using defaults: {e}")
            # Smart defaults based on content characteristics
            layout_type = "chart" if has_chart_data else "bullets"
            image_position = "right" if has_image else "none"

            return {
                "layout_type": layout_type,
                "content_emphasis": "balanced",
                "image_position": image_position,
                "font_size": "medium",
                "split_content": len(section_content) > 1500,
                "bullet_count_per_slide": 6,
                "reason": "Default layout based on content characteristics",
            }

    async def suggest_image_position(
        self,
        section_title: str,
        section_content: str,
        content_type: str,
        slide_dimensions: Tuple[float, float] = (13.333, 7.5),
    ) -> Dict[str, Any]:
        """
        Use AI to suggest optimal image positioning based on content.

        Args:
            section_title: Title of the section
            section_content: Content of the section
            content_type: Type of content (bullets, paragraph, chart, etc.)
            slide_dimensions: Width and height of the slide/page

        Returns:
            Dictionary with image positioning details
        """
        try:
            llm = await self._get_llm()

            slide_width, slide_height = slide_dimensions
            content_preview = section_content[:300] + "..." if len(section_content) > 300 else section_content

            prompt = f"""Determine optimal image positioning for this slide content.

Section Title: {section_title}
Content Type: {content_type}
Content Preview: {content_preview}
Slide Dimensions: {slide_width}" x {slide_height}"

Consider:
- Balance between text and image
- Visual hierarchy (what should be seen first)
- Reading flow (left-to-right, top-to-bottom)
- Amount of text content

Return ONLY a JSON object with positioning in inches:
{{
    "position": "right|left|top|bottom|background|corner",
    "x": 0.0,
    "y": 0.0,
    "width": 0.0,
    "height": 0.0,
    "text_area_adjustment": {{
        "width_reduction": 0.0,
        "x_offset": 0.0
    }},
    "overlay_text": false,
    "reason": "Brief explanation"
}}"""

            response = await llm.ainvoke(prompt)
            response_text = response.content if hasattr(response, 'content') else str(response)

            # Parse JSON response
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            if json_match:
                positioning = json.loads(json_match.group())
            else:
                positioning = json.loads(response_text)

            logger.debug(
                "AI image positioning complete",
                section=section_title,
                position=positioning.get("position"),
            )

            return positioning

        except Exception as e:
            logger.warning(f"AI image positioning failed, using defaults: {e}")
            # Default right-side positioning
            slide_width, slide_height = slide_dimensions
            return {
                "position": "right",
                "x": slide_width * 0.55,
                "y": 1.5,
                "width": slide_width * 0.4,
                "height": slide_height * 0.6,
                "text_area_adjustment": {
                    "width_reduction": slide_width * 0.45,
                    "x_offset": 0,
                },
                "overlay_text": False,
                "reason": "Default right-side image placement",
            }

    async def optimize_content_distribution(
        self,
        sections: List[Dict[str, Any]],
        output_format: str,
        max_items_per_page: int = 6,
    ) -> List[Dict[str, Any]]:
        """
        Use AI to optimize how content is distributed across pages/slides.

        Args:
            sections: List of section dictionaries with title and content
            output_format: Document format (pptx, docx, pdf)
            max_items_per_page: Maximum bullet points per page/slide

        Returns:
            List of optimized sections with split recommendations
        """
        try:
            llm = await self._get_llm()

            # Prepare sections summary
            sections_summary = []
            for i, section in enumerate(sections):
                content = section.get("content", "")
                bullet_count = content.count("- ") + content.count("• ")
                word_count = len(content.split())
                sections_summary.append({
                    "index": i,
                    "title": section.get("title", f"Section {i+1}"),
                    "bullet_count": bullet_count,
                    "word_count": word_count,
                })

            prompt = f"""Analyze these sections and recommend optimal content distribution.

Sections:
{json.dumps(sections_summary, indent=2)}

Output Format: {output_format}
Max Items Per Page/Slide: {max_items_per_page}

For each section, recommend:
1. Whether to split into multiple pages/slides
2. How many slides/pages needed
3. Grouping strategy (by topic, by importance, chronological)

Return ONLY a JSON array:
[
    {{
        "index": 0,
        "split": true|false,
        "pages_needed": 1,
        "items_per_page": [6, 4],
        "strategy": "topic|importance|chronological"
    }},
    ...
]"""

            response = await llm.ainvoke(prompt)
            response_text = response.content if hasattr(response, 'content') else str(response)

            # Parse JSON response
            json_match = re.search(r'\[.*\]', response_text, re.DOTALL)
            if json_match:
                distribution = json.loads(json_match.group())
            else:
                distribution = json.loads(response_text)

            # Apply recommendations to sections
            for rec in distribution:
                idx = rec.get("index", 0)
                if idx < len(sections):
                    sections[idx]["layout_recommendation"] = rec

            logger.info(
                "AI content distribution optimized",
                sections=len(sections),
                splits_recommended=sum(1 for r in distribution if r.get("split", False)),
            )

            return sections

        except Exception as e:
            logger.warning(f"AI content distribution failed, using defaults: {e}")
            # Apply simple heuristic-based defaults
            for section in sections:
                content = section.get("content", "")
                bullet_count = content.count("- ") + content.count("• ")
                section["layout_recommendation"] = {
                    "split": bullet_count > max_items_per_page,
                    "pages_needed": max(1, (bullet_count + max_items_per_page - 1) // max_items_per_page),
                    "items_per_page": [min(bullet_count, max_items_per_page)],
                    "strategy": "topic",
                }
            return sections

    async def verify_generated_document(
        self,
        file_path: str,
        output_format: "OutputFormat",
        expected_sections: int,
        job_title: str,
    ) -> Dict[str, Any]:
        """
        Verify the generated document is complete and properly formatted.

        Args:
            file_path: Path to the generated document
            output_format: Format of the document
            expected_sections: Expected number of sections
            job_title: Title of the document

        Returns:
            Dictionary with verification results and any issues found
        """
        from .models import OutputFormat

        issues = []
        warnings = []
        stats = {}

        try:
            if output_format == OutputFormat.PPTX:
                from pptx import Presentation
                from pptx.util import Pt

                prs = Presentation(file_path)

                # Count slides
                slide_count = len(prs.slides)
                stats["slide_count"] = slide_count

                # Check for empty slides
                empty_slides = []
                slides_with_notes = 0
                slides_with_images = 0
                slides_without_title = []

                for i, slide in enumerate(prs.slides):
                    has_content = False
                    has_image = False
                    has_title = False

                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text_frame.text.strip():
                            has_content = True
                            if not has_title:
                                try:
                                    for para in shape.text_frame.paragraphs:
                                        if para.runs and para.runs[0].font.size and para.runs[0].font.size >= Pt(18):
                                            has_title = True
                                            break
                                except Exception:
                                    pass
                        if shape.shape_type == 13:  # Picture
                            has_image = True
                            slides_with_images += 1

                    if not has_content:
                        empty_slides.append(i + 1)
                    if not has_title and i > 0:
                        slides_without_title.append(i + 1)

                    try:
                        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                            slides_with_notes += 1
                    except Exception:
                        pass

                stats["empty_slides"] = empty_slides
                stats["slides_with_notes"] = slides_with_notes
                stats["slides_with_images"] = slides_with_images
                stats["slides_without_title"] = slides_without_title

                if empty_slides:
                    issues.append(f"Empty slides detected: {empty_slides}")

                min_expected = expected_sections + 3
                if slide_count < min_expected:
                    issues.append(f"Expected at least {min_expected} slides, got {slide_count}")

                if slides_without_title:
                    warnings.append(f"Slides missing titles: {slides_without_title}")

            elif output_format == OutputFormat.DOCX:
                from docx import Document

                doc = Document(file_path)

                paragraph_count = len(doc.paragraphs)
                heading_count = sum(1 for p in doc.paragraphs if p.style.name.startswith('Heading'))
                empty_paragraphs = sum(1 for p in doc.paragraphs if not p.text.strip())

                stats["paragraph_count"] = paragraph_count
                stats["heading_count"] = heading_count
                stats["empty_paragraphs"] = empty_paragraphs

                total_text = " ".join(p.text for p in doc.paragraphs)
                if len(total_text.strip()) < 100:
                    issues.append("Document appears to have very little content")

                title_found = any(job_title.lower() in p.text.lower() for p in doc.paragraphs[:10])
                if not title_found:
                    issues.append("Document title not found on cover page")

                if heading_count < expected_sections:
                    issues.append(f"Expected {expected_sections} section headings, found {heading_count}")

                if empty_paragraphs > paragraph_count * 0.3:
                    warnings.append(f"Excessive empty paragraphs: {empty_paragraphs}")

            elif output_format == OutputFormat.PDF:
                file_size = os.path.getsize(file_path)
                stats["file_size_kb"] = file_size / 1024

                if file_size < 5000:
                    issues.append("PDF file seems too small, may be incomplete")

                try:
                    from PyPDF2 import PdfReader
                    reader = PdfReader(file_path)
                    page_count = len(reader.pages)
                    stats["page_count"] = page_count

                    min_expected = expected_sections + 2
                    if page_count < min_expected:
                        issues.append(f"Expected at least {min_expected} pages, got {page_count}")

                    total_text = ""
                    for page in reader.pages[:3]:
                        total_text += page.extract_text() or ""

                    if len(total_text.strip()) < 100:
                        issues.append("PDF appears to have very little text content")

                except ImportError:
                    warnings.append("PyPDF2 not installed, detailed PDF verification skipped")
                except Exception as e:
                    warnings.append(f"Could not read PDF for verification: {e}")

            elif output_format == OutputFormat.XLSX:
                import openpyxl

                wb = openpyxl.load_workbook(file_path)
                sheet_count = len(wb.sheetnames)
                stats["sheet_count"] = sheet_count

                empty_sheets = []
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    if sheet.max_row <= 1:
                        empty_sheets.append(sheet_name)

                if empty_sheets:
                    issues.append(f"Empty sheets detected: {empty_sheets}")

            else:
                file_size = os.path.getsize(file_path)
                stats["file_size_kb"] = file_size / 1024

                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()

                stats["character_count"] = len(content)
                stats["line_count"] = content.count('\n') + 1

                if len(content.strip()) < 100:
                    issues.append("Document appears to have very little content")

                # Check for title
                if output_format == OutputFormat.MARKDOWN:
                    if not content.strip().startswith('#'):
                        issues.append("Markdown document missing title heading")
                elif output_format == OutputFormat.HTML:
                    if '<h1' not in content.lower() and '<title' not in content.lower():
                        issues.append("HTML document missing title")

            verification_passed = len(issues) == 0
            verification_result = {
                "passed": verification_passed,
                "issues": issues,
                "warnings": warnings,
                "stats": stats,
                "file_path": file_path,
                "format": output_format.value,
            }

            if verification_passed:
                logger.info(
                    "Document verification passed",
                    format=output_format.value,
                    stats=stats,
                    warnings=len(warnings),
                )
            else:
                logger.warning(
                    "Document verification found issues",
                    format=output_format.value,
                    issues=issues,
                    warnings=warnings,
                )

            return verification_result

        except Exception as e:
            logger.error(f"Document verification failed: {e}")
            return {
                "passed": False,
                "issues": [f"Verification error: {str(e)}"],
                "warnings": [],
                "stats": {},
                "file_path": file_path,
                "format": output_format.value if output_format else "unknown",
            }

    async def repair_document(
        self,
        file_path: str,
        output_format: "OutputFormat",
        verification_result: Dict[str, Any],
        job: Optional["GenerationJob"] = None,
    ) -> Dict[str, Any]:
        """
        Repair issues found in document verification.

        Args:
            file_path: Path to the document to repair
            output_format: Format of the document
            verification_result: Results from verify_generated_document
            job: The generation job (for context)

        Returns:
            Dictionary with repair results
        """
        from .models import OutputFormat

        repairs_made = []
        repair_failed = []
        document_modified = False

        issues = verification_result.get("issues", [])
        stats = verification_result.get("stats", {})

        if not issues:
            return {
                "success": True,
                "repairs_made": [],
                "repair_failed": [],
                "document_modified": False,
                "message": "No issues to repair",
            }

        try:
            if output_format == OutputFormat.PPTX:
                from pptx import Presentation
                from pptx.util import Pt

                prs = Presentation(file_path)
                slide_count = len(prs.slides)

                # Repair: Remove empty slides
                empty_slides = stats.get("empty_slides", [])
                if empty_slides and len(empty_slides) < slide_count // 2:
                    for slide_idx in reversed(empty_slides):
                        if slide_idx > 1 and slide_idx < slide_count:
                            try:
                                slide_id = prs.slides._sldIdLst[slide_idx - 1].rId
                                prs.part.drop_rel(slide_id)
                                del prs.slides._sldIdLst[slide_idx - 1]
                                repairs_made.append(f"Removed empty slide {slide_idx}")
                                document_modified = True
                            except Exception as e:
                                repair_failed.append(f"Could not remove slide {slide_idx}: {e}")

                # Repair: Add missing notes
                if job and stats.get("slides_with_notes", 0) < slide_count - 2:
                    for i, slide in enumerate(prs.slides):
                        if i == 0:
                            continue
                        try:
                            notes_slide = slide.notes_slide
                            notes_frame = notes_slide.notes_text_frame
                            if not notes_frame.text.strip():
                                notes_frame.text = f"Slide {i + 1} of {slide_count}\nGenerated: {job.created_at.strftime('%Y-%m-%d %H:%M') if job.created_at else 'N/A'}"
                                repairs_made.append(f"Added notes to slide {i + 1}")
                                document_modified = True
                        except Exception as e:
                            repair_failed.append(f"Could not add notes to slide {i + 1}: {e}")

                if document_modified:
                    prs.save(file_path)
                    logger.info("PPTX document repaired", repairs=len(repairs_made))

            elif output_format == OutputFormat.DOCX:
                from docx import Document
                from docx.shared import Pt

                doc = Document(file_path)

                # Repair: Add missing title
                if any("title not found" in issue.lower() for issue in issues):
                    if job:
                        try:
                            first_para = doc.paragraphs[0] if doc.paragraphs else doc.add_paragraph()
                            if not first_para.text.strip():
                                run = first_para.add_run(job.title)
                                run.font.size = Pt(36)
                                run.font.bold = True
                                repairs_made.append("Added missing document title")
                                document_modified = True
                        except Exception as e:
                            repair_failed.append(f"Could not add title: {e}")

                # Repair: Remove excessive empty paragraphs
                if any("excessive empty" in issue.lower() or "excessive empty" in w.lower()
                       for issue in issues for w in verification_result.get("warnings", [])):
                    paragraphs_to_remove = []
                    consecutive_empty = 0
                    for para in doc.paragraphs:
                        if not para.text.strip():
                            consecutive_empty += 1
                            if consecutive_empty > 2:
                                paragraphs_to_remove.append(para)
                        else:
                            consecutive_empty = 0

                    removed_count = 0
                    for para in paragraphs_to_remove:
                        try:
                            p = para._element
                            p.getparent().remove(p)
                            removed_count += 1
                            document_modified = True
                        except Exception:
                            pass

                    if removed_count:
                        repairs_made.append(f"Removed {removed_count} excessive empty paragraphs")

                if document_modified:
                    doc.save(file_path)
                    logger.info("DOCX document repaired", repairs=len(repairs_made))

            elif output_format == OutputFormat.PDF:
                # PDF repair is limited - would require regeneration
                repair_failed.append("PDF documents cannot be directly repaired, regeneration required")

            elif output_format == OutputFormat.XLSX:
                import openpyxl

                wb = openpyxl.load_workbook(file_path)

                # Repair: Add headers to empty sheets
                empty_sheets = stats.get("empty_sheets", [])
                if empty_sheets:
                    for sheet_name in empty_sheets:
                        try:
                            sheet = wb[sheet_name]
                            if sheet.max_row <= 1:
                                sheet['A1'] = sheet_name
                                sheet['A1'].font = openpyxl.styles.Font(bold=True, size=14)
                                repairs_made.append(f"Added header to sheet '{sheet_name}'")
                                document_modified = True
                        except Exception as e:
                            repair_failed.append(f"Could not fix sheet '{sheet_name}': {e}")

                if document_modified:
                    wb.save(file_path)
                    logger.info("XLSX document repaired", repairs=len(repairs_made))

            elif output_format in (OutputFormat.MARKDOWN, OutputFormat.HTML):
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()

                # Repair: Add missing title
                if any("missing title" in issue.lower() for issue in issues) and job:
                    if output_format == OutputFormat.MARKDOWN:
                        if not content.strip().startswith('#'):
                            content = f"# {job.title}\n\n{content}"
                            repairs_made.append("Added missing title heading")
                            document_modified = True
                    else:  # HTML
                        if '<h1' not in content.lower():
                            # Insert after <body> or at beginning
                            if '<body' in content.lower():
                                idx = content.lower().find('<body')
                                end_idx = content.find('>', idx) + 1
                                content = content[:end_idx] + f"\n<h1>{job.title}</h1>\n" + content[end_idx:]
                            else:
                                content = f"<h1>{job.title}</h1>\n{content}"
                            repairs_made.append("Added missing title element")
                            document_modified = True

                if document_modified:
                    with open(file_path, 'w', encoding='utf-8') as f:
                        f.write(content)
                    logger.info(f"{output_format.value.upper()} document repaired", repairs=len(repairs_made))

            repair_result = {
                "success": len(repair_failed) == 0,
                "repairs_made": repairs_made,
                "repair_failed": repair_failed,
                "document_modified": document_modified,
                "message": f"Made {len(repairs_made)} repairs" if repairs_made else "No repairs possible",
            }

            logger.info(
                "Document repair completed",
                format=output_format.value,
                repairs_made=len(repairs_made),
                repair_failed=len(repair_failed),
            )

            return repair_result

        except Exception as e:
            logger.error(f"Document repair failed: {e}")
            return {
                "success": False,
                "repairs_made": repairs_made,
                "repair_failed": repair_failed + [f"Repair error: {str(e)}"],
                "document_modified": document_modified,
                "message": f"Repair failed: {str(e)}",
            }


# Singleton instance for easy access
_verifier_instance = None


def get_document_verifier() -> DocumentVerifier:
    """Get the singleton DocumentVerifier instance."""
    global _verifier_instance
    if _verifier_instance is None:
        _verifier_instance = DocumentVerifier()
    return _verifier_instance

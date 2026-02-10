"""
AIDocumentIndexer - Document Generation API Routes
===================================================

Endpoints for generating documents with human-in-the-loop workflow.
"""

from datetime import datetime
from typing import Optional, List
from uuid import UUID
import os
import shutil

from fastapi import APIRouter, Depends, HTTPException, Query, status, UploadFile, File
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field
from sqlalchemy.ext.asyncio import AsyncSession
import structlog

from backend.db.database import get_async_session, async_session_context
from backend.api.middleware.auth import AuthenticatedUser
from backend.services.generator import (
    DocumentGenerationService,
    GenerationJob,
    GenerationStatus,
    OutputFormat,
    DocumentOutline,
    Section,
    SourceReference,
    get_generation_service,
    THEMES,
    FONT_FAMILIES,
    LAYOUT_TEMPLATES,
    check_spelling,
)
from backend.services.image_generator import (
    ImageGeneratorService,
    ImageGeneratorConfig,
    ImageBackend,
    GeneratedImage,
    get_image_generator,
)

logger = structlog.get_logger(__name__)

router = APIRouter()


# =============================================================================
# Pydantic Models
# =============================================================================

class CreateJobRequest(BaseModel):
    """Request to create a new generation job."""
    title: str = Field(..., min_length=1, max_length=200)
    description: str = Field(..., min_length=10)
    output_format: str = Field(default="docx")
    output_language: str = Field(
        default="en",
        description="Language code for generated content (en, de, es, fr, it, pt, nl, pl, ru, zh, ja, ko, ar, hi)"
    )
    collection_filter: Optional[str] = None  # Single collection filter (legacy)
    collection_filters: Optional[List[str]] = None  # Multiple collection filters
    folder_id: Optional[str] = Field(default=None, description="Folder ID to scope query to")
    include_subfolders: bool = Field(default=True, description="Include documents in subfolders")
    metadata: Optional[dict] = None
    # Image generation - overrides admin setting if provided
    include_images: Optional[bool] = Field(
        default=None,
        description="Include images in generated document. If not set, uses admin setting."
    )
    # Auto chart/table generation - overrides admin setting if provided (PPTX only)
    auto_charts: Optional[bool] = Field(
        default=None,
        description="Auto-generate native charts and tables from numeric data. If not set, uses admin setting."
    )
    # Theme selection - defaults to business if not specified
    theme: Optional[str] = Field(
        default="business",
        description="Visual theme for the document. Options: business, creative, modern, nature"
    )
    # Enhanced theming options
    font_family: Optional[str] = Field(
        default=None,
        description="Font family key: modern, classic, professional, technical"
    )
    layout: Optional[str] = Field(
        default=None,
        description="Layout template key: standard, two_column, image_focused, minimal"
    )
    animations: Optional[bool] = Field(
        default=False,
        description="Enable slide animations (PPTX only)"
    )
    animation_speed: Optional[str] = Field(
        default="med",
        description="Animation speed: very_slow, slow, med, fast, very_fast, or 'custom' (PPTX only)"
    )
    animation_duration_ms: Optional[int] = Field(
        default=None,
        ge=200,
        le=5000,
        description="Custom animation duration in milliseconds (200-5000). Only used when animation_speed='custom'"
    )
    # Custom colors - override theme colors with user-selected hex values
    custom_colors: Optional[dict] = Field(
        default=None,
        description="Custom color overrides: {primary: '#hex', secondary: '#hex', accent: '#hex'}"
    )
    # Page/slide count control - None means auto (LLM decides)
    page_count: Optional[int] = Field(
        default=None,
        ge=1,
        le=20,
        description="Number of pages/slides (1-20). None = auto (LLM decides based on content)"
    )
    # Include sources page/slide/sheet
    include_sources: Optional[bool] = Field(
        default=None,
        description="Include sources page in generated document. If not set, uses admin setting."
    )
    # Style learning from existing documents
    use_existing_docs: bool = Field(
        default=False,
        description="Learn style and formatting from existing documents via RAG"
    )
    style_collection_filters: Optional[List[str]] = Field(
        default=None,
        description="Collections to learn style from (multi-select)"
    )
    style_folder_id: Optional[str] = Field(
        default=None,
        description="Folder to scope style learning"
    )
    include_style_subfolders: bool = Field(
        default=True,
        description="Include subfolders in style learning"
    )
    # AI Proofreading settings (uses CriticAgent for quality review)
    enable_critic_review: bool = Field(
        default=False,
        description="Enable AI proofreading - auto-reviews and fixes quality issues"
    )
    quality_threshold: float = Field(
        default=0.7,
        ge=0.6,
        le=0.9,
        description="Quality threshold (0.6-0.9). Content below this score will be auto-fixed."
    )
    fix_styling: bool = Field(
        default=True,
        description="Fix styling and formatting issues during proofreading"
    )
    fix_incomplete: bool = Field(
        default=True,
        description="Complete incomplete bullet points and sentences"
    )
    # Notes/Comments options (for PPTX speaker notes, DOCX comments, etc.)
    include_notes_explanation: bool = Field(
        default=False,
        description="Include AI explanation/reasoning in document notes (PPTX speaker notes)"
    )
    # Query Enhancement - controls query expansion and HyDE for source search
    enhance_query: Optional[bool] = Field(
        default=None,
        description="Enable query enhancement (expansion + HyDE) for source search. None = use admin default."
    )
    # PPTX Template - use an existing PPTX as a visual template
    template_pptx_id: Optional[str] = Field(
        default=None,
        description="ID of an uploaded PPTX template to use for styling/branding"
    )
    # Vision-based template analysis - uses vision LLM to analyze template slides
    enable_template_vision_analysis: Optional[bool] = Field(
        default=None,
        description="Enable vision-based template analysis (PPTX only). Uses vision LLM to analyze template slides for styling and layout. None = use system default."
    )
    template_vision_model: Optional[str] = Field(
        default=None,
        description="Vision model for template analysis. 'auto' uses the system default."
    )
    # Vision-based slide review - renders slides to images and uses vision LLM to check visual quality
    enable_vision_review: Optional[bool] = Field(
        default=None,
        description="Enable vision-based slide review (PPTX only). Renders slides to images and uses vision LLM to detect visual issues like text overflow, poor contrast, or layout problems. Resource-intensive - requires LibreOffice for rendering. None = use system default."
    )
    vision_review_model: Optional[str] = Field(
        default=None,
        description="Vision model for slide review. 'auto' uses the system default."
    )
    # LLM provider/model selection
    provider_id: Optional[str] = Field(
        default=None,
        description="LLM provider ID override for this job"
    )
    model: Optional[str] = Field(
        default=None,
        description="Model name override for this job"
    )
    # Phase 15 LLM Optimization - Advanced per-document overrides
    temperature_override: Optional[float] = Field(
        default=None,
        ge=0.0,
        le=1.0,
        description="Override content generation temperature for this document (0.0-1.0). None = use system default/optimized temperature."
    )
    # Dual Mode (RAG + General AI)
    dual_mode: bool = Field(
        default=False,
        description="Combine document knowledge with general AI for richer content"
    )
    dual_mode_blend: str = Field(
        default="merged",
        description="Dual mode blend strategy: merged or docs_first"
    )

    @property
    def effective_collection_filter(self) -> Optional[str]:
        """Get effective collection filter - prioritize collection_filters list."""
        if self.collection_filters:
            return ",".join(self.collection_filters)  # Join multiple as comma-separated
        return self.collection_filter


class OutlineModifications(BaseModel):
    """Modifications to apply to an outline."""
    title: Optional[str] = None
    sections: Optional[List[dict]] = None
    tone: Optional[str] = None
    theme: Optional[str] = None  # Allow changing theme after outline generation


class SectionFeedback(BaseModel):
    """Feedback for a section."""
    feedback: Optional[str] = None
    approved: bool = True


class SectionPlanApproval(BaseModel):
    """Approval/modification for a section plan (pre-generation review)."""
    section_id: str
    approved: bool = True  # False to skip this section
    title: Optional[str] = None  # Optional: edit the title
    description: Optional[str] = None  # Optional: edit the description


class ApproveSectionPlansRequest(BaseModel):
    """Request to approve section plans before generation."""
    section_approvals: Optional[List[SectionPlanApproval]] = None


class SourceReferenceResponse(BaseModel):
    """Source reference in response."""
    document_id: str
    document_name: str
    chunk_id: Optional[str]
    page_number: Optional[int]
    relevance_score: float
    snippet: str


class SectionResponse(BaseModel):
    """Section in response."""
    id: str
    title: str
    content: str
    order: int
    sources: List[SourceReferenceResponse]
    approved: bool
    feedback: Optional[str]
    # Pre-generation review fields
    description: Optional[str] = None
    generation_approved: bool = True
    skipped: bool = False
    # Generation metadata (content source, quality scores, etc.)
    metadata: Optional[dict] = None


class OutlineResponse(BaseModel):
    """Outline in response."""
    title: str
    description: str
    sections: List[dict]
    target_audience: Optional[str]
    tone: Optional[str]
    word_count_target: Optional[int]


class JobResponse(BaseModel):
    """Generation job response."""
    id: str
    user_id: str
    title: str
    description: str
    output_format: str
    status: str
    outline: Optional[OutlineResponse]
    sections: List[SectionResponse]
    sources_used_count: int
    created_at: datetime
    updated_at: datetime
    completed_at: Optional[datetime]
    has_output: bool
    error_message: Optional[str]
    # Generation settings used
    include_images: bool = False
    image_backend: Optional[str] = None


class JobListResponse(BaseModel):
    """List of generation jobs."""
    jobs: List[JobResponse]
    total: int


# =============================================================================
# Helper Functions
# =============================================================================

def job_to_response(job: GenerationJob) -> JobResponse:
    """Convert GenerationJob to response model."""
    outline_response = None
    if job.outline:
        outline_response = OutlineResponse(
            title=job.outline.title,
            description=job.outline.description,
            sections=job.outline.sections,
            target_audience=job.outline.target_audience,
            tone=job.outline.tone,
            word_count_target=job.outline.word_count_target,
        )

    sections_response = [
        SectionResponse(
            id=s.id,
            title=s.title,
            # Use rendered_content (what appears in output) for preview consistency
            content=s.rendered_content or s.revised_content or s.content,
            order=s.order,
            sources=[
                SourceReferenceResponse(
                    document_id=src.document_id,
                    document_name=src.document_name,
                    chunk_id=src.chunk_id,
                    page_number=src.page_number,
                    relevance_score=src.relevance_score,
                    snippet=src.snippet,
                )
                for src in s.sources
            ],
            approved=s.approved,
            feedback=s.feedback,
            # Pre-generation review fields
            description=getattr(s, 'description', None),
            generation_approved=getattr(s, 'generation_approved', True),
            skipped=getattr(s, 'skipped', False),
            metadata=getattr(s, 'metadata', None),
        )
        for s in job.sections
    ]

    return JobResponse(
        id=job.id,
        user_id=job.user_id,
        title=job.title,
        description=job.description,
        output_format=job.output_format.value,
        status=job.status.value,
        outline=outline_response,
        sections=sections_response,
        sources_used_count=len(job.sources_used),
        created_at=job.created_at,
        updated_at=job.updated_at,
        completed_at=job.completed_at,
        has_output=job.output_path is not None,
        error_message=job.error_message,
        # Image generation settings from job metadata
        include_images=job.metadata.get("include_images", False),
        image_backend=job.metadata.get("image_backend"),
    )


# =============================================================================
# API Endpoints
# =============================================================================

@router.post("/jobs", response_model=JobResponse, status_code=status.HTTP_201_CREATED)
async def create_generation_job(
    request: CreateJobRequest,
    user: AuthenticatedUser,
):
    """
    Create a new document generation job.

    This starts the human-in-the-loop workflow:
    1. Create job -> 2. Generate outline -> 3. Review/approve outline
    -> 4. Generate content -> 5. Review/approve sections -> 6. Download
    """
    logger.info(
        "Creating generation job",
        user_id=user.user_id,
        title=request.title,
        format=request.output_format,
    )

    # Validate output format
    try:
        output_format = OutputFormat(request.output_format.lower())
    except ValueError:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid output format: {request.output_format}. "
                   f"Valid formats: {[f.value for f in OutputFormat]}",
        )

    service = get_generation_service()

    # Build metadata with theme and page_count
    metadata = request.metadata or {}

    # Output language - store in metadata
    # "auto" means auto-detect from source documents
    valid_languages = {"auto", "en", "de", "es", "fr", "it", "pt", "nl", "pl", "ru", "zh", "ja", "ko", "ar", "hi"}
    if request.output_language and request.output_language not in valid_languages:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid output_language: {request.output_language}. Valid options: {valid_languages}",
        )
    metadata["output_language"] = request.output_language or "en"

    if request.theme:
        # Validate theme exists
        if request.theme not in THEMES:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid theme: {request.theme}. Valid themes: {list(THEMES.keys())}",
            )
        metadata["theme"] = request.theme

    # Enhanced theming options
    if request.font_family:
        if request.font_family not in FONT_FAMILIES:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid font_family: {request.font_family}. Valid options: {list(FONT_FAMILIES.keys())}",
            )
        metadata["font_family"] = request.font_family

    if request.layout:
        if request.layout not in LAYOUT_TEMPLATES:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid layout: {request.layout}. Valid options: {list(LAYOUT_TEMPLATES.keys())}",
            )
        metadata["layout"] = request.layout

    if request.animations is not None:
        metadata["animations"] = request.animations

    # Animation speed validation and storage
    if request.animation_speed:
        valid_speeds = {"very_slow", "slow", "med", "fast", "very_fast", "custom"}
        if request.animation_speed not in valid_speeds:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid animation_speed: {request.animation_speed}. Valid options: {valid_speeds}",
            )
        metadata["animation_speed"] = request.animation_speed

        # If custom speed, require animation_duration_ms
        if request.animation_speed == "custom":
            if not request.animation_duration_ms:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="animation_duration_ms is required when animation_speed is 'custom'",
                )
            metadata["animation_duration_ms"] = request.animation_duration_ms

    # Custom colors validation and storage
    if request.custom_colors:
        valid_color_keys = {"primary", "secondary", "accent", "text", "background"}
        import re
        hex_pattern = re.compile(r'^#[0-9A-Fa-f]{6}$')

        validated_colors = {}
        for key, value in request.custom_colors.items():
            if key not in valid_color_keys:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Invalid color key: {key}. Valid keys: {valid_color_keys}",
                )
            if not hex_pattern.match(value):
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Invalid hex color for {key}: {value}. Must be #RRGGBB format.",
                )
            validated_colors[key] = value

        metadata["custom_colors"] = validated_colors

    # Store page_count in metadata for outline generation
    if request.page_count is not None:
        metadata["page_count"] = request.page_count
        metadata["page_count_mode"] = "manual"
    else:
        metadata["page_count_mode"] = "auto"

    # Store include_sources preference (None means use admin setting)
    if request.include_sources is not None:
        metadata["include_sources"] = request.include_sources

    # Store auto_charts preference for PPTX (None means use admin setting)
    if request.auto_charts is not None:
        metadata["auto_charts"] = request.auto_charts

    # Store style learning settings if enabled
    if request.use_existing_docs:
        metadata["use_existing_docs"] = True
        metadata["style_collection_filters"] = request.style_collection_filters
        metadata["style_folder_id"] = request.style_folder_id
        metadata["include_style_subfolders"] = request.include_style_subfolders

    # Store AI proofreading settings
    if request.enable_critic_review:
        metadata["enable_critic_review"] = True
        metadata["quality_threshold"] = request.quality_threshold
        metadata["fix_styling"] = request.fix_styling
        metadata["fix_incomplete"] = request.fix_incomplete

    # Store notes/explanation settings
    if request.include_notes_explanation:
        metadata["include_notes_explanation"] = True

    # Store query enhancement preference (None means use admin setting)
    if request.enhance_query is not None:
        metadata["enhance_query"] = request.enhance_query

    # Store vision analysis preferences (for PPTX only) - per-document overrides
    if request.enable_template_vision_analysis is not None:
        metadata["enable_template_vision_analysis"] = request.enable_template_vision_analysis
    if request.template_vision_model:
        metadata["template_vision_model"] = request.template_vision_model
    if request.enable_vision_review is not None:
        metadata["enable_vision_review"] = request.enable_vision_review
    if request.vision_review_model:
        metadata["vision_review_model"] = request.vision_review_model

    # Store PPTX template if provided
    if request.template_pptx_id:
        # Look up template from uploaded documents
        from backend.db.models import Document
        from sqlalchemy import select

        template_path = None
        async with async_session_context() as session:
            query = select(Document).where(Document.id == request.template_pptx_id)
            result = await session.execute(query)
            doc = result.scalar_one_or_none()
            if doc and doc.file_path and os.path.exists(doc.file_path):
                template_path = doc.file_path

        if not template_path:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"PPTX template not found: {request.template_pptx_id}",
            )
        metadata["template_pptx_id"] = request.template_pptx_id
        metadata["template_pptx_path"] = template_path
        # Store original template filename for display in generated document metadata
        # Use original_filename (the actual uploaded name) over filename (UUID storage name)
        if doc:
            metadata["template_pptx_filename"] = doc.original_filename or doc.filename or os.path.basename(template_path)

    # Store LLM provider/model selection
    if request.provider_id:
        metadata["provider_id"] = request.provider_id
    if request.model:
        metadata["model"] = request.model

    # Store dual mode settings
    if request.dual_mode:
        metadata["dual_mode"] = True
        metadata["dual_mode_blend"] = request.dual_mode_blend

    # Store user email for notes/metadata display
    metadata["user_email"] = user.email

    job = await service.create_job(
        user_id=user.user_id,
        title=request.title,
        description=request.description,
        output_format=output_format,
        collection_filter=request.effective_collection_filter,
        folder_id=request.folder_id,
        include_subfolders=request.include_subfolders,
        metadata=metadata,
        include_images=request.include_images,  # Pass through to override admin setting
        temperature_override=request.temperature_override,  # Phase 15 LLM optimization
    )

    return job_to_response(job)


@router.get("/jobs", response_model=JobListResponse)
async def list_generation_jobs(
    user: AuthenticatedUser,
    status_filter: Optional[str] = Query(None, alias="status"),
):
    """
    List all generation jobs for the current user.
    """
    service = get_generation_service()

    # Parse status filter
    gen_status = None
    if status_filter:
        try:
            gen_status = GenerationStatus(status_filter)
        except ValueError:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid status: {status_filter}",
            )

    jobs = await service.list_jobs(
        user_id=user.user_id,
        status=gen_status,
    )

    return JobListResponse(
        jobs=[job_to_response(j) for j in jobs],
        total=len(jobs),
    )


@router.get("/jobs/{job_id}", response_model=JobResponse)
async def get_generation_job(
    job_id: str,
    user: AuthenticatedUser,
):
    """
    Get a specific generation job.
    """
    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id and not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    return job_to_response(job)


@router.post("/jobs/{job_id}/outline", response_model=JobResponse)
async def generate_outline(
    job_id: str,
    user: AuthenticatedUser,
    num_sections: Optional[int] = Query(default=None, ge=1, le=20),
):
    """
    Generate an outline for the document.

    The outline will be generated based on the job description and
    relevant sources from the knowledge base.

    num_sections priority:
    1. Query parameter (if provided) - highest priority, allows override
    2. Job metadata page_count (from CreateJobRequest)
    3. Auto mode (None) - LLM decides optimal count based on content
    """
    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    # Determine effective num_sections from priority order
    effective_num_sections = num_sections
    if effective_num_sections is None:
        # Fall back to job metadata page_count
        effective_num_sections = job.metadata.get("page_count")
    # If still None, auto mode - service will handle LLM-based determination

    logger.info(
        "Generating outline",
        job_id=job_id,
        num_sections=effective_num_sections,
        mode="auto" if effective_num_sections is None else "manual",
    )

    try:
        await service.generate_outline(job_id, num_sections=effective_num_sections)
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=str(e),
        )

    return job_to_response(job)


@router.post("/jobs/{job_id}/outline/approve", response_model=JobResponse)
async def approve_outline(
    job_id: str,
    user: AuthenticatedUser,
    modifications: Optional[OutlineModifications] = None,
):
    """
    Approve the outline to proceed with content generation.

    Optionally provide modifications to the outline before approval.
    """
    logger.info("Approving outline", job_id=job_id)

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    try:
        mods = modifications.model_dump(exclude_none=True) if modifications else None
        job = await service.approve_outline(job_id, modifications=mods)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=str(e),
        )

    return job_to_response(job)


@router.post("/jobs/{job_id}/sections/approve", response_model=JobResponse)
async def approve_section_plans(
    job_id: str,
    user: AuthenticatedUser,
    request: Optional[ApproveSectionPlansRequest] = None,
):
    """
    Approve section plans before content generation.

    This endpoint allows users to review and modify section plans before
    the actual content is generated. Users can:
    - Skip sections they don't want
    - Edit section titles and descriptions
    - Approve all sections to proceed with generation

    Requires status: sections_planning
    """
    logger.info("Approving section plans", job_id=job_id)

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    try:
        section_approvals = None
        if request and request.section_approvals:
            section_approvals = [
                {
                    "section_id": a.section_id,
                    "approved": a.approved,
                    "title": a.title,
                    "description": a.description,
                }
                for a in request.section_approvals
            ]
        job = await service.approve_section_plans(job_id, section_approvals=section_approvals)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=str(e),
        )

    return job_to_response(job)


@router.post("/jobs/{job_id}/generate", response_model=JobResponse)
async def generate_content(
    job_id: str,
    user: AuthenticatedUser,
):
    """
    Generate the document content.

    Requires an approved outline or approved section plans.
    """
    logger.info("Generating content", job_id=job_id)

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    try:
        job = await service.generate_content(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=str(e),
        )
    except Exception as e:
        logger.error("Content generation failed", job_id=job_id, error=str(e))
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Content generation failed",
        )

    return job_to_response(job)


@router.post("/jobs/{job_id}/sections/{section_id}/feedback", response_model=JobResponse)
async def provide_section_feedback(
    job_id: str,
    section_id: str,
    feedback: SectionFeedback,
    user: AuthenticatedUser,
):
    """
    Provide feedback on a section.

    If approved, marks the section as complete.
    If not approved, the section will need revision.
    """
    logger.info(
        "Section feedback",
        job_id=job_id,
        section_id=section_id,
        approved=feedback.approved,
    )

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    try:
        await service.approve_section(
            job_id=job_id,
            section_id=section_id,
            feedback=feedback.feedback,
            approved=feedback.approved,
        )
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=str(e),
        )

    return job_to_response(job)


@router.post("/jobs/{job_id}/sections/{section_id}/revise", response_model=JobResponse)
async def revise_section(
    job_id: str,
    section_id: str,
    user: AuthenticatedUser,
):
    """
    Revise a section based on provided feedback.
    """
    logger.info("Revising section", job_id=job_id, section_id=section_id)

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    try:
        await service.revise_section(job_id=job_id, section_id=section_id)
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=str(e),
        )

    return job_to_response(job)


@router.get("/jobs/{job_id}/download")
async def download_generated_document(
    job_id: str,
    user: AuthenticatedUser,
):
    """
    Download the generated document.

    Only available for completed jobs.
    """
    logger.info("Downloading document", job_id=job_id)

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id and not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    try:
        file_bytes, filename, content_type = await service.get_output_file(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=str(e),
        )

    return StreamingResponse(
        iter([file_bytes]),
        media_type=content_type,
        headers={
            "Content-Disposition": f'attachment; filename="{filename}"',
            "Content-Length": str(len(file_bytes)),
        },
    )


# =============================================================================
# Preview Endpoints
# =============================================================================

@router.get("/jobs/{job_id}/preview")
async def get_preview_metadata(
    job_id: str,
    user: AuthenticatedUser,
):
    """
    Get preview metadata for a generated document.

    Returns information about preview capabilities, page count, etc.
    """
    from backend.services.preview import preview_service

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id and not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    if job.status != "completed":
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Document not yet generated",
        )

    metadata = await preview_service.get_preview_metadata(job_id, job.output_format)
    return {"job_id": job_id, "format": job.output_format, **metadata}


@router.get("/jobs/{job_id}/preview/page/{page_num}")
async def get_preview_page(
    job_id: str,
    page_num: int,
    user: AuthenticatedUser,
):
    """
    Get a preview image for a specific page/slide.

    Returns a PNG image for PDF pages or PPTX slides.
    Returns HTML for DOCX.
    """
    from backend.services.preview import preview_service, PreviewError

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id and not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    if job.status != "completed":
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Document not yet generated",
        )

    format_lower = job.output_format.lower()

    try:
        if format_lower == "pdf":
            content, content_type = await preview_service.generate_pdf_preview(
                job_id, page=page_num
            )
            return StreamingResponse(
                iter([content]),
                media_type=content_type,
            )

        elif format_lower == "pptx":
            content, content_type = await preview_service.generate_pptx_slide_preview(
                job_id, slide_num=page_num
            )
            return StreamingResponse(
                iter([content]),
                media_type=content_type,
            )

        elif format_lower == "docx":
            html, content_type = await preview_service.generate_docx_preview(job_id)
            return StreamingResponse(
                iter([html.encode()]),
                media_type=content_type,
            )

        elif format_lower in ["md", "markdown", "html", "txt"]:
            content, content_type = await preview_service.generate_text_preview(
                job_id, format_lower
            )
            return StreamingResponse(
                iter([content.encode()]),
                media_type=content_type,
            )

        else:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Preview not supported for format: {format_lower}",
            )

    except PreviewError as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


@router.get("/jobs/{job_id}/preview/slides")
async def get_all_slide_previews(
    job_id: str,
    user: AuthenticatedUser,
):
    """
    Get preview images for all slides in a PPTX document.

    Returns a list of base64-encoded PNG images.
    """
    from backend.services.preview import preview_service, PreviewError

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id and not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    if job.status != "completed":
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Document not yet generated",
        )

    if job.output_format.lower() != "pptx":
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Slides preview only available for PPTX documents",
        )

    try:
        slides = await preview_service.generate_pptx_all_slides(job_id)
        return {"slides": slides, "count": len(slides)}
    except PreviewError as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


@router.get("/jobs/{job_id}/thumbnail")
async def get_thumbnail(
    job_id: str,
    user: AuthenticatedUser,
):
    """
    Get a thumbnail image for a generated document.

    Returns a small PNG thumbnail.
    """
    from backend.services.preview import preview_service, PreviewError

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id and not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    if job.status != "completed":
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Document not yet generated",
        )

    try:
        content, content_type = await preview_service.generate_thumbnail(
            job_id, job.output_format
        )
        return StreamingResponse(
            iter([content]),
            media_type=content_type,
        )
    except PreviewError as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


@router.delete("/jobs/{job_id}")
async def cancel_generation_job(
    job_id: str,
    user: AuthenticatedUser,
):
    """
    Cancel a generation job.
    """
    logger.info("Cancelling job", job_id=job_id)

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id and not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    try:
        await service.cancel_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=str(e),
        )

    return {"message": "Job cancelled", "job_id": job_id}


@router.get("/formats")
async def list_output_formats():
    """
    List all supported output formats.
    """
    return {
        "formats": [
            {"value": f.value, "name": f.name, "description": _get_format_description(f)}
            for f in OutputFormat
        ]
    }


def _get_format_description(format: OutputFormat) -> str:
    """Get description for output format."""
    descriptions = {
        OutputFormat.PPTX: "PowerPoint presentation (.pptx)",
        OutputFormat.DOCX: "Word document (.docx)",
        OutputFormat.PDF: "PDF document (.pdf)",
        OutputFormat.XLSX: "Excel spreadsheet (.xlsx)",
        OutputFormat.MARKDOWN: "Markdown text (.md)",
        OutputFormat.HTML: "HTML web page (.html)",
        OutputFormat.TXT: "Plain text (.txt)",
    }
    return descriptions.get(format, "")


# =============================================================================
# Image Generation Models
# =============================================================================

class ImageGenerateRequest(BaseModel):
    """Request to generate an image."""
    prompt: str = Field(..., min_length=5, max_length=1000)
    width: int = Field(default=800, ge=256, le=2048)
    height: int = Field(default=600, ge=256, le=2048)
    backend: Optional[str] = Field(None, description="Image generation backend: openai, stability, automatic1111, unsplash")
    model: Optional[str] = Field(None, description="Model to use (backend-specific, e.g., dall-e-3)")
    quality: Optional[str] = Field(None, description="Quality for DALL-E 3: standard or hd")
    style: Optional[str] = Field(None, description="Style for DALL-E 3: vivid or natural")


class ImageGenerateResponse(BaseModel):
    """Response from image generation."""
    success: bool
    path: Optional[str] = None
    image_base64: Optional[str] = None  # Base64-encoded image data
    width: int
    height: int
    prompt: str
    revised_prompt: Optional[str] = None  # DALL-E 3 may revise the prompt
    backend: str
    model: Optional[str] = None
    error: Optional[str] = None


class ImageConfigResponse(BaseModel):
    """Image generation configuration."""
    enabled: bool
    backend: str
    available_backends: List[str]
    sd_webui_available: bool


# =============================================================================
# Image Generation Endpoints
# =============================================================================

@router.post("/image", response_model=ImageGenerateResponse)
async def generate_image(
    request: ImageGenerateRequest,
    user: AuthenticatedUser,
):
    """
    Generate an image using AI.

    Supports multiple backends:
    - openai: OpenAI DALL-E 3/2 (requires OPENAI_API_KEY)
    - stability: Stability AI (requires STABILITY_API_KEY)
    - automatic1111: Local Stable Diffusion WebUI
    - unsplash: Placeholder images (free, always works)

    If no backend is specified, uses the configured default.
    """
    logger.info(
        "Generating image",
        user_id=user.user_id,
        prompt_preview=request.prompt[:50],
        backend=request.backend,
    )

    service = get_image_generator()

    # Check if image generation is enabled (unless explicitly requesting a backend)
    if not service.config.enabled and not request.backend:
        return ImageGenerateResponse(
            success=False,
            path=None,
            width=request.width,
            height=request.height,
            prompt=request.prompt,
            backend="disabled",
            error="Image generation is disabled. Enable with include_images=true in generation config, or specify a backend explicitly.",
        )

    try:
        import base64

        # Temporarily update config if backend/model specified
        original_backend = service.config.backend
        original_model = service.config.openai_model
        original_quality = service.config.openai_quality
        original_style = service.config.openai_style

        if request.backend:
            try:
                service.config.backend = ImageBackend(request.backend)
            except ValueError:
                return ImageGenerateResponse(
                    success=False,
                    path=None,
                    width=request.width,
                    height=request.height,
                    prompt=request.prompt,
                    backend=request.backend,
                    error=f"Invalid backend: {request.backend}. Valid options: openai, stability, automatic1111, unsplash",
                )

        if request.model:
            service.config.openai_model = request.model
        if request.quality:
            service.config.openai_quality = request.quality
        if request.style:
            service.config.openai_style = request.style

        # Generate image directly with the backend
        result = await service._generate_with_backend(
            prompt=request.prompt,
            width=request.width,
            height=request.height,
            backend=service.config.backend,
        )

        # Restore original config
        service.config.backend = original_backend
        service.config.openai_model = original_model
        service.config.openai_quality = original_quality
        service.config.openai_style = original_style

        if result and result.success:
            # Read image file and convert to base64
            image_base64 = None
            if result.path:
                try:
                    with open(result.path, "rb") as f:
                        image_base64 = base64.b64encode(f.read()).decode("utf-8")
                except Exception as e:
                    logger.warning(f"Could not read generated image file: {e}")

            return ImageGenerateResponse(
                success=True,
                path=result.path,
                image_base64=image_base64,
                width=result.width,
                height=result.height,
                prompt=result.prompt,
                revised_prompt=result.metadata.get("revised_prompt") if result.metadata else None,
                backend=result.backend.value,
                model=result.metadata.get("model") if result.metadata else None,
            )
        else:
            return ImageGenerateResponse(
                success=False,
                path=None,
                width=request.width,
                height=request.height,
                prompt=request.prompt,
                backend=result.backend.value if result else "unknown",
                error=result.error if result else "Unknown error",
            )

    except Exception as e:
        logger.error("Image generation failed", error=str(e))
        return ImageGenerateResponse(
            success=False,
            path=None,
            width=request.width,
            height=request.height,
            prompt=request.prompt,
            backend="error",
            error=str(e),
        )


@router.get("/image/config", response_model=ImageConfigResponse)
async def get_image_config(user: AuthenticatedUser):
    """
    Get image generation configuration and availability.
    """
    service = get_image_generator()

    # Check SD WebUI availability
    sd_available = await service.check_sd_webui_available()

    return ImageConfigResponse(
        enabled=service.config.enabled,
        backend=service.config.backend.value,
        available_backends=[b.value for b in ImageBackend],
        sd_webui_available=sd_available,
    )


@router.get("/image/backends")
async def get_image_backends(user: AuthenticatedUser):
    """
    Get list of available image generation backends with their configuration status.

    Returns which backends are configured and ready to use (have API keys, etc.)
    """
    service = get_image_generator()
    backends = service.get_available_backends()

    return {
        "backends": backends,
        "current_backend": service.config.backend.value,
        "enabled": service.config.enabled,
    }


# =============================================================================
# Theme Models and Endpoints
# =============================================================================

class ThemeInfo(BaseModel):
    """Information about a theme."""
    key: str
    name: str
    description: str
    primary: str
    secondary: str
    accent: str
    text: str


class ThemeSuggestionRequest(BaseModel):
    """Request for theme suggestions."""
    title: str = Field(..., min_length=1, max_length=200)
    description: str = Field(default="", description="Optional description for better suggestions")


class ThemeSuggestionResponse(BaseModel):
    """Response with theme suggestions."""
    recommended: str
    reason: str
    alternatives: List[str]
    themes: dict  # Full theme info for all themes
    # Enhanced suggestions - font, layout, animations
    font_family: Optional[str] = None
    layout: Optional[str] = None
    animations: Optional[bool] = None
    # Full details for UI display
    theme_details: Optional[dict] = None
    font_details: Optional[dict] = None
    layout_details: Optional[dict] = None
    available_fonts: Optional[dict] = None
    available_layouts: Optional[dict] = None


@router.get("/themes")
async def list_themes():
    """
    Get all available themes for document generation.

    Returns a list of themes with their colors and descriptions.
    """
    themes_list = []
    for key, theme in THEMES.items():
        themes_list.append(ThemeInfo(
            key=key,
            name=theme["name"],
            description=theme["description"],
            primary=theme["primary"],
            secondary=theme["secondary"],
            accent=theme["accent"],
            text=theme["text"],
        ))

    return {"themes": themes_list}


@router.post("/themes/suggest", response_model=ThemeSuggestionResponse)
async def suggest_theme(
    request: ThemeSuggestionRequest,
    user: AuthenticatedUser,
):
    """
    Get LLM-suggested themes, fonts, layouts based on document topic.

    The LLM analyzes the document title and description to recommend
    the most appropriate theme, font family, layout, and animations.
    User can still override with any settings.
    """
    logger.info(
        "Suggesting themes",
        user_id=user.user_id,
        title=request.title,
    )

    try:
        # Use the DocumentGenerationService for full theme suggestions
        service = get_generation_service()
        result = await service.suggest_theme(
            title=request.title,
            description=request.description,
            document_type="pptx",  # Default to PPTX for best suggestions
        )

        # Get alternatives (other themes not selected)
        alternatives = [k for k in THEMES.keys() if k != result.get("theme", "business")][:2]

        return ThemeSuggestionResponse(
            recommended=result.get("theme", "business"),
            reason=result.get("reason", "Recommended based on content analysis"),
            alternatives=alternatives,
            themes=THEMES,
            # Enhanced suggestions
            font_family=result.get("font_family"),
            layout=result.get("layout"),
            animations=result.get("animations"),
            # Full details for UI
            theme_details=result.get("theme_details"),
            font_details=result.get("font_details"),
            layout_details=result.get("layout_details"),
            available_fonts=FONT_FAMILIES,
            available_layouts=LAYOUT_TEMPLATES,
        )

    except Exception as e:
        logger.error("Theme suggestion failed", error=str(e))
        # Return default on error
        return ThemeSuggestionResponse(
            recommended="business",
            reason="Default professional theme (LLM unavailable)",
            alternatives=["modern", "creative"],
            themes=THEMES,
            font_family="modern",
            layout="standard",
            animations=False,
            theme_details=THEMES.get("business"),
            font_details=FONT_FAMILIES.get("modern"),
            layout_details=LAYOUT_TEMPLATES.get("standard"),
            available_fonts=FONT_FAMILIES,
            available_layouts=LAYOUT_TEMPLATES,
        )


# =============================================================================
# Spell Checking Models and Endpoints
# =============================================================================

class SpellCheckRequest(BaseModel):
    """Request to check spelling in text."""
    text: str = Field(..., min_length=1, max_length=50000)


class SpellIssue(BaseModel):
    """A spelling issue found in text."""
    word: str
    position: int
    suggestion: str
    context: str


class SpellCheckResponse(BaseModel):
    """Response from spell checking."""
    has_issues: bool
    issues: List[SpellIssue]


class SpellApplyRequest(BaseModel):
    """Request to apply spelling corrections."""
    job_id: str
    corrections: dict = Field(
        ...,
        description="Map of original word to corrected word"
    )


@router.post("/spell-check", response_model=SpellCheckResponse)
async def check_text_spelling(
    request: SpellCheckRequest,
    user: AuthenticatedUser,
):
    """
    Check spelling in provided text.

    Returns a list of potential spelling issues for user review.
    The user can then decide which corrections to apply.
    """
    logger.info(
        "Checking spelling",
        user_id=user.user_id,
        text_length=len(request.text),
    )

    result = check_spelling(request.text)

    return SpellCheckResponse(
        has_issues=result["has_issues"],
        issues=[
            SpellIssue(
                word=issue["word"],
                position=issue["position"],
                suggestion=issue["suggestion"],
                context=issue["context"],
            )
            for issue in result["issues"]
        ],
    )


@router.post("/jobs/{job_id}/spell-check", response_model=SpellCheckResponse)
async def check_job_spelling(
    job_id: str,
    user: AuthenticatedUser,
):
    """
    Check spelling in a generation job's content.

    Checks the title and all section titles and content for spelling issues.
    Returns a list of issues for user review before final document generation.
    """
    logger.info("Checking job spelling", job_id=job_id, user_id=user.user_id)

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id and not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    # Collect all text to check
    all_issues = []

    # Check title
    title_result = check_spelling(job.title)
    for issue in title_result["issues"]:
        issue["source"] = "title"
        all_issues.append(issue)

    # Check section titles and content
    for section in job.sections:
        # Check section title
        title_result = check_spelling(section.title)
        for issue in title_result["issues"]:
            issue["source"] = f"section_{section.id}_title"
            all_issues.append(issue)

        # Check section content
        content = section.revised_content or section.content
        content_result = check_spelling(content)
        for issue in content_result["issues"]:
            issue["source"] = f"section_{section.id}_content"
            all_issues.append(issue)

    # Limit total issues
    all_issues = all_issues[:50]

    return SpellCheckResponse(
        has_issues=len(all_issues) > 0,
        issues=[
            SpellIssue(
                word=issue["word"],
                position=issue["position"],
                suggestion=issue["suggestion"],
                context=issue["context"],
            )
            for issue in all_issues
        ],
    )


# =============================================================================
# PPTX Template Upload
# =============================================================================

class TemplateUploadResponse(BaseModel):
    """Response for template upload."""
    template_id: str
    filename: str
    slide_count: int
    master_layouts: List[str]
    theme_colors: Optional[dict] = None
    message: str


@router.post("/templates/pptx", response_model=TemplateUploadResponse)
async def upload_pptx_template(
    user: AuthenticatedUser,
    file: UploadFile = File(...),
) -> TemplateUploadResponse:
    """
    Upload a PPTX file to use as a visual template.

    The uploaded PPTX's slide master, colors, fonts, and layouts will be
    extracted and can be used when generating new presentations.
    """
    # Validate file type
    if not file.filename or not file.filename.lower().endswith('.pptx'):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Only .pptx files are supported as templates"
        )

    # Create templates directory
    templates_dir = os.path.join("storage", "templates", "pptx")
    os.makedirs(templates_dir, exist_ok=True)

    # Generate unique ID for template
    import uuid
    template_id = str(uuid.uuid4())
    template_path = os.path.join(templates_dir, f"{template_id}.pptx")

    # Save the file
    try:
        with open(template_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
    except Exception as e:
        logger.error("Failed to save template file", error=str(e))
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to save template file"
        )

    # Analyze the template
    try:
        from pptx import Presentation
        prs = Presentation(template_path)

        # Get slide count
        slide_count = len(prs.slides)

        # Get master layouts
        master_layouts = []
        for slide_master in prs.slide_masters:
            for layout in slide_master.slide_layouts:
                master_layouts.append(layout.name)

        # Extract theme colors if available
        theme_colors = None
        try:
            if prs.slide_masters and len(prs.slide_masters) > 0:
                master = prs.slide_masters[0]
                theme = master.theme_color_scheme if hasattr(master, 'theme_color_scheme') else None
                if theme:
                    theme_colors = {
                        "primary": str(theme.accent1) if hasattr(theme, 'accent1') else None,
                        "secondary": str(theme.accent2) if hasattr(theme, 'accent2') else None,
                    }
        except Exception as e:
            logger.debug("Theme extraction failed (optional)", error=str(e))

        logger.info(
            "PPTX template uploaded",
            template_id=template_id,
            filename=file.filename,
            slide_count=slide_count,
            layouts=len(master_layouts),
        )

        return TemplateUploadResponse(
            template_id=template_id,
            filename=file.filename or "template.pptx",
            slide_count=slide_count,
            master_layouts=master_layouts[:10],  # Limit to 10 layouts
            theme_colors=theme_colors,
            message=f"Template uploaded successfully with {slide_count} slides and {len(master_layouts)} layouts"
        )

    except Exception as e:
        # Clean up on failure
        if os.path.exists(template_path):
            os.remove(template_path)
        logger.error("Failed to analyze template", error=str(e))
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to analyze PPTX template: {str(e)}"
        )


@router.get("/templates/pptx")
async def list_pptx_templates(
    user: AuthenticatedUser,
) -> List[dict]:
    """List all uploaded PPTX documents that can be used as templates."""
    from backend.db.models import Document
    from sqlalchemy import select, or_

    templates = []

    # Query uploaded PPTX documents from database
    async with async_session_context() as session:
        query = select(Document).where(
            or_(
                Document.file_type == "pptx",
                Document.filename.ilike("%.pptx")
            )
        ).order_by(Document.created_at.desc()).limit(50)

        result = await session.execute(query)
        docs = result.scalars().all()

        for doc in docs:
            # Try to get slide count from the actual file
            slide_count = 0
            file_path = doc.file_path

            if file_path and os.path.exists(file_path):
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    slide_count = len(prs.slides)
                except Exception as e:
                    logger.warning(f"Failed to analyze PPTX {doc.filename}: {e}")

            # Get original filename from original_filename field or filename field
            original_name = doc.original_filename or doc.filename

            templates.append({
                "template_id": str(doc.id),
                "filename": original_name,
                "slide_count": slide_count,
                "created_at": doc.created_at.isoformat() if doc.created_at else datetime.now().isoformat(),
            })

    return templates


@router.delete("/templates/pptx/{template_id}")
async def delete_pptx_template(
    template_id: str,
    user: AuthenticatedUser,
) -> dict:
    """Delete a PPTX template."""
    template_path = os.path.join("storage", "templates", "pptx", f"{template_id}.pptx")

    if not os.path.exists(template_path):
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Template not found"
        )

    os.remove(template_path)
    return {"message": "Template deleted successfully"}


class SuggestTemplateRequest(BaseModel):
    """Request for AI template suggestions."""
    title: str = Field(..., description="Document title")
    description: str = Field(..., description="Document description/topic")


class TemplateSuggestion(BaseModel):
    """AI-suggested template with score and reason."""
    template_id: str
    filename: str
    slide_count: int
    score: float = Field(..., ge=0, le=100, description="Relevance score 0-100")
    reason: str = Field(..., description="Why this template fits")


class SuggestTemplateResponse(BaseModel):
    """Response with ranked template suggestions."""
    suggestions: List[TemplateSuggestion]
    message: str


@router.post("/templates/pptx/suggest", response_model=SuggestTemplateResponse)
async def suggest_pptx_templates(
    request: SuggestTemplateRequest,
    user: AuthenticatedUser,
) -> SuggestTemplateResponse:
    """
    Get AI-powered template suggestions based on document title and description.

    Analyzes uploaded PPTX documents and scores them for relevance to the
    new document topic, returning a ranked list with explanations.
    """
    from backend.services.llm import EnhancedLLMFactory
    from backend.db.models import Document
    from sqlalchemy import select, or_

    # Query uploaded PPTX documents from database
    templates = []
    async with async_session_context() as session:
        # Get PPTX documents that can serve as templates
        query = select(Document).where(
            or_(
                Document.file_type == "pptx",
                Document.filename.ilike("%.pptx")
            )
        ).order_by(Document.created_at.desc()).limit(20)

        result = await session.execute(query)
        docs = result.scalars().all()

        for doc in docs:
            # Try to get slide count from the actual file
            slide_count = 0
            layout_count = 0
            file_path = doc.file_path

            if file_path and os.path.exists(file_path):
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    slide_count = len(prs.slides)
                    layout_count = len(prs.slide_layouts)
                except Exception as e:
                    logger.warning(f"Failed to analyze PPTX {doc.filename}: {e}")

            # Get original filename from original_filename field or filename field
            original_name = doc.original_filename or doc.filename

            templates.append({
                "template_id": str(doc.id),
                "filename": original_name,
                "file_path": file_path,
                "slide_count": slide_count,
                "layout_count": layout_count,
                "title": doc.title or original_name,
            })

    if not templates:
        return SuggestTemplateResponse(
            suggestions=[],
            message="No PPTX documents found. Upload some presentations first to use as templates."
        )

    # Use LLM to score templates based on their titles/names
    llm, _ = await EnhancedLLMFactory.get_chat_model_for_operation(
        operation="content_generation",
        user_id=None,
    )

    # Build prompt for template scoring - include filenames which often describe the content
    template_info = "\n".join([
        f"- Template {i+1}: \"{t['title']}\" ({t['slide_count']} slides)"
        for i, t in enumerate(templates)
    ])

    prompt = f"""Analyze these existing PPTX presentations and score them as templates for a NEW document:

NEW Document Title: {request.title}
NEW Document Description: {request.description}

Available Presentations (to use as style/design templates):
{template_info}

Score each presentation as a potential TEMPLATE based on:
1. Topic similarity - Does the presentation topic match the new document?
2. Professional context - Is it appropriate for the same audience?
3. Design suitability - Would this presentation's style fit the new content?

Respond in JSON format:
{{
  "rankings": [
    {{"template_index": 0, "score": 85, "reason": "Similar business topic, professional style"}},
    {{"template_index": 1, "score": 45, "reason": "Different industry but good design"}}
  ]
}}

Include templates scoring above 30. Rank by score descending."""

    try:
        response = await llm.ainvoke(prompt)
        response_text = response.content if hasattr(response, 'content') else str(response)

        # Parse JSON from response
        import json
        import re

        # Extract JSON from response
        json_match = re.search(r'\{[\s\S]*\}', response_text)
        if json_match:
            rankings_data = json.loads(json_match.group())
            rankings = rankings_data.get("rankings", [])
        else:
            rankings = []

        # Build suggestions from rankings
        suggestions = []
        for ranking in rankings:
            idx = ranking.get("template_index", 0)
            if 0 <= idx < len(templates):
                template = templates[idx]
                suggestions.append(TemplateSuggestion(
                    template_id=template["template_id"],
                    filename=template["filename"],
                    slide_count=template["slide_count"],
                    score=min(100, max(0, ranking.get("score", 50))),
                    reason=ranking.get("reason", "AI-suggested template")
                ))

        # Sort by score descending
        suggestions.sort(key=lambda x: x.score, reverse=True)

        return SuggestTemplateResponse(
            suggestions=suggestions[:5],  # Top 5
            message=f"Found {len(suggestions)} matching templates from your uploaded presentations"
        )

    except Exception as e:
        logger.warning(f"AI template suggestion failed: {e}")
        # Fallback: return all templates with default score
        fallback_suggestions = [
            TemplateSuggestion(
                template_id=t["template_id"],
                filename=t["filename"],
                slide_count=t["slide_count"],
                score=50.0,
                reason="Available as template (AI scoring unavailable)"
            )
            for t in templates[:5]
        ]
        return SuggestTemplateResponse(
            suggestions=fallback_suggestions,
            message="AI scoring unavailable, showing uploaded presentations"
        )


# =============================================================================
# Document Verification and Repair Endpoints
# =============================================================================

class VerificationIssue(BaseModel):
    """A verification issue found in the document."""
    type: str
    severity: str  # error, warning, info
    message: str
    location: Optional[str] = None


class VerificationResult(BaseModel):
    """Document verification result."""
    passed: bool
    issues: List[VerificationIssue]
    checks_performed: List[str]
    document_stats: dict


class RepairResult(BaseModel):
    """Document repair result."""
    success: bool
    repairs_made: List[str]
    issues_remaining: List[VerificationIssue]
    message: str


class VerificationActionRequest(BaseModel):
    """Request to handle verification issues."""
    action: str = Field(
        ...,
        description="Action to take: 'proceed' (keep as-is), 'auto_repair' (let AI fix), 'regenerate' (start over)"
    )


@router.post("/jobs/{job_id}/verify", response_model=VerificationResult)
async def verify_document(
    job_id: str,
    user: AuthenticatedUser,
):
    """
    Verify a generated document for issues.

    Checks for:
    - Empty slides/sections
    - Missing content
    - Overflow issues
    - Format-specific problems

    Returns detailed verification results.
    """
    logger.info("Verifying document", job_id=job_id)

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id and not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    if job.status != GenerationStatus.COMPLETED:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Document must be completed before verification",
        )

    if not job.output_path:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No output file found for this job",
        )

    try:
        result = await service.verify_generated_document(
            file_path=job.output_path,
            output_format=job.output_format,
            expected_sections=len(job.sections),
            job_title=job.title,
        )

        return VerificationResult(
            passed=result["passed"],
            issues=[
                VerificationIssue(
                    type=issue.get("type", "unknown"),
                    severity=issue.get("severity", "warning"),
                    message=issue.get("message", "Unknown issue"),
                    location=issue.get("location"),
                )
                for issue in result.get("issues", [])
            ],
            checks_performed=result.get("checks_performed", []),
            document_stats=result.get("document_stats", {}),
        )

    except Exception as e:
        logger.error("Document verification failed", job_id=job_id, error=str(e))
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Verification failed: {str(e)}",
        )


@router.post("/jobs/{job_id}/repair", response_model=RepairResult)
async def repair_document(
    job_id: str,
    user: AuthenticatedUser,
):
    """
    Attempt to auto-repair issues found in a generated document.

    Repairs include:
    - Removing empty slides/sections
    - Adding missing speaker notes
    - Fixing formatting issues
    - Regenerating problematic sections

    Run /verify first to see what issues exist.
    """
    logger.info("Repairing document", job_id=job_id)

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id and not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    if job.status != GenerationStatus.COMPLETED:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Document must be completed before repair",
        )

    if not job.output_path:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No output file found for this job",
        )

    try:
        # First verify to get current issues
        verification_result = await service.verify_generated_document(
            file_path=job.output_path,
            output_format=job.output_format,
            expected_sections=len(job.sections),
            job_title=job.title,
        )

        if verification_result["passed"]:
            return RepairResult(
                success=True,
                repairs_made=[],
                issues_remaining=[],
                message="Document has no issues - no repairs needed",
            )

        # Attempt repair
        repair_result = await service.repair_document(
            file_path=job.output_path,
            output_format=job.output_format,
            verification_result=verification_result,
            job=job,
        )

        # Convert remaining issues
        remaining_issues = [
            VerificationIssue(
                type=issue.get("type", "unknown"),
                severity=issue.get("severity", "warning"),
                message=issue.get("message", "Unknown issue"),
                location=issue.get("location"),
            )
            for issue in repair_result.get("issues_remaining", [])
        ]

        return RepairResult(
            success=repair_result.get("success", False),
            repairs_made=repair_result.get("repairs_made", []),
            issues_remaining=remaining_issues,
            message=repair_result.get("message", "Repair completed"),
        )

    except Exception as e:
        logger.error("Document repair failed", job_id=job_id, error=str(e))
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Repair failed: {str(e)}",
        )


@router.post("/jobs/{job_id}/verification-action", response_model=JobResponse)
async def handle_verification_action(
    job_id: str,
    request: VerificationActionRequest,
    user: AuthenticatedUser,
):
    """
    Handle user's decision on verification issues.

    Actions:
    - proceed: Accept document as-is despite issues
    - auto_repair: Let AI automatically fix issues
    - regenerate: Mark job for regeneration

    This is the human-in-the-loop step after verification.
    """
    logger.info(
        "Handling verification action",
        job_id=job_id,
        action=request.action,
    )

    valid_actions = {"proceed", "auto_repair", "regenerate"}
    if request.action not in valid_actions:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid action: {request.action}. Valid options: {valid_actions}",
        )

    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id and not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    try:
        result = await service.handle_verification_action(
            job=job,
            action=request.action,
        )

        # Refresh job state
        job = await service.get_job(job_id)
        return job_to_response(job)

    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=str(e),
        )
    except Exception as e:
        logger.error("Verification action failed", job_id=job_id, error=str(e))
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Action failed: {str(e)}",
        )


@router.get("/jobs/{job_id}/verification-status")
async def get_verification_status(
    job_id: str,
    user: AuthenticatedUser,
):
    """
    Get the verification status for a completed document.

    Returns whether verification has been run and what the results were.
    """
    service = get_generation_service()

    try:
        job = await service.get_job(job_id)
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e),
        )

    # Verify ownership
    if job.user_id != user.user_id and not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this job",
        )

    # Check metadata for verification info
    verification_result = job.metadata.get("verification_result")
    requires_review = job.metadata.get("requires_manual_review", False)
    repair_mode = job.metadata.get("verification_repair_mode", "skip")

    return {
        "job_id": job_id,
        "status": job.status.value,
        "verified": verification_result is not None,
        "verification_passed": verification_result.get("passed") if verification_result else None,
        "requires_manual_review": requires_review,
        "repair_mode": repair_mode,
        "review_options": job.metadata.get("review_options") if requires_review else None,
    }

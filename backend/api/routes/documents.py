"""
AIDocumentIndexer - Documents API Routes
========================================

Endpoints for document management, listing, and retrieval.
All endpoints enforce permission checking based on user access tier.
"""

from datetime import datetime
from typing import Optional, List, Dict
from uuid import UUID
import io
import zipfile
import tempfile
import os

from fastapi import APIRouter, Depends, HTTPException, Query, Request, status, BackgroundTasks
from fastapi.responses import FileResponse, Response, StreamingResponse
from pathlib import Path
from pydantic import BaseModel, Field
from sqlalchemy import select, func, and_, or_, desc, asc, delete
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm import selectinload
import structlog

from backend.db.database import get_async_session
from backend.db.models import Document, Chunk, AccessTier, ProcessingStatus
from backend.api.middleware.auth import (
    get_user_context,
    require_admin,
    AuthenticatedUser,
    AdminUser,
    get_org_filter,
)
from backend.services.permissions import (
    UserContext,
    Permission,
    get_permission_service,
)
from backend.services.vectorstore import get_vector_store, SearchType
from backend.services.audit import AuditAction, get_audit_service
from backend.services.auto_tagger import AutoTaggerService

logger = structlog.get_logger(__name__)

router = APIRouter()


# =============================================================================
# Pydantic Models
# =============================================================================

class DocumentBase(BaseModel):
    """Base document model."""
    name: str = Field(..., min_length=1, max_length=500)
    file_type: str
    collection: Optional[str] = Field(None, max_length=100, pattern=r"^[a-zA-Z0-9_\-\s\.]+$")
    access_tier: int = Field(default=1, ge=1, le=100)


class DocumentCreate(DocumentBase):
    """Document creation request."""
    pass


class DocumentUpdate(BaseModel):
    """Document update request."""
    name: Optional[str] = Field(None, min_length=1, max_length=500)
    collection: Optional[str] = Field(None, max_length=100, pattern=r"^[a-zA-Z0-9_\-\s\.]+$")
    access_tier: Optional[int] = Field(default=None, ge=1, le=100)
    tags: Optional[List[str]] = None


class EnhancedMetadataResponse(BaseModel):
    """Enhanced metadata response model."""
    summary_short: Optional[str] = None
    summary_detailed: Optional[str] = None
    keywords: Optional[List[str]] = None
    topics: Optional[List[str]] = None
    entities: Optional[Dict[str, List[str]]] = None
    hypothetical_questions: Optional[List[str]] = None
    language: Optional[str] = None
    document_type: Optional[str] = None
    enhanced_at: Optional[str] = None
    model_used: Optional[str] = None


class DocumentResponse(BaseModel):
    """Document response model."""
    id: UUID
    name: str
    file_type: str
    file_size: int
    file_hash: str
    collection: Optional[str] = None
    access_tier: int
    access_tier_name: str
    status: str
    chunk_count: int
    word_count: Optional[int] = None
    page_count: Optional[int] = None
    created_at: datetime
    updated_at: datetime
    uploaded_by: Optional[UUID] = None
    tags: Optional[List[str]] = None
    enhanced_metadata: Optional[EnhancedMetadataResponse] = None
    is_enhanced: bool = False
    # Embedding status
    embedding_count: int = 0
    embedding_coverage: float = 0.0  # Percentage of chunks with embeddings
    has_all_embeddings: bool = False

    class Config:
        from_attributes = True


class DocumentListResponse(BaseModel):
    """Paginated document list response."""
    documents: List[DocumentResponse]
    total: int
    page: int
    page_size: int
    has_more: bool


class DocumentSearchRequest(BaseModel):
    """Document search request."""
    query: str
    collection: Optional[str] = None
    file_types: Optional[List[str]] = None
    min_tier: Optional[int] = None
    max_tier: Optional[int] = None
    limit: int = Field(default=20, ge=1, le=100)


class SearchResultItem(BaseModel):
    """Individual search result."""
    document_id: str
    document_name: str
    chunk_id: str
    content_preview: str
    relevance_score: float
    page_number: Optional[int] = None


class SearchResultResponse(BaseModel):
    """Search results response."""
    results: List[SearchResultItem]
    total_results: int
    query: str


class ChunkResponse(BaseModel):
    """Document chunk response."""
    id: UUID
    content: str
    chunk_index: int
    page_number: Optional[int]
    metadata: dict

    class Config:
        from_attributes = True


class CollectionInfo(BaseModel):
    """Collection summary information."""
    name: str
    document_count: int


class AutoTagRequest(BaseModel):
    """Auto-tag request for a single document."""
    max_tags: int = Field(default=5, ge=1, le=10, description="Maximum tags to generate")


class AutoTagResponse(BaseModel):
    """Auto-tag response."""
    document_id: str
    tags: List[str]
    collection: Optional[str] = None  # First tag as primary collection


class BulkAutoTagRequest(BaseModel):
    """Bulk auto-tag request."""
    document_ids: List[UUID]
    max_tags: int = Field(default=5, ge=1, le=10)


class BulkDownloadRequest(BaseModel):
    """Bulk download request."""
    document_ids: List[UUID] = Field(..., min_length=1, max_length=100)


class BulkAutoTagResponse(BaseModel):
    """Bulk auto-tag response."""
    status: str
    total: int
    processed: int = 0
    results: Optional[List[AutoTagResponse]] = None


class CollectionsResponse(BaseModel):
    """Collections list response."""
    collections: List[CollectionInfo]
    total_documents: int = 0  # Actual unique document count


# =============================================================================
# Helper Functions
# =============================================================================

def get_client_ip(request: Request) -> Optional[str]:
    """Extract client IP from request."""
    forwarded = request.headers.get("X-Forwarded-For")
    if forwarded:
        return forwarded.split(",")[0].strip()
    return request.client.host if request.client else None


def document_to_response(
    doc: Document,
    chunk_count: int = 0,
    embedding_count: int = 0,
) -> DocumentResponse:
    """Convert Document model to response."""
    # Parse enhanced metadata if present
    enhanced_metadata = None
    is_enhanced = False
    if doc.enhanced_metadata:
        is_enhanced = True
        enhanced_metadata = EnhancedMetadataResponse(
            summary_short=doc.enhanced_metadata.get("summary_short"),
            summary_detailed=doc.enhanced_metadata.get("summary_detailed"),
            keywords=doc.enhanced_metadata.get("keywords"),
            topics=doc.enhanced_metadata.get("topics"),
            entities=doc.enhanced_metadata.get("entities"),
            hypothetical_questions=doc.enhanced_metadata.get("hypothetical_questions"),
            language=doc.enhanced_metadata.get("language"),
            document_type=doc.enhanced_metadata.get("document_type"),
            enhanced_at=doc.enhanced_metadata.get("enhanced_at"),
            model_used=doc.enhanced_metadata.get("model_used"),
        )

    # Calculate embedding coverage
    embedding_coverage = (embedding_count / chunk_count * 100) if chunk_count > 0 else 0.0
    has_all_embeddings = embedding_count >= chunk_count if chunk_count > 0 else True

    return DocumentResponse(
        id=doc.id,
        name=doc.title or doc.original_filename,
        file_type=doc.file_type,
        file_size=doc.file_size,
        file_hash=doc.file_hash,
        collection=doc.tags[0] if doc.tags else None,  # Use first tag as collection
        access_tier=doc.access_tier.level if doc.access_tier else 1,
        access_tier_name=doc.access_tier.name if doc.access_tier else "Unknown",
        status=doc.processing_status.value,
        chunk_count=chunk_count,
        word_count=doc.word_count,
        page_count=doc.page_count,
        created_at=doc.created_at,
        updated_at=doc.updated_at,
        uploaded_by=doc.uploaded_by_id,
        tags=doc.tags,
        enhanced_metadata=enhanced_metadata,
        is_enhanced=is_enhanced,
        embedding_count=embedding_count,
        embedding_coverage=round(embedding_coverage, 1),
        has_all_embeddings=has_all_embeddings,
    )


# =============================================================================
# API Endpoints
# =============================================================================

@router.get("", response_model=DocumentListResponse)
async def list_documents(
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
    page: int = Query(default=1, ge=1),
    page_size: int = Query(default=20, ge=1, le=100),
    collection: Optional[str] = None,
    file_type: Optional[str] = None,
    status: Optional[str] = None,
    sort_by: str = Query(default="created_at"),
    sort_order: str = Query(default="desc", pattern="^(asc|desc)$"),
    # Date range filtering
    created_after: Optional[datetime] = Query(default=None, description="Filter documents created after this date (ISO 8601)"),
    created_before: Optional[datetime] = Query(default=None, description="Filter documents created before this date (ISO 8601)"),
    updated_after: Optional[datetime] = Query(default=None, description="Filter documents updated after this date (ISO 8601)"),
    updated_before: Optional[datetime] = Query(default=None, description="Filter documents updated before this date (ISO 8601)"),
    # Folder filtering
    folder_id: Optional[str] = Query(default=None, description="Filter by folder ID"),
    include_subfolders: bool = Query(default=True, description="Include documents from subfolders"),
):
    """
    List documents with pagination and filtering.

    Results are automatically filtered based on the user's access tier.
    Users can only see documents at or below their access tier level.
    """
    logger.info(
        "Listing documents",
        user_id=user.user_id,
        access_tier=user.access_tier_level,
        page=page,
        page_size=page_size,
        collection=collection,
    )

    # Build base query with access tier filtering
    # Exclude soft-deleted documents (those marked as failed with "Deleted by user" error)
    base_query = (
        select(Document)
        .join(AccessTier, Document.access_tier_id == AccessTier.id)
        .where(AccessTier.level <= user.access_tier_level)
        .where(
            ~(
                (Document.processing_status == ProcessingStatus.FAILED)
                & (Document.processing_error == "Deleted by user")
            )
        )
    )

    # PHASE 12 FIX: Apply organization filtering for multi-tenant isolation
    # Users can only see documents from their organization (or shared docs without org)
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        base_query = base_query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),  # Include legacy/shared docs
            )
        )

    # SECURITY FIX: Filter private documents - only owner or superadmin can see
    if not user.is_superadmin:
        base_query = base_query.where(
            or_(
                Document.is_private == False,
                and_(
                    Document.is_private == True,
                    Document.uploaded_by_id == UUID(user.user_id),
                ),
            )
        )

    # Apply filters
    if collection:
        base_query = base_query.where(Document.tags.contains([collection]))

    if file_type:
        base_query = base_query.where(Document.file_type == file_type)

    if status:
        try:
            status_enum = ProcessingStatus(status)
            base_query = base_query.where(Document.processing_status == status_enum)
        except ValueError:
            pass  # Invalid status, ignore filter

    # Apply date range filters
    if created_after:
        base_query = base_query.where(Document.created_at >= created_after)
    if created_before:
        base_query = base_query.where(Document.created_at <= created_before)
    if updated_after:
        base_query = base_query.where(Document.updated_at >= updated_after)
    if updated_before:
        base_query = base_query.where(Document.updated_at <= updated_before)

    # Apply folder filter
    if folder_id:
        import uuid as uuid_module
        try:
            folder_uuid = uuid_module.UUID(folder_id)
            if include_subfolders:
                # Get all documents in folder and subfolders
                from backend.services.folder_service import get_folder_service
                folder_service = get_folder_service()
                folder_doc_ids = await folder_service.get_folder_document_ids(
                    folder_id=folder_id,
                    include_subfolders=True,
                    user_tier_level=user.access_tier_level,
                )
                if folder_doc_ids:
                    doc_uuids = [uuid_module.UUID(d) for d in folder_doc_ids]
                    base_query = base_query.where(Document.id.in_(doc_uuids))
                else:
                    # No documents in folder, return empty
                    base_query = base_query.where(Document.id == None)
            else:
                # Only direct children
                base_query = base_query.where(Document.folder_id == folder_uuid)
        except ValueError:
            pass  # Invalid UUID, ignore filter

    # Count total before pagination
    count_query = select(func.count()).select_from(base_query.subquery())
    total_result = await db.execute(count_query)
    total = total_result.scalar() or 0

    # Apply sorting
    sort_column = getattr(Document, sort_by, Document.created_at)
    if sort_order == "desc":
        base_query = base_query.order_by(desc(sort_column))
    else:
        base_query = base_query.order_by(asc(sort_column))

    # Apply pagination
    offset = (page - 1) * page_size
    base_query = base_query.offset(offset).limit(page_size)

    # Include access tier relationship
    base_query = base_query.options(selectinload(Document.access_tier))

    # Execute query
    result = await db.execute(base_query)
    documents = result.scalars().all()

    # Get chunk counts and embedding counts for each document
    doc_ids = [doc.id for doc in documents]
    if doc_ids:
        # Query chunk counts
        chunk_counts_query = (
            select(Chunk.document_id, func.count(Chunk.id))
            .where(Chunk.document_id.in_(doc_ids))
            .group_by(Chunk.document_id)
        )
        chunk_result = await db.execute(chunk_counts_query)
        chunk_counts = {row[0]: row[1] for row in chunk_result.all()}

        # Query embedding counts (chunks stored in vector database)
        # Uses has_embedding flag instead of checking actual embedding column (saves storage)
        embedding_counts_query = (
            select(Chunk.document_id, func.count(Chunk.id))
            .where(
                Chunk.document_id.in_(doc_ids),
                Chunk.has_embedding == True,
            )
            .group_by(Chunk.document_id)
        )
        embedding_result = await db.execute(embedding_counts_query)
        embedding_counts = {row[0]: row[1] for row in embedding_result.all()}
    else:
        chunk_counts = {}
        embedding_counts = {}

    # Build response
    doc_responses = [
        document_to_response(
            doc,
            chunk_counts.get(doc.id, 0),
            embedding_counts.get(doc.id, 0),
        )
        for doc in documents
    ]

    return DocumentListResponse(
        documents=doc_responses,
        total=total,
        page=page,
        page_size=page_size,
        has_more=(offset + len(documents)) < total,
    )


@router.get("/{document_id}", response_model=DocumentResponse)
async def get_document(
    document_id: UUID,
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
):
    """
    Get a specific document by ID.

    Returns 404 if document doesn't exist or user doesn't have access.
    """
    logger.info(
        "Getting document",
        document_id=str(document_id),
        user_id=user.user_id,
    )

    # Query with access tier check and organization filtering
    query = (
        select(Document)
        .join(AccessTier, Document.access_tier_id == AccessTier.id)
        .where(
            and_(
                Document.id == document_id,
                AccessTier.level <= user.access_tier_level,
            )
        )
        .options(selectinload(Document.access_tier))
    )

    # PHASE 12 FIX: Apply organization filtering for multi-tenant isolation
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        query = query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),  # Include legacy/shared docs
            )
        )

    # SECURITY FIX: Filter private documents - only owner or superadmin can access
    if not user.is_superadmin:
        query = query.where(
            or_(
                Document.is_private == False,
                and_(
                    Document.is_private == True,
                    Document.uploaded_by_id == UUID(user.user_id),
                ),
            )
        )

    result = await db.execute(query)
    document = result.scalar_one_or_none()

    if not document:
        # Check if document exists but user doesn't have access
        exists_query = select(Document.id).where(Document.id == document_id)
        exists_result = await db.execute(exists_query)
        if exists_result.scalar_one_or_none():
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="You don't have access to this document",
            )
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document not found",
        )

    # Get chunk count
    chunk_count_query = (
        select(func.count(Chunk.id))
        .where(Chunk.document_id == document_id)
    )
    chunk_result = await db.execute(chunk_count_query)
    chunk_count = chunk_result.scalar() or 0

    # Get embedding count (chunks stored in vector database)
    embedding_count_query = (
        select(func.count(Chunk.id))
        .where(
            Chunk.document_id == document_id,
            Chunk.has_embedding == True,
        )
    )
    embedding_result = await db.execute(embedding_count_query)
    embedding_count = embedding_result.scalar() or 0

    return document_to_response(document, chunk_count, embedding_count)


@router.patch("/{document_id}", response_model=DocumentResponse)
async def update_document(
    document_id: UUID,
    update: DocumentUpdate,
    user: AuthenticatedUser,
    request: Request,
    db: AsyncSession = Depends(get_async_session),
):
    """
    Update document metadata.

    Users can only update documents they have write access to:
    - Document owner can always update
    - Users with tier >= 90 can update any document at their tier level or below
    - Admins can update any document
    """
    logger.info(
        "Updating document",
        document_id=str(document_id),
        user_id=user.user_id,
        update=update.model_dump(exclude_none=True),
    )

    # Get document with access tier
    query = (
        select(Document)
        .join(AccessTier, Document.access_tier_id == AccessTier.id)
        .where(Document.id == document_id)
        .options(selectinload(Document.access_tier))
    )

    # PHASE 12 FIX: Apply organization filtering for multi-tenant isolation
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        query = query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),  # Include legacy/shared docs
            )
        )

    result = await db.execute(query)
    document = result.scalar_one_or_none()

    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document not found",
        )

    # Check write permission
    permission_service = get_permission_service()
    has_write = await permission_service.check_document_access(
        user_context=user,
        document_id=str(document_id),
        permission=Permission.WRITE,
        session=db,
    )

    if not has_write:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have write access to this document",
        )

    # If updating access tier, ensure user can assign that tier
    if update.access_tier is not None:
        can_assign = await permission_service.can_assign_tier(
            user_context=user,
            target_tier_level=update.access_tier,
        )
        if not can_assign:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail=f"Cannot assign tier {update.access_tier}. Your tier: {user.access_tier_level}",
            )

        # Find the access tier by level (limit 1 in case of duplicate levels)
        tier_query = select(AccessTier).where(AccessTier.level == update.access_tier).limit(1)
        tier_result = await db.execute(tier_query)
        new_tier = tier_result.scalar_one_or_none()

        if new_tier:
            document.access_tier_id = new_tier.id

    # Apply updates
    if update.name is not None:
        document.title = update.name

    if update.collection is not None:
        # Store collection as first tag
        current_tags = document.tags or []
        if current_tags:
            current_tags[0] = update.collection
        else:
            current_tags = [update.collection]
        document.tags = current_tags

    if update.tags is not None:
        document.tags = update.tags

    await db.commit()
    await db.refresh(document)

    # Log the update
    audit_service = get_audit_service()
    await audit_service.log_document_action(
        action=AuditAction.DOCUMENT_UPDATE,
        user_id=user.user_id,
        document_id=str(document_id),
        document_name=document.original_filename,
        changes=update.model_dump(exclude_none=True),
        ip_address=get_client_ip(request),
        session=db,
    )

    # Get chunk count
    chunk_count_query = (
        select(func.count(Chunk.id))
        .where(Chunk.document_id == document_id)
    )
    chunk_result = await db.execute(chunk_count_query)
    chunk_count = chunk_result.scalar() or 0

    # Get embedding count (chunks stored in vector database)
    embedding_count_query = (
        select(func.count(Chunk.id))
        .where(
            Chunk.document_id == document_id,
            Chunk.has_embedding == True,
        )
    )
    embedding_result = await db.execute(embedding_count_query)
    embedding_count = embedding_result.scalar() or 0

    return document_to_response(document, chunk_count, embedding_count)


@router.delete("/{document_id}")
async def delete_document(
    document_id: UUID,
    user: AuthenticatedUser,
    request: Request,
    db: AsyncSession = Depends(get_async_session),
    hard_delete: bool = Query(default=False, description="Permanently delete (admin only)"),
):
    """
    Delete a document and all its chunks.

    - Soft delete by default (marks as deleted)
    - Hard delete requires admin privileges and explicit flag
    """
    logger.info(
        "Deleting document",
        document_id=str(document_id),
        user_id=user.user_id,
        hard_delete=hard_delete,
    )

    # Get document with organization filtering
    query = select(Document).where(Document.id == document_id)

    # PHASE 12 FIX: Apply organization filtering for multi-tenant isolation
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        query = query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),  # Include legacy/shared docs
            )
        )

    result = await db.execute(query)
    document = result.scalar_one_or_none()

    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document not found",
        )

    # Check delete permission
    permission_service = get_permission_service()
    has_delete = await permission_service.check_document_access(
        user_context=user,
        document_id=str(document_id),
        permission=Permission.DELETE,
        session=db,
    )

    if not has_delete:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have permission to delete this document",
        )

    # Store document name for audit log
    document_name = document.original_filename

    # Hard delete requires admin
    if hard_delete:
        if not user.is_admin():
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Hard delete requires admin privileges",
            )

        # Delete chunks from vector store first (ChromaDB/pgvector)
        try:
            vector_store = get_vector_store()
            if vector_store:
                await vector_store.delete_document_chunks(str(document_id))
                logger.info(
                    "Deleted chunks from vector store",
                    document_id=str(document_id),
                )
        except Exception as e:
            logger.warning(
                "Failed to delete chunks from vector store (continuing with DB delete)",
                document_id=str(document_id),
                error=str(e),
            )

        # Delete from Knowledge Graph (EntityMentions, EntityRelations)
        try:
            from backend.services.knowledge_graph import KnowledgeGraphService
            kg_service = KnowledgeGraphService(db)
            kg_result = await kg_service.delete_document_data(document_id)
            logger.info(
                "Deleted Knowledge Graph data",
                document_id=str(document_id),
                mentions_deleted=kg_result.get("mentions_deleted", 0),
                relations_deleted=kg_result.get("relations_deleted", 0),
            )
        except Exception as e:
            logger.warning(
                "Failed to delete Knowledge Graph data (continuing with DB delete)",
                document_id=str(document_id),
                error=str(e),
            )

        # Permanently delete document and chunks from database (cascade)
        await db.delete(document)
        await db.commit()

        # Log the deletion
        audit_service = get_audit_service()
        await audit_service.log_document_action(
            action=AuditAction.DOCUMENT_DELETE,
            user_id=user.user_id,
            document_id=str(document_id),
            document_name=document_name,
            changes={"hard_delete": True},
            ip_address=get_client_ip(request),
        )

        return {
            "message": "Document permanently deleted",
            "document_id": str(document_id),
        }

    # Soft delete - mark status as failed with error message
    document.processing_status = ProcessingStatus.FAILED
    document.processing_error = "Deleted by user"
    await db.commit()

    # Log the deletion
    audit_service = get_audit_service()
    await audit_service.log_document_action(
        action=AuditAction.DOCUMENT_DELETE,
        user_id=user.user_id,
        document_id=str(document_id),
        document_name=document_name,
        changes={"hard_delete": False, "soft_delete": True},
        ip_address=get_client_ip(request),
        session=db,
    )

    return {
        "message": "Document deleted successfully",
        "document_id": str(document_id),
    }


@router.get("/deleted/list")
async def list_deleted_documents(
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
    page: int = Query(default=1, ge=1),
    page_size: int = Query(default=25, ge=1, le=100),
):
    """
    List soft-deleted documents (admin only).

    Returns documents that were soft-deleted and can be restored.
    """
    # Admin only
    if not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Only admins can view deleted documents",
        )

    logger.info(
        "Listing deleted documents",
        user_id=user.user_id,
        page=page,
        page_size=page_size,
    )

    # Query soft-deleted documents
    base_query = (
        select(Document)
        .where(
            (Document.processing_status == ProcessingStatus.FAILED)
            & (Document.processing_error == "Deleted by user")
        )
    )

    # Count total
    count_query = select(func.count()).select_from(base_query.subquery())
    total_result = await db.execute(count_query)
    total = total_result.scalar() or 0

    # Apply pagination
    offset = (page - 1) * page_size
    query = base_query.order_by(desc(Document.updated_at)).offset(offset).limit(page_size)
    result = await db.execute(query)
    documents = result.scalars().all()

    return {
        "documents": [
            {
                "id": str(doc.id),
                "name": doc.original_filename or doc.filename,
                "file_type": doc.file_type,
                "file_size": doc.file_size,
                "deleted_at": doc.updated_at.isoformat() if doc.updated_at else None,
                "created_at": doc.created_at.isoformat() if doc.created_at else None,
            }
            for doc in documents
        ],
        "total": total,
        "page": page,
        "page_size": page_size,
        "total_pages": (total + page_size - 1) // page_size if page_size > 0 else 0,
    }


@router.post("/deleted/{document_id}/restore")
async def restore_deleted_document(
    document_id: UUID,
    user: AuthenticatedUser,
    request: Request,
    db: AsyncSession = Depends(get_async_session),
):
    """
    Restore a soft-deleted document (admin only).

    Restores the document to its previous state before deletion.
    """
    # Admin only
    if not user.is_admin():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Only admins can restore deleted documents",
        )

    logger.info(
        "Restoring deleted document",
        document_id=str(document_id),
        user_id=user.user_id,
    )

    # Find the soft-deleted document
    query = select(Document).where(
        (Document.id == document_id)
        & (Document.processing_status == ProcessingStatus.FAILED)
        & (Document.processing_error == "Deleted by user")
    )
    result = await db.execute(query)
    document = result.scalar_one_or_none()

    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Deleted document not found",
        )

    # Restore the document
    document.processing_status = ProcessingStatus.COMPLETED
    document.processing_error = None
    await db.commit()

    # Log the restoration
    audit_service = get_audit_service()
    await audit_service.log_document_action(
        action=AuditAction.DOCUMENT_UPDATE,
        user_id=user.user_id,
        document_id=str(document_id),
        document_name=document.original_filename,
        changes={"restored": True, "previous_status": "soft_deleted"},
        ip_address=get_client_ip(request),
        session=db,
    )

    return {
        "message": "Document restored successfully",
        "document_id": str(document_id),
        "name": document.original_filename or document.filename,
    }


@router.get("/{document_id}/chunks", response_model=List[ChunkResponse])
async def get_document_chunks(
    document_id: UUID,
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
    page: int = Query(default=1, ge=1),
    page_size: int = Query(default=50, ge=1, le=200),
):
    """
    Get all chunks for a document.

    Useful for debugging and viewing how documents were processed.
    """
    logger.info(
        "Getting document chunks",
        document_id=str(document_id),
        user_id=user.user_id,
        page=page,
        page_size=page_size,
    )

    # Check document access first
    doc_query = (
        select(Document.id)
        .join(AccessTier, Document.access_tier_id == AccessTier.id)
        .where(
            and_(
                Document.id == document_id,
                AccessTier.level <= user.access_tier_level,
            )
        )
    )

    # PHASE 12 FIX: Apply organization filtering for multi-tenant isolation
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        doc_query = doc_query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),  # Include legacy/shared docs
            )
        )

    doc_result = await db.execute(doc_query)
    if not doc_result.scalar_one_or_none():
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this document",
        )

    # Get chunks
    offset = (page - 1) * page_size
    chunks_query = (
        select(Chunk)
        .where(Chunk.document_id == document_id)
        .order_by(Chunk.chunk_index)
        .offset(offset)
        .limit(page_size)
    )

    result = await db.execute(chunks_query)
    chunks = result.scalars().all()

    return [
        ChunkResponse(
            id=chunk.id,
            content=chunk.content,
            chunk_index=chunk.chunk_index,
            page_number=chunk.page_number,
            metadata={
                "section_title": chunk.section_title,
                "token_count": chunk.token_count,
                "char_count": chunk.char_count,
            },
        )
        for chunk in chunks
    ]


@router.post("/search", response_model=SearchResultResponse)
async def search_documents(
    request: DocumentSearchRequest,
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
):
    """
    Search documents using semantic and keyword search.

    Combines vector similarity with full-text search for best results.
    Results are filtered based on user's access tier.
    """
    logger.info(
        "Searching documents",
        query=request.query,
        user_id=user.user_id,
        access_tier=user.access_tier_level,
        collection=request.collection,
        limit=request.limit,
    )

    # Use custom vector store for search
    vector_store = get_vector_store()

    # Perform hybrid search with access tier filtering
    results = await vector_store.search(
        query=request.query,
        search_type=SearchType.HYBRID,
        top_k=request.limit,
        access_tier_level=user.access_tier_level,
        session=db,
    )

    # Build response - batch fetch document info to avoid N+1 queries
    search_results = []

    # Collect all unique document IDs
    doc_ids = list({result.document_id for result in results})

    # Batch query for all document info in one query
    doc_info_lookup = {}
    if doc_ids:
        doc_info_query = (
            select(Document.id, Document.original_filename, Document.tags, Document.file_type)
            .where(Document.id.in_(doc_ids))
        )
        doc_info_result = await db.execute(doc_info_query)
        for row in doc_info_result.all():
            doc_info_lookup[row[0]] = {
                "name": row[1] or "Unknown",
                "tags": row[2] or [],
                "file_type": row[3],
            }

    for result in results:
        # Get document info from lookup (O(1) instead of O(n) queries)
        doc_info = doc_info_lookup.get(result.document_id, {})
        doc_name = doc_info.get("name", "Unknown")
        tags = doc_info.get("tags", [])
        file_type = doc_info.get("file_type")

        # Filter by collection if specified
        if request.collection and request.collection not in tags:
            continue

        # Filter by file type if specified
        if request.file_types and file_type not in request.file_types:
            continue

        search_results.append(
            SearchResultItem(
                document_id=str(result.document_id),
                document_name=doc_name,
                chunk_id=str(result.chunk_id),
                content_preview=result.content[:200] + "..." if len(result.content) > 200 else result.content,
                relevance_score=result.score,
                page_number=result.page_number,
            )
        )

    return SearchResultResponse(
        results=search_results[:request.limit],
        total_results=len(search_results),
        query=request.query,
    )


@router.post("/{document_id}/reprocess")
async def reprocess_document(
    document_id: UUID,
    background_tasks: BackgroundTasks,
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
):
    """
    Reprocess a document (re-extract text, regenerate embeddings).

    Useful after OCR improvements or when processing settings change.
    Requires write permission on the document.
    """
    logger.info(
        "Reprocessing document",
        document_id=str(document_id),
        user_id=user.user_id,
    )

    # Check write permission
    permission_service = get_permission_service()
    has_write = await permission_service.check_document_access(
        user_context=user,
        document_id=str(document_id),
        permission=Permission.WRITE,
        session=db,
    )

    if not has_write:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have write access to this document",
        )

    # Get document with eagerly loaded access_tier to avoid lazy loading issues
    query = (
        select(Document)
        .options(selectinload(Document.access_tier))
        .where(Document.id == document_id)
    )
    result = await db.execute(query)
    document = result.scalar_one_or_none()

    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document not found",
        )

    # Get file path
    file_path = document.file_path
    if not file_path or not Path(file_path).exists():
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document file not found on disk",
        )

    # Extract needed values before committing (avoid lazy loading in background task)
    doc_tags = list(document.tags) if document.tags else []
    doc_access_tier = document.access_tier.level if document.access_tier else 1

    # Reset processing status
    document.processing_status = ProcessingStatus.PENDING
    document.processing_error = None
    await db.commit()

    # Import and trigger background processing
    from backend.api.routes.upload import process_document_background, UploadOptions

    options = UploadOptions(
        collection=doc_tags[0] if doc_tags else None,
        access_tier=doc_access_tier,
    )

    background_tasks.add_task(
        process_document_background,
        file_id=document_id,
        file_path=file_path,
        options=options,
    )

    logger.info(
        "Document queued for reprocessing",
        document_id=str(document_id),
        status="processing",
    )

    return {
        "message": "Document queued for reprocessing",
        "document_id": str(document_id),
        "status": "processing",
    }


@router.get("/collections/list", response_model=CollectionsResponse)
async def list_collections(
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
):
    """
    List all collections (tags) the user has access to.

    Collections are derived from document tags, showing only
    collections from documents the user can access.
    """
    # Get all tags from accessible documents
    # Exclude soft-deleted documents (those marked as failed with "Deleted by user" error)
    # to match the filtering behavior of list_documents
    # Note: Empty arrays are filtered in Python (if tags:) since SQLite lacks array_length

    # PHASE 12 FIX: Apply organization filtering for multi-tenant isolation
    org_id = get_org_filter(user)

    query = (
        select(Document.tags)
        .join(AccessTier, Document.access_tier_id == AccessTier.id)
        .where(
            and_(
                AccessTier.level <= user.access_tier_level,
                Document.tags.isnot(None),
                # Exclude soft-deleted documents
                ~(
                    (Document.processing_status == ProcessingStatus.FAILED)
                    & (Document.processing_error == "Deleted by user")
                ),
            )
        )
    )

    if org_id and not user.is_superadmin:
        query = query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),  # Include legacy/shared docs
            )
        )

    result = await db.execute(query)
    all_tags = result.scalars().all()

    # Count documents per tag
    tag_counts: dict[str, int] = {}
    for tags in all_tags:
        if tags:
            for tag in tags:
                tag_counts[tag] = tag_counts.get(tag, 0) + 1

    # Debug logging
    logger.info(
        "list_collections query result",
        documents_with_tags=len(all_tags),
        unique_tags=len(tag_counts),
        tag_counts=tag_counts,
        user_access_tier=user.access_tier_level,
    )

    # Build response
    collections = [
        CollectionInfo(name=tag, document_count=count)
        for tag, count in sorted(tag_counts.items(), key=lambda x: -x[1])
    ]

    # Count untagged documents (null or empty tags)
    # Exclude soft-deleted documents to match list_documents behavior
    untagged_query = (
        select(func.count(Document.id))
        .join(AccessTier, Document.access_tier_id == AccessTier.id)
        .where(
            and_(
                AccessTier.level <= user.access_tier_level,
                or_(
                    Document.tags.is_(None),
                    Document.tags == [],
                ),
                # Exclude soft-deleted documents
                ~(
                    (Document.processing_status == ProcessingStatus.FAILED)
                    & (Document.processing_error == "Deleted by user")
                ),
            )
        )
    )

    # PHASE 12 FIX: Apply organization filtering
    if org_id and not user.is_superadmin:
        untagged_query = untagged_query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),
            )
        )

    untagged_result = await db.execute(untagged_query)
    untagged_count = untagged_result.scalar() or 0

    # Add "(Untagged)" pseudo-collection if there are untagged documents
    if untagged_count > 0:
        collections.append(CollectionInfo(name="(Untagged)", document_count=untagged_count))

    # Count total unique documents (with tags + untagged)
    total_query = (
        select(func.count(Document.id))
        .join(AccessTier, Document.access_tier_id == AccessTier.id)
        .where(
            and_(
                AccessTier.level <= user.access_tier_level,
                # Exclude soft-deleted documents
                ~(
                    (Document.processing_status == ProcessingStatus.FAILED)
                    & (Document.processing_error == "Deleted by user")
                ),
            )
        )
    )

    # PHASE 12 FIX: Apply organization filtering
    if org_id and not user.is_superadmin:
        total_query = total_query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),
            )
        )

    total_result = await db.execute(total_query)
    total_documents = total_result.scalar() or 0

    return CollectionsResponse(collections=collections, total_documents=total_documents)


@router.get("/{document_id}/download")
async def download_document(
    document_id: UUID,
    request: Request,
    db: AsyncSession = Depends(get_async_session),
    token: Optional[str] = Query(None, description="Auth token for direct browser downloads"),
    preview: bool = Query(False, description="If true, display inline instead of downloading"),
):
    """
    Download or preview the original document file.

    Args:
        preview: If true, returns file for inline viewing (Content-Disposition: inline).
                 If false (default), returns file for download (Content-Disposition: attachment).

    Supports both Authorization header and token query parameter for direct browser access.
    """
    from backend.api.middleware.auth import _get_dev_user_context, DEV_MODE

    # Try to get user from Authorization header first, then query parameter
    user = None
    auth_header = request.headers.get("Authorization")

    if auth_header and auth_header.startswith("Bearer "):
        header_token = auth_header.split(" ", 1)[1]
        if DEV_MODE and header_token == "dev-token":
            user = _get_dev_user_context()

    if not user and token:
        # Use token from query parameter
        if DEV_MODE and token == "dev-token":
            user = _get_dev_user_context()

    if not user:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Authentication required. Provide token as query parameter or Authorization header.",
        )

    logger.info(
        "Downloading document",
        document_id=str(document_id),
        user_id=user.user_id,
    )

    # Get document with organization filtering
    query = select(Document).where(Document.id == document_id)

    # PHASE 12 FIX: Apply organization filtering for multi-tenant isolation
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        query = query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),  # Include legacy/shared docs
            )
        )

    result = await db.execute(query)
    document = result.scalar_one_or_none()

    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document not found",
        )

    # Check read permission
    permission_service = get_permission_service()
    has_access = await permission_service.check_document_access(
        user_context=user,
        document_id=str(document_id),
        permission=Permission.READ,
        session=db,
    )

    if not has_access:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have permission to download this document",
        )

    # Check if file exists
    file_path = Path(document.file_path)
    if not file_path.exists():
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document file not found on disk",
        )

    # Get MIME type - infer from file extension if not stored
    mime_type = document.mime_type
    if not mime_type:
        # Common MIME type mappings
        ext = (document.file_type or file_path.suffix.lstrip('.')).lower()
        mime_mappings = {
            'pdf': 'application/pdf',
            'png': 'image/png',
            'jpg': 'image/jpeg',
            'jpeg': 'image/jpeg',
            'gif': 'image/gif',
            'webp': 'image/webp',
            'svg': 'image/svg+xml',
            'doc': 'application/msword',
            'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'ppt': 'application/vnd.ms-powerpoint',
            'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'xls': 'application/vnd.ms-excel',
            'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'txt': 'text/plain',
            'html': 'text/html',
            'json': 'application/json',
            'xml': 'application/xml',
            'csv': 'text/csv',
            'md': 'text/markdown',
        }
        mime_type = mime_mappings.get(ext, 'application/octet-stream')

    # Return file - if preview mode, display inline; otherwise download
    if preview:
        # For preview: don't set filename (displays inline in browser)
        return FileResponse(
            path=str(file_path),
            media_type=mime_type,
        )
    else:
        # For download: set filename (triggers download)
        return FileResponse(
            path=str(file_path),
            filename=document.original_filename,
            media_type=mime_type,
        )


@router.get("/{document_id}/preview/metadata")
async def get_uploaded_document_preview_metadata(
    document_id: UUID,
    request: Request,
    db: AsyncSession = Depends(get_async_session),
    token: Optional[str] = Query(None, description="Auth token for direct browser access"),
):
    """
    Get preview metadata for an uploaded document.

    Returns information about preview capabilities, page count, etc.
    """
    from backend.api.middleware.auth import _get_dev_user_context, DEV_MODE

    # Auth handling (same as download)
    user = None
    auth_header = request.headers.get("Authorization")
    if auth_header and auth_header.startswith("Bearer "):
        header_token = auth_header.split(" ", 1)[1]
        if DEV_MODE and header_token == "dev-token":
            user = _get_dev_user_context()
    if not user and token:
        if DEV_MODE and token == "dev-token":
            user = _get_dev_user_context()
    if not user:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Authentication required")

    # Get document
    query = select(Document).where(Document.id == document_id)
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        query = query.where(or_(Document.organization_id == org_id, Document.organization_id.is_(None)))
    result = await db.execute(query)
    document = result.scalar_one_or_none()

    if not document:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document not found")

    # Check permission
    permission_service = get_permission_service()
    has_access = await permission_service.check_document_access(
        user_context=user, document_id=str(document_id), permission=Permission.READ, session=db
    )
    if not has_access:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Access denied")

    file_path = Path(document.file_path)
    if not file_path.exists():
        return {"supported": False, "error": "File not found"}

    file_type = (document.file_type or file_path.suffix.lstrip('.')).lower()

    # Check preview support based on file type
    if file_type == "pdf":
        try:
            from pdf2image import pdfinfo_from_path
            info = pdfinfo_from_path(str(file_path))
            return {"supported": True, "type": "pdf", "page_count": info.get("Pages", 1), "format": "pdf"}
        except Exception:
            return {"supported": True, "type": "iframe", "format": "pdf"}  # Fallback to iframe

    elif file_type == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            return {"supported": True, "type": "slides", "page_count": len(prs.slides), "format": "pptx"}
        except Exception as e:
            return {"supported": False, "error": str(e)}

    elif file_type == "docx":
        try:
            import mammoth
            return {"supported": True, "type": "html", "page_count": 1, "format": "docx"}
        except ImportError:
            return {"supported": False, "error": "DOCX preview not available"}

    elif file_type in ["png", "jpg", "jpeg", "gif", "webp", "svg"]:
        return {"supported": True, "type": "image", "page_count": 1, "format": file_type}

    elif file_type in ["txt", "md", "html", "json", "csv", "xml"]:
        return {"supported": True, "type": "text", "page_count": 1, "format": file_type}

    else:
        return {"supported": False, "type": "unsupported", "format": file_type}


@router.get("/{document_id}/preview/page/{page_num}")
async def get_uploaded_document_preview_page(
    document_id: UUID,
    page_num: int,
    request: Request,
    db: AsyncSession = Depends(get_async_session),
    token: Optional[str] = Query(None, description="Auth token"),
    width: int = Query(1280, description="Image width"),
    height: int = Query(720, description="Image height"),
):
    """
    Get a rendered preview page/slide as an image.

    For PPTX: Renders slide as image
    For DOCX: Returns HTML content
    For PDF: Returns page as image (if pdf2image available)
    """
    from backend.api.middleware.auth import _get_dev_user_context, DEV_MODE
    import io
    from PIL import Image, ImageDraw, ImageFont

    # Auth handling
    user = None
    auth_header = request.headers.get("Authorization")
    if auth_header and auth_header.startswith("Bearer "):
        header_token = auth_header.split(" ", 1)[1]
        if DEV_MODE and header_token == "dev-token":
            user = _get_dev_user_context()
    if not user and token:
        if DEV_MODE and token == "dev-token":
            user = _get_dev_user_context()
    if not user:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Authentication required")

    # Get document
    query = select(Document).where(Document.id == document_id)
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        query = query.where(or_(Document.organization_id == org_id, Document.organization_id.is_(None)))
    result = await db.execute(query)
    document = result.scalar_one_or_none()

    if not document:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document not found")

    # Check permission
    permission_service = get_permission_service()
    has_access = await permission_service.check_document_access(
        user_context=user, document_id=str(document_id), permission=Permission.READ, session=db
    )
    if not has_access:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Access denied")

    file_path = Path(document.file_path)
    if not file_path.exists():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")

    file_type = (document.file_type or file_path.suffix.lstrip('.')).lower()

    # PPTX: Render slide as image using LibreOffice
    if file_type == "pptx":
        try:
            import subprocess
            import tempfile
            import hashlib

            # Get slide count first
            from pptx import Presentation
            prs = Presentation(str(file_path))
            slide_count = len(prs.slides)

            if page_num < 1 or page_num > slide_count:
                raise HTTPException(status_code=400, detail=f"Invalid slide number. Document has {slide_count} slides.")

            # Create a cache directory for converted PDFs
            cache_dir = Path(tempfile.gettempdir()) / "pptx_preview_cache"
            cache_dir.mkdir(exist_ok=True)

            # Generate cache key based on file path and modification time
            file_stat = file_path.stat()
            cache_key = hashlib.md5(f"{file_path}:{file_stat.st_mtime}".encode()).hexdigest()
            cached_pdf = cache_dir / f"{cache_key}.pdf"

            # Convert PPTX to PDF using LibreOffice if not cached
            if not cached_pdf.exists():
                # Find LibreOffice
                soffice_paths = [
                    "/opt/homebrew/bin/soffice",
                    "/usr/local/bin/soffice",
                    "/usr/bin/soffice",
                    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                ]
                soffice = None
                for path in soffice_paths:
                    if Path(path).exists():
                        soffice = path
                        break

                if not soffice:
                    raise HTTPException(status_code=500, detail="LibreOffice not found for PPTX rendering")

                # Convert to PDF
                result = subprocess.run(
                    [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(cache_dir), str(file_path)],
                    capture_output=True,
                    timeout=120,
                )

                if result.returncode != 0:
                    logger.error("LibreOffice conversion failed", stderr=result.stderr.decode())
                    raise HTTPException(status_code=500, detail="Failed to convert PPTX to PDF")

                # Rename the output file to our cache key
                output_pdf = cache_dir / f"{file_path.stem}.pdf"
                if output_pdf.exists():
                    output_pdf.rename(cached_pdf)
                else:
                    raise HTTPException(status_code=500, detail="PDF conversion produced no output")

            # Now use pdf2image to get the specific page
            try:
                from pdf2image import convert_from_path

                images = convert_from_path(
                    str(cached_pdf),
                    dpi=150,
                    first_page=page_num,
                    last_page=page_num,
                )

                if not images:
                    raise HTTPException(status_code=400, detail=f"Could not render page {page_num}")

                # Resize if needed
                img = images[0]
                if width and height:
                    img.thumbnail((width, height), Image.Resampling.LANCZOS)

                buffer = io.BytesIO()
                img.save(buffer, format="PNG")
                buffer.seek(0)

                return Response(content=buffer.getvalue(), media_type="image/png")

            except ImportError:
                raise HTTPException(status_code=500, detail="pdf2image not installed for slide rendering")

        except HTTPException:
            raise
        except subprocess.TimeoutExpired:
            raise HTTPException(status_code=500, detail="PPTX conversion timed out")
        except Exception as e:
            logger.error("PPTX preview error", error=str(e))
            raise HTTPException(status_code=500, detail=f"Failed to render slide: {e}")

    # DOCX: Return HTML
    elif file_type == "docx":
        try:
            import mammoth

            with open(file_path, "rb") as f:
                result = mammoth.convert_to_html(f)
                html = result.value

                styled_html = f"""<!DOCTYPE html>
<html><head><style>
body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; max-width: 800px; margin: 0 auto; padding: 20px; line-height: 1.6; color: #333; }}
h1, h2, h3 {{ color: #111; margin-top: 1.5em; }}
table {{ border-collapse: collapse; width: 100%; margin: 1em 0; }}
th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
img {{ max-width: 100%; height: auto; }}
</style></head><body>{html}</body></html>"""

                return Response(content=styled_html, media_type="text/html")
        except ImportError:
            raise HTTPException(status_code=500, detail="DOCX preview not available")
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to render DOCX: {e}")

    # PDF: Render page as image (if pdf2image available)
    elif file_type == "pdf":
        try:
            from pdf2image import convert_from_path

            images = convert_from_path(str(file_path), dpi=150, first_page=page_num, last_page=page_num)
            if not images:
                raise HTTPException(status_code=400, detail=f"No page {page_num}")

            buffer = io.BytesIO()
            images[0].save(buffer, format="PNG")
            buffer.seek(0)

            return Response(content=buffer.getvalue(), media_type="image/png")
        except ImportError:
            # Fallback: just return the PDF page (browser will handle)
            raise HTTPException(status_code=501, detail="PDF rendering not available, use iframe")
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to render PDF: {e}")

    # Text files: Return content
    elif file_type in ["txt", "md", "html", "json", "csv", "xml"]:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            content_types = {
                "txt": "text/plain", "md": "text/markdown", "html": "text/html",
                "json": "application/json", "csv": "text/csv", "xml": "application/xml"
            }
            return Response(content=content, media_type=content_types.get(file_type, "text/plain"))
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to read file: {e}")

    else:
        raise HTTPException(status_code=400, detail=f"Preview not supported for {file_type} files")


@router.get("/{document_id}/preview/slides")
async def get_all_uploaded_document_slides(
    document_id: UUID,
    request: Request,
    db: AsyncSession = Depends(get_async_session),
    token: Optional[str] = Query(None, description="Auth token"),
    width: int = Query(400, description="Thumbnail width"),
    height: int = Query(225, description="Thumbnail height"),
):
    """
    Get all slides from a PPTX as base64-encoded images.

    Returns a list of base64 PNG images for each slide.
    """
    from backend.api.middleware.auth import _get_dev_user_context, DEV_MODE
    import io
    import base64
    from PIL import Image, ImageDraw, ImageFont

    # Auth handling
    user = None
    auth_header = request.headers.get("Authorization")
    if auth_header and auth_header.startswith("Bearer "):
        header_token = auth_header.split(" ", 1)[1]
        if DEV_MODE and header_token == "dev-token":
            user = _get_dev_user_context()
    if not user and token:
        if DEV_MODE and token == "dev-token":
            user = _get_dev_user_context()
    if not user:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Authentication required")

    # Get document
    query = select(Document).where(Document.id == document_id)
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        query = query.where(or_(Document.organization_id == org_id, Document.organization_id.is_(None)))
    result = await db.execute(query)
    document = result.scalar_one_or_none()

    if not document:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document not found")

    # Check permission
    permission_service = get_permission_service()
    has_access = await permission_service.check_document_access(
        user_context=user, document_id=str(document_id), permission=Permission.READ, session=db
    )
    if not has_access:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Access denied")

    file_path = Path(document.file_path)
    file_type = (document.file_type or file_path.suffix.lstrip('.')).lower()

    if file_type != "pptx":
        raise HTTPException(status_code=400, detail="Slides endpoint only available for PPTX files")

    try:
        import subprocess
        import tempfile
        import hashlib
        from pptx import Presentation

        prs = Presentation(str(file_path))
        slide_count = len(prs.slides)

        # Create cache directory
        cache_dir = Path(tempfile.gettempdir()) / "pptx_preview_cache"
        cache_dir.mkdir(exist_ok=True)

        # Generate cache key
        file_stat = file_path.stat()
        cache_key = hashlib.md5(f"{file_path}:{file_stat.st_mtime}".encode()).hexdigest()
        cached_pdf = cache_dir / f"{cache_key}.pdf"

        # Convert PPTX to PDF if not cached
        if not cached_pdf.exists():
            soffice_paths = [
                "/opt/homebrew/bin/soffice",
                "/usr/local/bin/soffice",
                "/usr/bin/soffice",
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            ]
            soffice = None
            for path in soffice_paths:
                if Path(path).exists():
                    soffice = path
                    break

            if not soffice:
                raise HTTPException(status_code=500, detail="LibreOffice not found")

            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(cache_dir), str(file_path)],
                capture_output=True,
                timeout=120,
            )

            if result.returncode != 0:
                raise HTTPException(status_code=500, detail="Failed to convert PPTX")

            output_pdf = cache_dir / f"{file_path.stem}.pdf"
            if output_pdf.exists():
                output_pdf.rename(cached_pdf)

        # Generate thumbnails from PDF
        from pdf2image import convert_from_path

        all_images = convert_from_path(str(cached_pdf), dpi=72)  # Lower DPI for thumbnails
        slides = []

        for img in all_images:
            # Resize to thumbnail size
            img.thumbnail((width, height), Image.Resampling.LANCZOS)

            # Create a new image with exact dimensions (centered)
            thumb = Image.new("RGB", (width, height), (255, 255, 255))
            offset = ((width - img.width) // 2, (height - img.height) // 2)
            thumb.paste(img, offset)

            buffer = io.BytesIO()
            thumb.save(buffer, format="PNG")
            slides.append(base64.b64encode(buffer.getvalue()).decode())

        return {"slide_count": len(slides), "slides": slides}

    except HTTPException:
        raise
    except Exception as e:
        logger.error("PPTX slides error", error=str(e))
        raise HTTPException(status_code=500, detail=f"Failed to render slides: {e}")


@router.post("/bulk-download")
async def bulk_download_documents(
    request: BulkDownloadRequest,
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
):
    """
    Download multiple documents as a ZIP file.

    Returns a ZIP archive containing all requested documents.
    Only documents the user has read access to will be included.
    """
    logger.info(
        "Bulk downloading documents",
        user_id=user.user_id,
        document_count=len(request.document_ids),
    )

    if not request.document_ids:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No document IDs provided",
        )

    # Get documents the user has access to
    query = (
        select(Document)
        .join(AccessTier, Document.access_tier_id == AccessTier.id)
        .where(
            and_(
                Document.id.in_(request.document_ids),
                AccessTier.level <= user.access_tier_level,
            )
        )
    )

    # PHASE 12 FIX: Apply organization filtering for multi-tenant isolation
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        query = query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),  # Include legacy/shared docs
            )
        )

    result = await db.execute(query)
    documents = result.scalars().all()

    if not documents:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="No accessible documents found",
        )

    # Create ZIP file in memory
    zip_buffer = io.BytesIO()

    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        files_added = 0
        for doc in documents:
            file_path = Path(doc.file_path)
            if file_path.exists():
                # Use original filename, handle duplicates
                filename = doc.original_filename
                # Check for duplicates and add suffix if needed
                base_name = Path(filename).stem
                extension = Path(filename).suffix
                counter = 1
                while filename in [info.filename for info in zip_file.filelist]:
                    filename = f"{base_name}_{counter}{extension}"
                    counter += 1

                zip_file.write(file_path, filename)
                files_added += 1
            else:
                logger.warning(
                    "Document file not found on disk",
                    document_id=str(doc.id),
                    file_path=str(file_path),
                )

    if files_added == 0:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="No document files found on disk",
        )

    # Reset buffer position
    zip_buffer.seek(0)

    # Generate filename with timestamp
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    zip_filename = f"documents_{timestamp}.zip"

    logger.info(
        "Bulk download completed",
        user_id=user.user_id,
        files_added=files_added,
        zip_filename=zip_filename,
    )

    return StreamingResponse(
        zip_buffer,
        media_type="application/zip",
        headers={
            "Content-Disposition": f'attachment; filename="{zip_filename}"',
        },
    )


# =============================================================================
# Auto-Tagging Endpoints
# =============================================================================

@router.post("/{document_id}/auto-tag", response_model=AutoTagResponse)
async def auto_tag_document(
    document_id: UUID,
    request: AutoTagRequest,
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
):
    """
    Auto-generate tags for a document using LLM analysis.

    Analyzes document content and suggests relevant tags/categories.
    The first tag is set as the primary collection.
    """
    logger.info(
        "Auto-tagging document",
        document_id=str(document_id),
        user_id=user.user_id,
        max_tags=request.max_tags,
    )

    # Get document with chunks
    query = (
        select(Document)
        .options(selectinload(Document.chunks))
        .where(Document.id == document_id)
    )
    result = await db.execute(query)
    document = result.scalar_one_or_none()

    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document not found",
        )

    # Check write permission (modifying tags)
    permission_service = get_permission_service()
    has_access = await permission_service.check_document_access(
        user_context=user,
        document_id=str(document_id),
        permission=Permission.WRITE,
        session=db,
    )

    if not has_access:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have permission to modify this document",
        )

    # Get content sample from chunks
    chunks = sorted(document.chunks, key=lambda c: c.chunk_index)[:3]
    content_sample = "\n".join([chunk.content for chunk in chunks])

    if not content_sample:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Document has no indexed content for analysis",
        )

    # Get existing collections for context
    collections_query = (
        select(Document.tags)
        .join(AccessTier, Document.access_tier_id == AccessTier.id)
        .where(
            and_(
                AccessTier.level <= user.access_tier_level,
                Document.tags.isnot(None),
            )
        )
    )
    collections_result = await db.execute(collections_query)
    all_tags = collections_result.scalars().all()

    # Flatten and deduplicate collections
    existing_collections = list(set(
        tag for tags in all_tags if tags for tag in tags
    ))

    # Generate tags
    auto_tagger = AutoTaggerService()
    tags = await auto_tagger.generate_tags(
        document_name=document.original_filename or document.filename,
        content_sample=content_sample,
        existing_collections=existing_collections,
        max_tags=request.max_tags,
    )

    if not tags:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to generate tags. Please try again.",
        )

    # Merge auto-generated tags with existing user tags (preserve user tags)
    existing_tags = document.tags or []
    merged_tags = list(dict.fromkeys(existing_tags + tags))
    document.tags = merged_tags
    await db.commit()

    # Log audit
    audit_service = get_audit_service()
    await audit_service.log(
        action=AuditAction.DOCUMENT_UPDATE,
        user_id=user.user_id,
        resource_type="document",
        resource_id=str(document_id),
        details={"auto_tags": tags},
        session=db,
    )

    logger.info(
        "Document auto-tagged successfully",
        document_id=str(document_id),
        tags=tags,
    )

    return AutoTagResponse(
        document_id=str(document_id),
        tags=tags,
        collection=tags[0] if tags else None,
    )


@router.post("/bulk-auto-tag", response_model=BulkAutoTagResponse)
async def bulk_auto_tag_documents(
    request: BulkAutoTagRequest,
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
):
    """
    Auto-tag multiple documents using LLM analysis.

    Processes documents synchronously and returns all results.
    For large batches, consider implementing background processing.
    """
    logger.info(
        "Bulk auto-tagging documents",
        user_id=user.user_id,
        document_count=len(request.document_ids),
    )

    if len(request.document_ids) > 20:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Maximum 20 documents per bulk request. For more, use individual requests.",
        )

    # Get existing collections for context (once, not per document)
    collections_query = (
        select(Document.tags)
        .join(AccessTier, Document.access_tier_id == AccessTier.id)
        .where(
            and_(
                AccessTier.level <= user.access_tier_level,
                Document.tags.isnot(None),
            )
        )
    )
    collections_result = await db.execute(collections_query)
    all_tags = collections_result.scalars().all()
    existing_collections = list(set(
        tag for tags in all_tags if tags for tag in tags
    ))

    auto_tagger = AutoTaggerService()
    results = []
    processed = 0

    for doc_id in request.document_ids:
        try:
            # Get document with chunks
            query = (
                select(Document)
                .options(selectinload(Document.chunks))
                .where(Document.id == doc_id)
            )
            result = await db.execute(query)
            document = result.scalar_one_or_none()

            if not document:
                logger.warning(f"Document {doc_id} not found, skipping")
                continue

            # Check permission
            permission_service = get_permission_service()
            has_access = await permission_service.check_document_access(
                user_context=user,
                document_id=str(doc_id),
                permission=Permission.WRITE,
                session=db,
            )

            if not has_access:
                logger.warning(f"No write permission for document {doc_id}, skipping")
                continue

            # Get content sample
            chunks = sorted(document.chunks, key=lambda c: c.chunk_index)[:3]
            content_sample = "\n".join([chunk.content for chunk in chunks])

            if not content_sample:
                logger.warning(f"Document {doc_id} has no content, skipping")
                continue

            # Generate tags
            doc_name = document.title or document.original_filename or document.filename
            tags = await auto_tagger.generate_tags(
                document_name=doc_name,
                content_sample=content_sample,
                existing_collections=existing_collections,
                max_tags=request.max_tags,
            )

            if tags:
                # Merge auto-generated tags with existing user tags (preserve user tags)
                existing_tags = document.tags or []
                merged_tags = list(dict.fromkeys(existing_tags + tags))
                document.tags = merged_tags
                processed += 1
                results.append(AutoTagResponse(
                    document_id=str(doc_id),
                    tags=merged_tags,
                    collection=merged_tags[0] if merged_tags else None,
                ))

                # Add new tags to existing collections for next iterations
                for tag in tags:
                    if tag not in existing_collections:
                        existing_collections.append(tag)

        except Exception as e:
            logger.error(f"Failed to auto-tag document {doc_id}: {e}")
            continue

    await db.commit()

    logger.info(
        "Bulk auto-tagging completed",
        total=len(request.document_ids),
        processed=processed,
    )

    return BulkAutoTagResponse(
        status="completed",
        total=len(request.document_ids),
        processed=processed,
        results=results,
    )


# =============================================================================
# Document Name Extraction from File Metadata
# =============================================================================


class FixDocumentNamesRequest(BaseModel):
    """Request to fix document names from file metadata."""
    document_ids: Optional[List[UUID]] = None  # If None, fix all documents with "unknown" name


class FixDocumentNamesResponse(BaseModel):
    """Response after fixing document names."""
    status: str
    total: int
    fixed: int
    results: List[Dict[str, str]]


def _extract_name_from_pdf(file_path: str) -> Optional[str]:
    """Extract title from PDF metadata."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        if reader.metadata:
            # Try /Title first, then /Subject
            title = reader.metadata.get("/Title") or reader.metadata.get("/Subject")
            if title and title.strip() and title.strip().lower() != "untitled":
                return title.strip()
    except Exception as e:
        logger.debug("PDF metadata extraction failed", error=str(e))
    return None


def _extract_name_from_pptx(file_path: str) -> Optional[str]:
    """Extract title from PowerPoint metadata."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        if prs.core_properties.title:
            return prs.core_properties.title.strip()
        if prs.core_properties.subject:
            return prs.core_properties.subject.strip()
    except Exception as e:
        logger.debug("PPTX metadata extraction failed", error=str(e))
    return None


def _extract_name_from_file(file_path: str, file_type: str) -> Optional[str]:
    """Extract document name from file metadata based on type."""
    if file_type in ["pdf"]:
        return _extract_name_from_pdf(file_path)
    elif file_type in ["pptx", "ppt"]:
        return _extract_name_from_pptx(file_path)
    # For other types, we can't extract metadata - use filename
    return None


@router.post("/fix-names", response_model=FixDocumentNamesResponse)
async def fix_document_names(
    request: FixDocumentNamesRequest,
    user: AdminUser,  # Admin only
    db: AsyncSession = Depends(get_async_session),
):
    """
    Fix document names by extracting from file metadata.

    For documents with "unknown" as original_filename, attempts to:
    1. Extract title from PDF/PPTX metadata
    2. Fall back to using the file's base name with extension

    Admin only endpoint.
    """
    logger.info(
        "Fixing document names",
        user_id=user.user_id,
        document_ids=len(request.document_ids) if request.document_ids else "all",
    )

    # Build query for documents to fix
    if request.document_ids:
        query = (
            select(Document)
            .options(selectinload(Document.access_tier))
            .where(Document.id.in_(request.document_ids))
        )
    else:
        # Fix all documents with "unknown" name
        query = (
            select(Document)
            .options(selectinload(Document.access_tier))
            .where(
                or_(
                    Document.original_filename == "unknown",
                    Document.original_filename.is_(None),
                )
            )
        )

    result = await db.execute(query)
    documents = result.scalars().all()

    fixed = 0
    results = []

    for document in documents:
        old_name = document.original_filename or "unknown"
        new_name = None

        # Try to extract from file metadata
        if document.file_path and Path(document.file_path).exists():
            file_type = document.file_type or Path(document.file_path).suffix.lstrip(".")
            new_name = _extract_name_from_file(document.file_path, file_type.lower())

        # Fallback: use filename with proper formatting
        if not new_name:
            if document.filename:
                # Remove UUID prefix if present (e.g., "abc123.pdf" -> keep as is, or extract meaningful part)
                base_name = Path(document.filename).stem
                ext = Path(document.filename).suffix
                # If the stem looks like a UUID, create a readable name
                if len(base_name) == 36 and "-" in base_name:
                    new_name = f"Document{ext}"
                else:
                    new_name = document.filename

        if new_name and new_name != old_name:
            document.original_filename = new_name
            fixed += 1
            results.append({
                "document_id": str(document.id),
                "old_name": old_name,
                "new_name": new_name,
            })

    await db.commit()

    logger.info(
        "Document names fixed",
        total=len(documents),
        fixed=fixed,
    )

    return FixDocumentNamesResponse(
        status="completed",
        total=len(documents),
        fixed=fixed,
        results=results,
    )


# =============================================================================
# Document Move Endpoints
# =============================================================================

class MoveDocumentRequest(BaseModel):
    """Request to move a document to a folder."""
    folder_id: Optional[str] = Field(None, description="Target folder ID (null to remove from folder)")


class BulkMoveRequest(BaseModel):
    """Request to move multiple documents."""
    document_ids: List[str] = Field(..., min_length=1, description="List of document IDs to move")
    folder_id: Optional[str] = Field(None, description="Target folder ID (null to remove from folders)")


class MoveDocumentResponse(BaseModel):
    """Response for document move operation."""
    document_id: str
    folder_id: Optional[str]
    folder_name: Optional[str]
    message: str


class BulkMoveResponse(BaseModel):
    """Response for bulk move operation."""
    moved: int
    failed: int
    folder_id: Optional[str]
    folder_name: Optional[str]
    errors: List[Dict]


@router.post("/{document_id}/move", response_model=MoveDocumentResponse)
async def move_document(
    document_id: str,
    request: MoveDocumentRequest,
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
):
    """
    Move a document to a folder.

    - Set folder_id to move document into a folder
    - Set folder_id to null to remove document from its current folder
    """
    from backend.services.folder_service import get_folder_service
    import uuid as uuid_module

    # Get document with organization filtering
    query = (
        select(Document)
        .options(selectinload(Document.access_tier))
        .where(Document.id == document_id)
    )

    # PHASE 12 FIX: Apply organization filtering for multi-tenant isolation
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        query = query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),  # Include legacy/shared docs
            )
        )

    result = await db.execute(query)
    document = result.scalar_one_or_none()

    if not document:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document not found")

    # Check document access
    user_tier = user.access_tier_level
    doc_tier = document.access_tier.level if document.access_tier else 1
    if user_tier < doc_tier:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="No access to this document")

    folder_service = get_folder_service()
    folder_name = None

    if request.folder_id:
        # Check target folder access
        target_folder = await folder_service.get_folder_by_id(request.folder_id)
        if not target_folder:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Target folder not found")

        effective_tier = await folder_service.get_effective_tier_level(target_folder)
        if user_tier < effective_tier:
            raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="No access to target folder")

        document.folder_id = uuid_module.UUID(request.folder_id)
        folder_name = target_folder.name
    else:
        # Remove from folder
        document.folder_id = None

    await db.commit()

    logger.info(
        "Document moved",
        document_id=document_id,
        folder_id=request.folder_id,
    )

    return MoveDocumentResponse(
        document_id=document_id,
        folder_id=request.folder_id,
        folder_name=folder_name,
        message=f"Document moved to {'folder ' + folder_name if folder_name else 'root'}",
    )


@router.post("/bulk/move", response_model=BulkMoveResponse)
async def bulk_move_documents(
    request: BulkMoveRequest,
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
):
    """
    Move multiple documents to a folder.

    Moves all documents that the user has access to.
    Returns count of moved documents and any errors.
    """
    from backend.services.folder_service import get_folder_service
    import uuid as uuid_module

    folder_service = get_folder_service()
    user_tier = user.access_tier_level
    folder_name = None

    # Validate target folder
    if request.folder_id:
        target_folder = await folder_service.get_folder_by_id(request.folder_id)
        if not target_folder:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Target folder not found")

        effective_tier = await folder_service.get_effective_tier_level(target_folder)
        if user_tier < effective_tier:
            raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="No access to target folder")

        folder_name = target_folder.name

    # Get all documents with organization filtering
    query = (
        select(Document)
        .options(selectinload(Document.access_tier))
        .where(Document.id.in_(request.document_ids))
    )

    # PHASE 12 FIX: Apply organization filtering for multi-tenant isolation
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        query = query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),  # Include legacy/shared docs
            )
        )

    result = await db.execute(query)
    documents = result.scalars().all()

    moved = 0
    failed = 0
    errors = []

    for doc in documents:
        doc_tier = doc.access_tier.level if doc.access_tier else 1

        if user_tier < doc_tier:
            failed += 1
            errors.append({
                "document_id": str(doc.id),
                "error": "No access to this document",
            })
            continue

        if request.folder_id:
            doc.folder_id = uuid_module.UUID(request.folder_id)
        else:
            doc.folder_id = None

        moved += 1

    # Check for documents not found
    found_ids = {str(doc.id) for doc in documents}
    for doc_id in request.document_ids:
        if doc_id not in found_ids:
            failed += 1
            errors.append({
                "document_id": doc_id,
                "error": "Document not found",
            })

    await db.commit()

    logger.info(
        "Bulk move completed",
        moved=moved,
        failed=failed,
        folder_id=request.folder_id,
    )

    return BulkMoveResponse(
        moved=moved,
        failed=failed,
        folder_id=request.folder_id,
        folder_name=folder_name,
        errors=errors,
    )


# =============================================================================
# Bulk Delete
# =============================================================================

class BulkDeleteRequest(BaseModel):
    """Request to delete multiple documents."""
    document_ids: List[str] = Field(..., min_length=1, max_length=100, description="List of document IDs to delete")
    permanent: bool = Field(False, description="If true, permanently delete. If false, soft delete.")


class BulkDeleteResponse(BaseModel):
    """Response for bulk delete operation."""
    deleted: int
    failed: int
    permanent: bool
    errors: List[Dict]


@router.post("/bulk/delete", response_model=BulkDeleteResponse)
async def bulk_delete_documents(
    request: BulkDeleteRequest,
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
):
    """
    Delete multiple documents.

    Soft deletes by default (sets status to FAILED with "Deleted by user").
    Set permanent=true to permanently remove documents and their chunks.

    Requires access to each document. Returns count of deleted and failed.
    """
    from backend.db.models import ProcessingStatus
    import uuid as uuid_module

    user_tier = user.access_tier_level

    # Get all documents with organization filtering
    doc_uuids = []
    for doc_id in request.document_ids:
        try:
            doc_uuids.append(uuid_module.UUID(doc_id))
        except ValueError:
            pass

    query = (
        select(Document)
        .options(selectinload(Document.access_tier))
        .where(Document.id.in_(doc_uuids))
    )

    # PHASE 12 FIX: Apply organization filtering for multi-tenant isolation
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        query = query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),  # Include legacy/shared docs
            )
        )

    result = await db.execute(query)
    documents = result.scalars().all()

    deleted = 0
    failed = 0
    errors = []

    for doc in documents:
        doc_tier = doc.access_tier.level if doc.access_tier else 1

        if user_tier < doc_tier:
            failed += 1
            errors.append({
                "document_id": str(doc.id),
                "error": "No access to this document",
            })
            continue

        if request.permanent:
            # 1. Delete from vector store (ChromaDB)
            try:
                vector_store = get_vector_store()
                if vector_store:
                    await vector_store.delete_document_chunks(str(doc.id))
            except Exception as e:
                logger.warning(
                    "Failed to delete from vector store during bulk delete",
                    document_id=str(doc.id),
                    error=str(e),
                )

            # 2. Delete from Knowledge Graph
            try:
                from backend.services.knowledge_graph import KnowledgeGraphService
                kg_service = KnowledgeGraphService(db)
                await kg_service.delete_document_data(doc.id)
            except Exception as e:
                logger.warning(
                    "Failed to delete from Knowledge Graph during bulk delete",
                    document_id=str(doc.id),
                    error=str(e),
                )

            # 3. Delete chunks from SQLite
            await db.execute(
                delete(Chunk).where(Chunk.document_id == doc.id)
            )
            # 4. Delete document from SQLite
            await db.delete(doc)
        else:
            # Soft delete
            doc.processing_status = ProcessingStatus.FAILED
            doc.processing_error = "Deleted by user"

        deleted += 1

    # Check for documents not found
    found_ids = {str(doc.id) for doc in documents}
    for doc_id in request.document_ids:
        if doc_id not in found_ids:
            failed += 1
            errors.append({
                "document_id": doc_id,
                "error": "Document not found",
            })

    await db.commit()

    logger.info(
        "Bulk delete completed",
        deleted=deleted,
        failed=failed,
        permanent=request.permanent,
    )

    return BulkDeleteResponse(
        deleted=deleted,
        failed=failed,
        permanent=request.permanent,
        errors=errors,
    )


# =============================================================================
# Bulk Tag Update
# =============================================================================

class BulkTagRequest(BaseModel):
    """Request to update tags on multiple documents."""
    document_ids: List[str] = Field(..., min_length=1, max_length=100, description="List of document IDs")
    add_tags: List[str] = Field(default_factory=list, description="Tags to add to all documents")
    remove_tags: List[str] = Field(default_factory=list, description="Tags to remove from all documents")
    set_collection: Optional[str] = Field(None, description="Set collection for all documents")


class BulkTagResponse(BaseModel):
    """Response for bulk tag operation."""
    updated: int
    failed: int
    errors: List[Dict]


@router.post("/bulk/tags", response_model=BulkTagResponse)
async def bulk_update_tags(
    request: BulkTagRequest,
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
):
    """
    Update tags and/or collection on multiple documents.

    - add_tags: Tags to add to each document
    - remove_tags: Tags to remove from each document
    - set_collection: Set the collection field for all documents
    """
    import uuid as uuid_module
    import json

    user_tier = user.access_tier_level

    # Get all documents with organization filtering
    doc_uuids = []
    for doc_id in request.document_ids:
        try:
            doc_uuids.append(uuid_module.UUID(doc_id))
        except ValueError:
            pass

    query = (
        select(Document)
        .options(selectinload(Document.access_tier))
        .where(Document.id.in_(doc_uuids))
    )

    # PHASE 12 FIX: Apply organization filtering for multi-tenant isolation
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        query = query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),  # Include legacy/shared docs
            )
        )

    result = await db.execute(query)
    documents = result.scalars().all()

    updated = 0
    failed = 0
    errors = []

    for doc in documents:
        doc_tier = doc.access_tier.level if doc.access_tier else 1

        if user_tier < doc_tier:
            failed += 1
            errors.append({
                "document_id": str(doc.id),
                "error": "No access to this document",
            })
            continue

        # Update tags
        if request.add_tags or request.remove_tags:
            # Parse existing tags
            existing_tags = set()
            if doc.tags:
                try:
                    existing_tags = set(json.loads(doc.tags)) if isinstance(doc.tags, str) else set(doc.tags)
                except (json.JSONDecodeError, TypeError):
                    existing_tags = set()

            # Add new tags
            existing_tags.update(request.add_tags)

            # Remove tags
            existing_tags -= set(request.remove_tags)

            # Save back
            doc.tags = json.dumps(list(existing_tags))

        # Update collection
        if request.set_collection is not None:
            doc.collection = request.set_collection if request.set_collection else None

        updated += 1

    # Check for documents not found
    found_ids = {str(doc.id) for doc in documents}
    for doc_id in request.document_ids:
        if doc_id not in found_ids:
            failed += 1
            errors.append({
                "document_id": doc_id,
                "error": "Document not found",
            })

    await db.commit()

    logger.info(
        "Bulk tag update completed",
        updated=updated,
        failed=failed,
        add_tags=request.add_tags,
        remove_tags=request.remove_tags,
    )

    return BulkTagResponse(
        updated=updated,
        failed=failed,
        errors=errors,
    )


# =============================================================================
# Image Analysis Endpoints
# =============================================================================

class VisionStatusResponse(BaseModel):
    """Vision model status response."""
    available: bool
    provider: Optional[str] = None
    model: Optional[str] = None
    issues: List[str] = []
    recommendation: Optional[str] = None


class ImageReanalysisRequest(BaseModel):
    """Request to reanalyze images in a document."""
    force_reanalyze: bool = Field(
        default=False,
        description="Re-analyze even if already completed"
    )
    skip_duplicates: bool = Field(
        default=True,
        description="Skip images already analyzed elsewhere (uses cached captions)"
    )


class ImageAnalysisStatsResponse(BaseModel):
    """Image analysis statistics."""
    images_found: int = 0
    images_analyzed: int = 0
    images_skipped_small: int = 0
    images_skipped_duplicate: int = 0
    images_failed: int = 0
    cached_used: int = 0
    newly_analyzed: int = 0
    total_time_ms: int = 0


class ImageReanalysisResponse(BaseModel):
    """Response from image reanalysis."""
    success: bool
    document_id: str
    status: str
    stats: ImageAnalysisStatsResponse
    error: Optional[str] = None


class UnanalyzedDocumentsResponse(BaseModel):
    """Response listing documents with unanalyzed images."""
    documents: List[Dict]
    total: int
    page: int
    page_size: int


@router.get("/vision/status", response_model=VisionStatusResponse)
async def get_vision_status(
    user: AuthenticatedUser,
):
    """
    Check if vision model is configured and ready.

    Returns status of vision model availability including:
    - Which provider is configured (ollama, openai, anthropic)
    - Which model is being used
    - Any issues detected
    - Recommendations if not configured
    """
    from backend.services.image_analysis import get_image_analysis_service

    service = get_image_analysis_service()
    status = await service.check_vision_available()

    return VisionStatusResponse(
        available=status.available,
        provider=status.provider,
        model=status.model,
        issues=status.issues,
        recommendation=status.recommendation,
    )


@router.get("/unanalyzed-images", response_model=UnanalyzedDocumentsResponse)
async def list_documents_without_image_analysis(
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
    page: int = Query(default=1, ge=1),
    page_size: int = Query(default=20, ge=1, le=100),
):
    """
    List documents that have images but haven't been analyzed.

    Returns documents where:
    - images_extracted_count > 0
    - image_analysis_status is pending, skipped, or failed

    Respects user access tier and organization filtering.
    """
    # Base query for documents with unanalyzed images
    base_query = (
        select(Document)
        .join(AccessTier, Document.access_tier_id == AccessTier.id)
        .where(
            and_(
                AccessTier.level <= user.access_tier_level,
                Document.images_extracted_count > 0,
                Document.image_analysis_status.in_(["pending", "skipped", "failed"]),
            )
        )
    )

    # Apply organization filtering
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        base_query = base_query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),
            )
        )

    # Apply private document filtering
    if not user.is_superadmin:
        base_query = base_query.where(
            or_(
                Document.is_private == False,
                and_(
                    Document.is_private == True,
                    Document.uploaded_by_id == UUID(user.user_id),
                ),
            )
        )

    # Get total count
    count_query = select(func.count()).select_from(base_query.subquery())
    total_result = await db.execute(count_query)
    total = total_result.scalar() or 0

    # Apply pagination
    offset = (page - 1) * page_size
    query = base_query.order_by(desc(Document.created_at)).offset(offset).limit(page_size)

    result = await db.execute(query)
    documents = result.scalars().all()

    # Build response
    doc_list = []
    for doc in documents:
        doc_list.append({
            "id": str(doc.id),
            "name": doc.original_filename or doc.filename,
            "file_type": doc.file_type,
            "images_extracted_count": doc.images_extracted_count,
            "images_analyzed_count": doc.images_analyzed_count or 0,
            "image_analysis_status": doc.image_analysis_status,
            "image_analysis_error": doc.image_analysis_error,
            "created_at": doc.created_at.isoformat() if doc.created_at else None,
        })

    return UnanalyzedDocumentsResponse(
        documents=doc_list,
        total=total,
        page=page,
        page_size=page_size,
    )


@router.post("/{document_id}/reanalyze-images", response_model=ImageReanalysisResponse)
async def reanalyze_document_images(
    document_id: UUID,
    request: ImageReanalysisRequest,
    user: AuthenticatedUser,
    db: AsyncSession = Depends(get_async_session),
):
    """
    Reanalyze only images in a document (faster than full reprocess).

    This endpoint:
    - Extracts images from the stored file
    - Runs vision analysis on each image
    - Uses cached captions for duplicate images (if skip_duplicates=True)
    - Updates document image statistics

    Requires write permission on the document.
    """
    from backend.services.image_analysis import get_image_analysis_service

    logger.info(
        "Reanalyzing document images",
        document_id=str(document_id),
        user_id=user.user_id,
        force=request.force_reanalyze,
    )

    # Check write permission
    permission_service = get_permission_service()
    has_write = await permission_service.check_document_access(
        user_context=user,
        document_id=str(document_id),
        permission=Permission.WRITE,
        session=db,
    )

    if not has_write:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have write access to this document",
        )

    # Get document with organization filtering
    query = (
        select(Document)
        .options(selectinload(Document.access_tier))
        .where(Document.id == document_id)
    )

    # Apply organization filtering
    org_id = get_org_filter(user)
    if org_id and not user.is_superadmin:
        query = query.where(
            or_(
                Document.organization_id == org_id,
                Document.organization_id.is_(None),
            )
        )

    # Apply private document filtering
    if not user.is_superadmin:
        query = query.where(
            or_(
                Document.is_private == False,
                and_(
                    Document.is_private == True,
                    Document.uploaded_by_id == UUID(user.user_id),
                ),
            )
        )

    result = await db.execute(query)
    document = result.scalar_one_or_none()

    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document not found",
        )

    # Check vision availability
    service = get_image_analysis_service()
    vision_status = await service.check_vision_available()

    if not vision_status.available:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail=f"Vision model not configured. {vision_status.recommendation or ''}",
        )

    # Run image analysis
    analysis_result = await service.analyze_document_images(
        document_id=str(document_id),
        force=request.force_reanalyze,
        skip_duplicates=request.skip_duplicates,
        db=db,
    )

    logger.info(
        "Document image reanalysis completed",
        document_id=str(document_id),
        success=analysis_result.success,
        stats=analysis_result.stats.__dict__,
    )

    return ImageReanalysisResponse(
        success=analysis_result.success,
        document_id=str(document_id),
        status="completed" if analysis_result.success else "failed",
        stats=ImageAnalysisStatsResponse(
            images_found=analysis_result.stats.images_found,
            images_analyzed=analysis_result.stats.images_analyzed,
            images_skipped_small=analysis_result.stats.images_skipped_small,
            images_skipped_duplicate=analysis_result.stats.images_skipped_duplicate,
            images_failed=analysis_result.stats.images_failed,
            cached_used=analysis_result.stats.cached_used,
            newly_analyzed=analysis_result.stats.newly_analyzed,
            total_time_ms=analysis_result.stats.total_time_ms,
        ),
        error=analysis_result.error,
    )

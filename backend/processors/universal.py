"""
AIDocumentIndexer - Universal Document Processor
================================================

Processes all supported file types using appropriate extractors.
Supports PDF, Office documents, images, and more.
"""

import os
import mimetypes
from pathlib import Path
from typing import Optional, List, Dict, Any
from dataclasses import dataclass, field

import structlog

logger = structlog.get_logger(__name__)


@dataclass
class ExtractedImage:
    """Represents an extracted image from a document."""
    page_number: Optional[int] = None
    image_index: int = 0
    image_bytes: bytes = field(default_factory=bytes)
    extension: str = "png"  # png, jpg, etc.
    width: Optional[int] = None
    height: Optional[int] = None
    alt_text: Optional[str] = None  # If available from source
    source_type: str = "embedded"  # embedded, inline, background


@dataclass
class ExtractedContent:
    """Extracted content from a document."""
    text: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    pages: List[Dict[str, Any]] = field(default_factory=list)
    images: List[Dict[str, Any]] = field(default_factory=list)
    extracted_images: List[ExtractedImage] = field(default_factory=list)  # New: actual image data
    tables: List[Dict[str, Any]] = field(default_factory=list)
    word_count: int = 0
    page_count: int = 0
    language: str = "en"


@dataclass
class ImageOptimizationConfig:
    """Configuration for image optimization during processing."""
    # Zoom levels for OCR rendering
    # Higher zoom = better OCR quality for scanned documents
    default_zoom: float = 2.5  # Increased from 1.5 for better OCR quality
    large_pdf_zoom: float = 2.0  # Increased from 1.0 for better quality on large files

    # Size thresholds
    large_file_threshold_mb: float = 10.0  # Files > 10MB use large_pdf_zoom

    # Image optimization
    compress_before_ocr: bool = True  # Compress images before OCR
    max_image_dimension: int = 3000  # Increased from 2000 for better OCR
    convert_to_grayscale: bool = False  # Disabled - color can improve OCR accuracy


class UniversalProcessor:
    """
    Universal document processor supporting multiple file formats.

    Automatically selects the appropriate extractor based on file type.
    """

    # Mapping from Tesseract language codes to PaddleOCR language codes
    # Tesseract uses ISO 639-2 (3-letter), PaddleOCR uses various formats
    TESSERACT_TO_PADDLE_LANG = {
        "eng": "en",
        "deu": "de",
        "fra": "fr",
        "spa": "es",
        "ita": "it",
        "por": "pt",
        "nld": "nl",
        "pol": "pl",
        "rus": "ru",
        "jpn": "japan",
        "chi_sim": "ch",
        "chi_tra": "chinese_cht",
        "kor": "korean",
        "ara": "ar",
    }

    # File type to processor mapping
    PROCESSORS = {
        # PDF
        ".pdf": "_process_pdf",
        # Office documents
        ".docx": "_process_docx",
        ".doc": "_process_doc",
        ".pptx": "_process_pptx",
        ".ppt": "_process_ppt",
        ".xlsx": "_process_xlsx",
        ".xls": "_process_xls",
        # Open formats
        ".odt": "_process_odt",
        ".odp": "_process_odp",
        ".ods": "_process_ods",
        # Text files
        ".txt": "_process_text",
        ".md": "_process_text",
        ".rtf": "_process_rtf",
        ".html": "_process_html",
        ".htm": "_process_html",
        ".xml": "_process_text",
        ".json": "_process_text",
        # Images (OCR)
        ".png": "_process_image",
        ".jpg": "_process_image",
        ".jpeg": "_process_image",
        ".gif": "_process_image",
        ".webp": "_process_image",
        ".bmp": "_process_image",
        ".tiff": "_process_image",
        # Email
        ".eml": "_process_email",
        ".msg": "_process_msg",
        # CSV
        ".csv": "_process_csv",
    }

    def __init__(
        self,
        enable_ocr: bool = True,
        enable_image_analysis: bool = True,
        ocr_language: str = None,  # Will use OCR_LANGUAGE env var or default
        smart_image_handling: bool = True,
        optimization_config: Optional[ImageOptimizationConfig] = None,
        parallel_ocr: bool = True,  # Enable parallel OCR for multi-page PDFs
        max_ocr_workers: int = 4,  # Max parallel OCR workers
    ):
        """
        Initialize the universal processor.

        Args:
            enable_ocr: Enable OCR for scanned documents
            enable_image_analysis: Enable image description generation
            ocr_language: Language for OCR (uses OCR_LANGUAGE env var or defaults to 'eng')
            smart_image_handling: Optimize image quality based on content
            optimization_config: Image optimization configuration
            parallel_ocr: Enable parallel OCR for multi-page documents
            max_ocr_workers: Maximum number of parallel OCR workers
        """
        import os
        import hashlib
        self.enable_ocr = enable_ocr
        self.enable_image_analysis = enable_image_analysis
        # Use provided language, env var, or default to English
        self.ocr_language = ocr_language or os.getenv("OCR_LANGUAGE", "eng")
        self.smart_image_handling = smart_image_handling
        self.optimization_config = optimization_config or ImageOptimizationConfig()
        self.parallel_ocr = parallel_ocr
        self.max_ocr_workers = max_ocr_workers

        # Lazy-load heavy dependencies
        self._ocr_engine = None
        self._easyocr_engine = None
        self._pdf_processor = None

        # Image OCR cache for duplicate detection within a document
        # Key: image hash, Value: OCR result text
        self._image_ocr_cache: Dict[str, str] = {}

    def process(
        self,
        file_path: str,
        processing_mode: str = "full",
    ) -> ExtractedContent:
        """
        Process a document and extract content.

        Args:
            file_path: Path to the document
            processing_mode: "full", "ocr", or "basic"

        Returns:
            ExtractedContent: Extracted text and metadata
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        extension = path.suffix.lower()

        if extension not in self.PROCESSORS:
            raise ValueError(f"Unsupported file type: {extension}")

        processor_method = getattr(self, self.PROCESSORS[extension])

        logger.info(
            "Processing document",
            file_path=file_path,
            extension=extension,
            processing_mode=processing_mode,
        )

        # Clear OCR cache for new document (duplicate detection is per-document)
        self.clear_ocr_cache()

        try:
            content = processor_method(file_path, processing_mode)
            content.word_count = len(content.text.split())
            return content
        except Exception as e:
            logger.error(
                "Error processing document",
                file_path=file_path,
                error=str(e),
            )
            raise

    # =========================================================================
    # PDF Processing
    # =========================================================================

    def _process_pdf(self, file_path: str, mode: str) -> ExtractedContent:
        """Process PDF document using PyMuPDF with optional parallel OCR."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.warning("PyMuPDF not installed, using fallback")
            return self._process_pdf_fallback(file_path, mode)

        # Get file size for smart optimization
        file_size_mb = Path(file_path).stat().st_size / (1024 * 1024)

        doc = fitz.open(file_path)
        pages = []
        all_text = []
        images = []
        extracted_images: List[ExtractedImage] = []

        # Identify pages that need OCR
        pages_needing_ocr = []
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text()
            all_text.append(text)

            if self.enable_ocr and len(text.strip()) < 50:
                pages_needing_ocr.append((page_num - 1, page))  # 0-indexed

            pages.append({
                "page_number": page_num,
                "text": text,
                "has_images": len(page.get_images()) > 0,
            })

            # Extract images if not basic mode
            if mode != "basic":
                for img_index, img in enumerate(page.get_images()):
                    images.append({
                        "page_number": page_num,
                        "index": img_index,
                        "xref": img[0],
                    })

                    # Extract actual image bytes for multimodal processing
                    if self.enable_image_analysis:
                        try:
                            extracted_img = self._extract_pdf_image(doc, img[0], page_num, img_index)
                            if extracted_img:
                                extracted_images.append(extracted_img)
                        except Exception as e:
                            logger.debug(f"Failed to extract image xref={img[0]}: {e}")

        # Process OCR pages in parallel if there are multiple pages needing OCR
        if pages_needing_ocr:
            if len(pages_needing_ocr) > 1 and hasattr(self, 'parallel_ocr') and self.parallel_ocr:
                # Use parallel OCR for multiple pages
                ocr_results = self._ocr_pages_parallel(
                    file_path,
                    [idx for idx, _ in pages_needing_ocr],
                    file_size_mb
                )
                for idx, ocr_text in ocr_results.items():
                    if ocr_text:
                        all_text[idx] = ocr_text
                        pages[idx]["text"] = ocr_text
            else:
                # Sequential OCR (original behavior)
                for idx, page in pages_needing_ocr:
                    ocr_text = self._ocr_page(page, file_size_mb)
                    if ocr_text:
                        all_text[idx] = ocr_text
                        pages[idx]["text"] = ocr_text

        doc.close()

        return ExtractedContent(
            text="\n\n".join(all_text),
            metadata={
                "file_type": "pdf",
                "file_path": file_path,
                "file_size_mb": file_size_mb,
                "parallel_ocr": getattr(self, 'parallel_ocr', False),
                "images_extracted": len(extracted_images),
            },
            pages=pages,
            images=images,
            extracted_images=extracted_images,
            page_count=len(pages),
        )

    def _extract_pdf_image(
        self,
        doc,
        xref: int,
        page_number: int,
        image_index: int,
    ) -> Optional[ExtractedImage]:
        """
        Extract a single image from a PDF by xref.

        Args:
            doc: PyMuPDF document object
            xref: Image cross-reference number
            page_number: Page number where image is located
            image_index: Index of image on the page

        Returns:
            ExtractedImage or None if extraction fails
        """
        try:
            base_image = doc.extract_image(xref)
            if not base_image:
                return None

            image_bytes = base_image.get("image")
            if not image_bytes or len(image_bytes) < 100:  # Skip tiny images (likely artifacts)
                return None

            extension = base_image.get("ext", "png")
            width = base_image.get("width")
            height = base_image.get("height")

            # Skip very small images (likely icons/bullets)
            if width and height and (width < 50 or height < 50):
                return None

            return ExtractedImage(
                page_number=page_number,
                image_index=image_index,
                image_bytes=image_bytes,
                extension=extension,
                width=width,
                height=height,
                source_type="embedded",
            )

        except Exception as e:
            logger.debug(f"Failed to extract PDF image: {e}")
            return None

    def _process_pdf_fallback(self, file_path: str, mode: str) -> ExtractedContent:
        """Fallback PDF processing without PyMuPDF."""
        # Try pypdf as fallback
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
            return ExtractedContent(
                text=text,
                metadata={"file_type": "pdf", "file_path": file_path},
                page_count=len(reader.pages),
            )
        except ImportError:
            raise ImportError("Neither PyMuPDF nor pypdf is installed")

    # =========================================================================
    # Office Documents
    # =========================================================================

    def _process_docx(self, file_path: str, mode: str) -> ExtractedContent:
        """Process DOCX using python-docx."""
        try:
            from docx import Document
        except ImportError:
            return self._unstructured_fallback(file_path)

        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # Extract tables
        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append({"data": table_data})

        # Extract images if not basic mode
        extracted_images: List[ExtractedImage] = []
        if mode != "basic" and self.enable_image_analysis:
            extracted_images = self._extract_docx_images(doc)

        return ExtractedContent(
            text="\n\n".join(paragraphs),
            metadata={
                "file_type": "docx",
                "file_path": file_path,
                "images_extracted": len(extracted_images),
            },
            tables=tables,
            extracted_images=extracted_images,
            page_count=1,  # DOCX doesn't have page concept in API
        )

    def _extract_docx_images(self, doc) -> List[ExtractedImage]:
        """
        Extract images from a DOCX document.

        Args:
            doc: python-docx Document object

        Returns:
            List of ExtractedImage objects
        """
        extracted_images = []
        image_index = 0

        try:
            # Access the document's relationships to find embedded images
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        image_part = rel.target_part
                        image_bytes = image_part.blob

                        # Skip tiny images
                        if len(image_bytes) < 100:
                            continue

                        # Determine extension from content type
                        content_type = image_part.content_type
                        extension = content_type.split("/")[-1] if "/" in content_type else "png"
                        # Normalize extension
                        if extension == "jpeg":
                            extension = "jpg"

                        extracted_images.append(ExtractedImage(
                            page_number=None,  # DOCX doesn't have page concept
                            image_index=image_index,
                            image_bytes=image_bytes,
                            extension=extension,
                            source_type="embedded",
                        ))
                        image_index += 1

                    except Exception as e:
                        logger.debug(f"Failed to extract DOCX image: {e}")
                        continue

        except Exception as e:
            logger.warning(f"Error extracting DOCX images: {e}")

        logger.debug(f"Extracted {len(extracted_images)} images from DOCX")
        return extracted_images

    def _process_doc(self, file_path: str, mode: str) -> ExtractedContent:
        """Process legacy DOC files."""
        return self._unstructured_fallback(file_path)

    def _process_pptx(self, file_path: str, mode: str) -> ExtractedContent:
        """Process PPTX using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            return self._unstructured_fallback(file_path)

        prs = Presentation(file_path)
        slides = []
        all_text = []
        extracted_images: List[ExtractedImage] = []
        image_index = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                # Extract table text
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text for cell in row.cells]
                        slide_text.append(" | ".join(row_text))

                # Extract images if not basic mode
                if mode != "basic" and self.enable_image_analysis:
                    if hasattr(shape, "image"):
                        try:
                            img_blob = shape.image.blob
                            if img_blob and len(img_blob) > 100:  # Skip tiny images
                                extension = shape.image.ext or "png"
                                if extension == "jpeg":
                                    extension = "jpg"

                                extracted_images.append(ExtractedImage(
                                    page_number=slide_num,
                                    image_index=image_index,
                                    image_bytes=img_blob,
                                    extension=extension,
                                    source_type="embedded",
                                ))
                                image_index += 1
                        except Exception as e:
                            logger.debug(f"Failed to extract PPTX image: {e}")

            text = "\n".join(slide_text)
            slides.append({
                "page_number": slide_num,
                "text": text,
                "has_images": any(hasattr(s, "image") for s in slide.shapes),
            })
            all_text.append(f"[Slide {slide_num}]\n{text}")

        return ExtractedContent(
            text="\n\n".join(all_text),
            metadata={
                "file_type": "pptx",
                "file_path": file_path,
                "images_extracted": len(extracted_images),
            },
            pages=slides,
            extracted_images=extracted_images,
            page_count=len(slides),
        )

    def _process_ppt(self, file_path: str, mode: str) -> ExtractedContent:
        """Process legacy PPT files."""
        return self._unstructured_fallback(file_path)

    def _process_xlsx(self, file_path: str, mode: str) -> ExtractedContent:
        """Process XLSX using openpyxl."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            return self._unstructured_fallback(file_path)

        wb = load_workbook(file_path, read_only=True, data_only=True)
        all_text = []
        tables = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_data = []

            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(row_values):  # Skip empty rows
                    sheet_data.append(row_values)

            if sheet_data:
                all_text.append(f"[Sheet: {sheet_name}]")
                for row in sheet_data:
                    all_text.append(" | ".join(row))

                tables.append({
                    "sheet_name": sheet_name,
                    "data": sheet_data,
                })

        wb.close()

        return ExtractedContent(
            text="\n".join(all_text),
            metadata={"file_type": "xlsx", "file_path": file_path},
            tables=tables,
            page_count=len(wb.sheetnames),
        )

    def _process_xls(self, file_path: str, mode: str) -> ExtractedContent:
        """Process legacy XLS files."""
        return self._unstructured_fallback(file_path)

    # =========================================================================
    # Open Document Formats
    # =========================================================================

    def _process_odt(self, file_path: str, mode: str) -> ExtractedContent:
        """Process ODT files."""
        return self._unstructured_fallback(file_path)

    def _process_odp(self, file_path: str, mode: str) -> ExtractedContent:
        """Process ODP files."""
        return self._unstructured_fallback(file_path)

    def _process_ods(self, file_path: str, mode: str) -> ExtractedContent:
        """Process ODS files."""
        return self._unstructured_fallback(file_path)

    # =========================================================================
    # Text Files
    # =========================================================================

    def _process_text(self, file_path: str, mode: str) -> ExtractedContent:
        """Process plain text files."""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

        return ExtractedContent(
            text=text,
            metadata={"file_type": "text", "file_path": file_path},
            page_count=1,
        )

    def _process_rtf(self, file_path: str, mode: str) -> ExtractedContent:
        """Process RTF files."""
        try:
            from striprtf.striprtf import rtf_to_text
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                rtf_content = f.read()
            text = rtf_to_text(rtf_content)
            return ExtractedContent(
                text=text,
                metadata={"file_type": "rtf", "file_path": file_path},
                page_count=1,
            )
        except ImportError:
            return self._unstructured_fallback(file_path)

    def _process_html(self, file_path: str, mode: str) -> ExtractedContent:
        """Process HTML files."""
        try:
            from bs4 import BeautifulSoup
            import base64
            import re

            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f.read(), "html.parser")

            # Extract images if not basic mode
            extracted_images: List[ExtractedImage] = []
            if mode != "basic" and self.enable_image_analysis:
                image_index = 0
                for img in soup.find_all("img"):
                    src = img.get("src", "")
                    alt = img.get("alt", "")

                    # Handle base64-encoded inline images
                    if src.startswith("data:image/"):
                        try:
                            # Parse data:image/png;base64,... format
                            match = re.match(r"data:image/(\w+);base64,(.+)", src)
                            if match:
                                extension = match.group(1)
                                if extension == "jpeg":
                                    extension = "jpg"
                                image_data = base64.b64decode(match.group(2))

                                if len(image_data) > 100:  # Skip tiny images
                                    extracted_images.append(ExtractedImage(
                                        page_number=None,
                                        image_index=image_index,
                                        image_bytes=image_data,
                                        extension=extension,
                                        alt_text=alt if alt else None,
                                        source_type="inline",
                                    ))
                                    image_index += 1
                        except Exception as e:
                            logger.debug(f"Failed to extract HTML inline image: {e}")

            # Remove script and style elements
            for element in soup(["script", "style", "nav", "footer"]):
                element.decompose()

            text = soup.get_text(separator="\n", strip=True)
            return ExtractedContent(
                text=text,
                metadata={
                    "file_type": "html",
                    "file_path": file_path,
                    "images_extracted": len(extracted_images),
                },
                extracted_images=extracted_images,
                page_count=1,
            )
        except ImportError:
            # Fallback: basic text extraction
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            return ExtractedContent(
                text=text,
                metadata={"file_type": "html", "file_path": file_path},
                page_count=1,
            )

    # =========================================================================
    # Images (OCR)
    # =========================================================================

    def _process_image(self, file_path: str, mode: str) -> ExtractedContent:
        """Process image using OCR and optional VLM captioning."""
        if mode == "basic":
            return ExtractedContent(
                text="",
                metadata={"file_type": "image", "file_path": file_path, "skipped": True},
                page_count=1,
            )

        # OCR the image for text extraction
        text = self._ocr_image(file_path)

        # Also include the image for VLM captioning (like PDFs do)
        extracted_images: List[ExtractedImage] = []
        if self.enable_image_analysis:
            try:
                with open(file_path, "rb") as f:
                    image_bytes = f.read()

                # Determine extension from file path
                import os
                ext = os.path.splitext(file_path)[1].lower().lstrip(".")
                if ext == "jpeg":
                    ext = "jpg"

                if len(image_bytes) > 1000:  # Skip tiny images (< 1KB)
                    extracted_images.append(ExtractedImage(
                        page_number=1,
                        image_index=0,
                        image_bytes=image_bytes,
                        extension=ext or "png",
                        alt_text=None,
                        source_type="standalone",
                    ))
                    logger.debug(
                        "Image added for VLM captioning",
                        file_path=file_path,
                        size_kb=len(image_bytes) // 1024,
                    )
            except Exception as e:
                logger.warning("Failed to read image for VLM", error=str(e))

        return ExtractedContent(
            text=text,
            metadata={
                "file_type": "image",
                "file_path": file_path,
                "ocr_text_length": len(text),
                "images_extracted": len(extracted_images),
            },
            extracted_images=extracted_images,
            page_count=1,
        )

    def _convert_ocr_language(self, tesseract_lang: str) -> str:
        """
        Convert Tesseract language code to PaddleOCR format.

        Tesseract uses ISO 639-2 codes (e.g., 'eng', 'deu') while
        PaddleOCR uses various formats (e.g., 'en', 'de', 'german').

        Args:
            tesseract_lang: Language code in Tesseract format (e.g., 'deu+eng')

        Returns:
            PaddleOCR compatible language code
        """
        # Handle multi-language (e.g., "deu+eng" -> use first language)
        primary_lang = tesseract_lang.split("+")[0].strip()

        # Convert using mapping, fallback to original if not found
        paddle_lang = self.TESSERACT_TO_PADDLE_LANG.get(primary_lang, primary_lang)

        logger.debug(
            "Converted OCR language",
            tesseract=tesseract_lang,
            paddle=paddle_lang,
        )
        return paddle_lang

    def _ocr_image(self, image_path: str) -> str:
        """Perform OCR on an image file."""
        if not self.enable_ocr:
            return ""

        import os
        import time

        # Check settings for OCR provider (default to tesseract to avoid PaddleOCR memory issues)
        ocr_provider = os.getenv("OCR_PROVIDER", "tesseract").lower()

        fallback_used = False
        start_time = time.time()
        success = True
        error_message = None
        text = ""
        provider = ocr_provider

        # If provider is tesseract, skip PaddleOCR entirely
        if ocr_provider == "tesseract":
            try:
                text = self._ocr_tesseract(image_path)
                return text
            except Exception as e:
                logger.error("Tesseract OCR failed", error=str(e))
                return ""

        try:
            from paddleocr import PaddleOCR

            if self._ocr_engine is None:
                # Set model cache directory from environment or use project directory default
                default_paddle_home = os.path.join(os.getcwd(), "data", "paddle_models")
                paddle_home = os.getenv("PADDLEX_HOME", default_paddle_home)
                paddle_hub = os.getenv("PADDLE_HUB_HOME", os.path.join(paddle_home, "official_models"))

                # Ensure directories exist
                os.makedirs(paddle_home, exist_ok=True)
                os.makedirs(paddle_hub, exist_ok=True)

                # Set environment variables before PaddleOCR initialization
                os.environ.setdefault("PADDLEX_HOME", paddle_home)
                os.environ.setdefault("PADDLE_HUB_HOME", paddle_hub)

                # Use HuggingFace mirror if specified
                if os.getenv("PADDLE_PDX_MODEL_SOURCE"):
                    os.environ["PADDLE_PDX_MODEL_SOURCE"] = os.getenv("PADDLE_PDX_MODEL_SOURCE")

                # Convert Tesseract language code to PaddleOCR format
                paddle_lang = self._convert_ocr_language(self.ocr_language)

                # Initialize PaddleOCR
                # Note: use_textline_orientation replaces deprecated use_angle_cls
                # Disable doc preprocessing to avoid UVDoc model issues
                self._ocr_engine = PaddleOCR(
                    use_textline_orientation=True,
                    use_doc_orientation_classify=False,
                    use_doc_unwarping=False,
                    lang=paddle_lang,
                )

                logger.info(
                    "PaddleOCR initialized",
                    language=paddle_lang,
                    cache_dir=paddle_home,
                )

            # Note: cls parameter removed - text orientation is configured via use_textline_orientation in init
            result = self._ocr_engine.ocr(image_path)

            if result and result[0]:
                text_lines = [line[1][0] for line in result[0]]
                text = "\n".join(text_lines)
            else:
                text = ""

        except ImportError:
            logger.warning("PaddleOCR not installed, trying Tesseract")
            provider = "tesseract"
            fallback_used = True
            text = self._ocr_tesseract(image_path)
        except Exception as e:
            logger.warning("PaddleOCR failed, falling back to Tesseract", error=str(e))
            provider = "tesseract"
            fallback_used = True
            error_message = str(e)
            try:
                text = self._ocr_tesseract(image_path)
            except Exception as fallback_error:
                success = False
                error_message = f"PaddleOCR: {str(e)}, Tesseract: {str(fallback_error)}"
        finally:
            # Record metrics (async operation, fire-and-forget)
            processing_time_ms = int((time.time() - start_time) * 1000)
            self._record_ocr_metrics(
                provider=provider,
                processing_time_ms=processing_time_ms,
                success=success,
                character_count=len(text) if text else 0,
                error_message=error_message,
                fallback_used=fallback_used,
            )

        return text

    def _record_ocr_metrics(
        self,
        provider: str,
        processing_time_ms: int,
        success: bool,
        character_count: int = 0,
        error_message: str = None,
        fallback_used: bool = False,
    ):
        """
        Record OCR metrics for analytics.

        This is a fire-and-forget async operation that doesn't block OCR processing.
        """
        try:
            import asyncio
            from backend.db.database import get_async_session_context
            from backend.services.ocr_metrics import OCRMetricsService

            # Convert language code to standard format
            paddle_lang = self._convert_ocr_language(self.ocr_language)

            async def record_metrics():
                try:
                    async with get_async_session_context() as session:
                        metrics_service = OCRMetricsService(session)
                        await metrics_service.record_ocr_operation(
                            provider=provider,
                            language=paddle_lang,
                            processing_time_ms=processing_time_ms,
                            success=success,
                            character_count=character_count,
                            error_message=error_message,
                            fallback_used=fallback_used,
                        )
                except Exception as e:
                    logger.debug("Failed to record OCR metrics", error=str(e))

            # Run in background without blocking
            try:
                loop = asyncio.get_running_loop()
                loop.create_task(record_metrics())
            except RuntimeError:
                # No event loop running, skip metrics
                logger.debug("No event loop running, skipping OCR metrics")

        except Exception as e:
            # Metrics recording should never break OCR processing
            logger.debug("Failed to record OCR metrics", error=str(e))

    def _ocr_tesseract(self, image_path: str) -> str:
        """Fallback OCR using Tesseract."""
        try:
            import os
            import pytesseract
            from PIL import Image

            # Ensure TESSDATA_PREFIX is set for pytesseract to find language files
            tessdata_prefix = os.getenv("TESSDATA_PREFIX")
            if tessdata_prefix and os.path.isdir(tessdata_prefix):
                os.environ["TESSDATA_PREFIX"] = tessdata_prefix
                logger.debug("Using TESSDATA_PREFIX", path=tessdata_prefix)

            img = Image.open(image_path)
            text = pytesseract.image_to_string(img, lang=self.ocr_language)
            return text
        except ImportError:
            logger.warning("Neither PaddleOCR nor pytesseract installed")
            return ""
        except Exception as e:
            logger.error("Tesseract OCR failed", error=str(e), language=self.ocr_language)
            return ""

    def _ocr_easyocr(self, image_path: str) -> str:
        """OCR using EasyOCR."""
        try:
            import easyocr

            if self._easyocr_engine is None:
                # Map common language codes to EasyOCR format
                lang_map = {
                    "en": "en",
                    "de": "de",
                    "fr": "fr",
                    "es": "es",
                    "zh": "ch_sim",
                    "ja": "ja",
                    "ko": "ko",
                    "ar": "ar",
                }

                # Convert Tesseract language code to EasyOCR format
                easyocr_lang = lang_map.get(self.ocr_language, "en")

                # Initialize EasyOCR Reader
                # Use GPU if available (will fall back to CPU automatically)
                self._easyocr_engine = easyocr.Reader([easyocr_lang], gpu=True)

                logger.info("EasyOCR initialized", language=easyocr_lang)

            # Perform OCR
            result = self._easyocr_engine.readtext(image_path)

            # Extract text from results
            if result:
                text_lines = [detection[1] for detection in result]
                text = "\n".join(text_lines)
            else:
                text = ""

            return text

        except ImportError:
            logger.warning("EasyOCR not installed")
            return ""
        except Exception as e:
            logger.error("EasyOCR failed", error=str(e), language=self.ocr_language)
            return ""

    def _compute_image_hash(self, img_data: bytes) -> str:
        """
        Compute a hash of image data for duplicate detection.

        Uses SHA-256 for fast, reliable hashing.
        """
        import hashlib
        return hashlib.sha256(img_data).hexdigest()

    def _get_cached_ocr(self, img_data: bytes) -> Optional[str]:
        """
        Check if OCR result for this image is already cached.

        Args:
            img_data: Raw image bytes

        Returns:
            Cached OCR text if found, None otherwise
        """
        img_hash = self._compute_image_hash(img_data)
        cached = self._image_ocr_cache.get(img_hash)
        if cached is not None:
            logger.debug(
                "OCR cache hit - reusing result for duplicate image",
                hash=img_hash[:16],
            )
        return cached

    def _cache_ocr_result(self, img_data: bytes, ocr_text: str) -> None:
        """
        Cache OCR result for an image.

        Args:
            img_data: Raw image bytes
            ocr_text: OCR result text
        """
        img_hash = self._compute_image_hash(img_data)
        self._image_ocr_cache[img_hash] = ocr_text
        logger.debug(
            "Cached OCR result for image",
            hash=img_hash[:16],
            text_length=len(ocr_text),
        )

    def clear_ocr_cache(self) -> None:
        """Clear the image OCR cache (call between documents)."""
        cache_size = len(self._image_ocr_cache)
        self._image_ocr_cache.clear()
        if cache_size > 0:
            logger.debug("Cleared OCR cache", cached_images=cache_size)

    def _optimize_image_for_ocr(self, img_data: bytes) -> bytes:
        """
        Optimize image before OCR to reduce processing time.

        - Resizes if larger than max dimension
        - Converts to grayscale for faster OCR
        - Compresses PNG output
        """
        try:
            from PIL import Image
            import io

            img = Image.open(io.BytesIO(img_data))

            # Resize if too large
            max_dim = self.optimization_config.max_image_dimension
            if max(img.size) > max_dim:
                ratio = max_dim / max(img.size)
                new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
                img = img.resize(new_size, Image.Resampling.LANCZOS)
                logger.debug(
                    "Resized image for OCR",
                    original_size=img.size,
                    new_size=new_size,
                )

            # Convert to grayscale for faster OCR
            if self.optimization_config.convert_to_grayscale:
                img = img.convert('L')

            # Compress output
            output = io.BytesIO()
            img.save(output, format='PNG', optimize=True)
            return output.getvalue()

        except Exception as e:
            logger.warning("Image optimization failed, using original", error=str(e))
            return img_data

    def _ocr_page(self, page, file_size_mb: float = 0) -> str:
        """
        OCR a PDF page with smart optimization and duplicate detection.

        Uses adaptive zoom based on file size and applies image optimization
        when smart_image_handling is enabled. Caches results to avoid
        re-OCR'ing duplicate images within the same document.
        """
        try:
            import fitz  # Import here since it's conditional

            # Calculate appropriate zoom based on file size
            if self.smart_image_handling:
                if file_size_mb > self.optimization_config.large_file_threshold_mb:
                    zoom = self.optimization_config.large_pdf_zoom
                    logger.debug(
                        "Using reduced zoom for large file",
                        file_size_mb=file_size_mb,
                        zoom=zoom,
                    )
                else:
                    zoom = self.optimization_config.default_zoom
            else:
                zoom = 2.0  # Original default

            # Render page to image
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
            img_data = pix.tobytes("png")

            # Apply image optimization if smart handling is enabled
            if self.smart_image_handling and self.optimization_config.compress_before_ocr:
                img_data = self._optimize_image_for_ocr(img_data)

            # Check cache for duplicate image
            cached_text = self._get_cached_ocr(img_data)
            if cached_text is not None:
                return cached_text

            # Save temp file for OCR
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
                f.write(img_data)
                temp_path = f.name

            text = self._ocr_image(temp_path)

            # Clean up
            os.unlink(temp_path)

            # Cache result for future duplicate detection
            self._cache_ocr_result(img_data, text)

            return text
        except Exception as e:
            logger.error("Page OCR failed", error=str(e))
            return ""

    def _ocr_pages_parallel(
        self,
        file_path: str,
        page_indices: List[int],
        file_size_mb: float = 0,
    ) -> Dict[int, str]:
        """
        OCR multiple PDF pages in parallel using ThreadPoolExecutor.

        Args:
            file_path: Path to the PDF file
            page_indices: List of page indices (0-based) to OCR
            file_size_mb: File size for optimization decisions

        Returns:
            Dict mapping page index to OCR text
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed
        import tempfile

        logger.info(
            "Starting parallel OCR",
            file_path=file_path,
            pages=len(page_indices),
            max_workers=self.max_ocr_workers,
        )

        results: Dict[int, str] = {}

        def ocr_single_page(page_idx: int) -> tuple:
            """OCR a single page (runs in thread pool) with duplicate detection."""
            try:
                import fitz

                # Calculate zoom
                if self.smart_image_handling:
                    if file_size_mb > self.optimization_config.large_file_threshold_mb:
                        zoom = self.optimization_config.large_pdf_zoom
                    else:
                        zoom = self.optimization_config.default_zoom
                else:
                    zoom = 2.0

                # Open document in this thread
                doc = fitz.open(file_path)
                page = doc[page_idx]

                # Render page to image
                pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
                img_data = pix.tobytes("png")

                # Apply image optimization
                if self.smart_image_handling and self.optimization_config.compress_before_ocr:
                    img_data = self._optimize_image_for_ocr(img_data)

                doc.close()

                # Check cache for duplicate image (thread-safe read)
                cached_text = self._get_cached_ocr(img_data)
                if cached_text is not None:
                    return (page_idx, cached_text)

                # Save temp file for OCR
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
                    f.write(img_data)
                    temp_path = f.name

                # Run OCR
                text = self._ocr_image(temp_path)

                # Clean up temp file
                os.unlink(temp_path)

                # Cache result (may have race condition but acceptable for performance)
                self._cache_ocr_result(img_data, text)

                return (page_idx, text)

            except Exception as e:
                logger.error(
                    "Parallel OCR failed for page",
                    page_idx=page_idx,
                    error=str(e),
                )
                return (page_idx, "")

        # Use ThreadPoolExecutor for parallel processing
        with ThreadPoolExecutor(max_workers=self.max_ocr_workers) as executor:
            # Submit all OCR tasks
            future_to_page = {
                executor.submit(ocr_single_page, idx): idx
                for idx in page_indices
            }

            # Collect results as they complete
            for future in as_completed(future_to_page):
                page_idx = future_to_page[future]
                try:
                    idx, text = future.result()
                    results[idx] = text
                except Exception as e:
                    logger.error(
                        "Parallel OCR task failed",
                        page_idx=page_idx,
                        error=str(e),
                    )
                    results[page_idx] = ""

        logger.info(
            "Parallel OCR complete",
            pages_processed=len(results),
            successful=sum(1 for v in results.values() if v),
        )

        return results

    # =========================================================================
    # Email
    # =========================================================================

    def _process_email(self, file_path: str, mode: str) -> ExtractedContent:
        """Process EML email files."""
        import email
        from email import policy

        with open(file_path, "rb") as f:
            msg = email.message_from_binary_file(f, policy=policy.default)

        parts = []
        parts.append(f"From: {msg.get('From', '')}")
        parts.append(f"To: {msg.get('To', '')}")
        parts.append(f"Subject: {msg.get('Subject', '')}")
        parts.append(f"Date: {msg.get('Date', '')}")
        parts.append("")

        # Get body
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    parts.append(part.get_content())
        else:
            parts.append(msg.get_content())

        return ExtractedContent(
            text="\n".join(parts),
            metadata={
                "file_type": "email",
                "file_path": file_path,
                "subject": msg.get("Subject", ""),
            },
            page_count=1,
        )

    def _process_msg(self, file_path: str, mode: str) -> ExtractedContent:
        """Process MSG Outlook files."""
        try:
            import extract_msg
            msg = extract_msg.Message(file_path)

            parts = []
            parts.append(f"From: {msg.sender}")
            parts.append(f"To: {msg.to}")
            parts.append(f"Subject: {msg.subject}")
            parts.append(f"Date: {msg.date}")
            parts.append("")
            parts.append(msg.body or "")

            return ExtractedContent(
                text="\n".join(parts),
                metadata={
                    "file_type": "msg",
                    "file_path": file_path,
                    "subject": msg.subject,
                },
                page_count=1,
            )
        except ImportError:
            return self._unstructured_fallback(file_path)

    # =========================================================================
    # CSV
    # =========================================================================

    def _process_csv(self, file_path: str, mode: str) -> ExtractedContent:
        """Process CSV files."""
        import csv

        rows = []
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" | ".join(row))

        return ExtractedContent(
            text="\n".join(rows),
            metadata={"file_type": "csv", "file_path": file_path},
            tables=[{"data": rows}],
            page_count=1,
        )

    # =========================================================================
    # Fallback
    # =========================================================================

    def _unstructured_fallback(self, file_path: str) -> ExtractedContent:
        """
        Fallback to Unstructured.io for unsupported or complex formats.
        """
        try:
            from unstructured.partition.auto import partition

            elements = partition(filename=file_path)
            text = "\n\n".join(str(el) for el in elements)

            return ExtractedContent(
                text=text,
                metadata={"file_type": "unstructured", "file_path": file_path},
                page_count=1,
            )
        except ImportError:
            logger.error("Unstructured not installed for fallback")
            raise ImportError(
                f"Cannot process {file_path}. "
                "Install 'unstructured' package for this file type."
            )
